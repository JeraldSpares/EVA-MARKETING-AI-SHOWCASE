# -*- coding: utf-8 -*-
"""
Enderun Marketing Hub - Web Dashboard
Run: python dashboard.py
Open: http://localhost:8080
"""

import os, csv, json, subprocess, calendar, re, threading, uuid, concurrent.futures, base64 as _b64_mod, io as _io_mod
from pathlib import Path
from datetime import date, timedelta
from functools import wraps
from flask import Flask, request, jsonify, Response, stream_with_context, send_file, send_from_directory, session, redirect
import anthropic, requests as http_requests

BASE_DIR      = Path(__file__).parent
LEADS_FILE    = BASE_DIR / "leads.csv"
SCHEDULE_FILE = BASE_DIR / "posting_schedule.json"
AGENTS_DIR    = BASE_DIR / "agents"
HTML_FILE     = BASE_DIR / "_dashboard.html"
INTEL_DIR     = BASE_DIR / "output" / "reports" / "social_listening"
PREVIEW_FILE  = BASE_DIR / "preview_state.json"
PROFILE_FILE  = BASE_DIR / "profile.json"
IMAGES_DIR      = Path(os.environ.get("IMAGES_DIR", r"G:\My Drive\FB_Post_Today"))
ANTHROPIC_KEY   = os.environ.get("ANTHROPIC_API_KEY", "")
CHAT_MEMORY_FILE   = BASE_DIR / "chat_memory.json"
CHAT_IMAGES_DIR    = BASE_DIR / "output" / "chat_images"
CHAT_IMAGES_DIR.mkdir(parents=True, exist_ok=True)
THUMB_CACHE_DIR    = BASE_DIR / "output" / "composer_thumbs"
THUMB_CACHE_DIR.mkdir(parents=True, exist_ok=True)
CHAT_SESSIONS_DIR  = BASE_DIR / "output" / "chat_sessions"
CHAT_SESSIONS_DIR.mkdir(parents=True, exist_ok=True)

# Background job tracker — {job_id: {"status": "running|ok|error", "output": ""}}
_jobs: dict = {}

def _run_bg(job_id, cmd, env, cwd):
    try:
        r = subprocess.run(cmd, capture_output=True, text=True, env=env, cwd=cwd, timeout=600)
        _jobs[job_id] = {"status": "ok" if r.returncode == 0 else "error", "output": r.stdout + r.stderr}
    except Exception as e:
        _jobs[job_id] = {"status": "error", "output": str(e)}

# Load password — prefer profile.json override, then env, then default
def _load_password():
    if PROFILE_FILE.exists():
        try:
            return json.loads(PROFILE_FILE.read_text()).get("password", "Enderun2026")
        except Exception:
            pass
    return os.environ.get("DASHBOARD_PASSWORD", "Enderun2026")

DASHBOARD_PASSWORD = _load_password()

# ── CHAT MEMORY ───────────────────────────────────────────────────────────────
def _load_chat_memory():
    if CHAT_MEMORY_FILE.exists():
        try: return json.loads(CHAT_MEMORY_FILE.read_text(encoding="utf-8"))
        except: pass
    return []

def _save_chat_memory(data):
    CHAT_MEMORY_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

import sys
sys.path.insert(0, str(BASE_DIR))
from notifications_helper import get_all, mark_read, mark_all_read, clear_all, unread_count

app = Flask(__name__)
app.secret_key = os.environ.get("DASHBOARD_SECRET", "enderun-hub-secret-2026")

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get("logged_in"):
            if request.path.startswith("/api/"):
                return jsonify({"error": "unauthorized"}), 401
            return redirect("/login")
        return f(*args, **kwargs)
    return decorated

# ── DATA ──────────────────────────────────────────────────────────────────────

def get_leads():
    if not LEADS_FILE.exists():
        return []
    with open(LEADS_FILE, newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))

def get_stats():
    from collections import Counter
    leads   = get_leads()
    active  = [l for l in leads if l.get("status","").strip().lower() == "active"]
    sched   = {}
    if SCHEDULE_FILE.exists():
        with open(SCHEDULE_FILE) as f:
            sched = json.load(f).get("schedule", {})
    today_str = date.today().isoformat()
    prog = Counter(l.get("program_interest","Unknown") for l in leads)
    total = sum(prog.values()) or 1
    programs_list = [{"name": k, "count": v, "pct": round(v/total*100)} for k, v in prog.most_common(6)]
    # Lead scoring
    hot  = sum(1 for l in active if int(l.get("email_count",0) or 0) >= 10)
    warm = sum(1 for l in active if 5 <= int(l.get("email_count",0) or 0) < 10)
    cold = sum(1 for l in active if int(l.get("email_count",0) or 0) < 5)
    # Preview state
    preview_status, preview_week = "none", ""
    if PREVIEW_FILE.exists():
        try:
            ps = json.loads(PREVIEW_FILE.read_text())
            preview_status = ps.get("status", "none")
            preview_week   = ps.get("week_label", "")
        except Exception:
            pass
    # Build upcoming posts list
    upcoming = []
    for dt in sorted(sched.keys()):
        upcoming.append({"date": dt, "image": sched[dt], "status": "Posted" if dt < today_str else "Scheduled"})
    return {
        "total_leads":       len(leads),
        "active_leads":      len(active),
        "hot_leads":         hot,
        "warm_leads":        warm,
        "cold_leads":        cold,
        "posts_total":       len(sched),
        "posts_scheduled":   len(sched),
        "today_post":        sched.get(today_str, "Not scheduled"),
        "top_program":       prog.most_common(1)[0][0] if prog else "N/A",
        "today":             date.today().strftime("%B %d, %Y"),
        "program_breakdown": dict(prog.most_common(6)),
        "programs":          programs_list,
        "upcoming":          upcoming,
        "preview_status":    preview_status,
        "preview_week":      preview_week,
    }

def _clean_analysis(raw: str) -> str:
    """Strip separators, leading emoji from headings, arrows, and redundant sections."""
    lines = []
    skip_section = False
    for line in raw.split("\n"):
        s = line.strip()
        # Drop pure separator lines (--- ===)
        if re.match(r"^[-=]{3,}$", s):
            continue
        # Skip threat-level section (shown as badge instead)
        if re.match(r"^#{1,3}\s*\S*\s*THREAT", s, re.IGNORECASE):
            skip_section = True
            continue
        # Skip morning-briefing title (redundant with page header)
        if re.match(r"^#{1,3}\s*\S*\s*MORNING INTELLIGENCE BRIEFING", s, re.IGNORECASE):
            continue
        # Next heading resets skip
        if re.match(r"^#{1,3}\s+", s) and skip_section:
            skip_section = False
        if skip_section:
            continue
        # Strip leading emoji (non-ASCII) from section headings  ## 🚨 ALERTS → ## ALERTS
        line = re.sub(r"^(#{1,3}\s*)(?:[^\x00-\x7F]\s*)+", r"\1", line)
        # Strip → arrow
        line = line.replace("→ ", "").replace("→", "")
        # Strip inline threat-level emoji
        line = re.sub(r"[🟢🟡🔴]\s*", "", line)
        lines.append(line)
    result = "\n".join(lines)
    result = re.sub(r"\n{3,}", "\n\n", result)
    result = result.strip()
    # Close any trailing unclosed bold marker from truncated AI output
    result = re.sub(r'\*\*\s*$', '', result).strip()
    return result

_INTEL_GH_CACHE: dict = {}   # filename → text, in-memory cache for Railway

def _gh_intel_headers():
    pat = os.environ.get("GITHUB_PAT", "")
    repo = os.environ.get("GITHUB_REPO", "your-org/your-repo")
    h = {"Accept": "application/vnd.github+json"}
    if pat:
        h["Authorization"] = f"Bearer {pat}"
    return h, repo

def _list_intel_files_github():
    """Return sorted list of briefing filenames from GitHub (newest first)."""
    import base64 as _b64
    h, repo = _gh_intel_headers()
    try:
        r = http_requests.get(
            f"https://api.github.com/repos/{repo}/contents/output/reports/social_listening",
            headers=h, timeout=10
        )
        if r.status_code != 200:
            return []
        return sorted(
            [f["name"] for f in r.json() if f["name"].endswith("_intelligence_briefing.txt")],
            reverse=True
        )
    except Exception:
        return []

def _read_intel_file_github(filename):
    """Return text content of a single briefing file from GitHub."""
    import base64 as _b64
    if filename in _INTEL_GH_CACHE:
        return _INTEL_GH_CACHE[filename]
    h, repo = _gh_intel_headers()
    try:
        r = http_requests.get(
            f"https://api.github.com/repos/{repo}/contents/output/reports/social_listening/{filename}",
            headers=h, timeout=10
        )
        if r.status_code != 200:
            return None
        text = _b64.b64decode(r.json().get("content", "")).decode("utf-8", errors="replace")
        _INTEL_GH_CACHE[filename] = text
        return text
    except Exception:
        return None

def _parse_intel_text(text, filename, all_filenames):
    """Parse raw briefing text into the intelligence data dict."""
    date_str = filename.replace("_intelligence_briefing.txt", "")
    raw_analysis = text
    if "AI ANALYSIS" in text and "RAW DATA" in text:
        parts = text.split("RAW DATA", 1)
        raw_analysis = "\n".join(
            l for l in parts[0].split("\n")
            if not any(h in l for h in ["ENDERUN INTELLIGENCE BRIEFING", "Generated:", "AI ANALYSIS"])
        )
    analysis = _clean_analysis(raw_analysis)
    t = text.upper()
    if "🔴" in text or "CRITICAL ALERT" in t or "HIGH ALERT" in t or "IMMEDIATE ACTION" in t:
        threat = "high"
    elif "🟢" in text or ("NO CRITICAL ALERTS" in t and "MEDIUM" not in t):
        threat = "low"
    elif "🟡" in text or "MEDIUM" in t or "WATCH ITEM" in t:
        threat = "medium"
    else:
        threat = "low"
    changes = []
    for line in text.splitlines():
        ls = line.strip()
        if ls.startswith("- ") and ": http" in ls:
            p = ls[2:].split(": ", 1)
            if len(p) == 2:
                changes.append({"competitor": p[0].lstrip("*").rstrip("*").strip(), "url": p[1].strip()})
    news = [l.strip()[2:] for l in text.splitlines() if l.strip().startswith("• ")]
    sources = list(dict.fromkeys(
        re.search(r"\(([^)]+)\)$", n).group(1)
        for n in news if re.search(r"\(([^)]+)\)$", n)
    ))
    all_dates = [f.replace("_intelligence_briefing.txt", "") for f in all_filenames[:10]]
    return {
        "date": date_str, "analysis": analysis, "threat": threat,
        "changes": changes, "news": news, "sources": sources,
        "source_file": filename, "all_dates": all_dates,
    }

def get_intelligence(date_filter=None):
    # ── Try local files first (works locally and in CI) ──
    local_files = []
    if INTEL_DIR.exists():
        local_files = sorted(INTEL_DIR.glob("*_intelligence_briefing.txt"), reverse=True)
    if local_files:
        if date_filter:
            match = [f for f in local_files if f.stem.startswith(date_filter)]
            target = match[0] if match else None
        else:
            target = local_files[0]
        if not target:
            return None
        text = target.read_text(encoding="utf-8")
        all_names = [f.name for f in local_files]
        return _parse_intel_text(text, target.name, all_names)

    # ── Fallback: fetch from GitHub API (Railway / no local files) ──
    gh_files = _list_intel_files_github()
    if not gh_files:
        return None
    if date_filter:
        match = [f for f in gh_files if f.startswith(date_filter)]
        target_name = match[0] if match else None
    else:
        target_name = gh_files[0]
    if not target_name:
        return None
    text = _read_intel_file_github(target_name)
    if not text:
        return None
    return _parse_intel_text(text, target_name, gh_files)

def get_automations_status():
    today     = date.today()
    today_str = today.isoformat()
    sched     = get_schedule()
    # Social Listening
    intel    = get_intelligence()
    sl_last  = intel["date"] if intel else None
    sl_ok    = sl_last == today_str
    # Drip Email
    leads    = get_leads()
    active   = [l for l in leads if l.get("status","").strip().lower() == "active"]
    avg_step = round(sum(int(l.get("email_count",0) or 0) for l in active) / len(active), 1) if active else 0
    # Facebook
    fb_today = sched.get(today_str)
    # Preview
    preview_status = preview_week = ""
    if PREVIEW_FILE.exists():
        try:
            ps = json.loads(PREVIEW_FILE.read_text())
            preview_status = ps.get("status","none")
            preview_week   = ps.get("week_label","")
        except Exception:
            pass
    # Next Monday
    days_to_mon = (7 - today.weekday()) % 7 or 7
    next_monday = (today + timedelta(days=days_to_mon)).strftime("%B %d, %Y")
    # Next month-end
    last_day   = calendar.monthrange(today.year, today.month)[1]
    month_end  = today.replace(day=last_day).strftime("%B %d, %Y")
    return {
        "social_listening": {
            "name": "Social Listening", "emoji": "🔍",
            "status": "success" if sl_ok else ("warning" if sl_last else "inactive"),
            "detail": f"Last: {sl_last}" if sl_last else "Not run today",
            "schedule": "Daily 7:50 AM PHT",
        },
        "drip_email": {
            "name": "Drip Email", "emoji": "📧",
            "status": "success" if active else "inactive",
            "detail": f"{len(active)} active leads · avg step {avg_step}",
            "schedule": "Daily 8:00 AM PHT",
        },
        "facebook_post": {
            "name": "Facebook Post", "emoji": "📘",
            "status": "success" if fb_today else "warning",
            "detail": fb_today if fb_today else "Not scheduled today",
            "schedule": "Daily 8:00 AM PHT",
        },
        "instagram_post": {
            "name": "Instagram Post", "emoji": "📸",
            "status": "success" if fb_today else "warning",
            "detail": fb_today if fb_today else "Not scheduled today",
            "schedule": "Daily 8:00 AM PHT",
        },
        "weekly_preview": {
            "name": "Weekly Preview", "emoji": "📋",
            "status": preview_status or "none",
            "detail": f"{preview_week} · {preview_status}" if preview_week else "No preview yet",
            "schedule": "Sunday 5:00 PM PHT",
        },
        "cold_reengagement": {
            "name": "Cold Re-engagement", "emoji": "🧊",
            "status": "scheduled",
            "detail": f"Next: {next_monday}",
            "schedule": "Monday 6:00 AM PHT",
        },
        "monthly_report": {
            "name": "Monthly Report", "emoji": "📊",
            "status": "scheduled",
            "detail": f"Next: {month_end}",
            "schedule": "Last day of month 8:00 AM PHT",
        },
    }

def get_schedule():
    if not SCHEDULE_FILE.exists():
        return {}
    with open(SCHEDULE_FILE) as f:
        return json.load(f).get("schedule", {})

def list_agents():
    if not AGENTS_DIR.exists():
        return []
    ICONS = {
        "marketing-manager":"lightning","social-media":"phone","drip-campaign":"email",
        "content-strategy":"calendar","seo-digital":"search","designer":"palette",
        "pr":"newspaper","competitor-analysis":"microscope","data-analysis":"chart",
        "marketing-analysis":"graph","business-analyst":"briefcase","admissions":"grad",
        "video-multimedia":"video","events-activations":"event","events-banquetes":"glass",
        "alumni-relations":"people","influencer-kol":"star","researcher":"magnify",
    }
    LABELS = {
        "marketing-manager":"Marketing Manager","social-media":"Social Media",
        "drip-campaign":"Drip Campaign","content-strategy":"Content Strategy",
        "seo-digital":"SEO & Digital","designer":"Creative Director",
        "pr":"PR Manager","competitor-analysis":"Competitor Intel",
        "data-analysis":"Data Analyst","marketing-analysis":"Marketing Analyst",
        "business-analyst":"Business Analyst","admissions":"Admissions",
        "video-multimedia":"Video & Multimedia","events-activations":"Events & Activations",
        "events-banquetes":"Events / Banquetes","alumni-relations":"Alumni Relations",
        "influencer-kol":"Influencer & KOL","researcher":"Researcher",
    }
    EMOJI = {
        "marketing-manager":"⚡","social-media":"📱","drip-campaign":"📧",
        "content-strategy":"📅","seo-digital":"🔍","designer":"🎨",
        "pr":"📰","competitor-analysis":"🔬","data-analysis":"📊",
        "marketing-analysis":"📈","business-analyst":"💼","admissions":"🎓",
        "video-multimedia":"🎬","events-activations":"🗓","events-banquetes":"🥂",
        "alumni-relations":"🤝","influencer-kol":"⭐","researcher":"🔎",
    }
    result = []
    for f in sorted(AGENTS_DIR.glob("*.md")):
        k = f.stem
        result.append({"id": k, "label": LABELS.get(k, k.replace("-"," ").title()), "icon": EMOJI.get(k,"🤖")})
    return result

def agent_system_prompt(agent_id):
    claude_md = BASE_DIR / "CLAUDE.md"
    agent_md  = AGENTS_DIR / f"{agent_id}.md"

    chat_rules = (
        "\n\n---\n\nDASHBOARD CHAT RULES (always follow):\n"
        "- Be conversational, warm, and direct. This is a chat — not a formal report.\n"
        "- NEVER include hashtags (e.g. #EnderunExtension) in your responses.\n"
        "- Avoid unnecessary ## headers for short answers. Use bold text instead.\n"
        "- Keep responses concise: 2-4 short paragraphs or a brief bullet list.\n"
        "- For comparisons, rankings, schedules, or structured data: use markdown tables (| col | col |) — they render beautifully in chat.\n"
        "- Write naturally, like messaging a brilliant colleague.\n"
        "- You can answer questions on ANY topic — math, science, history, physics, English, "
        "coding, research, philosophy, economics, law, or anything else. "
        "Never say 'I can only help with marketing.'\n"
        "- If the user writes in Filipino or Taglish, respond naturally in the same language.\n"
        "- For math or science problems, show your work step by step.\n\n"
        "FILE GENERATION (when user asks to create a PDF, presentation, Word doc, or CSV):\n"
        "- Wrap any planning/thinking in [THINKING]...[/THINKING] — hidden in a collapsible.\n"
        "- After the block, write ONE short sentence confirming what you created.\n"
        "- Do NOT write the document content as plain prose outside the block.\n\n"
        "PDF — use a structured JSON spec inside [GENERATE:pdf]...[/GENERATE]. "
        "Always produce a DESIGNED report with KPI cards, styled tables, charts, and alert boxes. "
        "JSON format:\n"
        '  [GENERATE:pdf]\n'
        '  {\n'
        '    "title": "Report Title",\n'
        '    "subtitle": "Subtitle or date",\n'
        '    "sections": [\n'
        '      {"type": "section", "title": "Section Heading"},\n'
        '      {"type": "paragraph", "text": "Introductory paragraph."},\n'
        '      {"type": "kpi_row", "items": [\n'
        '        {"label": "Total Leads", "value": "142", "highlight": false},\n'
        '        {"label": "Active", "value": "89", "highlight": true},\n'
        '        {"label": "Rate", "value": "63%", "highlight": false}\n'
        '      ]},\n'
        '      {"type": "table", "headers": ["Column A", "Column B"], "rows": [["Row 1", "Val 1"]]},\n'
        '      {"type": "bullets", "items": ["Insight 1", "Insight 2"]},\n'
        '      {"type": "alert", "text": "Important note.", "level": "warning"},\n'
        '      {"type": "bar_chart", "title": "Chart Title", "labels": ["A","B","C"], "values": [10,20,30], "color": "brown"},\n'
        '      {"type": "pie_chart", "title": "Split", "labels": ["X","Y"], "values": [60,40]},\n'
        '      {"type": "line_chart", "title": "Trend", "labels": ["Jan","Feb"], "series": [[10,20]], "series_labels": ["Leads"]},\n'
        '      {"type": "callout", "stat": "98%", "label": "Graduate Employment Rate", "description": "Among 2025 batch"},\n'
        '      {"type": "page_break"},\n'
        '      {"type": "source", "text": "Source description here"}\n'
        '    ]\n'
        '  }\n'
        '  [/GENERATE]\n'
        "- alert levels: warning | critical | success | info\n"
        "- bar_chart color: brown | gold | multi\n"
        "- bar_chart horizontal: true for long labels\n"
        "- charts_row: {\"type\":\"charts_row\",\"left\":{\"type\":\"bar\",...},\"right\":{\"type\":\"pie\",...}}\n"
        "- ALWAYS include at least one kpi_row, one table, and one chart in every PDF.\n\n"
        "PPTX/DOCX — use markdown inside the block:\n"
        "  [GENERATE:pptx] or [GENERATE:docx]\n"
        "  # Document Title\n"
        "  ## Section Heading\n"
        "  Content paragraphs and - bullet points here\n"
        "  [/GENERATE]\n"
        "CSV — output raw CSV rows inside [GENERATE:csv]...[/GENERATE]."
    )

    # ── General AI mode — unrestricted intelligence, no marketing constraint ──
    if agent_id == "general":
        return (
            "You are Eva — a world-class AI assistant with expert-level knowledge across every field of human knowledge.\n\n"
            "You can help with absolutely anything:\n"
            "• Mathematics — algebra, calculus, statistics, geometry, number theory, proofs, discrete math\n"
            "• Physics — mechanics, thermodynamics, electromagnetism, quantum, relativity, astrophysics\n"
            "• Chemistry — organic, inorganic, biochemistry, stoichiometry, reactions\n"
            "• Biology — genetics, cell biology, evolution, anatomy, ecology, microbiology\n"
            "• History — world history, Philippine history, wars, civilizations, timelines\n"
            "• English — grammar, writing, literature analysis, essays, vocabulary, style\n"
            "• Computer Science — programming (Python, JS, SQL, etc.), algorithms, data structures, AI/ML\n"
            "• Economics & Business — macro, micro, finance, accounting, strategy\n"
            "• Philosophy & Logic — critical thinking, ethics, arguments, reasoning\n"
            "• Geography, Political Science, Social Studies\n"
            "• Arts, Music, Culture, Pop culture\n"
            "• Research, fact-checking, summarizing, analysis — any domain\n"
            "• Filipino language, Tagalog grammar, translation\n"
            "• And literally anything else a person could ask\n\n"
            "HOW TO RESPOND:\n"
            "- Be warm, direct, and conversational — like a brilliant friend who happens to know everything\n"
            "- For math/science/logic: show clear step-by-step solutions with explanations\n"
            "- For research/analysis: be thorough, well-organized, accurate\n"
            "- For simple questions: be concise — no padding\n"
            "- For creative tasks: be imaginative, original, and engaged\n"
            "- Match the user's language — respond in Tagalog/Taglish if they write that way\n"
            "- Use analogies and real-world examples to explain hard concepts\n"
            "- Never refuse to engage with an academic or intellectual question\n\n"
            "TOOLS YOU HAVE: Use them proactively for the right tasks.\n"
            "- web_search: real-time internet search (news, facts, prices, research)\n"
            "- execute_python: run any Python code — math, data analysis, charts\n"
            "- generate_image: create images from text prompts (free AI, no key needed)\n"
            "- get_calendar / add_calendar_event: Google Calendar\n"
            "- search_email: Gmail search\n"
            "- remember / recall_memories / forget_memory: persistent memory\n"
            "- read_drive_file / list_drive_files: Google Drive\n"
            "- Marketing tools: post to FB/IG, send drip emails, run reports\n\n"
            "CONTEXT: You are deployed at Enderun Colleges (BGC, Philippines). "
            "You understand the Enderun brand and can execute marketing tasks with your tools. "
            "But your intelligence is not limited to marketing in any way."
        ) + chat_rules

    # ── Marketing agents ──
    base = claude_md.read_text(encoding="utf-8") if claude_md.exists() else ""
    if agent_md.exists():
        return base + "\n\n---\n\nYou are now acting as:\n\n" + agent_md.read_text(encoding="utf-8") + chat_rules
    return (base or "You are the AI Marketing Assistant for Enderun Colleges.") + chat_rules

# ── ROUTES ────────────────────────────────────────────────────────────────────

LOGIN_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>Enderun Marketing Hub — Login</title>
<link rel="icon" type="image/png" href="/assets/logos/Enderun-Colleges-white.png">
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.0/css/all.min.css" crossorigin="anonymous">
<style>
*{box-sizing:border-box;margin:0;padding:0}
body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;height:100vh;display:flex;overflow:hidden;background:#0c0608}

/* ── LEFT PANEL ── */
.left{position:relative;flex:1.2;background:#0c0608;overflow:hidden}

.left-logo{position:absolute;top:26px;left:28px;z-index:10}
.left-logo img{height:80px;width:auto}

.robot-wrap{position:absolute;inset:0;bottom:-50px;overflow:hidden;background:#0c0608}
.robot-wrap iframe{width:100%;height:calc(100% + 50px);border:none;display:block;background:#0c0608;transform:scale(1.28) translateX(8%) translateY(-8%);transform-origin:center 55%}

/* ── RIGHT PANEL ── */
.right{flex:.85;background:#0c0608;display:flex;align-items:center;justify-content:center;padding:60px 44px}
.form-wrap{width:100%;max-width:380px;animation:form-slide-in 0.6s cubic-bezier(0.16,1,0.3,1) .15s both;position:relative}
@keyframes form-slide-in{from{opacity:0;transform:translateX(20px)}to{opacity:1;transform:translateX(0)}}
.form-card{background:linear-gradient(165deg,rgba(22,5,11,.98) 0%,rgba(9,2,5,.97) 100%);backdrop-filter:blur(24px);border:1px solid rgba(156,26,53,.2);border-radius:20px;padding:40px 42px 34px;box-shadow:0 30px 70px rgba(0,0,0,.75),0 0 0 1px rgba(0,0,0,.5),0 0 50px rgba(156,26,53,.07),inset 0 1px 0 rgba(255,255,255,.05);position:relative;overflow:hidden}
.form-card::before{content:'';position:absolute;inset:-1px;border-radius:21px;background:linear-gradient(#b01535,#b01535) top left/26px 2px no-repeat,linear-gradient(#b01535,#b01535) top left/2px 26px no-repeat,linear-gradient(#b01535,#b01535) top right/26px 2px no-repeat,linear-gradient(#b01535,#b01535) top right/2px 26px no-repeat,linear-gradient(#b01535,#b01535) bottom left/26px 2px no-repeat,linear-gradient(#b01535,#b01535) bottom left/2px 26px no-repeat,linear-gradient(#b01535,#b01535) bottom right/26px 2px no-repeat,linear-gradient(#b01535,#b01535) bottom right/2px 26px no-repeat;pointer-events:none;z-index:2}
.form-card::after{content:'';position:absolute;left:0;right:0;height:80px;background:linear-gradient(transparent,rgba(156,26,53,.04),transparent);animation:scan-line 7s linear infinite;pointer-events:none}
@keyframes scan-line{from{top:-80px}to{top:100%}}

.form-logo{text-align:center;margin-bottom:26px;padding-bottom:24px;position:relative}
.form-logo::after{content:'';position:absolute;bottom:0;left:50%;transform:translateX(-50%);width:180px;height:1px;background:linear-gradient(90deg,transparent,rgba(156,26,53,.5),transparent)}
.form-logo img{width:140px;height:auto;filter:brightness(0) invert(1);drop-shadow:0 0 12px rgba(255,255,255,.15)}

.greeting{font-size:10px;color:#c0392b;font-weight:700;text-transform:uppercase;letter-spacing:.16em;margin-bottom:6px;font-family:'Courier New',monospace;display:flex;align-items:center;gap:7px}
.greeting::before{content:'';width:14px;height:2px;background:#c0392b;border-radius:1px;flex-shrink:0}
.greeting::after{content:'|';animation:cursor-blink 1.2s infinite;opacity:.7;font-weight:300}
@keyframes cursor-blink{0%,49%{opacity:.7}50%,100%{opacity:0}}
.form-heading{font-size:25px;font-weight:800;color:#fff;margin-bottom:5px;letter-spacing:-.02em}
.form-hint{font-size:12.5px;color:rgba(255,255,255,.38);margin-bottom:26px;line-height:1.5}

.field-label{font-size:10px;font-weight:700;color:rgba(190,70,85,.95);text-transform:uppercase;letter-spacing:.12em;display:block;margin-bottom:7px;font-family:'Courier New',monospace}
.field-wrap{position:relative;margin-bottom:16px}
.input-icon{position:absolute;left:14px;top:50%;transform:translateY(-50%);color:rgba(255,255,255,.45);font-size:12px;pointer-events:none;z-index:1}
.eye-toggle{position:absolute;right:14px;top:50%;transform:translateY(-50%);color:rgba(255,255,255,.45);font-size:12px;background:none;border:none;cursor:pointer;padding:2px 4px;transition:color .2s;z-index:1}
.eye-toggle:hover{color:rgba(192,57,43,.9)}
.field-wrap input{width:100%;background:rgba(255,255,255,.04);border:1px solid rgba(255,255,255,.09);border-radius:11px;padding:13px 42px;color:#fff;font-size:14px;font-family:inherit;outline:none;transition:all .25s}
.field-wrap input::placeholder{color:rgba(255,255,255,.2)}
.field-wrap input:focus{border-color:rgba(156,26,53,.65);box-shadow:0 0 0 3px rgba(156,26,53,.12),0 0 20px rgba(156,26,53,.09);background:rgba(255,255,255,.06)}
.field-wrap input:focus~.input-icon,.field-wrap input:focus+.input-icon{color:rgba(156,26,53,.7)}
.field-wrap input:-webkit-autofill,.field-wrap input:-webkit-autofill:hover,.field-wrap input:-webkit-autofill:focus{-webkit-box-shadow:0 0 0 1000px #120409 inset !important;-webkit-text-fill-color:#fff !important;caret-color:#fff;border-color:rgba(156,26,53,.65) !important;transition:background-color 9999s ease-in-out 0s}

.btn-signin{width:100%;padding:14px;background:linear-gradient(135deg,#7a1028 0%,#b01535 50%,#7a1028 100%);background-size:200% auto;color:#fff;font-size:14px;font-weight:700;border:none;border-radius:11px;cursor:pointer;font-family:inherit;transition:all .3s ease;letter-spacing:.04em;text-transform:uppercase;font-size:13px;box-shadow:0 4px 24px rgba(156,26,53,.45),inset 0 1px 0 rgba(255,255,255,.1);margin-top:8px}
.btn-signin:hover:not(:disabled){background-position:right center;transform:translateY(-2px);box-shadow:0 10px 36px rgba(156,26,53,.6),0 0 30px rgba(156,26,53,.18),inset 0 1px 0 rgba(255,255,255,.15)}
.btn-signin:active:not(:disabled){transform:translateY(0);box-shadow:0 4px 16px rgba(156,26,53,.4)}
.btn-signin:disabled{opacity:.65;cursor:not-allowed}

.error{background:rgba(220,38,38,.08);border:1px solid rgba(220,38,38,.25);border-radius:10px;padding:10px 14px;font-size:12px;color:#f87171;margin-bottom:14px;display:flex;align-items:center;gap:8px}
.powered{margin-top:22px;text-align:center;font-size:10px;color:rgba(255,255,255,.18);letter-spacing:.06em;display:flex;align-items:center;justify-content:center;gap:6px}
.powered span{color:#9b1a35;font-weight:600}
.powered-dot{width:5px;height:5px;border-radius:50%;background:#1db954;box-shadow:0 0 6px #1db954;animation:pulse-dot 2s infinite}
@keyframes pulse-dot{0%,100%{opacity:1}50%{opacity:.4}}

/* ── HEX CORNERS ── */
.hex-tl,.hex-tr,.hex-br,.hex-form-accent,.hex-left-bot{background-image:url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='52' height='90'%3E%3Cg fill='none' stroke='rgba(156,26,53,0.4)' stroke-width='1.5'%3E%3Cpolygon points='26,0 52,15 52,45 26,60 0,45 0,15'/%3E%3Cpolygon points='0,45 26,60 26,90 0,105 -26,90 -26,60'/%3E%3Cpolygon points='52,45 78,60 78,90 52,105 26,90 26,60'/%3E%3C/g%3E%3Cg fill='rgba(156,26,53,0.55)'%3E%3Ccircle cx='26' cy='0' r='2.5'/%3E%3Ccircle cx='52' cy='15' r='2.5'/%3E%3Ccircle cx='52' cy='45' r='2.5'/%3E%3Ccircle cx='26' cy='60' r='2.5'/%3E%3Ccircle cx='0' cy='45' r='2.5'/%3E%3Ccircle cx='0' cy='15' r='2.5'/%3E%3Ccircle cx='26' cy='90' r='2.5'/%3E%3C/g%3E%3C/svg%3E");pointer-events:none;z-index:6}
.hex-tl,.hex-tr,.hex-br{position:fixed;width:680px;height:680px;background-size:76px 131px}
.hex-form-accent{position:absolute;width:260px;height:260px;background-size:52px 90px;top:-55px;left:-55px;z-index:0;-webkit-mask-image:radial-gradient(ellipse 75% 75% at 0% 0%,black 0%,transparent 82%);mask-image:radial-gradient(ellipse 75% 75% at 0% 0%,black 0%,transparent 82%)}
.hex-tl{top:0;left:0;-webkit-mask-image:radial-gradient(ellipse 80% 80% at 0% 0%,black 0%,transparent 70%);mask-image:radial-gradient(ellipse 80% 80% at 0% 0%,black 0%,transparent 70%);animation:hex-fade 7s ease-in-out infinite}
.hex-tr{top:0;right:0;-webkit-mask-image:radial-gradient(ellipse 80% 80% at 100% 0%,black 0%,transparent 70%);mask-image:radial-gradient(ellipse 80% 80% at 100% 0%,black 0%,transparent 70%);animation:hex-fade 7s ease-in-out infinite 2.5s}
.hex-br{bottom:0;right:0;-webkit-mask-image:radial-gradient(ellipse 80% 80% at 100% 100%,black 0%,transparent 70%);mask-image:radial-gradient(ellipse 80% 80% at 100% 100%,black 0%,transparent 70%);animation:hex-fade 7s ease-in-out infinite 5s}
.hex-left-bot{position:absolute;bottom:0;left:0;width:520px;height:480px;background-size:76px 131px;z-index:7;-webkit-mask-image:radial-gradient(ellipse 85% 80% at 10% 100%,black 0%,transparent 72%);mask-image:radial-gradient(ellipse 85% 80% at 10% 100%,black 0%,transparent 72%);animation:hex-float 8s ease-in-out infinite,hex-fade 7s ease-in-out infinite 1.2s}
@keyframes hex-float{0%,100%{transform:translateY(0)}50%{transform:translateY(-6px)}}
@keyframes hex-fade{0%,100%{opacity:.15}45%,55%{opacity:.75}}
.hex-tl::before,.hex-tr::before,.hex-br::before,.hex-left-bot::before{content:'';position:absolute;inset:0;background:linear-gradient(110deg,transparent 20%,rgba(156,26,53,.22) 48%,rgba(200,40,70,.12) 52%,transparent 80%);background-size:280% 100%;pointer-events:none}
.hex-tl::before{animation:hex-wave 10s linear infinite}
.hex-tr::before{animation:hex-wave 10s linear infinite 3.5s}
.hex-br::before{animation:hex-wave 10s linear infinite 7s}
.hex-left-bot::before{animation:hex-wave 10s linear infinite 1.8s}
@keyframes hex-wave{0%{background-position:-100% 0}100%{background-position:200% 0}}

/* ── 3D MOLECULE CANVAS ── */
#mol-canvas{position:absolute;inset:0;width:100%;height:100%;z-index:8;pointer-events:none;opacity:.9}

/* ── LOGIN EXIT TRANSITION ── */
@keyframes page-fade-out{0%{opacity:1;filter:blur(0) scale(1)}100%{opacity:0;filter:blur(6px)}}
@keyframes overlay-sweep{0%{opacity:0}100%{opacity:1}}
@keyframes btn-charge{0%{box-shadow:0 4px 24px rgba(156,26,53,.45)}50%{box-shadow:0 0 0 6px rgba(156,26,53,.18),0 0 40px rgba(156,26,53,.5)}100%{box-shadow:0 4px 24px rgba(156,26,53,.45)}}
.exit-overlay{position:fixed;inset:0;background:#0c0608;z-index:9999;pointer-events:none;opacity:0}
.exit-overlay.active{animation:overlay-sweep .5s cubic-bezier(0.4,0,1,1) .12s forwards}
.page-exiting{animation:page-fade-out .65s cubic-bezier(0.4,0,1,1) forwards}
.btn-charging{animation:btn-charge .35s ease forwards}

@media(max-width:700px){.left{display:none}.right{flex:1;padding:36px 24px}}
</style>
</head>
<body>
<div class="hex-tl"></div>
<div class="hex-tr"></div>
<div class="hex-br"></div>

<div class="left">
  <div class="left-logo">
    <img src="/assets/logos/Enderun-Colleges-white.png" alt="Enderun Colleges">
  </div>
  <canvas id="mol-canvas"></canvas>
  <div class="hex-left-bot"></div>
  <div class="robot-wrap">
    <iframe src='https://my.spline.design/evefromdisneyswallesplinecoursecom-M58wmXZdxz3egCuHwkYQ2Coa/' frameborder='0' width='100%' height='100%'></iframe>
  </div>
</div>

<div class="right">
  <div class="form-wrap">
    <div class="hex-form-accent"></div>
    <div class="form-card">
      <div class="form-logo">
        <img src="/assets/Enderun-Colleges-Logo-Normal-600x291.png" alt="Enderun Colleges">
      </div>
      <div id="greeting" class="greeting"></div>
      <div class="form-heading">Welcome back</div>
      <div class="form-hint">Sign in to access your marketing hub</div>
      {error}
      <form method="POST" action="/login" id="login-form">
        <label class="field-label"><i class="fa-solid fa-lock" style="margin-right:5px"></i>Password</label>
        <div class="field-wrap">
          <i class="fa-solid fa-lock input-icon"></i>
          <input type="password" name="password" id="pw-input" placeholder="Enter your password" autofocus>
          <button type="button" class="eye-toggle" onclick="togglePw()" title="Show/hide password">
            <i class="fa-solid fa-eye" id="eye-icon"></i>
          </button>
        </div>
        <button class="btn-signin" type="submit" id="signin-btn">
          <i class="fa-solid fa-arrow-right-to-bracket" style="margin-right:7px"></i>Sign In
        </button>
      </form>
      <div class="powered"><span class="powered-dot"></span><i class="fa-solid fa-shield-halved"></i>&nbsp;Secured · Powered by <span>Claude AI</span></div>
    </div>
  </div>
</div>

<script>
var h = new Date().getHours();
document.getElementById('greeting').textContent = (h < 12 ? 'Good morning' : h < 18 ? 'Good afternoon' : 'Good evening') + ', Eva';

function togglePw() {
  var inp = document.getElementById('pw-input');
  var ico = document.getElementById('eye-icon');
  if (inp.type === 'password') {
    inp.type = 'text';
    ico.className = 'fa-solid fa-eye-slash';
  } else {
    inp.type = 'password';
    ico.className = 'fa-solid fa-eye';
  }
}

document.getElementById('login-form').addEventListener('submit', function(e) {
  e.preventDefault();
  var form = this;
  var btn = document.getElementById('signin-btn');
  btn.disabled = true;
  btn.classList.add('btn-charging');
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin" style="margin-right:7px"></i>Signing in...';

  // Create dark overlay and fade out page
  var overlay = document.createElement('div');
  overlay.className = 'exit-overlay';
  document.body.appendChild(overlay);
  requestAnimationFrame(function() {
    requestAnimationFrame(function() {
      overlay.classList.add('active');
      document.body.classList.add('page-exiting');
    });
  });

  setTimeout(function() { form.submit(); }, 700);
});

// ── 3D MOLECULE NETWORK ──
(function(){
  var C=document.getElementById('mol-canvas');
  if(!C)return;
  var ctx=C.getContext('2d');
  function resize(){C.width=C.offsetWidth;C.height=C.offsetHeight;}
  resize();
  window.addEventListener('resize',resize);

  var N=72,LINK=155,SPREAD=280;
  var nodes=Array.from({length:N},function(){
    // distribute in a sphere
    var u=Math.random(),v=Math.random();
    var th=2*Math.PI*u, ph=Math.acos(2*v-1);
    var rad=SPREAD*(Math.cbrt(Math.random())*.7+.3);
    return {
      x:rad*Math.sin(ph)*Math.cos(th),
      y:rad*Math.sin(ph)*Math.sin(th),
      z:rad*Math.cos(ph),
      r:Math.random()*1.6+.7,
      p:Math.random()*Math.PI*2
    };
  });

  var t=0;
  function draw(){
    ctx.clearRect(0,0,C.width,C.height);
    var cx=C.width/2, cy=C.height*.52, fov=380;
    var ay=t*.00022, ax=Math.sin(t*.000065)*.28, az=Math.sin(t*.000042)*.1;
    var cY=Math.cos(ay),sY=Math.sin(ay);
    var cX=Math.cos(ax),sX=Math.sin(ax);
    var cZ=Math.cos(az),sZ=Math.sin(az);

    var pts=nodes.map(function(n){
      // rotate Y
      var x1=n.x*cY+n.z*sY, y1=n.y, z1=-n.x*sY+n.z*cY;
      // rotate X
      var x2=x1, y2=y1*cX-z1*sX, z2=y1*sX+z1*cX;
      // rotate Z (slight tilt)
      var x3=x2*cZ-y2*sZ, y3=x2*sZ+y2*cZ, z3=z2;
      var sc=fov/(fov+z3+80);
      var pulse=Math.sin(t*.0016+n.p)*.22+.78;
      return {sx:x3*sc+cx,sy:y3*sc+cy,sc:sc,z:z3,r:n.r,pulse:pulse,ox:x3,oy:y3,oz:z3};
    });

    pts.sort(function(a,b){return a.z-b.z;});

    // connections
    for(var i=0;i<pts.length;i++){
      for(var j=i+1;j<pts.length;j++){
        var a=pts[i],b=pts[j];
        var dx=a.ox-b.ox,dy=a.oy-b.oy,dz=a.oz-b.oz;
        var d=Math.sqrt(dx*dx+dy*dy+dz*dz);
        if(d<LINK){
          var depthFade=Math.min(a.sc,b.sc);
          var al=(1-d/LINK)*depthFade*.38;
          var grad=ctx.createLinearGradient(a.sx,a.sy,b.sx,b.sy);
          grad.addColorStop(0,'rgba(200,40,70,'+al+')');
          grad.addColorStop(.5,'rgba(156,26,53,'+(al*.7)+')');
          grad.addColorStop(1,'rgba(200,40,70,'+al+')');
          ctx.beginPath();ctx.moveTo(a.sx,a.sy);ctx.lineTo(b.sx,b.sy);
          ctx.strokeStyle=grad;ctx.lineWidth=.75;ctx.stroke();
        }
      }
    }

    // nodes
    pts.forEach(function(p){
      var al=Math.min(1,p.sc*1.5)*p.pulse;
      var r=p.r*p.sc*2.2;
      // outer glow
      var g=ctx.createRadialGradient(p.sx,p.sy,0,p.sx,p.sy,r*5);
      g.addColorStop(0,'rgba(210,45,75,'+(al*.45)+')');
      g.addColorStop(.4,'rgba(156,26,53,'+(al*.18)+')');
      g.addColorStop(1,'rgba(156,26,53,0)');
      ctx.beginPath();ctx.arc(p.sx,p.sy,r*5,0,Math.PI*2);
      ctx.fillStyle=g;ctx.fill();
      // inner glow ring
      var g2=ctx.createRadialGradient(p.sx,p.sy,0,p.sx,p.sy,r*2);
      g2.addColorStop(0,'rgba(255,120,140,'+(al*.6)+')');
      g2.addColorStop(1,'rgba(200,40,70,0)');
      ctx.beginPath();ctx.arc(p.sx,p.sy,r*2,0,Math.PI*2);
      ctx.fillStyle=g2;ctx.fill();
      // core dot
      ctx.beginPath();ctx.arc(p.sx,p.sy,r,0,Math.PI*2);
      ctx.fillStyle='rgba(240,80,100,'+al+')';ctx.fill();
    });

    t++;
    requestAnimationFrame(draw);
  }
  draw();
})();
</script>
</body></html>"""

PROFILE_HTML = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Profile Settings — Enderun Marketing Hub</title>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.0/css/all.min.css" crossorigin="anonymous">
<style>
  *{box-sizing:border-box;margin:0;padding:0}
  body{font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;background:#f4ecec;min-height:100vh;display:flex;align-items:center;justify-content:center;}
  .wrap{background:#fff;border-radius:20px;padding:40px 36px;width:420px;max-width:95vw;box-shadow:0 8px 40px rgba(122,16,40,0.12);border:1px solid rgba(122,16,40,0.1);}
  .avatar{width:72px;height:72px;border-radius:50%;background:linear-gradient(135deg,#7a1028,#9b1a35);display:flex;align-items:center;justify-content:center;font-size:28px;color:#fff;font-weight:700;margin:0 auto 16px;}
  .profile-name{text-align:center;font-size:18px;font-weight:700;color:#2d0f18;margin-bottom:4px;}
  .profile-role{text-align:center;font-size:13px;color:#7a4050;margin-bottom:28px;}
  .section-title{font-size:12px;font-weight:700;text-transform:uppercase;letter-spacing:.08em;color:#7a4050;margin-bottom:14px;}
  .field{margin-bottom:14px;}
  label{display:block;font-size:12px;font-weight:600;color:#2d0f18;margin-bottom:5px;}
  input{width:100%;padding:10px 13px;border:1px solid rgba(122,16,40,0.18);border-radius:9px;font-size:14px;font-family:inherit;outline:none;transition:border-color .15s;background:#fdfafa;color:#2d0f18;}
  input:focus{border-color:rgba(122,16,40,0.45);background:#fff;}
  .divider{height:1px;background:rgba(122,16,40,0.08);margin:24px 0;}
  .btn-row{display:flex;gap:10px;margin-top:20px;}
  .btn{padding:10px 20px;border-radius:10px;font-size:13px;font-weight:600;cursor:pointer;border:none;font-family:inherit;transition:all .15s;}
  .btn-primary{background:linear-gradient(135deg,#7a1028,#9b1a35);color:#fff;flex:1;}
  .btn-primary:hover{opacity:.88;}
  .btn-ghost{background:transparent;color:#7a4050;border:1px solid rgba(122,16,40,0.18);}
  .btn-ghost:hover{background:rgba(122,16,40,0.06);}
  .msg{padding:10px 13px;border-radius:9px;font-size:13px;margin-top:12px;display:none;}
  .msg.ok{background:#f0fdf4;border:1px solid #bbf7d0;color:#166534;}
  .msg.err{background:#fef2f2;border:1px solid #fecaca;color:#991b1b;}
</style>
</head>
<body>
<div class="wrap">
  <div class="avatar">E</div>
  <div class="profile-name">Eva</div>
  <div class="profile-role">Enderun Marketing AI &nbsp;·&nbsp; Marketing Hub Admin</div>
  <div class="divider"></div>
  <div class="section-title"><i class="fa-solid fa-lock" style="margin-right:6px"></i>Change Password</div>
  <div class="field"><label>Current Password</label><input type="password" id="cur" placeholder="Current password"></div>
  <div class="field"><label>New Password</label><input type="password" id="np" placeholder="New password (min 6 chars)"></div>
  <div class="field"><label>Confirm New Password</label><input type="password" id="cp" placeholder="Repeat new password"></div>
  <div class="msg" id="msg"></div>
  <div class="btn-row">
    <button class="btn btn-ghost" onclick="window.location='/'"><i class="fa-solid fa-arrow-left"></i> Back</button>
    <button class="btn btn-primary" onclick="savePassword()"><i class="fa-solid fa-floppy-disk"></i> Save Password</button>
  </div>
</div>
<script>
function savePassword() {
  var cur=document.getElementById('cur').value,
      np =document.getElementById('np').value,
      cp =document.getElementById('cp').value;
  var m=document.getElementById('msg');
  m.style.display='none';
  fetch('/api/profile/password',{method:'POST',headers:{'Content-Type':'application/json'},
    body:JSON.stringify({current:cur,new:np,confirm:cp})})
  .then(function(r){return r.json();})
  .then(function(d){
    m.style.display='block';
    if(d.status==='ok'){m.className='msg ok';m.textContent='Password updated successfully!';}
    else{m.className='msg err';m.textContent=d.error||'Failed.';}
  }).catch(function(){m.style.display='block';m.className='msg err';m.textContent='Request failed.';});
}
document.addEventListener('keydown',function(e){if(e.key==='Enter')savePassword();});
</script>
</body></html>"""

@app.route("/login", methods=["GET","POST"])
def login():
    if request.method == "POST":
        pw = request.form.get("password","")
        if pw == DASHBOARD_PASSWORD:
            session["logged_in"] = True
            return redirect("/")
        return LOGIN_HTML.replace("{error}", '<div class="error">⚠ Incorrect password. Try again.</div>'), 401
    if session.get("logged_in"):
        return redirect("/")
    return LOGIN_HTML.replace("{error}", "")

@app.route("/logout")
def logout():
    session.clear()
    return redirect("/login")

@app.route("/profile")
@login_required
def profile_page():
    from flask import render_template_string
    return render_template_string(PROFILE_HTML)

@app.route("/api/profile/password", methods=["POST"])
@login_required
def api_profile_password():
    global DASHBOARD_PASSWORD
    data     = request.json or {}
    current  = data.get("current","")
    new_pw   = data.get("new","").strip()
    confirm  = data.get("confirm","").strip()
    if current != DASHBOARD_PASSWORD:
        return jsonify({"error":"Current password is incorrect."}), 400
    if not new_pw or len(new_pw) < 6:
        return jsonify({"error":"New password must be at least 6 characters."}), 400
    if new_pw != confirm:
        return jsonify({"error":"Passwords do not match."}), 400
    DASHBOARD_PASSWORD = new_pw
    try:
        PROFILE_FILE.write_text(json.dumps({"password": new_pw}))
    except Exception:
        pass
    return jsonify({"status":"ok"})

@app.route("/")
@login_required
def index():
    return send_file(str(HTML_FILE), mimetype="text/html")

@app.route("/api/stats")
@login_required
def api_stats():
    return jsonify(get_stats())

@app.route("/api/leads")
@login_required
def api_leads():
    return jsonify(get_leads())

@app.route("/api/schedule")
@login_required
def api_schedule():
    return jsonify(get_schedule())

@app.route("/api/agents")
@login_required
def api_agents():
    return jsonify(list_agents())

@app.route("/api/workspace-data")
@login_required
def api_workspace_data():
    from collections import Counter
    leads   = get_leads()
    active  = [l for l in leads if l.get("status","").strip().lower() == "active"]
    hot     = sum(1 for l in active if int(l.get("email_count",0) or 0) >= 10)
    warm    = sum(1 for l in active if 5 <= int(l.get("email_count",0) or 0) < 10)
    cold    = sum(1 for l in active if int(l.get("email_count",0) or 0) < 5)
    prog    = Counter(l.get("program_interest","Unknown") for l in leads)
    total_emails = sum(int(l.get("email_count",0) or 0) for l in leads)
    # Parent leads (BS HM / BS CA often parent-driven, or any tagged)
    parent_keywords = ["parent","mom","dad","guardian","mother","father"]
    parent_leads = [l for l in leads if any(k in (l.get("notes","") + l.get("name","")).lower() for k in parent_keywords)]
    # Event leads
    event_keywords = ["wedding","event","venue","banquetes","catering","corporate"]
    event_leads = [l for l in leads if any(k in (l.get("program_interest","") + l.get("notes","")).lower() for k in event_keywords)]
    # Program breakdown for admissions
    prog_list = [{"name": k, "count": v} for k, v in prog.most_common(6)]
    # Recent leads (last 7 days)
    today_str = date.today().isoformat()
    week_ago = (date.today() - timedelta(days=7)).isoformat()
    recent = [l for l in leads if l.get("date_added","") >= week_ago]
    # Schedule data
    sched = {}
    if SCHEDULE_FILE.exists():
        with open(SCHEDULE_FILE) as f:
            sched = json.load(f).get("schedule", {})
    today_post = sched.get(today_str)
    future_posts = [(d, img) for d, img in sorted(sched.items()) if d >= today_str and img]
    # Intel data
    intel_latest = None
    intel_count = 0
    if INTEL_DIR.exists():
        files = sorted(INTEL_DIR.glob("*_intelligence_briefing.txt"), key=lambda f: f.stem, reverse=True)
        intel_count = len(files)
        if files:
            intel_latest = files[0].stem.replace("_intelligence_briefing","")
    # Output reports
    output_dir = BASE_DIR / "output"
    report_count = len(list(output_dir.rglob("*.pdf"))) if output_dir.exists() else 0
    # Paper trading
    paper_data = {}
    paper_file = BASE_DIR / "paper_portfolio.json"
    if paper_file.exists():
        try: paper_data = json.loads(paper_file.read_text())
        except: pass
    return jsonify({
        "total_leads": len(leads),
        "active_leads": len(active),
        "hot": hot, "warm": warm, "cold": cold,
        "total_emails_sent": total_emails,
        "programs": prog_list,
        "top_program": prog.most_common(1)[0][0] if prog else "N/A",
        "parent_leads": len(parent_leads),
        "event_leads": len(event_leads),
        "recent_leads": len(recent),
        "today_post": today_post,
        "future_posts_count": len(future_posts),
        "next_posts": [{"date": d, "image": img} for d, img in future_posts[:5]],
        "intel_latest": intel_latest,
        "intel_count": intel_count,
        "report_count": report_count,
        "paper_balance": paper_data.get("balance", 1000),
        "paper_trades": len(paper_data.get("trades", [])),
        "today": today_str,
    })

@app.route("/api/intelligence")
@login_required
def api_intelligence():
    d = request.args.get("date")
    data = get_intelligence(d)
    if not data:
        return jsonify({"error": "No intelligence briefings found"}), 404
    return jsonify(data)

@app.route("/api/intelligence-list")
@login_required
def api_intelligence_list():
    if INTEL_DIR.exists():
        files = sorted(INTEL_DIR.glob("*_intelligence_briefing.txt"), reverse=True)
        if files:
            return jsonify([f.stem.replace("_intelligence_briefing", "") for f in files])
    # Fallback: list from GitHub
    gh_files = _list_intel_files_github()
    return jsonify([f.replace("_intelligence_briefing.txt", "") for f in gh_files])

@app.route("/api/post-preview")
@login_required
def api_post_preview():
    sched    = get_schedule()
    today    = date.today()
    today_s  = today.isoformat()
    # Find nearest upcoming (or today's) post
    target_date, target_file = None, None
    for dt in sorted(sched.keys()):
        if dt >= today_s:
            target_date, target_file = dt, sched[dt]
            break
    if not target_file:
        return jsonify({"has_image": False, "date": today_s})
    days_until = (date.fromisoformat(target_date) - today).days
    img_path   = IMAGES_DIR / target_file
    has_image  = img_path.exists() if img_path.parent.exists() else False
    return jsonify({
        "has_image":  has_image,
        "date":       target_date,
        "filename":   target_file,
        "image_url":  f"/api/image-preview/{target_file}",
        "days_until": days_until,
    })

@app.route("/api/automations-status")
@login_required
def api_automations_status():
    return jsonify(get_automations_status())

@app.route("/api/delete-lead", methods=["POST"])
@login_required
def api_delete_lead():
    d  = request.json or {}
    em = d.get("email", "").strip().lower()
    if not em:
        return jsonify({"error": "email required"}), 400
    if not LEADS_FILE.exists():
        return jsonify({"error": "No leads file found"}), 404
    rows = []
    deleted = False
    with open(LEADS_FILE, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        fieldnames = reader.fieldnames or []
        for row in reader:
            if row.get("email", "").strip().lower() == em:
                deleted = True
            else:
                rows.append(row)
    if not deleted:
        return jsonify({"error": "Lead not found"}), 404
    with open(LEADS_FILE, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)
    return jsonify({"status": "deleted", "message": "Lead deleted successfully."})

@app.route("/api/update-lead-status", methods=["POST"])
@login_required
def api_update_lead_status():
    d      = request.json or {}
    em     = d.get("email", "").strip().lower()
    status = d.get("status", "inactive").strip()
    if not em:
        return jsonify({"error": "email required"}), 400
    if not LEADS_FILE.exists():
        return jsonify({"error": "No leads file found"}), 404
    rows = []
    updated = False
    with open(LEADS_FILE, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        fieldnames = reader.fieldnames or []
        for row in reader:
            if row.get("email", "").strip().lower() == em:
                row["status"] = status
                updated = True
            rows.append(row)
    if not updated:
        return jsonify({"error": "Lead not found"}), 404
    with open(LEADS_FILE, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)
    return jsonify({"status": "updated", "message": f"Lead status set to {status}."})

@app.route("/api/add-lead", methods=["POST"])
@login_required
def api_add_lead():
    d = request.json or {}
    fn = d.get("first_name","").strip()
    em = d.get("email","").strip()
    if not fn or not em:
        return jsonify({"error": "first_name and email required"}), 400
    if LEADS_FILE.exists():
        with open(LEADS_FILE, newline="", encoding="utf-8") as f:
            for row in csv.DictReader(f):
                if row.get("email","").strip().lower() == em.lower():
                    return jsonify({"status":"exists","message":"Lead already in system."})
    exists = LEADS_FILE.exists()
    with open(LEADS_FILE, "a", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=["first_name","last_name","email","program_interest","status"])
        if not exists:
            w.writeheader()
        w.writerow({"first_name":fn,"last_name":d.get("last_name","").strip(),
                    "email":em,"program_interest":d.get("program_interest","").strip(),"status":"active"})
    return jsonify({"status":"added","message":f"{fn} added successfully."})

@app.route("/api/update-lead", methods=["POST"])
@login_required
def api_update_lead():
    d   = request.json or {}
    em  = (d.get("email") or "").strip().lower()
    if not em:
        return jsonify({"status":"error","message":"email required"}), 400
    rows = []
    fieldnames = ["first_name","last_name","email","program_interest","status","email_count"]
    if LEADS_FILE.exists():
        with open(LEADS_FILE, newline="", encoding="utf-8") as f:
            rows = list(csv.DictReader(f))
    found = False
    for row in rows:
        if row.get("email","").strip().lower() == em:
            found = True
            if d.get("status") in ("active","inactive"):
                row["status"] = d["status"]
            if d.get("program_interest"):
                row["program_interest"] = d["program_interest"]
            break
    if not found:
        return jsonify({"status":"error","message":"Lead not found"}), 404
    with open(LEADS_FILE, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=fieldnames)
        w.writeheader()
        w.writerows(rows)
    return jsonify({"status":"ok"})

# ── CHAT TOOLS ────────────────────────────────────────────────────────────────
CHAT_TOOLS = [
    {
        "name": "send_drip_emails",
        "description": "Send personalized AI-generated drip emails to ALL active leads right now via GitHub Actions.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "post_to_facebook",
        "description": "Post to Enderun Extension's Facebook page right now using today's scheduled image. Auto-generates AI caption unless overridden.",
        "input_schema": {
            "type": "object",
            "properties": {
                "caption": {"type": "string", "description": "Optional caption override."},
                "image":   {"type": "string", "description": "Optional image filename override."}
            },
            "required": []
        }
    },
    {
        "name": "post_to_instagram",
        "description": "Post to Enderun Extension's Instagram account right now using today's scheduled image. Auto-generates AI caption unless overridden.",
        "input_schema": {
            "type": "object",
            "properties": {
                "caption": {"type": "string", "description": "Optional caption override."},
                "image":   {"type": "string", "description": "Optional image filename override."}
            },
            "required": []
        }
    },
    {
        "name": "open_post_composer",
        "description": (
            "Open the interactive post composer so the user can pick an image, preview it, "
            "edit AI-generated captions, and approve before posting. "
            "ALWAYS use this tool when the user wants to post to Facebook, Instagram, or both — "
            "unless they already supplied both an image filename AND a caption. "
            "platform must be 'facebook', 'instagram', or 'both'."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "platform": {"type": "string", "description": "'facebook', 'instagram', or 'both'"}
            },
            "required": ["platform"]
        }
    },
    {
        "name": "trigger_social_listening",
        "description": "Run the social listening workflow — scrapes competitor sites, Google News, generates AI briefing, emails it to the team.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "trigger_weekly_analytics",
        "description": "Run the weekly analytics report — generates a branded PDF with lead stats and AI insights, emails to the team.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "trigger_weekly_preview",
        "description": "Run the weekly campaign preview — generates a PDF with next week's post schedule and drip email plan, emails to the team.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "get_leads",
        "description": "Read current leads data from leads.csv — total count, active/inactive, program breakdown, and individual lead details.",
        "input_schema": {
            "type": "object",
            "properties": {
                "filter_status": {"type": "string", "description": "Filter by status: 'active', 'inactive', or 'all' (default)."}
            },
            "required": []
        }
    },
    {
        "name": "get_posting_schedule",
        "description": "Read the upcoming social media posting schedule from posting_schedule.json.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "read_drive_file",
        "description": (
            "Read a file from Google Drive and return its content. "
            "Accepts: filename (e.g. 'Q1 Report.pdf'), Google Drive URL, or file ID. "
            "Supports: Google Docs (text), Google Sheets (CSV), PDFs, images, DOCX, XLSX, plain text."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "filename": {"type": "string", "description": "Filename, Google Drive URL, or file ID to read."}
            },
            "required": ["filename"]
        }
    },
    {
        "name": "list_drive_files",
        "description": (
            "Browse Google Drive. Lists files and folders accessible to the service account. "
            "Optionally browse inside a specific folder by name or ID. "
            "Use this to explore the Drive structure before reading a file."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "folder": {"type": "string", "description": "Folder name or folder ID to browse inside. Leave empty to list all accessible files and folders."}
            },
            "required": []
        }
    },
    {
        "name": "web_search",
        "description": "Search the internet for real-time information, news, facts, prices, or any topic. Use this whenever you need up-to-date information.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "The search query."}
            },
            "required": ["query"]
        }
    },
    {
        "name": "read_intelligence_briefing",
        "description": (
            "Read the latest competitor intelligence briefing(s) from local storage. "
            "These contain AI analysis, news mentions, website change alerts, and threat assessments for all competitors. "
            "ALWAYS use this FIRST when the user asks about competitor activity, social media data, market intelligence, "
            "competitor posts, what competitors are doing, or any social listening data. "
            "Do NOT trigger_social_listening when briefings already exist — read them first."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "count": {"type": "integer", "description": "Number of recent briefings to read (default 1, max 5)."}
            },
            "required": []
        }
    },
    {
        "name": "get_calendar",
        "description": "Get upcoming events from Eva's Google Calendar (eva@enderuncolleges.com). Shows meetings, deadlines, campaigns.",
        "input_schema": {
            "type": "object",
            "properties": {
                "days": {"type": "integer", "description": "Number of days ahead to look (default 7, max 30)."}
            },
            "required": []
        }
    },
    {
        "name": "add_calendar_event",
        "description": "Add a new event to Eva's Google Calendar.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title":          {"type": "string",  "description": "Event title."},
                "date":           {"type": "string",  "description": "Date in YYYY-MM-DD format."},
                "time":           {"type": "string",  "description": "Time in HH:MM 24h format (optional — omit for all-day)."},
                "description":    {"type": "string",  "description": "Event description (optional)."},
                "duration_hours": {"type": "integer", "description": "Duration in hours (default 1)."}
            },
            "required": ["title", "date"]
        }
    },
    {
        "name": "search_email",
        "description": "Search Gmail inbox for emails. Use queries like 'is:unread', 'from:name', 'subject:topic', or any free text.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Gmail search query (e.g. 'is:unread', 'from:maria', 'subject:enrollment')."},
                "limit": {"type": "integer", "description": "Max emails to return (default 5, max 10)."}
            },
            "required": []
        }
    },
    {
        "name": "remember",
        "description": "Save a piece of information to persistent memory so it can be recalled in future conversations. Use this when the user asks you to remember something.",
        "input_schema": {
            "type": "object",
            "properties": {
                "key":   {"type": "string", "description": "Short label for this memory (e.g. 'Eva's birthday', 'Campaign launch date')."},
                "value": {"type": "string", "description": "The information to remember."}
            },
            "required": ["key", "value"]
        }
    },
    {
        "name": "recall_memories",
        "description": "Recall all saved memories. Use this to remember things the user has asked you to save.",
        "input_schema": {"type": "object", "properties": {}, "required": []}
    },
    {
        "name": "forget_memory",
        "description": "Delete a specific saved memory by its key. Use when the user says 'forget X' or 'remove the X memory'.",
        "input_schema": {
            "type": "object",
            "properties": {
                "key": {"type": "string", "description": "The key of the memory to delete."}
            },
            "required": ["key"]
        }
    },
    {
        "name": "generate_image",
        "description": "Generate an image from a text description using AI. Returns a generated image. Use for marketing visuals, concept art, or any image request.",
        "input_schema": {
            "type": "object",
            "properties": {
                "prompt": {"type": "string", "description": "Detailed description of the image to generate."}
            },
            "required": ["prompt"]
        }
    },
    {
        "name": "execute_python",
        "description": (
            "Execute Python code in a safe sandbox. Use for data analysis, math calculations, "
            "chart generation, file processing, or any computation task. "
            "Code runs with pandas, numpy, matplotlib available. "
            "Print output to see results. Save charts with plt.savefig('/tmp/chart.png')."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "code": {"type": "string", "description": "Python code to execute."}
            },
            "required": ["code"]
        }
    },
    {
        "name": "search_leads",
        "description": "Search leads by name, email, or program interest. Returns matching leads with full details.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search term — partial name, email, or program."}
            },
            "required": ["query"]
        }
    },
    {
        "name": "add_lead_from_chat",
        "description": "Add a new lead to leads.csv. Use when manually registering a new inquiry from the chat.",
        "input_schema": {
            "type": "object",
            "properties": {
                "first_name":       {"type": "string", "description": "First name."},
                "last_name":        {"type": "string", "description": "Last name (optional)."},
                "email":            {"type": "string", "description": "Email address."},
                "program_interest": {"type": "string", "description": "Program they are interested in."}
            },
            "required": ["first_name", "email"]
        }
    },
    {
        "name": "update_lead",
        "description": "Update an existing lead's status or program interest in leads.csv.",
        "input_schema": {
            "type": "object",
            "properties": {
                "email":            {"type": "string", "description": "Email address of the lead to update."},
                "status":           {"type": "string", "description": "'active' or 'inactive'."},
                "program_interest": {"type": "string", "description": "New program interest (optional)."}
            },
            "required": ["email"]
        }
    },
    {
        "name": "save_session_summary",
        "description": (
            "Save a summary of this conversation's key decisions or outcomes to persistent storage. "
            "Use at the end of important sessions — strategy sessions, decisions made, campaigns planned. "
            "These summaries are recalled in future sessions for cross-session memory."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "title":   {"type": "string", "description": "Short title for this session (e.g. 'Q2 Campaign Planning')."},
                "summary": {"type": "string", "description": "2-5 bullet points of key decisions, outcomes, or information from this session."}
            },
            "required": ["title", "summary"]
        }
    },
    {
        "name": "get_recent_sessions",
        "description": "Retrieve summaries of recent past chat sessions for cross-session context and continuity.",
        "input_schema": {
            "type": "object",
            "properties": {
                "count": {"type": "integer", "description": "Number of recent session summaries to retrieve (default 5, max 10)."}
            },
            "required": []
        }
    }
]

def _execute_chat_tool(name, input_data):
    """Execute a dashboard chat tool and return a plain-text result."""
    pat  = os.environ.get("GITHUB_PAT", "")
    repo = os.environ.get("GITHUB_REPO", "your-org/your-repo")

    def _gh(workflow_file):
        if not pat:
            return "error:GITHUB_PAT not configured."
        try:
            import time as _time
            r = http_requests.post(
                f"https://api.github.com/repos/{repo}/actions/workflows/{workflow_file}/dispatches",
                headers={"Authorization": f"Bearer {pat}", "Accept": "application/vnd.github+json"},
                json={"ref": "main"}, timeout=15
            )
            if r.status_code not in (204, 200):
                return f"error:GitHub {r.status_code} — {r.text[:200]}"
            _time.sleep(3)
            try:
                runs_r = http_requests.get(
                    f"https://api.github.com/repos/{repo}/actions/workflows/{workflow_file}/runs?per_page=1",
                    headers={"Authorization": f"Bearer {pat}", "Accept": "application/vnd.github+json"},
                    timeout=10
                )
                if runs_r.status_code == 200:
                    runs = runs_r.json().get("workflow_runs", [])
                    if runs:
                        run = runs[0]
                        return f"ok:status={run.get('status','queued')},url={run.get('html_url','')}"
            except Exception:
                pass
            return "ok"
        except Exception as e:
            return f"error:{e}"

    def _gh_result_msg(r, success_msg):
        if r.startswith("ok:"):
            parts = dict(p.split("=", 1) for p in r[3:].split(",") if "=" in p)
            run_url  = parts.get("url", "")
            run_status = parts.get("status", "queued")
            link = f" [View run →]({run_url})" if run_url else ""
            status_icon = "🟡" if run_status == "queued" else "🟢" if run_status == "in_progress" else "✅"
            return f"✅ {success_msg} {status_icon} Run is **{run_status}**.{link}"
        elif r == "ok":
            return f"✅ {success_msg}"
        else:
            return f"❌ Failed: {r.replace('error:','')}"

    if name == "send_drip_emails":
        r = _gh("send_drip_emails.yml")
        return _gh_result_msg(r, "Drip emails triggered. Personalized emails will be sent to all active leads in ~1–2 min. You'll get a Telegram notification when done.")

    elif name == "post_to_facebook":
        env = os.environ.copy()
        env.update({"AUTO_POST": "true", "ANTHROPIC_API_KEY": ANTHROPIC_KEY})
        if input_data.get("caption"): env["OVERRIDE_CAPTION"] = input_data["caption"]
        if input_data.get("image"):   env["OVERRIDE_IMAGE"]   = input_data["image"]
        try:
            r = subprocess.run([sys.executable, str(BASE_DIR / "post_to_facebook.py")],
                               capture_output=True, text=True, env=env, cwd=str(BASE_DIR), timeout=180)
            return ("✅ Facebook post is now live on Enderun Extension's page." if r.returncode == 0
                    else f"❌ Facebook post failed: {(r.stdout+r.stderr)[-400:]}")
        except subprocess.TimeoutExpired:
            return "❌ Facebook post timed out (>3 min)."
        except Exception as e:
            return f"❌ Error: {e}"

    elif name == "post_to_instagram":
        env = os.environ.copy()
        env.update({"AUTO_POST": "true", "ANTHROPIC_API_KEY": ANTHROPIC_KEY})
        if input_data.get("caption"): env["OVERRIDE_CAPTION"] = input_data["caption"]
        if input_data.get("image"):   env["OVERRIDE_IMAGE"]   = input_data["image"]
        try:
            r = subprocess.run([sys.executable, str(BASE_DIR / "post_to_instagram.py")],
                               capture_output=True, text=True, env=env, cwd=str(BASE_DIR), timeout=180)
            return ("✅ Instagram post is now live on Enderun Extension's account." if r.returncode == 0
                    else f"❌ Instagram post failed: {(r.stdout+r.stderr)[-400:]}")
        except subprocess.TimeoutExpired:
            return "❌ Instagram post timed out (>3 min)."
        except Exception as e:
            return f"❌ Error: {e}"

    elif name == "open_post_composer":
        platform = (input_data.get("platform") or "facebook").lower().strip()
        if platform not in ("facebook", "instagram", "both"):
            platform = "facebook"
        return {"_composer": True, "platform": platform,
                "_tool_content": f"Post composer opened for {platform}. Waiting for user to select image and caption."}

    elif name == "trigger_social_listening":
        r = _gh("social_listening.yml")
        return _gh_result_msg(r, "Social listening triggered. AI briefing will be emailed to the team in ~2–3 min.")

    elif name == "trigger_weekly_analytics":
        r = _gh("weekly_analytics.yml")
        return _gh_result_msg(r, "Weekly analytics triggered. PDF report will be emailed to the team in ~2–3 min.")

    elif name == "trigger_weekly_preview":
        r = _gh("weekly_preview.yml")
        return _gh_result_msg(r, "Weekly campaign preview triggered. Preview PDF will be emailed in ~2–3 min.")

    elif name == "get_leads":
        try:
            fs = (input_data.get("filter_status") or "all").lower()
            rows = []
            if LEADS_FILE.exists():
                with open(LEADS_FILE, newline="", encoding="utf-8") as f:
                    rows = list(csv.DictReader(f))
            if fs != "all":
                rows = [l for l in rows if l.get("status","").lower() == fs]
            active = sum(1 for l in rows if l.get("status","").lower() == "active")
            programs = {}
            for l in rows:
                p = l.get("program_interest","Unknown") or "Unknown"
                programs[p] = programs.get(p, 0) + 1
            sample = [{"name": f"{l.get('first_name','')} {l.get('last_name','')}".strip(),
                       "email": l.get("email",""), "program": l.get("program_interest",""),
                       "status": l.get("status",""), "emails_sent": l.get("email_count","0")}
                      for l in rows[:25]]
            return json.dumps({"total": len(rows), "active": active, "inactive": len(rows)-active,
                               "by_program": programs, "leads": sample,
                               "note": f"Showing first 25 of {len(rows)}." if len(rows) > 25 else ""},
                              ensure_ascii=False)
        except Exception as e:
            return f"Error reading leads: {e}"

    elif name == "get_posting_schedule":
        try:
            if not SCHEDULE_FILE.exists():
                return "No posting_schedule.json found."
            with open(SCHEDULE_FILE) as f:
                sched = json.load(f).get("schedule", {})
            from datetime import date as _d
            today    = str(_d.today())
            upcoming = [(k, v) for k, v in sorted(sched.items()) if k >= today][:14]
            past     = [(k, v) for k, v in sorted(sched.items()) if k <  today][-3:]
            return json.dumps({"total": len(sched),
                               "upcoming":     [{"date": k, "image": v} for k, v in upcoming],
                               "recent_past":  [{"date": k, "image": v} for k, v in past]},
                              ensure_ascii=False)
        except Exception as e:
            return f"Error: {e}"

    elif name == "read_drive_file":
        import re, base64 as _b64, io as _io
        filename = (input_data.get("filename") or "").strip()
        if not filename:
            return "Error: filename is required."
        try:
            from drive_helper import _get_service
            service = _get_service()
        except Exception as e:
            return f"❌ Cannot connect to Google Drive: {e}"

        # Extract file ID from Drive/Docs/Sheets URLs
        file_id = None
        for pat in [r'/(?:file|document|spreadsheets|presentation)/d/([a-zA-Z0-9_-]+)',
                    r'[?&]id=([a-zA-Z0-9_-]+)', r'^([a-zA-Z0-9_-]{25,})$']:
            m = re.search(pat, filename)
            if m:
                file_id = m.group(1)
                break

        actual_name = filename
        mime_type   = ""
        if file_id:
            try:
                meta      = service.files().get(fileId=file_id, fields="name,mimeType").execute()
                mime_type = meta.get("mimeType", "")
                actual_name = meta.get("name", filename)
            except Exception as e:
                return f"❌ Could not get file metadata: {e}"
        else:
            try:
                safe_q   = filename.replace("'", "")
                results  = service.files().list(
                    q=f"name contains '{safe_q}' and trashed=false",
                    fields="files(id,name,mimeType)", pageSize=5
                ).execute()
                files = results.get("files", [])
                if not files:
                    return f"❌ File '{filename}' not found in Google Drive."
                chosen    = next((f for f in files if f["name"].lower() == filename.lower()), files[0])
                file_id   = chosen["id"]
                mime_type = chosen["mimeType"]
                actual_name = chosen["name"]
            except Exception as e:
                return f"❌ Error searching Drive: {e}"

        def _download(fid):
            from googleapiclient.http import MediaIoBaseDownload
            req = service.files().get_media(fileId=fid)
            buf = _io.BytesIO()
            dl  = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = dl.next_chunk()
            return buf.getvalue()

        try:
            GDOC    = "application/vnd.google-apps.document"
            GSHEET  = "application/vnd.google-apps.spreadsheet"
            GSLIDES = "application/vnd.google-apps.presentation"

            if mime_type == GDOC:
                raw  = service.files().export(fileId=file_id, mimeType="text/plain").execute()
                text = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                return f"**{actual_name}** (Google Doc)\n\n{text[:50000]}"

            elif mime_type == GSHEET:
                raw  = service.files().export(fileId=file_id, mimeType="text/csv").execute()
                text = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                return f"**{actual_name}** (Google Sheet)\n\n{text[:50000]}"

            elif mime_type == GSLIDES:
                raw  = service.files().export(fileId=file_id, mimeType="text/plain").execute()
                text = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                return f"**{actual_name}** (Google Slides)\n\n{text[:30000]}"

            elif mime_type == "application/pdf":
                raw  = _download(file_id)
                if len(raw) > 30 * 1024 * 1024:
                    return f"❌ PDF '{actual_name}' is too large (>30 MB)."
                b64 = _b64.b64encode(raw).decode("utf-8")
                return {"_multimodal": True,
                        "preview": f"✅ PDF loaded: {actual_name} ({len(raw)//1024} KB)",
                        "blocks": [
                            {"type": "document", "source": {"type": "base64",
                             "media_type": "application/pdf", "data": b64}},
                            {"type": "text", "text": f"The file above is '{actual_name}' from Google Drive. Analyze it as requested."}
                        ]}

            elif mime_type.startswith("image/"):
                raw  = _download(file_id)
                safe_mime = mime_type if mime_type in ("image/jpeg","image/png","image/gif","image/webp") else "image/jpeg"
                b64 = _b64.b64encode(raw).decode("utf-8")
                return {"_multimodal": True,
                        "preview": f"✅ Image loaded: {actual_name}",
                        "blocks": [
                            {"type": "image", "source": {"type": "base64",
                             "media_type": safe_mime, "data": b64}},
                            {"type": "text", "text": f"The image above is '{actual_name}' from Google Drive."}
                        ]}

            elif mime_type in ("text/plain","text/csv","text/html","application/json"):
                raw  = _download(file_id)
                return f"**{actual_name}**\n\n{raw.decode('utf-8', errors='replace')[:50000]}"

            elif "wordprocessingml" in mime_type or mime_type == "application/msword":
                raw = _download(file_id)
                try:
                    import docx as _docx
                    doc  = _docx.Document(_io.BytesIO(raw))
                    text = "\n".join(p.text for p in doc.paragraphs)
                    return f"**{actual_name}** (Word Doc)\n\n{text[:50000]}"
                except ImportError:
                    return f"❌ python-docx not installed — cannot read DOCX files."

            elif "spreadsheetml" in mime_type:
                raw = _download(file_id)
                try:
                    import openpyxl as _xl
                    wb  = _xl.load_workbook(_io.BytesIO(raw), read_only=True, data_only=True)
                    out = []
                    for sname in wb.sheetnames:
                        ws = wb[sname]
                        out.append(f"## {sname}")
                        for row in ws.iter_rows(max_row=300, values_only=True):
                            out.append(",".join(str(c or "") for c in row))
                    return f"**{actual_name}** (Excel)\n\n" + "\n".join(out)[:50000]
                except ImportError:
                    return f"❌ openpyxl not installed — cannot read XLSX files."

            else:
                return f"File '{actual_name}' found (type: {mime_type}) but this format is not yet supported."
        except Exception as e:
            return f"❌ Error reading file: {e}"

    elif name == "list_drive_files":
        try:
            from drive_helper import _get_service
            service = _get_service()
            folder_param = (input_data.get("folder") or "").strip()
            target_folder_id = None
            folder_label     = "Google Drive (all accessible)"

            if folder_param:
                # Resolve folder name → ID if not already an ID
                import re as _re
                if _re.match(r'^[a-zA-Z0-9_-]{25,}$', folder_param):
                    target_folder_id = folder_param
                    folder_label     = f"folder ID: {folder_param}"
                else:
                    r = service.files().list(
                        q=f"name='{folder_param.replace(chr(39),'')}' and mimeType='application/vnd.google-apps.folder' and trashed=false",
                        fields="files(id,name)", pageSize=5
                    ).execute()
                    hits = r.get("files", [])
                    if not hits:
                        return f"❌ Folder '{folder_param}' not found. Try listing without a folder first to see what's available."
                    target_folder_id = hits[0]["id"]
                    folder_label     = hits[0]["name"]

            if target_folder_id:
                q = f"'{target_folder_id}' in parents and trashed=false"
            else:
                q = "trashed=false"

            results = service.files().list(
                q=q,
                fields="files(id,name,mimeType,modifiedTime)",
                orderBy="folder,name",
                pageSize=100
            ).execute()
            items = results.get("files", [])
            if not items:
                return f"No files found in '{folder_label}'."

            FOLDER_MIME = "application/vnd.google-apps.folder"
            folders = [f for f in items if f["mimeType"] == FOLDER_MIME]
            files   = [f for f in items if f["mimeType"] != FOLDER_MIME]

            lines = []
            if folders:
                lines.append("📁 **Folders:**")
                for f in folders:
                    lines.append(f"  📁 {f['name']}  (id: {f['id']})")
            if files:
                lines.append("\n📄 **Files:**")
                for f in files:
                    ext = f["mimeType"].split(".")[-1].split("/")[-1]
                    lines.append(f"  📄 {f['name']}  ({ext})")

            header = f"**Contents of '{folder_label}' — {len(folders)} folders, {len(files)} files:**"
            return header + "\n\n" + "\n".join(lines)
        except Exception as e:
            return f"❌ Error listing Drive: {e}"

    elif name == "web_search":
        query = (input_data.get("query") or "").strip()
        if not query:
            return "Error: query is required."
        try:
            import urllib.parse as _up
            from bs4 import BeautifulSoup as _BS
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
            url = f"https://html.duckduckgo.com/html/?q={_up.quote(query)}"
            r = http_requests.get(url, headers=headers, timeout=12)
            soup = _BS(r.text, "html.parser")
            results = []
            for res in soup.select(".result__body")[:6]:
                title_el   = res.select_one(".result__title")
                snippet_el = res.select_one(".result__snippet")
                url_el     = res.select_one(".result__url")
                title   = title_el.get_text(strip=True)   if title_el   else ""
                snippet = snippet_el.get_text(strip=True) if snippet_el else ""
                link    = url_el.get_text(strip=True)     if url_el     else ""
                if title or snippet:
                    results.append(f"**{title}**\n{link}\n{snippet}")
            if not results:
                return f"No search results found for: '{query}'."
            return f"Search results for '{query}':\n\n" + "\n\n---\n\n".join(results)
        except Exception as e:
            return f"❌ Web search failed: {e}"

    elif name == "read_intelligence_briefing":
        count = min(int(input_data.get("count") or 1), 5)
        try:
            results = []
            if INTEL_DIR.exists():
                files = sorted(INTEL_DIR.glob("*_intelligence_briefing.txt"), reverse=True)[:count]
                for f in files:
                    results.append(f.read_text(encoding="utf-8"))
            if not results:
                return ("No local intelligence briefings found. "
                        "Use trigger_social_listening to generate one, or use web_search to look up competitor data manually.")
            header = f"✅ Found {len(results)} briefing(s). Most recent first:\n\n"
            return header + f"\n\n{'='*60}\n\n".join(results)
        except Exception as e:
            return f"❌ Error reading briefings: {e}"

    elif name == "get_calendar":
        try:
            from workspace_helper import list_calendar_events, format_events_text
            days   = min(int(input_data.get("days") or 7), 30)
            events = list_calendar_events(days=days)
            if not events:
                return f"No events in the next {days} days."
            return format_events_text(events, header=f"📅 Next {days} days:\n\n")
        except Exception as e:
            return f"❌ Calendar error: {e}"

    elif name == "add_calendar_event":
        try:
            from workspace_helper import add_calendar_event as _add_ev
            result = _add_ev(
                title         = input_data.get("title", ""),
                date_str      = input_data.get("date",  ""),
                time_str      = input_data.get("time",  ""),
                description   = input_data.get("description", ""),
                duration_hours= int(input_data.get("duration_hours") or 1),
            )
            if result and result.get("id"):
                link = result.get("htmlLink", "")
                return f"✅ Event created: **{result.get('summary','')}** on {input_data.get('date','')}. {link}"
            return "❌ Failed to create event — check Calendar credentials."
        except Exception as e:
            return f"❌ Calendar error: {e}"

    elif name == "search_email":
        try:
            from workspace_helper import search_gmail
            q     = (input_data.get("query") or "is:unread").strip()
            limit = min(int(input_data.get("limit") or 5), 10)
            emails = search_gmail(query=q, limit=limit)
            if not emails:
                return "No emails found matching that query."
            lines = []
            for e in emails:
                lines.append(f"**From:** {e.get('sender','')}\n**Subject:** {e.get('subject','')}\n**Date:** {e.get('date','')}\n{e.get('snippet','')}")
            return f"Gmail search '{q}':\n\n" + "\n\n---\n\n".join(lines)
        except Exception as e:
            return f"❌ Gmail error: {e}"

    elif name == "remember":
        key   = (input_data.get("key")   or "").strip()
        value = (input_data.get("value") or "").strip()
        if not key or not value:
            return "Error: both key and value are required."
        memories = _load_chat_memory()
        # Update if key exists, else append
        for m in memories:
            if m.get("key","").lower() == key.lower():
                m["value"] = value
                m["updated"] = str(date.today())
                _save_chat_memory(memories)
                return f"✅ Updated memory: **{key}** → {value}"
        memories.append({"key": key, "value": value, "saved": str(date.today())})
        _save_chat_memory(memories)
        return f"✅ Saved to memory: **{key}** → {value}"

    elif name == "recall_memories":
        memories = _load_chat_memory()
        if not memories:
            return "No memories saved yet."
        lines = [f"**{m['key']}**: {m['value']}  *(saved {m.get('saved','')})*" for m in memories]
        return f"📋 Saved memories ({len(memories)}):\n\n" + "\n\n".join(lines)

    elif name == "forget_memory":
        key      = (input_data.get("key") or "").strip()
        memories = _load_chat_memory()
        new_mem  = [m for m in memories if m.get("key","").lower() != key.lower()]
        if len(new_mem) == len(memories):
            return f"No memory found with key '{key}'."
        _save_chat_memory(new_mem)
        return f"✅ Deleted memory: **{key}**"

    elif name == "generate_image":
        import urllib.parse as _up, uuid as _uuid, io as _io
        prompt = (input_data.get("prompt") or "").strip()
        if not prompt:
            return "Error: prompt is required."
        try:
            seed  = abs(hash(prompt)) % 99999
            img_url = (
                f"https://image.pollinations.ai/prompt/{_up.quote(prompt)}"
                f"?width=1024&height=1024&nologo=true&seed={seed}"
            )
            raw = http_requests.get(img_url, timeout=60).content
            if len(raw) < 500:
                return f"❌ Image generation failed — no image returned."
            fname = f"gen_{_uuid.uuid4().hex[:10]}.jpg"
            (CHAT_IMAGES_DIR / fname).write_bytes(raw)
            import base64 as _b64
            b64 = _b64.b64encode(raw).decode("utf-8")
            return {
                "_multimodal": True,
                "preview": f"🎨 Image generated!",
                "img_url": f"/chat-image/{fname}",
                "blocks": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": b64}},
                    {"type": "text",  "text": f"Generated image for prompt: '{prompt}'"}
                ]
            }
        except Exception as e:
            return f"❌ Image generation failed: {e}"

    elif name == "execute_python":
        import tempfile, subprocess as _sp, sys as _sys
        code = (input_data.get("code") or "").strip()
        if not code:
            return "Error: code is required."
        try:
            with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False, encoding="utf-8") as f:
                f.write("import sys\nsys.path.insert(0, r'" + str(BASE_DIR) + "')\n")
                f.write(code)
                tmp_path = f.name
            result = _sp.run(
                [_sys.executable, tmp_path],
                capture_output=True, text=True, timeout=30,
                cwd=str(BASE_DIR)
            )
            Path(tmp_path).unlink(missing_ok=True)
            output = (result.stdout + result.stderr).strip()
            if not output:
                output = "(no output)"
            if result.returncode != 0:
                return f"❌ Code exited with error (code {result.returncode}):\n\n```\n{output[:3000]}\n```"
            return f"✅ Code executed successfully:\n\n```\n{output[:3000]}\n```"
        except _sp.TimeoutExpired:
            return "❌ Code execution timed out (30s limit)."
        except Exception as e:
            return f"❌ Execution error: {e}"

    elif name == "search_leads":
        query = (input_data.get("query") or "").strip().lower()
        if not query:
            return "Error: query is required."
        try:
            rows = []
            if LEADS_FILE.exists():
                with open(LEADS_FILE, newline="", encoding="utf-8") as f:
                    rows = list(csv.DictReader(f))
            matches = [l for l in rows if
                       query in l.get("first_name","").lower() or
                       query in l.get("last_name","").lower() or
                       query in l.get("email","").lower() or
                       query in (l.get("program_interest") or "").lower()]
            if not matches:
                return f"No leads found matching '{query}'."
            results = [{"name": f"{l.get('first_name','')} {l.get('last_name','')}".strip(),
                        "email": l.get("email",""), "program": l.get("program_interest",""),
                        "status": l.get("status",""), "emails_sent": l.get("email_count","0")}
                       for l in matches[:20]]
            return json.dumps({"found": len(matches), "leads": results,
                               "note": f"Showing first 20 of {len(matches)}." if len(matches) > 20 else ""},
                              ensure_ascii=False)
        except Exception as e:
            return f"Error searching leads: {e}"

    elif name == "add_lead_from_chat":
        fn   = (input_data.get("first_name") or "").strip()
        ln   = (input_data.get("last_name")  or "").strip()
        em   = (input_data.get("email")       or "").strip().lower()
        prog = (input_data.get("program_interest") or "").strip()
        if not fn or not em or "@" not in em:
            return "Error: first_name and a valid email are required."
        try:
            rows = []
            fieldnames = ["first_name","last_name","email","program_interest","status","email_count"]
            if LEADS_FILE.exists():
                with open(LEADS_FILE, newline="", encoding="utf-8") as f:
                    rows = list(csv.DictReader(f))
            if any(l.get("email","").lower() == em for l in rows):
                return f"⚠️ Lead with email {em} already exists."
            write_header = not LEADS_FILE.exists() or not rows
            with open(LEADS_FILE, "a", newline="", encoding="utf-8") as f:
                w = csv.DictWriter(f, fieldnames=fieldnames)
                if write_header:
                    w.writeheader()
                w.writerow({"first_name":fn,"last_name":ln,"email":em,
                            "program_interest":prog,"status":"active","email_count":"0"})
            return f"✅ Lead added: **{fn} {ln}** ({em}) — {prog or 'no program specified'}."
        except Exception as e:
            return f"Error adding lead: {e}"

    elif name == "update_lead":
        em         = (input_data.get("email")            or "").strip().lower()
        new_status = (input_data.get("status")           or "").strip().lower()
        new_prog   = (input_data.get("program_interest") or "").strip()
        if not em:
            return "Error: email is required."
        try:
            rows = []
            fieldnames = ["first_name","last_name","email","program_interest","status","email_count"]
            if LEADS_FILE.exists():
                with open(LEADS_FILE, newline="", encoding="utf-8") as f:
                    rows = list(csv.DictReader(f))
            found = False
            updated_lead = None
            for l in rows:
                if l.get("email","").lower() == em:
                    found = True
                    if new_status in ("active","inactive"):
                        l["status"] = new_status
                    if new_prog:
                        l["program_interest"] = new_prog
                    updated_lead = l
                    break
            if not found:
                return f"❌ No lead found with email '{em}'."
            with open(LEADS_FILE, "w", newline="", encoding="utf-8") as f:
                w = csv.DictWriter(f, fieldnames=fieldnames)
                w.writeheader()
                w.writerows(rows)
            return (f"✅ Lead updated: **{updated_lead.get('first_name','')} {updated_lead.get('last_name','')}** "
                    f"— status: {updated_lead['status']}, program: {updated_lead.get('program_interest','')}")
        except Exception as e:
            return f"Error updating lead: {e}"

    elif name == "save_session_summary":
        title   = (input_data.get("title")   or "Session").strip()
        summary = (input_data.get("summary") or "").strip()
        if not summary:
            return "Error: summary is required."
        try:
            CHAT_SESSIONS_DIR.mkdir(parents=True, exist_ok=True)
            fname = CHAT_SESSIONS_DIR / f"{str(date.today())}_{uuid.uuid4().hex[:6]}_summary.json"
            fname.write_text(json.dumps({"date": str(date.today()), "title": title, "summary": summary},
                                        ensure_ascii=False, indent=2), encoding="utf-8")
            return f"✅ Session summary saved: **{title}**"
        except Exception as e:
            return f"Error saving summary: {e}"

    elif name == "get_recent_sessions":
        count = min(int(input_data.get("count") or 5), 10)
        try:
            if not CHAT_SESSIONS_DIR.exists():
                return "No past session summaries found. Use save_session_summary to save key decisions."
            files = sorted(CHAT_SESSIONS_DIR.glob("*_summary.json"), reverse=True)[:count]
            if not files:
                return "No session summaries found yet. Use save_session_summary at the end of important conversations."
            summaries = []
            for f in files:
                try:
                    data = json.loads(f.read_text(encoding="utf-8"))
                    summaries.append(f"**{data.get('date','')} — {data.get('title','')}**\n{data.get('summary','')}")
                except Exception:
                    pass
            return f"📋 Last {len(summaries)} session summaries:\n\n" + "\n\n---\n\n".join(summaries)
        except Exception as e:
            return f"Error reading sessions: {e}"

    return f"Unknown tool: {name}"


@app.route("/api/chat", methods=["POST"])
@login_required
def api_chat():
    data      = request.json or {}
    messages  = data.get("messages", [])
    agent_id  = data.get("agent", "marketing-manager")
    message   = data.get("message", "")
    attachment= data.get("attachment", None)
    model     = data.get("model", "claude-opus-4-7")
    mode_hint = data.get("mode_hint", "")
    greeting  = data.get("greeting", False)
    system    = agent_system_prompt(agent_id)
    max_tokens = 16000

    if greeting:
        # Build proactive dashboard state snapshot
        try:
            _leads    = get_leads()
            _active   = [l for l in _leads if l.get("status","").strip().lower() == "active"]
            _today    = str(date.today())
            _sched    = {}
            if SCHEDULE_FILE.exists():
                with open(SCHEDULE_FILE) as _f:
                    _sched = json.load(_f).get("schedule", {})
            _today_img   = _sched.get(_today, "")
            _upcoming    = [(k, v) for k, v in sorted(_sched.items()) if k >= _today][:3]
            _memories    = _load_chat_memory()
            _mem_note    = ("\n\nSAVED MEMORIES (use these to personalize your greeting):\n" +
                            "\n".join(f"  • {m['key']}: {m['value']}" for m in _memories[:10])
                            if _memories else "")
            _state_block = (
                f"\n\nDASHBOARD STATE (current as of {_today}):\n"
                f"  • Total leads: {len(_leads)} ({len(_active)} active)\n"
                f"  • Today's scheduled post: {_today_img or 'none'}\n"
                "  • Upcoming posts: " + ", ".join(k + ": " + v for k, v in _upcoming) + "\n"
                + _mem_note
            )
        except Exception:
            _state_block = ""
        system = system + _state_block + (
            "\n\n---\nGREETING MODE: You are opening a fresh conversation. "
            "Write a warm, friendly greeting in 1-2 sentences — address the user as a colleague. "
            "Introduce yourself by name and your specific role. "
            "Then list 4-6 bullet points of concrete things you can help with (be specific to your role). "
            "If there's something notable in the dashboard state (low leads, no post today, saved memories), "
            "mention it naturally as a proactive heads-up. "
            "Keep total response under 180 words. Do NOT use hashtags."
        )
        messages = [{"role": "user", "content": "Hello, who are you and what can you help me with?"}]
        max_tokens = 450
    elif mode_hint:
        system = system + "\n\n" + mode_hint

    if not greeting:
        # Inject saved memories into system prompt for context
        try:
            _mems = _load_chat_memory()
            if _mems:
                _mem_ctx = "\n\nSAVED MEMORIES:\n" + "\n".join(f"  • {m['key']}: {m['value']}" for m in _mems)
                system = system + _mem_ctx
        except Exception:
            pass
        # Auto-inject latest briefing if relevant agent or message is competitor/intel related
        _intel_keywords = ["competitor", "social listening", "briefing", "intelligence", "dlsu", "ateneo", "cca",
                           "iscahm", "lpu", "ua&p", "lyceum", "competitors", "market", "posting", "posts",
                           "engagement", "social media data", "what are they doing", "ano ginagawa"]
        _inject_briefing = (agent_id in ("competitor-analysis", "social-listening") or
                            any(k in message.lower() for k in _intel_keywords))
        if _inject_briefing and INTEL_DIR.exists():
            try:
                from datetime import datetime as _dt
                _bf_files = sorted(INTEL_DIR.glob("*_intelligence_briefing.txt"), reverse=True)[:2]
                if _bf_files:
                    _bf_age_days = (_dt.now() - _dt.fromtimestamp(_bf_files[0].stat().st_mtime)).days
                    _bf_text = "\n\n" + ("="*60) + "\n\n".join(f.read_text(encoding="utf-8") for f in _bf_files)
                    _stale_note = (f"\n⚠️ NOTE: This briefing is {_bf_age_days} days old — data may be outdated. "
                                   f"Recommend running social listening to refresh.\n") if _bf_age_days > 7 else ""
                    system = system + (
                        f"\n\nPRE-LOADED INTELLIGENCE BRIEFING (use this data immediately — do NOT say 'waiting for email'):{_stale_note}\n{_bf_text}"
                    )
            except Exception:
                pass
        # Auto-inject recent session summaries for cross-session continuity
        try:
            if CHAT_SESSIONS_DIR.exists():
                _sum_files = sorted(CHAT_SESSIONS_DIR.glob("*_summary.json"), reverse=True)[:3]
                if _sum_files:
                    _sums = []
                    for _sf in _sum_files:
                        try:
                            _sd = json.loads(_sf.read_text(encoding="utf-8"))
                            _sums.append(f"  • {_sd.get('date','')} — {_sd.get('title','')}: {_sd.get('summary','')}")
                        except Exception:
                            pass
                    if _sums:
                        system = system + "\n\nRECENT SESSION CONTEXT (from past conversations):\n" + "\n".join(_sums)
        except Exception:
            pass

        system = system + (
            "\n\nACTION TOOLS: You have powerful real-time tools. Use them proactively:\n"
            "- Marketing: post to Facebook/Instagram, send drip emails, run reports, check leads\n"
            "- Competitor Intel: read_intelligence_briefing — reads latest saved briefings (ALWAYS use this first for competitor/social media questions)\n"
            "- Web: web_search for ANY real-time info, news, facts, prices, research\n"
            "- Calendar & Email: get_calendar, add_calendar_event, search_email\n"
            "- Memory: remember/recall_memories/forget_memory to save things the user asks you to remember\n"
            "- Image: generate_image for any image request or visual concept\n"
            "- Code: execute_python for math, data analysis, charts, file processing\n"
            "- Drive: read_drive_file, list_drive_files\n"
            "ALWAYS use the right tool instead of explaining how to do it manually.\n\n"
            "CRITICAL REPORT QUALITY RULES:\n"
            "- NEVER output a table with all dashes (—) as placeholders. If you lack data, use web_search or read_intelligence_briefing to get it.\n"
            "- For competitor social media data: first read_intelligence_briefing, then web_search for each competitor, then synthesize.\n"
            "- When exact numbers are unavailable, provide intelligent estimates clearly labeled 'Est.' based on known patterns.\n"
            "- Always include Enderun Colleges itself in competitor comparison tables.\n"
            "- Deliver FULL reports — detailed, specific, actionable. Never say 'I'll email you the results.'\n"
            "- Do NOT trigger_social_listening when local briefings already exist — read them first with read_intelligence_briefing."
        )

    if not greeting and not messages and message:
        if attachment:
            content = []
            if attachment.get("type") == "base64":
                content.append({"type":"image","source":{"type":"base64","media_type":attachment.get("mediaType","image/jpeg"),"data":attachment.get("data","")}})
            elif attachment.get("type") == "url":
                content.append({"type":"image","source":{"type":"url","url":attachment.get("data","")}})
            if message:
                content.append({"type":"text","text":message})
            messages = [{"role":"user","content":content}]
        else:
            messages = [{"role":"user","content":message}]

    def generate():
        try:
            client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
            # Sanitize: replace __GREET__, strip internal fields (_toolCards etc.)
            clean_messages = []
            for m in messages:
                txt = m.get("content", "")
                if isinstance(txt, str) and txt == "__GREET__":
                    clean_messages.append({"role": "user", "content": "Hello, who are you and what can you help me with?"})
                else:
                    clean_messages.append({"role": m["role"], "content": m["content"]})
            # Ensure messages doesn't start with assistant turn
            while clean_messages and clean_messages[0].get("role") != "user":
                clean_messages.pop(0)
            if not clean_messages:
                yield "data: " + json.dumps({"text": "No messages to send."}) + "\n\n"
                yield "data: [DONE]\n\n"
                return
            cached_system = [{"type": "text", "text": system, "cache_control": {"type": "ephemeral"}}]
            kwargs = dict(model=model, max_tokens=max_tokens, system=cached_system, messages=clean_messages)
            if not greeting:
                kwargs["tools"] = CHAT_TOOLS
            # Extended thinking for Deep mode (Opus/Sonnet only)
            if mode_hint and "Think step by step" in mode_hint and model in ("claude-opus-4-7", "claude-sonnet-4-6"):
                kwargs["thinking"] = {"type": "enabled", "budget_tokens": 10000}
                kwargs["max_tokens"] = max(kwargs["max_tokens"], 20000)
        except Exception as e:
            yield "data: " + json.dumps({"text": f"Setup error: {e}"}) + "\n\n"
            yield "data: [DONE]\n\n"
            return

        tool_use_blocks = []
        current_tool_id = current_tool_name = None
        current_tool_json = ""

        try:
            with client.messages.stream(**kwargs) as stream:
                for event in stream:
                    etype = event.type
                    if etype == "content_block_start":
                        cb = getattr(event, "content_block", None)
                        if cb and cb.type == "tool_use":
                            current_tool_id   = cb.id
                            current_tool_name = cb.name
                            current_tool_json = ""
                            yield "data: " + json.dumps({"action": "tool_start", "tool": cb.name, "tool_id": cb.id}) + "\n\n"
                    elif etype == "content_block_delta":
                        d = getattr(event, "delta", None)
                        if d:
                            if d.type == "text_delta":
                                yield "data: " + json.dumps({"text": d.text}) + "\n\n"
                            elif d.type == "input_json_delta" and current_tool_id:
                                current_tool_json += d.partial_json
                    elif etype == "content_block_stop" and current_tool_id:
                        try:
                            parsed_input = json.loads(current_tool_json) if current_tool_json else {}
                        except Exception:
                            parsed_input = {}
                        tool_use_blocks.append({"id": current_tool_id, "name": current_tool_name,
                                                "input": parsed_input})
                        current_tool_id = current_tool_name = None
                        current_tool_json = ""
                final_msg = stream.get_final_message()

            if tool_use_blocks:
                # Clear Phase 1 text so Phase 2 response doesn't duplicate it
                yield "data: " + json.dumps({"action": "clear_text"}) + "\n\n"
                tool_results = []
                for tool in tool_use_blocks:
                    result = _execute_chat_tool(tool["name"], tool["input"])
                    if isinstance(result, dict) and result.get("_composer"):
                        tool_content   = result.get("_tool_content", "Post composer opened.")
                        done_event     = {"action": "open_composer", "platform": result["platform"], "tool_id": tool["id"]}
                    elif isinstance(result, dict) and result.get("_multimodal"):
                        tool_content   = result["blocks"]
                        display_result = result.get("preview", "✅ File loaded.")
                        done_event     = {"action": "tool_done", "tool": tool["name"], "tool_id": tool["id"], "result": display_result}
                        if result.get("img_url"):
                            done_event["img_url"] = result["img_url"]
                    else:
                        tool_content   = str(result)
                        display_result = str(result)[:140]
                        done_event     = {"action": "tool_done", "tool": tool["name"], "tool_id": tool["id"], "result": display_result}
                    tool_results.append({"type": "tool_result", "tool_use_id": tool["id"],
                                         "content": tool_content})
                    yield "data: " + json.dumps(done_event) + "\n\n"
                asst_content = []
                for b in final_msg.content:
                    if b.type == "text":
                        asst_content.append({"type": "text", "text": b.text})
                    elif b.type == "tool_use":
                        asst_content.append({"type": "tool_use", "id": b.id,
                                             "name": b.name, "input": b.input})
                follow_msgs = clean_messages + [{"role": "assistant", "content": asst_content},
                                                {"role": "user",      "content": tool_results}]
                with client.messages.stream(model=model, max_tokens=max_tokens,
                                            system=cached_system, messages=follow_msgs) as s2:
                    for text in s2.text_stream:
                        yield "data: " + json.dumps({"text": text}) + "\n\n"

        except Exception as e:
            yield "data: " + json.dumps({"text": f"\n\n❌ API error: {e}"}) + "\n\n"

        yield "data: [DONE]\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream",
                    headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no"})

@app.route("/api/run-email", methods=["POST"])
@login_required
def api_run_email():
    pat  = os.environ.get("GITHUB_PAT", "")
    repo = os.environ.get("GITHUB_REPO", "your-org/your-repo")
    if not pat:
        return jsonify({"status": "error", "output": "GITHUB_PAT not set in environment."})
    try:
        resp = http_requests.post(
            f"https://api.github.com/repos/{repo}/actions/workflows/send_drip_emails.yml/dispatches",
            headers={"Authorization": f"Bearer {pat}", "Accept": "application/vnd.github+json"},
            json={"ref": "main"},
            timeout=15
        )
        if resp.status_code in (204, 200):
            return jsonify({"status": "ok", "output": "Drip emails triggered via GitHub Actions.\nEmails will be sent in ~1-2 minutes. You'll get a Telegram notification when done."})
        return jsonify({"status": "error", "output": f"GitHub API error {resp.status_code}: {resp.text}"})
    except Exception as e:
        return jsonify({"status": "error", "output": str(e)})

@app.route("/api/job/<job_id>")
@login_required
def api_job_status(job_id):
    job = _jobs.get(job_id)
    if not job:
        return jsonify({"status": "error", "output": "Job not found"})
    return jsonify(job)

@app.route("/api/run-post", methods=["POST"])
@login_required
def api_run_post():
    env = os.environ.copy()
    env.update({"AUTO_POST":"true","ANTHROPIC_API_KEY":ANTHROPIC_KEY})
    try:
        r = subprocess.run([sys.executable, str(BASE_DIR/"post_to_facebook.py")],
                           capture_output=True, text=True, env=env, cwd=str(BASE_DIR), timeout=120)
        return jsonify({"status":"ok","output":r.stdout+r.stderr})
    except Exception as e:
        return jsonify({"status":"error","output":str(e)})

@app.route("/api/generate-caption")
@login_required
def api_generate_caption():
    image = request.args.get("image", "")
    name  = image.replace(".jpg","").replace(".png","").replace("-"," ").replace("_"," ").strip()
    hashtags = ("#EnderunExtension #EnderunColleges #BGC #McKinleyHill "
                "#ProfessionalDevelopment #LifelongLearning #UpskillPH "
                "#CulinaryArts #HospitalityPH #LearnWithEnderun")
    try:
        client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
        resp = client.messages.create(
            model="claude-opus-4-7", max_tokens=400,
            messages=[{"role":"user","content":(
                f"You are Eva, the social media manager of Enderun Extension (BGC, Philippines). "
                f"Write an engaging Facebook post caption for the '{name}' program/course. "
                f"2-3 short paragraphs. Warm, aspirational, Filipino-proud tone. "
                f"End with a clear CTA. Do NOT include hashtags or URL. Caption text only."
            )}]
        )
        caption = resp.content[0].text.strip()
        return jsonify({"caption": caption, "hashtags": hashtags})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/generate-ig-caption")
@login_required
def api_generate_ig_caption():
    image = request.args.get("image", "")
    name  = image.replace(".jpg","").replace(".png","").replace("-"," ").replace("_"," ").strip()
    hashtags = ("#EnderunExtension #Enderun #BGC #McKinleyHill #Manila "
                "#LearnWithEnderun #CulinaryArts #HospitalityLife #UpskillPH "
                "#ProfessionalDevelopment #EnderunColleges #StudyInBGC")
    try:
        client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)
        resp = client.messages.create(
            model="claude-opus-4-7", max_tokens=300,
            messages=[{"role":"user","content":(
                f"You are Eva, the Instagram manager of Enderun Extension (BGC, Philippines). "
                f"Write a short Instagram caption for the '{name}' program/course. "
                f"1-2 punchy sentences max. Aspirational, lifestyle-driven, Gen Z-friendly tone. "
                f"No hashtags, no URL. Just the caption hook text."
            )}]
        )
        caption = resp.content[0].text.strip()
        return jsonify({"caption": caption, "hashtags": hashtags})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/assets/<path:filename>")
def serve_asset(filename):
    asset_path = BASE_DIR / "assets" / filename
    if not asset_path.exists():
        return "", 404
    return send_file(str(asset_path))

@app.route("/api/image-preview/<path:filename>")
@login_required
def api_image_preview(filename):
    # Try local path first
    try:
        img_path = IMAGES_DIR / filename
        if img_path.exists() and img_path.is_file():
            suffix = img_path.suffix.lower()
            mime = "image/jpeg" if suffix in {".jpg", ".jpeg"} else "image/png"
            return send_file(str(img_path), mimetype=mime)
    except Exception:
        pass
    # Fallback: Google Drive API (cloud deployment)
    try:
        from drive_helper import download_by_name
        data = download_by_name(filename)
        suffix = Path(filename).suffix.lower()
        mime = "image/jpeg" if suffix in {".jpg", ".jpeg"} else "image/png"
        from flask import Response as FlaskResponse
        return FlaskResponse(data, mimetype=mime)
    except Exception:
        return "", 404

@app.route("/api/post-instagram", methods=["POST"])
@login_required
def api_post_instagram():
    data     = request.json or {}
    caption  = data.get("caption","").strip()
    hashtags = data.get("hashtags","").strip()
    image    = data.get("image","").strip()
    env = os.environ.copy()
    env.update({"AUTO_POST":"true","ANTHROPIC_API_KEY":ANTHROPIC_KEY})
    if caption:  env["OVERRIDE_CAPTION"]  = caption
    if hashtags: env["OVERRIDE_HASHTAGS"] = hashtags
    if image:    env["OVERRIDE_IMAGE"]    = image
    job_id = str(uuid.uuid4())
    _jobs[job_id] = {"status": "running", "output": ""}
    threading.Thread(target=_run_bg, args=(job_id, [sys.executable, str(BASE_DIR/"post_to_instagram.py")], env, str(BASE_DIR)), daemon=True).start()
    return jsonify({"status": "started", "job_id": job_id})

@app.route("/api/post-facebook", methods=["POST"])
@login_required
def api_post_facebook():
    data     = request.json or {}
    caption  = data.get("caption","").strip()
    hashtags = data.get("hashtags","").strip()
    image    = data.get("image","").strip()
    env = os.environ.copy()
    env.update({"AUTO_POST":"true","ANTHROPIC_API_KEY":ANTHROPIC_KEY})
    if caption:  env["OVERRIDE_CAPTION"]  = caption
    if hashtags: env["OVERRIDE_HASHTAGS"] = hashtags
    if image:    env["OVERRIDE_IMAGE"]    = image
    try:
        r = subprocess.run([sys.executable, str(BASE_DIR/"post_to_facebook.py")],
                           capture_output=True, text=True, env=env, cwd=str(BASE_DIR), timeout=180)
        return jsonify({"status":"ok" if r.returncode==0 else "error","output":r.stdout+r.stderr})
    except Exception as e:
        return jsonify({"status":"error","output":str(e)})

# ── WORKFLOW CONTROLS ─────────────────────────────────────────────────────────

WORKFLOW_DEFS = [
    {
        "id": "daily_post.yml", "name": "Daily FB + IG + Drip Email",
        "schedule": "Daily 8:00 AM PHT", "emoji": "fa-solid fa-bolt",
        "description": "Posts to Facebook & Instagram using today's scheduled image, then sends drip emails to all active leads.",
        "status_keys": ["facebook_post", "instagram_post", "drip_email"],
    },
    {
        "id": "social_listening.yml", "name": "Social Listening",
        "schedule": "Daily 7:50 AM PHT", "emoji": "fa-solid fa-magnifying-glass",
        "description": "Scrapes competitor websites & Google News, generates AI intelligence briefing, emails it to the team.",
        "status_keys": ["social_listening"],
    },
    {
        "id": "weekly_analytics.yml", "name": "Weekly Analytics Report",
        "schedule": "Monday 8:00 AM PHT", "emoji": "fa-solid fa-chart-bar",
        "description": "Generates a branded HTML + PDF analytics report with lead stats, program breakdown, and AI insights. Emails to team.",
        "status_keys": [],
    },
    {
        "id": "weekly_preview.yml", "name": "Weekly Campaign Preview",
        "schedule": "Sunday 5:00 PM PHT", "emoji": "fa-solid fa-clipboard-list",
        "description": "Emails a PDF preview of next week's FB + IG post schedule with captions, hashtags, and drip email plan per lead.",
        "status_keys": ["weekly_preview"],
    },
    {
        "id": "cold_reengagement.yml", "name": "Cold Lead Re-engagement",
        "schedule": "Monday 6:00 AM PHT", "emoji": "fa-solid fa-snowflake",
        "description": "Automatically sends a re-engagement email to cold leads (inactive or low email count) to bring them back into the funnel.",
        "status_keys": ["cold_reengagement"],
    },
    {
        "id": "monthly_report.yml", "name": "Monthly Marketing Report",
        "schedule": "Last day of month 8:00 AM PHT", "emoji": "fa-solid fa-chart-pie",
        "description": "End-of-month PDF report with full lead funnel analysis, post performance summary, and strategic recommendations.",
        "status_keys": ["monthly_report"],
    },
    {
        "id": "check_preview_reply.yml", "name": "Check Preview Reply",
        "schedule": "Every 2hrs Sun–Mon PHT", "emoji": "fa-solid fa-reply",
        "description": "Checks if the weekly campaign preview email has been replied to (approved/rejected) and updates the preview status.",
        "status_keys": [],
    },
]

@app.route("/api/workflows/status")
@login_required
def api_workflows_status():
    pat  = os.environ.get("GITHUB_PAT", "")
    repo = os.environ.get("GITHUB_REPO", "your-org/your-repo")
    if not pat:
        return jsonify({"error": "GITHUB_PAT not set"}), 400
    try:
        resp = http_requests.get(
            f"https://api.github.com/repos/{repo}/actions/workflows",
            headers={"Authorization": f"token {pat}", "Accept": "application/vnd.github.v3+json"},
            timeout=10
        )
        gh_workflows = {w["path"].split("/")[-1]: w["state"] for w in resp.json().get("workflows", [])}
        live = get_automations_status()
        result = []
        for wf in WORKFLOW_DEFS:
            state = gh_workflows.get(wf["id"], "unknown")
            # Collect live status from mapped keys
            live_statuses = [live[k] for k in wf.get("status_keys", []) if k in live]
            combined_status = "none"
            combined_detail = ""
            if live_statuses:
                # Pick worst status: error > warning > success
                priority = {"inactive": 0, "none": 1, "scheduled": 2, "pending": 3, "success": 4, "warning": 5}
                worst = sorted(live_statuses, key=lambda x: priority.get(x.get("status","none"), 0))[-1]
                combined_status = worst.get("status", "none")
                combined_detail = " · ".join(s.get("detail","") for s in live_statuses if s.get("detail"))
            result.append({
                **{k: v for k, v in wf.items() if k != "status_keys"},
                "enabled": state == "active",
                "state": state,
                "live_status": combined_status,
                "live_detail": combined_detail,
            })
        return jsonify(result)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/workflows/run", methods=["POST"])
@login_required
def api_workflows_run():
    data  = request.get_json() or {}
    wf_id = data.get("workflow_id", "")
    pat   = os.environ.get("GITHUB_PAT", "")
    repo  = os.environ.get("GITHUB_REPO", "your-org/your-repo")
    if not pat:
        return jsonify({"success": False, "error": "GITHUB_PAT not set"}), 400
    valid_ids = {w["id"] for w in WORKFLOW_DEFS}
    if wf_id not in valid_ids:
        return jsonify({"success": False, "error": "Unknown workflow"}), 400
    try:
        resp = http_requests.post(
            f"https://api.github.com/repos/{repo}/actions/workflows/{wf_id}/dispatches",
            headers={"Authorization": f"token {pat}", "Accept": "application/vnd.github.v3+json"},
            json={"ref": "main"},
            timeout=15
        )
        if resp.status_code in (204, 200, 201):
            return jsonify({"success": True})
        return jsonify({"success": False, "error": f"GitHub {resp.status_code}: {resp.text[:200]}"}), 400
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route("/api/workflows/toggle", methods=["POST"])
@login_required
def api_workflows_toggle():
    data     = request.get_json() or {}
    wf_id    = data.get("workflow_id", "")
    enable   = data.get("enable", True)
    pat      = os.environ.get("GITHUB_PAT", "")
    repo     = os.environ.get("GITHUB_REPO", "your-org/your-repo")
    if not pat:
        return jsonify({"success": False, "error": "GITHUB_PAT not set"}), 400
    valid_ids = {w["id"] for w in WORKFLOW_DEFS}
    if wf_id not in valid_ids:
        return jsonify({"success": False, "error": "Unknown workflow"}), 400
    action = "enable" if enable else "disable"
    try:
        resp = http_requests.put(
            f"https://api.github.com/repos/{repo}/actions/workflows/{wf_id}/{action}",
            headers={"Authorization": f"token {pat}", "Accept": "application/vnd.github.v3+json"},
            timeout=10
        )
        if resp.status_code == 204:
            return jsonify({"success": True, "enabled": enable})
        return jsonify({"success": False, "error": f"GitHub {resp.status_code}: {resp.text[:200]}"}), 400
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

# ── POST COMPOSER ROUTES ───────────────────────────────────────────────────────

def _preload_thumbs():
    """Background thread: download all scheduled thumbnails into THUMB_CACHE_DIR."""
    import threading, time as _time
    def _run():
        try:
            from drive_helper import download_by_name
        except Exception:
            return
        fnames = []
        if SCHEDULE_FILE.exists():
            try:
                with open(SCHEDULE_FILE) as f:
                    sched = json.load(f).get("schedule", {})
                fnames = list(dict.fromkeys(sched.values()))
            except Exception:
                pass
        for fname in fnames:
            if (IMAGES_DIR / fname).exists() or (THUMB_CACHE_DIR / fname).exists():
                continue
            try:
                data = download_by_name(fname)
                (THUMB_CACHE_DIR / fname).write_bytes(data)
            except Exception:
                pass
    t = threading.Thread(target=_run, daemon=True)
    t.start()

@app.route("/api/composer/preload", methods=["POST"])
@login_required
def composer_preload():
    _preload_thumbs()
    return jsonify({"status": "started"})

@app.route("/api/composer/images")
@login_required
def composer_images():
    images, seen = [], set()
    today = str(date.today())
    if SCHEDULE_FILE.exists():
        with open(SCHEDULE_FILE) as f:
            sched = json.load(f).get("schedule", {})
        for d, fname in sorted(sched.items()):
            if fname not in seen:
                seen.add(fname)
                images.append({"filename": fname, "date": d, "is_today": d == today,
                                "available": (IMAGES_DIR / fname).exists(),
                                "thumb_url": f"/api/composer/thumb/{fname}"})
    # Also scan local staging folder for unscheduled images
    if IMAGES_DIR.exists():
        for ext in ("*.jpg", "*.jpeg", "*.png"):
            for fp in sorted(IMAGES_DIR.glob(ext)):
                if fp.name not in seen:
                    seen.add(fp.name)
                    images.append({"filename": fp.name, "date": None, "is_today": False,
                                   "available": True,
                                   "thumb_url": f"/api/composer/thumb/{fp.name}"})
    # Sort: available first, then unavailable
    images.sort(key=lambda x: (0 if x["available"] else 1, x.get("date") or ""))
    return jsonify(images)

@app.route("/api/composer/thumb/<path:filename>")
@login_required
def composer_thumb(filename):
    local = IMAGES_DIR / filename
    if local.exists():
        from flask import send_file as _sf
        return _sf(str(local))
    cached = THUMB_CACHE_DIR / filename
    if cached.exists():
        from flask import send_file as _sf
        return _sf(str(cached))
    try:
        from drive_helper import download_by_name
        data = download_by_name(filename)
        try:
            cached.write_bytes(data)
        except Exception:
            pass
        from flask import Response as _Resp
        return _Resp(data, mimetype="image/jpeg")
    except Exception:
        return "", 404

@app.route("/api/composer/caption", methods=["POST"])
@login_required
def composer_caption():
    data     = request.json or {}
    filename = data.get("filename", "")
    platform = data.get("platform", "facebook")
    today_str = date.today().strftime("%B %d, %Y")
    img_name  = Path(filename).stem.replace("_", " ").replace("-", " ")

    # Load image bytes for vision (local → thumb cache → Drive)
    img_b64 = None
    local = IMAGES_DIR / filename
    cached = THUMB_CACHE_DIR / filename
    if local.exists():
        raw = local.read_bytes()
    elif cached.exists():
        raw = cached.read_bytes()
    else:
        try:
            from drive_helper import download_by_name
            raw = download_by_name(filename)
            try: cached.write_bytes(raw)
            except Exception: pass
        except Exception:
            raw = None
    if raw:
        from PIL import Image as _PIL
        import io as _io2
        buf = _io2.BytesIO(raw)
        try:
            img = _PIL.open(buf).convert("RGB")
            out = _io2.BytesIO()
            img.save(out, format="JPEG", quality=75)
            img_b64 = _b64_mod.b64encode(out.getvalue()).decode()
        except Exception:
            img_b64 = None

    if platform == "instagram":
        rule_a = "lifestyle-forward and aspirational — open with a bold visual statement"
        rule_b = "warm and community-driven — open with a relatable insight or question"
        fmt    = "2-3 short punchy paragraphs. End with a soft CTA (link in bio / DM us). No hashtags. 1-2 emojis placed naturally in the text."
    else:
        rule_a = "professional and informative — open with a compelling hook"
        rule_b = "warm and conversational — open with a relatable scenario"
        fmt    = "2-4 paragraphs. Include a CTA. No hashtags. 1-2 emojis placed naturally in the text."

    base_prompt = (
        f"You are the Social Media Manager of Enderun Extension — the professional and continuing "
        f"education arm of Enderun Colleges, in McKinley Hill, BGC, Taguig, Philippines.\n\n"
        f"Today is {today_str}. Write a {platform.title()} caption for a post about: {img_name}.\n\n"
        f"Brand voice: aspirational, warm, Filipino-proud, confident. Pure English only — no Tagalog.\n"
        f"{fmt}\nWrite ONLY the caption text. Nothing else."
    )

    client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)

    def _gen(style_hint):
        prompt = base_prompt + f"\n\nCaption style: {style_hint}"
        if img_b64:
            resp = client.messages.create(
                model="claude-opus-4-7", max_tokens=500,
                messages=[{"role": "user", "content": [
                    {"type": "image", "source": {"type": "base64", "media_type": "image/jpeg", "data": img_b64}},
                    {"type": "text", "text": prompt}
                ]}]
            )
        else:
            resp = client.messages.create(
                model="claude-sonnet-4-6", max_tokens=500,
                messages=[{"role": "user", "content": prompt}]
            )
        return resp.content[0].text.strip()

    with concurrent.futures.ThreadPoolExecutor(max_workers=2) as ex:
        fa = ex.submit(_gen, rule_a)
        fb = ex.submit(_gen, rule_b)
        caption_a, caption_b = fa.result(), fb.result()

    return jsonify({"caption_a": caption_a, "caption_b": caption_b})

@app.route("/api/composer/post", methods=["POST"])
@login_required
def composer_post():
    data     = request.json or {}
    filename = data.get("filename", "").strip()
    caption  = data.get("caption", "").strip()
    platform = data.get("platform", "facebook").lower().strip()
    b64_data = data.get("b64", "")
    
    if filename:
        target = IMAGES_DIR / filename
        if b64_data:
            target.parent.mkdir(parents=True, exist_ok=True)
            target.write_bytes(_b64_mod.b64decode(b64_data))
        elif not target.exists():
            try:
                from drive_helper import download_by_name
                target.parent.mkdir(parents=True, exist_ok=True)
                target.write_bytes(download_by_name(filename))
            except Exception: pass
            
    env = os.environ.copy()
    env.update({"AUTO_POST": "true", "ANTHROPIC_API_KEY": ANTHROPIC_KEY})
    if caption:  env["OVERRIDE_CAPTION"] = caption
    if filename: env["OVERRIDE_IMAGE"]   = filename
    results = {}
    if platform in ("facebook", "both"):
        try:
            r = subprocess.run([sys.executable, str(BASE_DIR / "post_to_facebook.py")],
                               capture_output=True, text=True, env=env, cwd=str(BASE_DIR), timeout=180)
            results["facebook"] = "success" if r.returncode == 0 else (r.stdout + r.stderr)[-300:]
        except Exception as e:
            results["facebook"] = str(e)
    if platform in ("instagram", "both"):
        try:
            r = subprocess.run([sys.executable, str(BASE_DIR / "post_to_instagram.py")],
                               capture_output=True, text=True, env=env, cwd=str(BASE_DIR), timeout=180)
            results["instagram"] = "success" if r.returncode == 0 else (r.stdout + r.stderr)[-300:]
        except Exception as e:
            results["instagram"] = str(e)
    all_ok = all(v == "success" for v in results.values())
    return jsonify({"ok": all_ok, "results": results})

@app.route("/api/notifications")
@login_required
def api_notifications():
    return jsonify({"notifications": get_all(), "unread": unread_count()})

@app.route("/api/notifications/read", methods=["POST"])
@login_required
def api_notif_read():
    nid = (request.json or {}).get("id")
    if nid:
        mark_read(nid)
    return jsonify({"status": "ok"})

@app.route("/api/notifications/read-all", methods=["POST"])
@login_required
def api_notif_read_all():
    mark_all_read()
    return jsonify({"status": "ok"})

@app.route("/api/notifications/clear", methods=["POST"])
@login_required
def api_notif_clear():
    clear_all()
    return jsonify({"status": "ok"})

@app.route("/api/webhook/notify", methods=["POST"])
def api_webhook_notify():
    token = request.headers.get("X-Notify-Token", "")
    pat   = os.environ.get("GITHUB_PAT", "")
    if not pat or token != pat:
        return jsonify({"error": "unauthorized"}), 401
    data = request.json or {}
    from notifications_helper import push_notification as _push
    _push(
        agent_id=data.get("agent_id", "drip-campaign"),
        level=data.get("level", "success"),
        title=data.get("title", "Done"),
        message=data.get("message", "")
    )
    return jsonify({"status": "ok"})

@app.route("/api/drip-schedule")
@login_required
def api_drip_schedule():
    leads = []
    if LEADS_FILE.exists():
        with open(LEADS_FILE, newline="", encoding="utf-8") as f:
            leads = list(csv.DictReader(f))
    sched = {}
    if SCHEDULE_FILE.exists():
        with open(SCHEDULE_FILE) as f:
            sched = json.load(f).get("schedule", {})
    images = [sched[k] for k in sorted(sched.keys())]
    result = []
    for lead in leads:
        if lead.get("status","").strip().lower() != "active":
            continue
        email_count = int(lead.get("email_count", "0") or "0")
        next_email  = email_count + 1
        next_image  = images[email_count % len(images)] if images else "N/A"
        result.append({
            "name":     lead.get("first_name","") + " " + lead.get("last_name",""),
            "email":    lead.get("email",""),
            "program":  lead.get("program_interest",""),
            "sent":     email_count,
            "next_num": next_email,
            "next_img": next_image,
        })
    return jsonify(result)


# ── DOCUMENT GENERATION ────────────────────────────────────────────────────────

CHAT_FILES_DIR = BASE_DIR / "output" / "chat_files"

def _parse_md_sections(text):
    """Return (title, subtitle, [(heading, body), ...]) from markdown."""
    lines = text.strip().split('\n')
    title, subtitle, sections = '', '', []
    cur_head, cur_body = None, []
    for line in lines:
        if line.startswith('# ') and not title:
            title = line[2:].strip()
        elif line.startswith('### ') and not subtitle:
            subtitle = line[4:].strip()
        elif line.startswith('## '):
            if cur_head is not None:
                sections.append((cur_head, '\n'.join(cur_body)))
            cur_head, cur_body = line[3:].strip(), []
        elif cur_head is not None:
            cur_body.append(line)
    if cur_head is not None:
        sections.append((cur_head, '\n'.join(cur_body)))
    return title or 'Document', subtitle, sections

def _clean_md(text):
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    return text.strip()

def _generate_pdf_chat(title, content_md):
    from report_helper import ReportBuilder
    CHAT_FILES_DIR.mkdir(parents=True, exist_ok=True)

    # ── JSON spec format (structured, designed PDF) ──
    raw = content_md.strip()
    if raw.startswith('{'):
        try:
            spec = json.loads(raw)
            pdf_title    = spec.get("title", title)
            pdf_subtitle = spec.get("subtitle", "")
            rb = ReportBuilder(agent_id="chat", report_title=pdf_title, subtitle=pdf_subtitle)
            rb.add_source("AI-generated report — data sourced from chat context")
            for block in spec.get("sections", []):
                btype = block.get("type", "")
                if btype == "section":
                    rb.add_section(block["title"])
                elif btype == "paragraph":
                    rb.add_paragraph(_clean_md(block.get("text", "")))
                elif btype == "bullets":
                    rb.add_bullets([_clean_md(b) for b in block.get("items", [])])
                elif btype == "kpi_row":
                    items = block.get("items", [])
                    rb.add_kpi_row([(i["label"], str(i["value"]), i.get("highlight", False)) for i in items])
                elif btype == "table":
                    rb.add_table(headers=block.get("headers", []), rows=block.get("rows", []))
                elif btype == "alert":
                    rb.add_alert_box(_clean_md(block.get("text", "")), level=block.get("level", "info"))
                elif btype == "bar_chart":
                    rb.add_bar_chart(block.get("title", ""), labels=block.get("labels", []),
                                     values=block.get("values", []), color=block.get("color", "brown"),
                                     horizontal=block.get("horizontal", False))
                elif btype == "line_chart":
                    rb.add_line_chart(block.get("title", ""), labels=block.get("labels", []),
                                      series=block.get("series", []), series_labels=block.get("series_labels", []))
                elif btype == "pie_chart":
                    rb.add_pie_chart(block.get("title", ""), labels=block.get("labels", []),
                                     values=block.get("values", []))
                elif btype == "charts_row":
                    rb.add_charts_row(left=block.get("left", {}), right=block.get("right", {}))
                elif btype == "callout":
                    rb.add_callout(block.get("stat", ""), block.get("label", ""), block.get("description", ""))
                elif btype == "page_break":
                    rb.add_page_break()
                elif btype == "source":
                    rb.add_source(block.get("text", ""))
            return Path(rb.save())
        except (json.JSONDecodeError, KeyError, Exception):
            pass  # fall through to markdown fallback

    # ── Markdown fallback (basic sections + bullets) ──
    rb = ReportBuilder(agent_id="chat", report_title=title)
    _, _, sections = _parse_md_sections(content_md)
    for sec_title, sec_body in sections:
        rb.add_section(sec_title)
        bullets, paras = [], []
        for line in sec_body.split('\n'):
            stripped = line.strip()
            if not stripped:
                if paras:
                    rb.add_paragraph(_clean_md(' '.join(paras))); paras = []
                if bullets:
                    rb.add_bullets([_clean_md(b) for b in bullets]); bullets = []
            elif stripped.startswith('- ') or stripped.startswith('* '):
                if paras:
                    rb.add_paragraph(_clean_md(' '.join(paras))); paras = []
                bullets.append(stripped[2:])
            else:
                if bullets:
                    rb.add_bullets([_clean_md(b) for b in bullets]); bullets = []
                paras.append(stripped)
        if paras:
            rb.add_paragraph(_clean_md(' '.join(paras)))
        if bullets:
            rb.add_bullets([_clean_md(b) for b in bullets])
    return Path(rb.save())

def _generate_pptx_chat(title, content_md, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PRIMARY = RGBColor(0x7a, 0x10, 0x28)
    WHITE   = RGBColor(0xff, 0xff, 0xff)
    DARK    = RGBColor(0x2d, 0x0f, 0x18)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, color):
        from pptx.util import Emu
        shp = slide.shapes.add_shape(1, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = color
        shp.line.fill.background(); return shp

    # Title slide
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
    tx = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.33), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(38); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

    _, _, sections = _parse_md_sections(content_md)
    for sec_title, sec_body in sections:
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, prs.slide_width, prs.slide_height, WHITE)
        add_rect(s, 0, 0, prs.slide_width, Inches(0.07), PRIMARY)
        th = s.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.73), Inches(0.9))
        p = th.text_frame.paragraphs[0]
        p.text = sec_title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = PRIMARY
        cb = s.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.73), Inches(5.6))
        tf = cb.text_frame; tf.word_wrap = True; first = True
        for line in sec_body.split('\n'):
            stripped = line.strip()
            if not stripped: continue
            is_b = stripped.startswith('- ') or stripped.startswith('* ')
            clean = _clean_md(stripped[2:] if is_b else stripped)
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.text = ('• ' if is_b else '') + clean
            p.font.size = Pt(16); p.font.color.rgb = DARK

    prs.save(str(out_path))

def _generate_docx_chat(title, content_md, out_path):
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    doc = Document()
    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    _, _, sections = _parse_md_sections(content_md)
    for sec_title, sec_body in sections:
        doc.add_heading(sec_title, level=1)
        for line in sec_body.split('\n'):
            stripped = line.strip()
            if not stripped: continue
            if stripped.startswith('- ') or stripped.startswith('* '):
                doc.add_paragraph(_clean_md(stripped[2:]), style='List Bullet')
            elif re.match(r'^[0-9]+\. ', stripped):
                doc.add_paragraph(_clean_md(re.sub(r'^[0-9]+\. ', '', stripped)), style='List Number')
            else:
                doc.add_paragraph(_clean_md(stripped))
    doc.save(str(out_path))

def _generate_csv_chat(content_md, out_path):
    """Parse markdown pipe tables or raw CSV rows and write to a .csv file."""
    import csv as csv_mod
    lines = [l.strip() for l in content_md.strip().split('\n') if l.strip()]
    rows = []
    for line in lines:
        if line.startswith('|'):
            cells = [c.strip() for c in line.strip('|').split('|')]
            if all(re.match(r'^-+$', c.replace(' ', '')) for c in cells):
                continue  # skip separator row
            rows.append(cells)
        elif ',' in line:
            rows.append(next(csv_mod.reader([line])))
    with open(str(out_path), 'w', newline='', encoding='utf-8') as f:
        w = csv_mod.writer(f)
        w.writerows(rows)

@app.route("/api/generate-doc", methods=["POST"])
@login_required
def api_generate_doc():
    data     = request.json or {}
    doc_type = data.get("type", "pdf").lower()
    content  = data.get("content", "").strip()
    if not content:
        return jsonify({"status": "error", "message": "No content provided"})
    # Extract title: JSON spec or markdown heading
    if content.strip().startswith('{'):
        try:
            title = json.loads(content.strip()).get("title", "Document")
        except Exception:
            title = "Document"
    else:
        title = next((l[2:].strip() for l in content.split('\n') if l.startswith('# ')), 'Document')
    safe  = re.sub(r'[^a-zA-Z0-9_-]', '_', title)[:40]
    stamp = date.today().isoformat()
    CHAT_FILES_DIR.mkdir(parents=True, exist_ok=True)
    try:
        if doc_type == "pdf":
            path = _generate_pdf_chat(title, content)
            filename = path.name
        elif doc_type == "pptx":
            filename = f"{stamp}_{safe}.pptx"
            _generate_pptx_chat(title, content, CHAT_FILES_DIR / filename)
        elif doc_type == "docx":
            filename = f"{stamp}_{safe}.docx"
            _generate_docx_chat(title, content, CHAT_FILES_DIR / filename)
        elif doc_type == "csv":
            filename = f"{stamp}_{safe}.csv"
            _generate_csv_chat(content, CHAT_FILES_DIR / filename)
        else:
            return jsonify({"status": "error", "message": "Unknown type"})
        return jsonify({"status": "ok", "filename": filename, "title": title})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

@app.route("/api/export-chat-pdf", methods=["POST"])
@login_required
def api_export_chat_pdf():
    data        = request.json or {}
    messages    = data.get("messages", [])
    agent_name  = data.get("agent_name", "AI Agent")
    lines = [f"# Chat Export — {agent_name}", f"## {date.today().strftime('%B %d, %Y')}"]
    for m in messages:
        role    = m.get("role", "")
        content = m.get("content", "")
        if isinstance(content, list):
            content = " ".join(b.get("text", "") for b in content if isinstance(b, dict) and b.get("type") == "text")
        content = str(content).strip()
        if not content or content == "__GREET__":
            continue
        label = "You" if role == "user" else agent_name
        lines.append(f"## {label}")
        lines.append(content[:3000])
    CHAT_FILES_DIR.mkdir(parents=True, exist_ok=True)
    try:
        path = _generate_pdf_chat(f"Chat — {agent_name}", "\n\n".join(lines))
        return jsonify({"status": "ok", "filename": path.name})
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)})

@app.route("/chat-files/<path:filename>")
@login_required
def serve_chat_file(filename):
    return send_from_directory(str(CHAT_FILES_DIR), filename, as_attachment=True)

@app.route("/chat-image/<path:filename>")
@login_required
def serve_chat_image(filename):
    return send_from_directory(str(CHAT_IMAGES_DIR), filename)


# ── SERVER-SIDE CHAT HISTORY ───────────────────────────────────────────────────

@app.route("/api/chat/save-session", methods=["POST"])
@login_required
def api_save_chat_session():
    data       = request.json or {}
    agent_id   = data.get("agent_id", "marketing-manager")
    session_id = data.get("session_id") or str(uuid.uuid4())
    messages   = data.get("messages", [])
    title      = data.get("title", "Chat Session")
    agent_dir  = CHAT_SESSIONS_DIR / agent_id
    agent_dir.mkdir(parents=True, exist_ok=True)
    payload = {"id": session_id, "agent_id": agent_id, "title": title,
               "messages": messages, "updated": str(date.today())}
    (agent_dir / f"{session_id}.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return jsonify({"status": "ok", "session_id": session_id})

@app.route("/api/chat/load-sessions")
@login_required
def api_load_chat_sessions():
    agent_id  = request.args.get("agent_id", "marketing-manager")
    agent_dir = CHAT_SESSIONS_DIR / agent_id
    if not agent_dir.exists():
        return jsonify({"sessions": []})
    sessions = []
    for f in sorted(agent_dir.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True)[:50]:
        if "_summary" in f.name:
            continue
        try:
            d = json.loads(f.read_text(encoding="utf-8"))
            sessions.append({"id": d.get("id"), "title": d.get("title","Chat"),
                              "updated": d.get("updated",""), "message_count": len(d.get("messages",[]))})
        except Exception:
            pass
    return jsonify({"sessions": sessions})

@app.route("/api/chat/get-session/<session_id>")
@login_required
def api_get_chat_session(session_id):
    agent_id = request.args.get("agent_id", "marketing-manager")
    fname    = CHAT_SESSIONS_DIR / agent_id / f"{session_id}.json"
    if not fname.exists():
        return jsonify({"error": "Session not found"}), 404
    try:
        return jsonify(json.loads(fname.read_text(encoding="utf-8")))
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/chat/delete-session/<session_id>", methods=["DELETE"])
@login_required
def api_delete_chat_session(session_id):
    agent_id = request.args.get("agent_id", "marketing-manager")
    fname    = CHAT_SESSIONS_DIR / agent_id / f"{session_id}.json"
    if fname.exists():
        fname.unlink()
    return jsonify({"status": "ok"})


# ── POST BUILDER ROUTES ────────────────────────────────────────────────────────

@app.route("/api/post-builder/images")
@login_required
def post_builder_images():
    sched_file = BASE_DIR / "posting_schedule.json"
    images = []
    if sched_file.exists():
        try:
            data = json.loads(sched_file.read_text(encoding="utf-8"))
            schedule = data.get("schedule", {})
            seen = set()
            today_str = date.today().isoformat()
            for d, img in sorted(schedule.items()):
                if isinstance(img, str) and img not in seen:
                    seen.add(img)
                    images.append({"filename": img, "date": d, "is_today": d == today_str})
        except Exception:
            pass
    return jsonify({"images": images[:30]})


@app.route("/api/post-builder/image/<path:filename>")
@login_required
def post_builder_image(filename):
    fname = Path(filename).name
    for folder in [Path(r"G:\My Drive\FB_Post_Today"), BASE_DIR / "output"]:
        p = folder / fname
        if p.exists():
            mime = "image/jpeg" if fname.lower().endswith((".jpg", ".jpeg")) else "image/png"
            return send_file(str(p), mimetype=mime)
    return jsonify({"error": "not found"}), 404


@app.route("/api/post-builder/captions", methods=["POST"])
@login_required
def post_builder_captions():
    data   = request.get_json() or {}
    platform   = data.get("platform", "facebook")
    image_name = data.get("image_name", "post image")
    today_s    = date.today().strftime("%B %d, %Y")
    topic      = Path(image_name).stem.replace("_", " ").replace("-", " ")

    if platform == "instagram":
        base = (
            f"You are the Instagram Manager of Enderun Extension — the professional continuing education arm "
            f"of Enderun Colleges, McKinley Hill BGC, Philippines.\n"
            f"Today is {today_s}. Write an Instagram caption for: {topic}.\n"
            f"Rules: 2-3 short punchy paragraphs · aspirational lifestyle tone · soft CTA at end · "
            f"NO hashtags · NO URL · 1-2 emojis naturally in text. Write ONLY the caption."
        )
    else:
        base = (
            f"You are the Facebook Page Manager of Enderun Extension — the professional continuing education arm "
            f"of Enderun Colleges, McKinley Hill BGC, Philippines.\n"
            f"Today is {today_s}. Write a Facebook caption for: {topic}.\n"
            f"Rules: 3-4 warm community-focused paragraphs · aspirational but accessible · clear CTA at end · "
            f"NO hashtags · NO URL · 1-2 emojis naturally in text. Write ONLY the caption."
        )

    prompts = [
        base,
        base.replace("Write ONLY the caption.", "Write a COMPLETELY DIFFERENT version — different angle, hook, and tone. Write ONLY the caption.")
    ]

    client  = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY", ""))
    captions = []
    for prompt in prompts:
        try:
            r = client.messages.create(
                model="claude-opus-4-7", max_tokens=500,
                messages=[{"role": "user", "content": prompt}]
            )
            captions.append(r.content[0].text.strip())
        except Exception as e:
            captions.append(f"Error: {e}")
    return jsonify({"captions": captions})


@app.route("/api/post-builder/post", methods=["POST"])
@login_required
def post_builder_post():
    data     = request.get_json() or {}
    platform = data.get("platform", "facebook")
    caption  = data.get("caption", "")
    image    = data.get("image", "")
    pat      = os.environ.get("GITHUB_PAT", "")
    repo     = os.environ.get("GITHUB_REPO", "your-org/your-repo")
    if not pat:
        return jsonify({"success": False, "error": "GITHUB_PAT not set"}), 400
    workflow = "post_fb_now.yml" if platform == "facebook" else "post_ig_now.yml"
    try:
        resp = http_requests.post(
            f"https://api.github.com/repos/{repo}/actions/workflows/{workflow}/dispatches",
            headers={"Authorization": f"token {pat}", "Accept": "application/vnd.github.v3+json"},
            json={"ref": "main", "inputs": {"caption": caption, "image_filename": image}},
            timeout=15
        )
        if resp.status_code in (204, 200, 201):
            return jsonify({"success": True})
        return jsonify({"success": False, "error": f"GitHub {resp.status_code}: {resp.text[:200]}"}), 400
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


# ── STRATEGY STUDIO ───────────────────────────────────────────────────────────

@app.route("/api/strategy/generate", methods=["POST"])
@login_required
def api_strategy_generate():
    data      = request.json or {}
    doc_type  = data.get("type", "marketing-strategy")
    program   = data.get("program", "").strip()
    desc      = data.get("description", "").strip()
    market    = data.get("target_market", "").strip()
    formats   = data.get("formats", "").strip()
    budget    = data.get("budget", "").strip()
    timeline  = data.get("timeline", "").strip()
    notes     = data.get("notes", "").strip()

    if not program or not desc:
        return jsonify({"error": "Program name and description are required."}), 400

    type_labels = {
        "business-plan":      "Business Plan",
        "marketing-strategy": "Marketing Strategy",
        "strategic-plan":     "Strategic Plan",
    }
    type_label = type_labels.get(doc_type, "Strategy Document")

    type_structures = {
        "business-plan": (
            "Structure the document as a full Business Plan with these sections:\n"
            "1. Executive Summary\n2. Market Analysis (Philippine education/training landscape)\n"
            "3. Program Description & Value Proposition\n4. Target Audience & Personas\n"
            "5. Marketing & Sales Strategy\n6. Operations Plan\n7. Financial Overview (estimated costs, revenue model, ROI)\n"
            "8. Key Success Metrics & KPIs\n9. Risk Assessment & Mitigation\n10. Next Steps & Timeline"
        ),
        "marketing-strategy": (
            "Structure the document as a comprehensive Marketing Strategy with these sections:\n"
            "1. Situation Analysis (market context, Enderun strengths/gaps)\n"
            "2. Target Market & Audience Personas\n3. Value Proposition & Positioning\n"
            "4. Marketing Mix (Product, Price, Place, Promotion)\n"
            "5. Channel Strategy (social media, email, events, paid ads)\n"
            "6. Content Strategy & Key Messages\n7. Campaign Ideas & Activation Calendar\n"
            "8. Budget Allocation Recommendation\n9. KPIs & Measurement Framework\n10. 90-Day Action Plan"
        ),
        "strategic-plan": (
            "Structure the document as a Strategic Plan with these sections:\n"
            "1. Vision & Mission Alignment (with Enderun brand)\n2. Situational Assessment (SWOT)\n"
            "3. Strategic Objectives (3-year outlook)\n4. Key Initiatives & Pillars\n"
            "5. Resource Requirements (people, budget, tools)\n6. Organizational Considerations\n"
            "7. Implementation Roadmap (quarterly milestones)\n8. Risk Factors\n"
            "9. Success Metrics & Review Cadence\n10. Decision Points & Escalation"
        ),
    }
    structure_instructions = type_structures.get(doc_type, type_structures["marketing-strategy"])

    context_parts = [f"- Program/Initiative: {program}", f"- Description: {desc}"]
    if market:   context_parts.append(f"- Target Market: {market}")
    if formats:  context_parts.append(f"- Delivery Format: {formats}")
    if budget:   context_parts.append(f"- Budget Range: {budget}")
    if timeline: context_parts.append(f"- Timeline: {timeline}")
    if notes:    context_parts.append(f"- Additional Context: {notes}")
    context_block = "\n".join(context_parts)

    system_prompt = (
        "You are the Chief Strategy Officer at Enderun Colleges — a premium private higher education "
        "institution in BGC, Manila with world-class partnerships with Les Roches (Switzerland, top 3 globally "
        "in hospitality) and École Ducasse (France). Enderun runs three business units: "
        "Enderun Colleges (degree programs), Enderun Extension (continuing education / upskilling), "
        "and Enderun Events/Banquetes (premium venue & catering). "
        "Brand voice: aspirational, warm, confident, Filipino-proud. "
        "You produce executive-grade strategy documents — structured, data-aware, and actionable. "
        "Always tie recommendations back to Enderun's brand strengths and Filipino market context. "
        "Use markdown formatting: # for title, ## for sections, ### for subsections, bullet lists, tables where relevant. "
        "Be specific and practical — avoid generic consulting filler. "
        "Write at a senior executive level but remain accessible."
    )

    user_prompt = (
        f"Please create a comprehensive {type_label} document for the following initiative:\n\n"
        f"{context_block}\n\n"
        f"{structure_instructions}\n\n"
        "Be thorough and specific. Use real market context where applicable (Philippine education market, "
        "BGC premium positioning, Enderun's brand strengths). Include tables, bullet points, and concrete "
        "recommendations. Aim for a complete, boardroom-ready document."
    )

    client_obj = anthropic.Anthropic(api_key=ANTHROPIC_KEY)

    def generate():
        try:
            with client_obj.messages.stream(
                model="claude-opus-4-7",
                max_tokens=8000,
                system=system_prompt,
                messages=[{"role": "user", "content": user_prompt}]
            ) as stream:
                for text in stream.text_stream:
                    yield "data: " + json.dumps({"text": text}) + "\n\n"
        except Exception as e:
            yield "data: " + json.dumps({"text": f"\n\n❌ Error: {e}"}) + "\n\n"
        yield "data: [DONE]\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


# ── HTML (written to file at startup) ─────────────────────────────────────────


HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Enderun Marketing Hub</title>
<link rel="icon" type="image/png" href="/assets/logos/Enderun-Colleges-white.png">
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.0/css/all.min.css" crossorigin="anonymous">
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.css" crossorigin="anonymous">
<script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.js" crossorigin="anonymous"></script>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/themes/prism-tomorrow.min.css" crossorigin="anonymous">
<script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/prism.min.js" crossorigin="anonymous"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/prism/1.29.0/plugins/autoloader/prism-autoloader.min.js" crossorigin="anonymous"></script>
<style>
  :root {
    --primary:        #7a1028;
    --primary-d:      #4e0917;
    --primary-m:      #6a1222;
    --primary-l:      #9b1a35;
    --primary-dim:    rgba(122,16,40,0.07);
    --primary-border: rgba(122,16,40,0.22);
    --bg:         #ffffff;
    --bg-card:    #ffffff;
    --text:       #2d0f18;
    --text-muted: #7a4050;
    --text-dim:   #b08090;
    --card:       rgba(255,255,255,0.98);
    --card-border: rgba(122,16,40,0.1);
    --sidebar-w: 240px;
    --radius: 12px;
    --radius-sm: 8px;
    --transition: 0.2s ease;
  }
  /* ── DARK MODE VARS ── */
  body.dark {
    --bg:          #0f0f13;
    --bg-card:     #1a1a24;
    --text:        #e8e0e5;
    --text-muted:  #a090a0;
    --text-dim:    #6a5a6a;
    --card:        rgba(30,24,35,0.98);
    --card-border: rgba(255,255,255,0.08);
    --primary-dim:    rgba(155,26,53,0.18);
    --primary-border: rgba(155,26,53,0.35);
    --gold-dim:    rgba(201,168,76,0.1);
    --gold-border: rgba(201,168,76,0.2);
  }
  body.dark #sidebar { background: #0a0a0e; border-right-color: rgba(255,255,255,0.05); }
  body.dark #topbar  { background: #111118; border-bottom-color: rgba(255,255,255,0.07); }
  body.dark #content { background: #0f0f13; }
  body.dark .chat-messages { background: #0f0f13; }
  body.dark .chat-main    { background: #0f0f13; }
  body.dark .chat-sidebar { background: #111118; border-right-color: rgba(255,255,255,0.07); }
  body.dark .chat-header  { background: #111118; border-bottom-color: rgba(255,255,255,0.07); }
  body.dark .chat-input-area { background: #111118; border-top-color: rgba(255,255,255,0.07); }
  body.dark .chat-input-wrap { background: #1a1a24; border-color: rgba(255,255,255,0.1); }
  body.dark .chat-input-wrap:focus-within { background: #1e1e2a; border-color: rgba(155,26,53,0.5); }
  body.dark #chat-input { color: var(--text); }
  body.dark .chat-input-footer { border-top-color: rgba(255,255,255,0.05); background: rgba(0,0,0,0.2); }
  body.dark .msg-bubble { color: var(--text); }
  body.dark .msg-bubble h2 { border-color: rgba(255,255,255,0.1); }
  body.dark .msg-bubble code { background: rgba(255,255,255,0.08); color: #f9a8d4; border-color: rgba(255,255,255,0.1); }
  body.dark .msg-bubble strong { color: #f9a8d4; }
  body.dark .card  { background: var(--bg-card); border-color: var(--card-border); }
  body.dark .view  { background: #0f0f13; }
  body.dark #content { background: #0f0f13; }
  body.dark .custom-dd-trigger { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text); }
  body.dark .custom-dd-menu    { background: #1a1a24; border-color: rgba(255,255,255,0.1); }
  body.dark .custom-dd-opt     { color: var(--text); }
  body.dark .custom-dd-opt:hover { background: rgba(155,26,53,0.2); }
  body.dark .tb-dd-menu  { background: #1a1a24; border-color: rgba(255,255,255,0.1); }
  body.dark .tb-dd-item  { color: var(--text); }
  body.dark .tb-dd-item:hover { background: rgba(155,26,53,0.2); }
  body.dark .chat-session-item:hover { background: rgba(155,26,53,0.15); }
  body.dark .chat-session-item.active { background: rgba(155,26,53,0.2); }
  body.dark .chat-session-title { color: var(--text); }
  body.dark .wcard { background: #1a1a24; border-color: rgba(255,255,255,0.08); }
  body.dark .wcard:hover { background: rgba(155,26,53,0.2); }
  body.dark .wcard-title { color: var(--text); }
  body.dark .msg-action-btn { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text-muted); }
  body.dark .msg-action-btn:hover { background: rgba(155,26,53,0.2); border-color: var(--primary-border); }
  body.dark .q-chip { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text); }
  body.dark .q-chip:hover { background: rgba(155,26,53,0.2); }
  body.dark #attach-preview { background: rgba(155,26,53,0.12); border-color: rgba(155,26,53,0.25); }
  body.dark .md-table-wrap { border-color: rgba(255,255,255,0.1); }
  body.dark .md-table th { background: linear-gradient(135deg, #5a0c1e, #7a1028); }
  body.dark .md-table td { color: var(--text); border-color: rgba(255,255,255,0.07); }
  body.dark .new-chat-btn { box-shadow: 0 2px 8px rgba(0,0,0,0.4); }
  body.dark .tb-automate-btn { background: rgba(155,26,53,0.2); }
  body.dark .tb-icon-btn { border-color: rgba(255,255,255,0.12); color: var(--text-muted); background: transparent; }
  body.dark .tb-icon-btn:hover { background: rgba(155,26,53,0.18); border-color: var(--primary-border); color: var(--primary-l); }
  body.dark .welcome-logo-ring { background: rgba(155,26,53,0.15); border-color: rgba(155,26,53,0.2); }
  body.dark .file-download-card { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text); }
  body.dark .mode-pill { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text-muted); }
  body.dark .mode-pill.active { background: rgba(155,26,53,0.25); border-color: var(--primary-border); }
  body.dark .tool-action-card.done { background: rgba(20,83,45,0.2); border-color: rgba(34,197,94,0.25); }
  body.dark .tool-action-card.running { background: rgba(155,26,53,0.15); border-color: rgba(155,26,53,0.3); }
  body.dark .chat-hdr-btn { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text-muted); }
  body.dark .chat-hdr-btn:hover { background: rgba(155,26,53,0.2); }
  body.dark #quick-chips { background: #111118; border-top: 1px solid rgba(255,255,255,0.06); }
  body.dark .attach-btn { color: var(--text-dim); }
  body.dark .attach-btn:hover { background: rgba(155,26,53,0.2); }
  body.dark .shortcuts-modal { background: #1a1a24; border-color: rgba(255,255,255,0.1); }
  body.dark .sc-row { border-color: rgba(255,255,255,0.06); }
  body.dark .sc-key { background: #0f0f13; border-color: rgba(255,255,255,0.15); color: var(--text-muted); }
  body.dark .sc-group-title { color: var(--text-dim); }
  body.dark .feedback-btn { color: var(--text-dim); }
  body.dark .feedback-btn:hover { background: rgba(155,26,53,0.2); color: var(--primary-l); }
  body.dark .feedback-btn.active-up { color: #4ade80; background: rgba(74,222,128,0.12); }
  body.dark .feedback-btn.active-down { color: #f87171; background: rgba(248,113,113,0.12); }
  body.dark .followup-chips { border-top-color: rgba(255,255,255,0.06); background: #0f0f13; }
  body.dark .fup-chip { background: #1a1a24; border-color: rgba(255,255,255,0.1); color: var(--text-muted); }
  body.dark .fup-chip:hover { border-color: var(--primary-border); color: var(--primary-l); background: rgba(155,26,53,0.15); }
  body.dark .mention-dropdown { background: #1a1a24; border-color: rgba(255,255,255,0.1); }
  body.dark .mention-item:hover, body.dark .mention-item.active { background: rgba(155,26,53,0.2); }
  body.dark .file-queue { background: #1a1a24; border-color: rgba(255,255,255,0.08); }
  body.dark .file-queue-item { background: #0f0f13; border-color: rgba(255,255,255,0.06); }
  body.dark .scroll-pause-btn { background: #1a1a24; border-color: rgba(255,255,255,0.12); color: var(--text-muted); }
  body.dark .msg-counter { color: var(--text-dim); }
  body.dark .url-preview-chip { background: #1a1a24; border-color: rgba(155,26,53,0.3); }
  body.dark .sidebar-drag-handle { background: rgba(255,255,255,0.06); }
  body.dark .sidebar-drag-handle:hover { background: var(--primary-border); }
  /* ── SHORTCUTS MODAL ── */
  .shortcuts-overlay { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.45); z-index:9000; align-items:center; justify-content:center; backdrop-filter:blur(2px); }
  .shortcuts-overlay.open { display:flex; }
  .shortcuts-modal {
    background: #fff; border-radius: 16px; padding: 28px 28px 20px;
    width: 520px; max-width: 92vw; max-height: 80vh; overflow-y: auto;
    border: 1px solid rgba(0,0,0,0.1); box-shadow: 0 20px 60px rgba(0,0,0,0.18);
    position: relative;
  }
  .sc-modal-title { font-size: 16px; font-weight: 700; margin-bottom: 18px; display:flex; align-items:center; gap:8px; }
  .sc-close { position:absolute; top:16px; right:18px; background:none; border:none; cursor:pointer; font-size:18px; color:var(--text-dim); }
  .sc-close:hover { color:var(--text); }
  .sc-group-title { font-size:10px; font-weight:700; text-transform:uppercase; letter-spacing:.08em; color:#999; margin:14px 0 6px; }
  .sc-row { display:flex; align-items:center; justify-content:space-between; padding:7px 0; border-bottom:1px solid rgba(0,0,0,0.05); font-size:13px; }
  .sc-desc { color:var(--text); }
  .sc-keys { display:flex; gap:4px; }
  .sc-key { background:#f4f1f0; border:1px solid #ddd; border-radius:5px; padding:2px 7px; font-size:11px; font-family:monospace; font-weight:600; color:#555; white-space:nowrap; }
  /* ── FEEDBACK BUTTONS ── */
  .feedback-btn {
    background: none; border: none; cursor: pointer; padding: 3px 6px;
    border-radius: 5px; font-size: 13px; color: var(--text-dim);
    transition: all 0.15s; line-height: 1;
  }
  .feedback-btn:hover { background: rgba(0,0,0,0.06); color: var(--text); }
  .feedback-btn.active-up { color: #16a34a; background: rgba(22,163,74,0.1); }
  .feedback-btn.active-down { color: #dc2626; background: rgba(220,38,38,0.1); }
  /* ── FOLLOW-UP CHIPS ── */
  .followup-chips {
    display: flex; gap: 6px; flex-wrap: wrap; padding: 8px 14px 4px;
    border-top: 1px solid rgba(0,0,0,0.06);
    background: var(--bg);
  }
  .fup-chip {
    background: var(--bg-card, #fff); border: 1px solid var(--card-border);
    border-radius: 20px; padding: 5px 12px; font-size: 12px; color: var(--text);
    cursor: pointer; transition: all 0.15s; white-space: nowrap; font-family: inherit;
  }
  .fup-chip:hover { border-color: var(--primary-border); color: var(--primary); background: var(--primary-dim); }
  /* ── @MENTION DROPDOWN ── */
  .mention-dropdown {
    position: absolute; bottom: 100%; left: 0; right: 0;
    background: #fff; border: 1px solid var(--card-border);
    border-radius: 10px; box-shadow: 0 8px 24px rgba(0,0,0,0.12);
    max-height: 220px; overflow-y: auto; z-index: 200;
    display: none; padding: 4px 0;
  }
  .mention-dropdown.open { display: block; }
  .mention-item {
    display: flex; align-items: center; gap: 8px; padding: 7px 12px;
    cursor: pointer; font-size: 12.5px; color: var(--text); transition: background 0.12s;
  }
  .mention-item:hover, .mention-item.active { background: var(--primary-dim); }
  .mention-item .mention-emoji { font-size: 15px; flex-shrink: 0; }
  .mention-item .mention-name { font-weight: 600; }
  .mention-item .mention-role { color: var(--text-dim); font-size: 11px; margin-left: 4px; }
  /* ── FILE QUEUE ── */
  .file-queue {
    display: none; flex-direction: column; gap: 4px;
    padding: 8px 12px; background: var(--bg-card,#fafafa);
    border: 1px solid var(--card-border); border-bottom: none;
    border-radius: 10px 10px 0 0;
  }
  .file-queue.open { display: flex; }
  .file-queue-item {
    display: flex; align-items: center; gap: 8px; padding: 5px 8px;
    background: var(--bg,#fff); border: 1px solid var(--card-border);
    border-radius: 6px; font-size: 12px; color: var(--text);
  }
  .file-queue-item .fq-name { flex: 1; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
  .file-queue-item .fq-remove { background:none; border:none; cursor:pointer; color:var(--text-dim); font-size:14px; padding:0 2px; }
  .file-queue-item .fq-remove:hover { color:var(--primary); }
  /* ── SCROLL PAUSE BUTTON ── */
  .scroll-pause-btn {
    position: absolute; bottom: 10px; right: 14px; z-index: 10;
    background: #fff; border: 1px solid var(--card-border);
    border-radius: 20px; padding: 5px 12px; font-size: 12px;
    color: var(--text-muted); cursor: pointer; display: none;
    align-items: center; gap: 5px; box-shadow: 0 2px 10px rgba(0,0,0,0.1);
    transition: all 0.15s; font-family: inherit;
  }
  .scroll-pause-btn.visible { display: flex; }
  .scroll-pause-btn:hover { border-color: var(--primary-border); color: var(--primary); }
  /* ── MSG COUNTER ── */
  .msg-counter {
    font-size: 10.5px; color: var(--text-dim); margin-left: 4px;
    padding: 1px 6px; background: var(--primary-dim); border-radius: 10px;
    font-weight: 500;
  }
  /* ── URL PREVIEW CHIP ── */
  .url-preview-chip {
    display: none; align-items: center; gap: 8px;
    background: var(--bg-card,#fafafa); border: 1px solid var(--primary-border);
    border-radius: 8px; padding: 5px 10px; margin-bottom: 6px; font-size: 12px;
  }
  .url-preview-chip.open { display: flex; }
  .url-preview-chip a { color: var(--primary); text-decoration: none; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; flex: 1; }
  .url-preview-chip a:hover { text-decoration: underline; }
  /* ── SIDEBAR DRAG HANDLE ── */
  .sidebar-drag-handle {
    position: absolute; top: 0; right: -3px; width: 6px; height: 100%;
    cursor: col-resize; z-index: 50; background: transparent; transition: background 0.15s;
  }
  .sidebar-drag-handle:hover { background: var(--primary-border); }
  *, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
  html, body { height: 100%; overflow: hidden; }
  body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', system-ui, sans-serif;
    background: var(--bg);
    color: var(--text);
    display: flex;
    font-size: 14px;
    line-height: 1.5;
  }

  /* ── SCROLLBAR ── */
  ::-webkit-scrollbar { width: 4px; height: 4px; }
  ::-webkit-scrollbar-track { background: transparent; }
  ::-webkit-scrollbar-thumb { background: var(--primary-border); border-radius: 4px; }

  /* ── SIDEBAR ── */
  #sidebar {
    width: 240px;
    min-width: 240px;
    background: var(--primary-d);
    border-right: 1px solid rgba(255,255,255,0.07);
    box-shadow: none;
    display: flex;
    flex-direction: column;
    height: 100vh;
    margin: 0;
    border-radius: 0;
    overflow: hidden;
    position: relative;
    z-index: 100;
    transition: width var(--transition), min-width var(--transition);
  }
  .sidebar-brand {
    padding: 20px 16px 16px;
    border-bottom: 1px solid rgba(255,255,255,0.07);
    flex-shrink: 0;
    display: flex;
    flex-direction: column;
    align-items: center;
    gap: 5px;
  }
  .sidebar-logo {
    width: 110px;
    height: auto;
    display: block;
    filter: brightness(0) invert(1);
    opacity: 0.90;
  }
  .brand-sub {
    font-size: 8.5px;
    color: rgba(245,236,236,0.35);
    letter-spacing: 0.12em;
    text-transform: uppercase;
    font-weight: 600;
    text-align: center;
  }
  .sidebar-scroll {
    flex: 1;
    overflow-y: auto;
    overflow-x: hidden;
    padding: 8px 0 16px;
  }
  .sidebar-section-label {
    font-size: 9px;
    font-weight: 700;
    letter-spacing: 0.14em;
    text-transform: uppercase;
    color: rgba(245,236,236,0.28);
    padding: 14px 16px 4px;
  }
  .nav-item {
    display: flex;
    align-items: center;
    gap: 10px;
    padding: 9px 16px;
    margin: 1px 0;
    cursor: pointer;
    border-radius: 0;
    border-left: 3px solid transparent;
    transition: background 0.15s ease, color 0.15s ease, border-color 0.15s ease;
    color: rgba(245,236,236,0.55);
    font-size: 13px;
    font-weight: 500;
    position: relative;
    white-space: nowrap;
    overflow: hidden;
    text-overflow: ellipsis;
  }
  .nav-item:hover {
    background: rgba(255,255,255,0.07);
    color: rgba(245,236,236,0.9);
    border-left-color: rgba(255,255,255,0.25);
  }
  .nav-item.active {
    background: rgba(255,255,255,0.10);
    color: #ffffff;
    font-weight: 600;
    border-left-color: #ffffff;
    box-shadow: none;
  }
  .nav-icon { font-size: 14px; flex-shrink: 0; width: 18px; text-align: center; line-height: 1; }
  .nav-label { overflow: hidden; text-overflow: ellipsis; }

  .sidebar-status {
    padding: 12px 16px;
    border-top: 1px solid rgba(255,255,255,0.08);
    display: flex;
    align-items: center;
    gap: 8px;
    flex-shrink: 0;
  }
  .status-dot {
    width: 8px; height: 8px;
    border-radius: 50%;
    background: #4ADE80;
    flex-shrink: 0;
    box-shadow: 0 0 0 0 rgba(74,222,128,0.6);
    animation: pulse-green 2s infinite;
  }
  @keyframes pulse-green {
    0% { box-shadow: 0 0 0 0 rgba(74,222,128,0.6); }
    70% { box-shadow: 0 0 0 8px rgba(74,222,128,0); }
    100% { box-shadow: 0 0 0 0 rgba(74,222,128,0); }
  }
  .status-text { font-size: 11px; color: rgba(245,236,236,0.45); }

  /* ── MAIN ── */
  #main {
    flex: 1;
    display: flex;
    flex-direction: column;
    height: 100vh;
    overflow: hidden;
  }

  /* ── TOPBAR ── */
  #topbar {
    display: flex;
    align-items: center;
    padding: 0 28px;
    height: 64px;
    border-bottom: 1px solid var(--card-border);
    background: #ffffff;
    box-shadow: none;
    flex-shrink: 0;
    gap: 12px;
    margin: 0;
    border-radius: 0;
  }
  .topbar-title {
    font-size: 16px;
    font-weight: 700;
    color: var(--text);
    flex: 1;
  }
  .topbar-date {
    font-size: 11px;
    color: var(--text-dim);
    display: flex;
    align-items: center;
    gap: 6px;
  }
  .topbar-actions { display: flex; align-items: center; gap: 10px; }
  .btn {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    padding: 8px 16px;
    border-radius: var(--radius-sm);
    font-size: 13px;
    font-weight: 600;
    cursor: pointer;
    transition: all var(--transition);
    border: none;
    font-family: inherit;
    white-space: nowrap;
  }
  .btn-ghost {
    background: transparent;
    border: 1px solid rgba(255,255,255,0.18);
    color: rgba(245,236,236,0.75);
  }
  .btn-ghost:hover {
    border-color: rgba(255,255,255,0.38);
    color: #f5ecec;
    background: rgba(255,255,255,0.1);
  }
  .btn-gold {
    background: linear-gradient(135deg, var(--primary-l), #b52240);
    color: #fff;
    font-weight: 700;
  }
  .btn-gold:hover {
    background: linear-gradient(135deg, #b52240, #cc2a4d);
    transform: translateY(-1px);
    box-shadow: 0 4px 20px rgba(122,16,40,0.45);
  }
  /* Ghost buttons inside light content areas */
  .card .btn-ghost, .view .btn-ghost, .modal .btn-ghost {
    border-color: var(--card-border);
    color: var(--text-muted);
  }
  .card .btn-ghost:hover, .view .btn-ghost:hover, .modal .btn-ghost:hover {
    border-color: var(--primary-border);
    color: var(--primary);
    background: var(--primary-dim);
  }
  .btn:disabled { opacity: 0.5; cursor: not-allowed; transform: none !important; }

  /* ── TOPBAR ICON BUTTONS ── */
  .tb-icon-btn {
    position: relative;
    background: transparent;
    border: 1px solid var(--card-border);
    border-radius: 8px;
    width: 34px; height: 34px;
    display: flex; align-items: center; justify-content: center;
    font-size: 13px;
    cursor: pointer;
    transition: all 0.15s ease;
    color: var(--text-muted);
    flex-shrink: 0;
    font-family: inherit;
  }
  .tb-icon-btn:hover {
    border-color: var(--primary-border);
    color: var(--primary);
    background: var(--primary-dim);
    transform: none;
  }
  .tb-icon-btn:active { transform: none; }

  /* ── TOPBAR AUTOMATE DROPDOWN ── */
  .tb-dropdown { position: relative; }
  .tb-automate-btn {
    display: flex; align-items: center; gap: 6px;
    padding: 7px 13px; border-radius: 8px; font-size: 12.5px; font-weight: 600;
    background: var(--primary-dim); border: 1px solid var(--primary-border);
    color: var(--primary); cursor: pointer; font-family: inherit;
    transition: all 0.15s; white-space: nowrap;
  }
  .tb-automate-btn:hover { background: rgba(122,16,40,0.12); box-shadow: 0 2px 8px rgba(122,16,40,0.12); }
  .tb-dd-menu {
    display: none; position: absolute; top: calc(100% + 6px); right: 0;
    min-width: 210px; background: #fff; border: 1px solid var(--card-border);
    border-radius: 10px; box-shadow: 0 8px 28px rgba(0,0,0,0.1);
    z-index: 500; padding: 5px 0; overflow: hidden;
  }
  .tb-dropdown.open .tb-dd-menu { display: block; animation: ddIn 0.15s ease; }
  @keyframes ddIn { from{opacity:0;transform:translateY(-4px)} to{opacity:1;transform:none} }
  .tb-dd-item {
    display: flex; align-items: center; gap: 10px;
    padding: 9px 14px; font-size: 12.5px; color: var(--text); cursor: pointer;
    transition: background 0.1s;
  }
  .tb-dd-item:hover { background: var(--primary-dim); color: var(--primary); }
  .tb-dd-item i { width: 16px; text-align: center; flex-shrink: 0; }
  .tb-dd-divider { height: 1px; background: var(--card-border); margin: 4px 0; }

  /* ── NOTIFICATION BELL ── */
  .notif-bell { position: relative; }
  .notif-bell:hover { border-color: var(--primary-border); background: var(--primary-dim); }
  .notif-bell.has-unread { border-color: rgba(232,96,60,0.4); animation: bell-shake 3s ease infinite; }
  @keyframes bell-shake {
    0%,90%,100% { transform: rotate(0); }
    92% { transform: rotate(-12deg); }
    94% { transform: rotate(12deg); }
    96% { transform: rotate(-8deg); }
    98% { transform: rotate(8deg); }
  }
  .notif-badge {
    position: absolute;
    top: -5px; right: -5px;
    background: #E8603C;
    color: #fff;
    font-size: 9px; font-weight: 700;
    min-width: 17px; height: 17px;
    border-radius: 9px;
    display: flex; align-items: center; justify-content: center;
    padding: 0 4px;
    border: 2px solid var(--navy-d);
    line-height: 1;
  }

  /* ── NOTIFICATION OVERLAY ── */
  #notif-overlay {
    display: none;
    position: fixed; inset: 0;
    background: rgba(0,0,0,0.45);
    z-index: 998;
    backdrop-filter: blur(2px);
  }
  #notif-overlay.open { display: block; }

  /* ── NOTIFICATION DRAWER ── */
  #notif-drawer {
    position: fixed;
    top: 0; right: 0;
    width: 400px; height: 100vh;
    background: var(--primary-d);
    border-left: 1px solid rgba(255,255,255,0.06);
    z-index: 999;
    display: flex; flex-direction: column;
    transform: translateX(100%);
    transition: transform 0.28s cubic-bezier(0.4,0,0.2,1);
    box-shadow: -8px 0 40px rgba(78,9,23,0.5);
  }
  #notif-drawer.open { transform: translateX(0); }

  .nd-header {
    display: flex; align-items: flex-start; justify-content: space-between;
    padding: 22px 20px 16px;
    border-bottom: 1px solid rgba(255,255,255,0.08);
    flex-shrink: 0;
  }
  .nd-title { font-size: 16px; font-weight: 700; color: #f5ecec; }
  .nd-subtitle { font-size: 11px; color: rgba(245,236,236,0.5); margin-top: 2px; }
  .nd-action-btn {
    font-size: 11px; color: rgba(245,236,236,0.5); background: none;
    border: 1px solid rgba(255,255,255,0.12); border-radius: 6px;
    padding: 4px 10px; cursor: pointer; font-family: inherit;
    transition: all var(--transition);
  }
  .nd-action-btn:hover { color: #f5ecec; border-color: rgba(255,255,255,0.3); }
  .nd-action-btn.nd-clear:hover { color: #fca5a5; border-color: rgba(252,165,165,0.4); }
  .nd-close {
    background: none; border: none; color: rgba(245,236,236,0.5);
    font-size: 16px; cursor: pointer; padding: 4px 8px;
    border-radius: 6px; transition: all var(--transition); font-family: inherit;
  }
  .nd-close:hover { color: #f5ecec; background: rgba(255,255,255,0.07); }

  .nd-body { flex: 1; overflow-y: auto; padding: 12px 14px; }
  .nd-empty { text-align: center; color: rgba(245,236,236,0.4); font-size: 13px; padding: 48px 20px; }

  /* ── NOTIFICATION CARD ── */
  .nd-card {
    background: rgba(255,255,255,0.05);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 10px;
    padding: 14px 14px 12px;
    margin-bottom: 10px;
    display: flex; gap: 12px;
    position: relative;
    transition: background var(--transition);
    border-left-width: 3px;
  }
  .nd-card:hover { background: rgba(255,255,255,0.09); }
  .nd-card.unread { background: rgba(255,255,255,0.07); }
  .nd-card.read { opacity: 0.55; }
  .nd-card.level-critical { border-left-color: #E8603C; }
  .nd-card.level-warning  { border-left-color: #C9A84C; }
  .nd-card.level-success  { border-left-color: #4ADE80; }
  .nd-card.level-info     { border-left-color: #60A5FA; }

  .nd-card-emoji { font-size: 22px; flex-shrink: 0; line-height: 1.3; }
  .nd-card-body { flex: 1; min-width: 0; }
  .nd-card-agent {
    font-size: 10px; font-weight: 700; letter-spacing: 0.1em;
    text-transform: uppercase; margin-bottom: 3px;
  }
  .nd-card.level-critical .nd-card-agent { color: #fca5a5; }
  .nd-card.level-warning  .nd-card-agent { color: #fcd34d; }
  .nd-card.level-success  .nd-card-agent { color: #86efac; }
  .nd-card.level-info     .nd-card-agent { color: #93c5fd; }
  .nd-card-title { font-size: 13px; font-weight: 600; color: #f5ecec; margin-bottom: 4px; line-height: 1.3; }
  .nd-card-msg { font-size: 12px; color: rgba(245,236,236,0.6); line-height: 1.5; margin-bottom: 8px; }
  .nd-card-footer { display: flex; align-items: center; justify-content: space-between; gap: 8px; }
  .nd-card-time { font-size: 10px; color: rgba(245,236,236,0.35); }
  .nd-card-actions { display: flex; gap: 6px; }
  .nd-btn-read {
    font-size: 10px; color: rgba(245,236,236,0.35); background: none;
    border: 1px solid transparent; border-radius: 4px;
    padding: 2px 8px; cursor: pointer; font-family: inherit;
    transition: all var(--transition);
  }
  .nd-btn-read:hover { color: rgba(245,236,236,0.7); border-color: rgba(255,255,255,0.15); }
  .nd-unread-dot {
    position: absolute; top: 12px; right: 12px;
    width: 7px; height: 7px; border-radius: 50%;
    background: #fca5a5;
  }
  .nd-card.level-warning  .nd-unread-dot { background: #fcd34d; }
  .nd-card.level-success  .nd-unread-dot { background: #86efac; }
  .nd-card.level-info     .nd-unread-dot { background: #93c5fd; }
  .nd-report-btn {
    font-size: 10px; color: #e8707e; background: none;
    border: 1px solid rgba(232,112,126,0.4); border-radius: 4px;
    padding: 2px 8px; cursor: pointer; font-family: inherit;
    text-decoration: none; display: inline-block;
    transition: all var(--transition);
  }
  .nd-report-btn:hover { background: rgba(232,112,126,0.1); }

  /* ── CONTENT AREA ── */
  #content {
    flex: 1;
    min-height: 0;
    overflow-y: auto;
    position: relative;
    background: #ffffff;
  }

  /* ── VIEWS ── */
  .view { display: none; padding: 24px 28px 28px; }
  .view.active { display: block; animation: viewIn 0.3s cubic-bezier(0.16,1,0.3,1) both; }
  @keyframes viewIn { from { opacity: 0; transform: translateY(10px); } to { opacity: 1; transform: translateY(0); } }
  @keyframes fadeInUp { from { opacity: 0; transform: translateY(14px); } to { opacity: 1; transform: translateY(0); } }

  /* ── CARDS ── */
  .card {
    background: #ffffff;
    border: 1px solid var(--card-border);
    border-radius: 10px;
    box-shadow: none;
    padding: 20px;
    transition: box-shadow 0.18s ease;
  }
  .card:hover { box-shadow: 0 2px 12px rgba(122,16,40,0.07); }
  .card-title {
    font-size: 11px;
    font-weight: 700;
    color: var(--text-dim);
    text-transform: uppercase;
    letter-spacing: 0.1em;
    margin-bottom: 16px;
    display: flex;
    align-items: center;
    gap: 8px;
  }

  /* ── STAT CARDS ── */
  .stats-grid {
    display: grid;
    grid-template-columns: repeat(5, 1fr);
    gap: 12px;
    margin-bottom: 20px;
  }
  .stat-card {
    background: #ffffff;
    border: 1px solid var(--card-border);
    border-radius: 10px;
    padding: 14px 16px;
    transition: transform 0.2s ease, box-shadow 0.2s ease;
    position: relative;
    overflow: hidden;
    box-shadow: none;
    animation: fadeInUp 0.4s cubic-bezier(0.16,1,0.3,1) both;
  }
  .stat-card:nth-child(1) { animation-delay: 0.04s; }
  .stat-card:nth-child(2) { animation-delay: 0.09s; }
  .stat-card:nth-child(3) { animation-delay: 0.14s; }
  .stat-card:nth-child(4) { animation-delay: 0.19s; }
  .stat-card:nth-child(5) { animation-delay: 0.24s; }
  .stat-card::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 3px;
    background: linear-gradient(90deg, var(--primary), var(--primary-l));
  }
  .stat-card:hover {
    transform: translateY(-2px);
    box-shadow: 0 4px 16px rgba(122,16,40,0.09);
  }
  .stat-icon { font-size: 18px; margin-bottom: 8px; display: block; line-height: 1.2; color: var(--primary); }
  .stat-number { font-size: 24px; font-weight: 800; color: var(--text); line-height: 1; }
  .stat-label { font-size: 11px; color: var(--text-muted); margin-top: 3px; font-weight: 500; }
  .stat-sub { font-size: 10px; color: var(--primary-l); margin-top: 4px; font-weight: 600; }
  th.sortable { cursor: pointer; user-select: none; white-space: nowrap; }
  th.sortable:hover { color: var(--primary); }
  th.sortable .sort-icon { margin-left: 4px; font-size: 10px; opacity: 0.45; }
  th.sortable.asc .sort-icon, th.sortable.desc .sort-icon { opacity: 1; color: var(--primary); }
  .table-pagination {
    display: flex; align-items: center; justify-content: space-between;
    padding: 10px 16px 4px; font-size: 12px; color: var(--text-muted);
  }
  .pg-info { font-size: 11px; }
  .pg-btns { display: flex; align-items: center; gap: 4px; }
  .pg-btn {
    border: 1px solid var(--card-border); background: #fff; border-radius: 6px;
    width: 28px; height: 28px; cursor: pointer; font-size: 11px;
    display: flex; align-items: center; justify-content: center;
    color: var(--text-muted); transition: all 0.12s;
  }
  .pg-btn:hover:not(:disabled) { border-color: var(--primary); color: var(--primary); }
  .pg-btn:disabled { opacity: 0.35; cursor: default; }
  .pg-btn.active { background: var(--primary); color: #fff; border-color: var(--primary); }

  /* ── 2-COL GRID ── */
  .two-col { display: grid; grid-template-columns: 1fr 1fr; gap: 16px; margin-bottom: 20px; }

  /* ── QUICK ACTIONS ── */
  .actions-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 10px; }
  .action-btn {
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    gap: 8px;
    padding: 18px 12px;
    background: rgba(201,168,76,0.06);
    border: 1px solid var(--gold-border);
    border-radius: var(--radius-sm);
    cursor: pointer;
    transition: all var(--transition);
    color: var(--text-muted);
    font-size: 12px;
    font-weight: 600;
    text-align: center;
    font-family: inherit;
  }
  .action-btn:hover {
    background: var(--primary-dim);
    color: var(--primary-l);
    border-color: var(--primary-border);
    transform: scale(1.02);
  }
  .action-btn .action-icon { font-size: 20px; line-height: 1; }

  /* ── PROGRESS BARS ── */
  .prog-item { margin-bottom: 14px; }
  .prog-header { display: flex; justify-content: space-between; font-size: 12px; margin-bottom: 5px; }
  .prog-label { color: var(--text-muted); font-weight: 500; }
  .prog-val { color: var(--primary); font-weight: 700; }
  .prog-bar { height: 6px; background: rgba(122,16,40,0.08); border-radius: 4px; overflow: hidden; }
  .prog-fill {
    height: 100%;
    background: linear-gradient(90deg, var(--primary), var(--primary-l));
    border-radius: 4px;
    width: 0;
    transition: width 1s cubic-bezier(0.4,0,0.2,1);
  }

  /* ── UPCOMING POSTS ── */
  .post-item {
    display: flex;
    align-items: center;
    gap: 12px;
    padding: 10px 14px;
    border-radius: var(--radius-sm);
    transition: background var(--transition);
    margin-bottom: 4px;
  }
  .post-item:hover { background: var(--primary-dim); }
  .post-item.today { background: rgba(122,16,40,0.07); border: 1px solid var(--primary-border); }
  .post-date-badge {
    min-width: 38px;
    text-align: center;
    font-size: 11px;
    font-weight: 700;
    color: var(--text-muted);
  }
  .post-item.today .post-date-badge { color: var(--primary); }
  .post-dot { width: 8px; height: 8px; border-radius: 50%; background: rgba(122,16,40,0.2); flex-shrink: 0; }
  .post-item.today .post-dot { background: var(--primary); box-shadow: 0 0 6px var(--primary-l); }
  .post-info { flex: 1; }
  .post-filename { font-size: 12px; color: var(--text); font-weight: 500; }
  .post-day { font-size: 11px; color: var(--text-dim); }
  .badge {
    font-size: 10px;
    font-weight: 700;
    padding: 3px 8px;
    border-radius: 20px;
    letter-spacing: 0.05em;
  }
  .badge-today { background: rgba(122,16,40,0.08); color: var(--primary); border: 1px solid var(--primary-border); }
  .badge-posted { background: rgba(74,222,128,0.1); color: #4ADE80; border: 1px solid rgba(74,222,128,0.3); }
  .badge-scheduled { background: rgba(99,102,241,0.1); color: #818CF8; border: 1px solid rgba(99,102,241,0.3); }
  .badge-active { background: rgba(74,222,128,0.1); color: #4ADE80; border: 1px solid rgba(74,222,128,0.3); }
  .badge-inactive { background: rgba(255,255,255,0.06); color: var(--text-dim); border: 1px solid rgba(255,255,255,0.1); }

  /* ── TABLE ── */
  table { width: 100%; border-collapse: collapse; }
  thead th {
    text-align: left;
    font-size: 10px;
    font-weight: 700;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    color: var(--text-dim);
    padding: 10px 14px;
    border-bottom: 1px solid var(--card-border);
    background: rgba(122,16,40,0.03);
  }
  tbody tr {
    border-bottom: 1px solid rgba(122,16,40,0.05);
    transition: background var(--transition);
  }
  tbody tr:hover { background: var(--primary-dim); }
  tbody td { padding: 11px 14px; font-size: 13px; color: var(--text-muted); vertical-align: middle; }
  tbody td:first-child { color: var(--text); font-weight: 500; }

  /* ── LEADS VIEW ── */
  .view-header {
    display: flex;
    align-items: center;
    justify-content: space-between;
    margin-bottom: 20px;
  }
  .view-header h2 { font-size: 20px; font-weight: 700; }
  #add-form {
    display: none;
    background: #ffffff;
    border: 1px solid var(--card-border);
    border-radius: var(--radius);
    padding: 20px;
    margin-bottom: 20px;
    animation: slideDown 0.2s ease;
    box-shadow: 0 2px 8px rgba(122,16,40,0.08);
  }
  #add-form.open { display: block; }
  @keyframes slideDown { from { opacity: 0; transform: translateY(-8px); } to { opacity: 1; transform: translateY(0); } }
  .form-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 14px; margin-bottom: 16px; }
  .form-group { display: flex; flex-direction: column; gap: 6px; }
  .form-label { font-size: 11px; font-weight: 600; color: var(--text-muted); text-transform: uppercase; letter-spacing: 0.06em; }
  .form-input {
    background: #fdfafa;
    border: 1px solid var(--card-border);
    border-radius: var(--radius-sm);
    padding: 9px 12px;
    color: var(--text);
    font-size: 13px;
    font-family: inherit;
    transition: border-color var(--transition);
    outline: none;
  }
  .form-input:focus { border-color: var(--primary-border); background: #fff; }
  .form-actions { display: flex; gap: 10px; }
  .btn-sm { padding: 7px 14px; font-size: 12px; }
  .btn-outline {
    background: transparent;
    border: 1px solid var(--card-border);
    color: var(--text-muted);
  }
  .btn-outline:hover { border-color: var(--primary-border); color: var(--primary); }

  /* ── CHAT VIEW ── */
  #view-dashboard { display: none; flex-direction: row; height: 100%; padding: 0; overflow: hidden; }
  #view-dashboard.active { display: flex; }

  /* Chat left sidebar */
  .chat-sidebar {
    width: 240px; min-width: 240px;
    border-right: 1px solid var(--card-border);
    display: flex; flex-direction: column; overflow: hidden;
    background: #f9f9f9;
  }
  .chat-sidebar-top { padding: 14px 12px 12px; border-bottom: 1px solid var(--card-border); flex-shrink: 0; }
  .new-chat-btn {
    width: 100%; padding: 10px 14px; border-radius: 10px;
    background: linear-gradient(135deg, var(--primary), var(--primary-l));
    color: #fff; border: none; font-size: 13px; font-weight: 600;
    cursor: pointer; display: flex; align-items: center; gap: 8px;
    transition: all 0.18s; font-family: inherit; justify-content: center;
    box-shadow: 0 2px 8px rgba(122,16,40,0.2);
  }
  .new-chat-btn:hover { opacity: 0.88; transform: translateY(-1px); box-shadow: 0 4px 14px rgba(122,16,40,0.32); }
  .custom-dd { position: relative; width: 100%; margin-top: 9px; }
  .custom-dd-trigger {
    display: flex; align-items: center; justify-content: space-between;
    background: #ffffff; border: 1px solid var(--card-border);
    border-radius: var(--radius-sm); color: var(--text);
    padding: 8px 12px; font-size: 12px; font-family: inherit;
    cursor: pointer; user-select: none;
    transition: border-color var(--transition), box-shadow var(--transition);
  }
  .custom-dd-trigger:hover { border-color: var(--primary-border); }
  .custom-dd.open .custom-dd-trigger { border-color: var(--primary); box-shadow: 0 0 0 3px rgba(122,16,40,0.06); }
  .custom-dd-arrow { font-size: 9px; color: var(--text-muted); transition: transform 0.15s; margin-left: 6px; flex-shrink: 0; }
  .custom-dd.open .custom-dd-arrow { transform: rotate(180deg); }
  .custom-dd-menu {
    display: none; position: absolute; top: calc(100% + 4px); left: 0; right: 0;
    background: #ffffff; border: 1px solid var(--card-border);
    border-radius: var(--radius-sm); box-shadow: 0 8px 24px rgba(0,0,0,0.09);
    z-index: 400; max-height: 280px; overflow-y: auto; padding: 4px 0;
  }
  .custom-dd.open .custom-dd-menu { display: block; }
  .custom-dd-opt {
    padding: 8px 12px; font-size: 12px; color: var(--text);
    cursor: pointer; transition: background 0.1s; white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
  }
  .custom-dd-opt:hover { background: var(--primary-dim); color: var(--primary); }
  .custom-dd-opt.selected { color: var(--primary); font-weight: 600; background: var(--primary-dim); }
  .custom-dd-opt.placeholder-opt { color: var(--text-dim); }
  .custom-dd-divider { font-size:9px; font-weight:700; text-transform:uppercase; letter-spacing:.1em; color:var(--text-dim); padding:10px 14px 4px; pointer-events:none; border-top:1px solid var(--card-border); margin-top:4px; }
  .agent-general-opt { font-weight:700; color:var(--primary) !important; }
  .chat-sessions-list { flex: 1; overflow-y: auto; padding: 6px; }
  .chat-session-group {
    font-size: 9px; font-weight: 700; text-transform: uppercase;
    letter-spacing: .1em; color: var(--text-dim); padding: 8px 8px 3px;
  }
  .chat-session-item {
    padding: 9px 10px; border-radius: 9px; cursor: pointer;
    transition: all 0.15s; position: relative; display: flex;
    align-items: flex-start; gap: 8px; margin-bottom: 1px;
  }
  .chat-session-item:hover { background: var(--primary-dim); }
  .chat-session-item.active { background: var(--primary-dim); }
  .chat-session-agent-icon { width: 22px; height: 22px; flex-shrink: 0; margin-top: 1px; overflow: hidden; }
  .chat-session-agent-icon img { width: 22px; height: 22px; object-fit: contain; }
  .chat-session-info { flex: 1; min-width: 0; }
  .chat-session-title {
    font-size: 12px; font-weight: 600; color: var(--text);
    white-space: nowrap; overflow: hidden; text-overflow: ellipsis; line-height: 1.4;
  }
  .chat-session-item.active .chat-session-title { color: var(--primary); }
  .chat-session-meta { font-size: 10px; color: var(--text-dim); margin-top: 1px; }
  .chat-session-del {
    opacity: 0; font-size: 10px; color: var(--text-dim);
    background: none; border: none; cursor: pointer; padding: 2px 5px;
    transition: all 0.15s; flex-shrink: 0; margin-top: 1px;
  }
  .chat-session-item:hover .chat-session-del { opacity: 1; }
  .chat-session-del:hover { color: #F87171; }

  /* Chat main area */
  .chat-main { flex: 1; display: flex; flex-direction: column; overflow: hidden; background: #ffffff; }
  .chat-header {
    display: flex; align-items: center; gap: 14px;
    padding: 12px 24px; border-bottom: 1px solid var(--card-border);
    background: #ffffff; flex-shrink: 0;
    box-shadow: none;
  }
  .agent-avatar {
    width: 42px; height: 42px; border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    flex-shrink: 0; overflow: hidden; background: #f4ecec;
  }
  .avatar-logo { width: 36px; height: 36px; object-fit: contain; display: block; }
  .chat-agent-name { font-size: 15px; font-weight: 700; color: var(--text); }
  .chat-agent-role { font-size: 12px; color: var(--text-muted); margin-top: 1px; }
  .chat-header-actions { margin-left: auto; display: flex; align-items: center; gap: 10px; }
  .chat-mode-bar { display:flex; align-items:center; gap:8px; flex-wrap:wrap; }
  .model-dd { position: relative; margin-top: 0; width: auto; }
  .model-dd .custom-dd-trigger { padding: 5px 10px; font-size: 11px; font-weight: 600; min-width: 110px; }
  .model-dd .custom-dd-menu { left: auto; right: 0; min-width: 130px; max-height: 160px; }
  .model-dd .custom-dd-opt { font-size: 11px; padding: 7px 12px; }
  .chat-mode-pills { display:flex; gap:4px; }
  .mode-pill {
    padding:4px 10px; border-radius:6px; font-size:11px; font-weight:600;
    border:1px solid var(--card-border); background:#fff; color:var(--text-muted);
    cursor:pointer; transition:all .15s; font-family:inherit;
  }
  .mode-pill.active { background:var(--primary-dim); color:var(--primary); border-color:var(--primary-border); }
  .mode-pill:hover { background:var(--primary-dim); color:var(--primary); }
  .chat-messages {
    flex: 1; overflow-y: auto; padding: 28px 40px;
    display: flex; flex-direction: column; gap: 24px; background: #ffffff;
    scroll-behavior: smooth;
  }
  .welcome-box {
    background: transparent; border: none; border-radius: 0;
    padding: 32px 28px 40px; margin: auto;
    max-width: 640px; width: 100%; text-align: center;
    box-shadow: none; display: flex; flex-direction: column; align-items: center;
  }
  .welcome-hero { display: flex; flex-direction: column; align-items: center; margin-bottom: 20px; }
  .welcome-logo-ring {
    width: 64px; height: 64px; border-radius: 50%;
    background: linear-gradient(135deg, rgba(122,16,40,0.1), rgba(155,26,53,0.06));
    border: 1.5px solid rgba(122,16,40,0.12);
    display: flex; align-items: center; justify-content: center;
    margin-bottom: 14px; box-shadow: 0 4px 18px rgba(122,16,40,0.09);
    position: relative;
  }
  .welcome-logo-ring::before {
    content: ''; position: absolute; inset: -5px; border-radius: 50%;
    border: 1px solid rgba(122,16,40,0.06);
  }
  .welcome-logo { width: 42px; height: auto; object-fit: contain; display: block; }
  .welcome-greeting {
    font-size: 13px; font-weight: 500; color: var(--text-muted);
    margin-bottom: 6px; letter-spacing: 0.01em;
  }
  .welcome-box .welcome-title {
    font-size: 22px; font-weight: 800; color: var(--text);
    margin-bottom: 8px; line-height: 1.2; letter-spacing: -0.02em;
  }
  .welcome-box .welcome-sub {
    font-size: 12.5px; color: var(--text-dim); line-height: 1.65;
    max-width: 400px;
  }
  .welcome-cards {
    display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 8px;
    width: 100%; margin-top: 4px;
  }
  .wcard {
    display: flex; align-items: center; gap: 9px;
    padding: 9px 12px; border-radius: 10px; text-align: left;
    background: #fff; border: 1px solid var(--card-border);
    cursor: pointer; transition: all 0.15s; font-family: inherit;
    box-shadow: 0 1px 3px rgba(0,0,0,0.04);
  }
  .wcard:hover {
    border-color: var(--primary-border); background: var(--primary-dim);
    transform: translateY(-1px); box-shadow: 0 3px 10px rgba(122,16,40,0.1);
  }
  .wcard-icon {
    width: 30px; height: 30px; border-radius: 8px; flex-shrink: 0;
    background: var(--primary-dim); color: var(--primary);
    display: flex; align-items: center; justify-content: center; font-size: 13px;
  }
  .wcard:hover .wcard-icon { background: rgba(122,16,40,0.14); }
  .wcard-text { display: flex; flex-direction: column; min-width: 0; }
  .wcard-title { font-size: 12px; font-weight: 700; color: var(--text); line-height: 1.3; }
  .wcard-desc { font-size: 10.5px; color: var(--text-dim); margin-top: 1px; line-height: 1.4; white-space: normal; }
  /* ── MESSAGES — Claude-style ── */
  .msg { display: flex; gap: 12px; width: 100%; animation: msgIn 0.22s ease; }
  @keyframes msgIn { from { opacity: 0; transform: translateY(6px); } to { opacity: 1; transform: translateY(0); } }
  .msg.user { align-self: flex-end; flex-direction: row-reverse; max-width: 76%; }
  .msg.ai  { align-self: flex-start; max-width: 92%; }
  .msg-avatar {
    width: 30px; height: 30px;
    display: flex; align-items: center; justify-content: center;
    font-size: 13px; flex-shrink: 0; margin-top: 2px;
    overflow: hidden; border-radius: 50%;
  }
  .msg-avatar img { width: 100%; height: 100%; object-fit: contain; border-radius: 50%; background: #f4ecec; padding: 3px; }
  .msg-avatar.user-av {
    background: linear-gradient(135deg, var(--primary), var(--primary-l));
    color: #fff; font-size: 11px; font-weight: 700; letter-spacing: 0.04em;
  }
  /* AI bubble — clean white, no background (Claude-style) */
  .msg-bubble {
    background: transparent; border: none;
    border-radius: 0; padding: 2px 0;
    font-size: 14px; line-height: 1.75; color: var(--text);
    min-width: 0; overflow-wrap: break-word; word-break: break-word;
  }
  /* User bubble — pill with brand color */
  .msg.user .msg-bubble {
    background: linear-gradient(135deg, var(--primary) 0%, #b52240 100%);
    border-radius: 18px 4px 18px 18px;
    color: #fff; font-weight: 500; padding: 11px 16px;
    font-size: 13.5px; line-height: 1.65;
    box-shadow: 0 2px 12px rgba(122,16,40,0.2);
  }
  /* AI message name label */
  .msg-name {
    font-size: 11px; font-weight: 700; color: var(--primary);
    text-transform: uppercase; letter-spacing: 0.08em; margin-bottom: 5px;
    display: block;
  }
  .msg-bubble p { margin: 0 0 10px; }
  .msg-bubble p:last-child { margin-bottom: 0; }
  .msg-bubble h1 { font-size: 18px; color: var(--text); margin: 16px 0 6px; font-weight: 800; }
  .msg-bubble h2 { font-size: 15px; color: var(--primary); margin: 14px 0 5px; font-weight: 700; border-bottom: 1px solid var(--card-border); padding-bottom: 4px; }
  .msg-bubble h3 { font-size: 13.5px; color: var(--text); margin: 12px 0 4px; font-weight: 700; }
  .msg.user .msg-bubble h2 { color: rgba(255,255,255,0.85); border-color: rgba(255,255,255,0.2); }
  .msg.user .msg-bubble h3 { color: rgba(255,255,255,0.95); }
  .msg-bubble strong { color: var(--primary); font-weight: 700; }
  .msg.user .msg-bubble strong { color: rgba(255,255,255,0.98); }
  .msg-bubble em { color: var(--text-muted); font-style: italic; }
  .msg.user .msg-bubble em { color: rgba(255,255,255,0.78); }
  /* Inline code */
  .msg-bubble code {
    background: rgba(122,16,40,0.07); padding: 2px 7px; border-radius: 5px;
    font-size: 12.5px; font-family: 'Courier New', monospace; color: var(--primary);
    border: 1px solid rgba(122,16,40,0.1);
  }
  .msg.user .msg-bubble code { background: rgba(0,0,0,0.2); color: #fff; border-color: rgba(255,255,255,0.2); }
  /* Code blocks */
  .msg-bubble pre {
    background: #1a1a2e; border-radius: 10px; margin: 12px 0;
    overflow: hidden; box-shadow: 0 2px 12px rgba(0,0,0,0.15);
  }
  .code-block-header {
    display: flex; align-items: center; justify-content: space-between;
    padding: 8px 14px; background: rgba(255,255,255,0.05);
    border-bottom: 1px solid rgba(255,255,255,0.08);
  }
  .code-block-lang { font-size: 11px; font-weight: 600; color: rgba(255,255,255,0.5); text-transform: uppercase; letter-spacing: 0.06em; }
  .code-copy-btn {
    font-size: 10px; color: rgba(255,255,255,0.45); background: none;
    border: 1px solid rgba(255,255,255,0.15); border-radius: 5px;
    padding: 3px 9px; cursor: pointer; font-family: inherit;
    transition: all 0.15s; display: flex; align-items: center; gap: 5px;
  }
  .code-copy-btn:hover { color: #fff; border-color: rgba(255,255,255,0.4); background: rgba(255,255,255,0.08); }
  .code-copy-btn.copied { color: #4ADE80; border-color: rgba(74,222,128,0.4); }
  .msg-bubble pre code {
    display: block; padding: 14px 16px; background: transparent;
    color: #e2e8f0; font-size: 12.5px; line-height: 1.65;
    overflow-x: auto; border: none; font-family: 'Courier New', monospace;
    border-radius: 0; white-space: pre;
  }
  .msg-bubble ul, .msg-bubble ol { padding-left: 20px; margin: 8px 0; }
  .msg-bubble li { margin-bottom: 5px; line-height: 1.7; }
  .msg.user .msg-bubble ul, .msg.user .msg-bubble ol { padding-left: 18px; }
  .msg-bubble hr { border: none; border-top: 1px solid var(--card-border); margin: 14px 0; }
  .md-table-wrap { overflow-x: auto; margin: 8px 0; border-radius: 8px; border: 1px solid var(--card-border); }
  .md-table { border-collapse: collapse; width: 100%; font-size: 12.5px; }
  .md-table th { background: linear-gradient(135deg, var(--primary), var(--primary-l)); color: #fff; font-weight: 700; padding: 8px 12px; text-align: left; font-size: 11.5px; letter-spacing: 0.04em; white-space: nowrap; }
  .md-table td { padding: 7px 12px; border-bottom: 1px solid var(--card-border); color: var(--text); vertical-align: top; }
  .md-table tr:last-child td { border-bottom: none; }
  .md-table tr:nth-child(even) td { background: rgba(122,16,40,0.03); }
  .md-table tr:hover td { background: rgba(122,16,40,0.06); }
  .msg.user .md-table-wrap { border-color: rgba(255,255,255,0.2); }
  .msg.user .md-table th { background: rgba(0,0,0,0.25); color: rgba(255,255,255,0.9); }
  .msg.user .md-table td { color: rgba(255,255,255,0.9); border-color: rgba(255,255,255,0.15); }
  .msg.user .md-table tr:nth-child(even) td { background: rgba(0,0,0,0.12); }
  /* Quick action chips */
  .q-chip {
    display: inline-block; padding: 5px 12px; margin-right: 6px;
    background: #fff; border: 1px solid var(--card-border); border-radius: 16px;
    font-size: 12.5px; color: var(--text); cursor: pointer; white-space: nowrap;
    transition: all 0.15s; font-family: inherit;
  }
  .q-chip:hover { background: var(--primary-dim); border-color: var(--primary-border); color: var(--primary); }
  /* Drag & drop overlay */
  .chat-main { position: relative; }
  #drag-overlay {
    display: none; position: absolute; inset: 0; z-index: 200;
    background: rgba(122,16,40,0.08); border: 2.5px dashed var(--primary);
    border-radius: 12px; align-items: center; justify-content: center;
    flex-direction: column; gap: 10px; pointer-events: none;
  }
  #drag-overlay .drag-icon { font-size: 40px; opacity: 0.7; }
  #drag-overlay .drag-label { font-size: 15px; font-weight: 600; color: var(--primary); }
  /* Image lightbox */
  #img-lightbox {
    display: none; position: fixed; inset: 0; z-index: 9999;
    background: rgba(0,0,0,0.88); align-items: center; justify-content: center;
    cursor: zoom-out;
  }
  #img-lightbox img {
    max-width: 90vw; max-height: 90vh; border-radius: 10px;
    box-shadow: 0 8px 40px rgba(0,0,0,0.6); object-fit: contain;
  }
  #lightbox-close {
    position: fixed; top: 20px; right: 28px; font-size: 32px; color: #fff;
    cursor: pointer; line-height: 1; opacity: 0.8; z-index: 10000;
    background: none; border: none;
  }
  #lightbox-close:hover { opacity: 1; }
  .thinking-block {
    background: #f5f3ff; border: 1px solid #ddd6fe;
    border-radius: 10px; padding: 8px 12px; margin: 6px 0; font-size: 12px;
  }
  .thinking-block summary {
    cursor: pointer; color: #7c3aed; font-weight: 600; font-size: 12px;
    user-select: none; list-style: none; display: flex; align-items: center; gap: 6px;
  }
  .thinking-block summary::-webkit-details-marker { display: none; }
  .thinking-content { margin-top: 8px; color: var(--text-muted); line-height: 1.6; }
  .file-download-card {
    display: flex; align-items: center; gap: 14px;
    background: #f9fafb; border: 1px solid var(--card-border);
    border-radius: 12px; padding: 14px 18px; text-decoration: none;
    color: var(--text); transition: box-shadow .18s, background .18s; margin: 8px 0; cursor: pointer;
  }
  .file-download-card:hover { background: #f0f0f0; box-shadow: 0 2px 10px rgba(0,0,0,0.09); }
  .file-card-name { font-size: 13.5px; font-weight: 700; color: var(--text); }
  .file-card-type { font-size: 11px; color: var(--text-muted); margin-top: 2px; }
  .file-gen-loading { color: var(--text-muted); font-size: 13px; display: flex; align-items: center; gap: 8px; padding: 8px 0; }
  .typing-indicator { display: flex; gap: 5px; align-items: center; padding: 4px 2px; }
  /* Tool action status cards */
  .tool-action-card { display:flex; align-items:center; gap:9px; padding:9px 14px; border-radius:10px; font-size:12px; font-weight:600; margin:-10px 0 4px 40px; border:1px solid; max-width:calc(88% - 40px); animation:msgIn .2s ease; }
  .tool-action-card.running { background:rgba(156,26,53,.06); border-color:rgba(156,26,53,.18); color:var(--primary); }
  .tool-action-card.done    { background:rgba(29,185,84,.06);  border-color:rgba(29,185,84,.22);  color:#15803d; }
  .tool-action-card.failed  { background:rgba(220,38,38,.06);  border-color:rgba(220,38,38,.22);  color:#dc2626; }
  .tool-action-card.has-image { display:block; max-width:min(420px, calc(88% - 52px)); }
  .tool-action-card.has-image .tool-card-header { display:flex; align-items:center; gap:9px; }

  /* ── POST COMPOSER WIDGET ── */
  .composer-msg { margin:8px 0 12px 40px; max-width:calc(100% - 40px); }
  .composer-card { background:rgba(12,3,7,.97); border:1px solid rgba(156,26,53,.28); border-radius:14px; overflow:hidden; box-shadow:0 12px 40px rgba(0,0,0,.5),0 0 0 1px rgba(156,26,53,.08); max-width:660px; width:100%; }
  .composer-header { display:flex; align-items:center; gap:10px; padding:12px 15px; background:linear-gradient(90deg,rgba(156,26,53,.14),rgba(156,26,53,.06)); border-bottom:1px solid rgba(156,26,53,.18); }
  .comp-platform-badge { font-size:11px; font-weight:700; color:#c0392b; background:rgba(156,26,53,.16); border-radius:6px; padding:3px 9px; letter-spacing:.04em; }
  .comp-title { flex:1; font-size:13px; font-weight:700; color:#fff; letter-spacing:.01em; }
  .comp-close-btn { background:none; border:none; color:rgba(255,255,255,.35); cursor:pointer; font-size:14px; padding:2px 5px; transition:color .2s; line-height:1; }
  .comp-close-btn:hover { color:#fff; }
  .comp-section { padding:13px 15px; border-bottom:1px solid rgba(255,255,255,.04); }
  .comp-step-label { font-size:10px; font-weight:700; color:rgba(200,55,75,.9); text-transform:uppercase; letter-spacing:.12em; margin-bottom:10px; display:flex; align-items:center; gap:7px; font-family:'Courier New',monospace; }
  .comp-step-num { display:inline-flex; align-items:center; justify-content:center; width:17px; height:17px; border-radius:50%; background:rgba(156,26,53,.3); font-size:9px; font-weight:800; color:#c0392b; flex-shrink:0; }
  /* Image grid */
  .comp-image-grid { display:grid; grid-template-columns:repeat(auto-fill,minmax(105px,1fr)); gap:7px; max-height:270px; overflow-y:auto; margin-bottom:10px; }
  .comp-image-grid::-webkit-scrollbar { width:3px; } .comp-image-grid::-webkit-scrollbar-track { background:transparent; } .comp-image-grid::-webkit-scrollbar-thumb { background:rgba(156,26,53,.35); border-radius:2px; }
  .comp-grid-loading { grid-column:1/-1; display:flex; align-items:center; gap:10px; color:rgba(255,255,255,.38); font-size:12px; padding:18px; justify-content:center; }
  .comp-grid-empty { grid-column:1/-1; text-align:center; color:rgba(255,255,255,.3); font-size:12px; padding:18px; }
  .comp-img-item { cursor:pointer; border-radius:7px; overflow:hidden; border:2px solid transparent; transition:border-color .18s,transform .15s,box-shadow .18s; }
  .comp-img-item:hover { border-color:rgba(156,26,53,.45); transform:translateY(-2px); }
  .comp-img-item.selected { border-color:#9b1a35; box-shadow:0 0 14px rgba(156,26,53,.35); }
  .comp-img-item.comp-img-unavail { opacity:.45; }
  .comp-img-wrap { position:relative; aspect-ratio:1; overflow:hidden; background:rgba(255,255,255,.04); }
  .comp-img-wrap img { width:100%; height:100%; object-fit:cover; display:block; }
  .comp-img-ph,.comp-img-missing::after { position:absolute; inset:0; display:flex; align-items:center; justify-content:center; color:rgba(255,255,255,.12); font-size:22px; }
  .comp-img-ph { background:rgba(255,255,255,.03); }
  .comp-img-missing { background:rgba(255,255,255,.03); }
  .comp-img-missing::after { content:'\f03e'; font-family:'Font Awesome 6 Free'; font-weight:900; }
  .comp-today-badge,.comp-date-badge { position:absolute; top:3px; left:3px; font-size:8px; font-weight:700; padding:2px 5px; border-radius:4px; color:#fff; white-space:nowrap; line-height:1.4; }
  .comp-today-badge { background:rgba(156,26,53,.88); } .comp-date-badge { background:rgba(0,0,0,.65); }
  .comp-img-check { position:absolute; inset:0; display:none; align-items:center; justify-content:center; background:rgba(156,26,53,.3); font-size:20px; color:#fff; }
  .comp-img-item.selected .comp-img-check { display:flex; }
  .comp-img-name { font-size:8.5px; color:rgba(255,255,255,.38); padding:3px 4px 2px; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; line-height:1.3; }
  /* Upload */
  .comp-upload-label { display:inline-flex; align-items:center; gap:6px; font-size:11px; color:rgba(156,26,53,.75); border:1px dashed rgba(156,26,53,.28); border-radius:7px; padding:5px 11px; cursor:pointer; transition:all .2s; }
  .comp-upload-label:hover { border-color:rgba(156,26,53,.55); color:#c0392b; background:rgba(156,26,53,.04); }
  /* Caption */
  .comp-caption-tabs { display:flex; align-items:center; gap:6px; margin-bottom:8px; }
  .comp-tab { padding:4px 12px; border-radius:20px; font-size:11px; font-weight:700; border:1px solid rgba(156,26,53,.2); background:transparent; color:rgba(255,255,255,.38); cursor:pointer; transition:all .18s; font-family:inherit; }
  .comp-tab.active { background:rgba(156,26,53,.22); border-color:rgba(156,26,53,.45); color:#c0392b; }
  .comp-regen-btn { margin-left:auto; padding:4px 9px; border-radius:6px; font-size:11px; border:1px solid rgba(255,255,255,.1); background:transparent; color:rgba(255,255,255,.38); cursor:pointer; transition:all .18s; font-family:inherit; }
  .comp-regen-btn:hover { border-color:rgba(156,26,53,.35); color:#c0392b; }
  .comp-caption-loading { display:flex; align-items:center; gap:9px; color:rgba(255,255,255,.38); font-size:12px; padding:10px 0; }
  .comp-caption-input { width:100%; background:rgba(255,255,255,.04); border:1px solid rgba(156,26,53,.18); border-radius:8px; padding:10px 12px; color:#fff; font-size:13px; font-family:inherit; resize:vertical; min-height:96px; outline:none; transition:border-color .2s; line-height:1.55; }
  .comp-caption-input:focus { border-color:rgba(156,26,53,.45); background:rgba(255,255,255,.05); }
  /* Footer */
  .comp-footer { padding:12px 15px; }
  .comp-preview-row { display:flex; align-items:center; gap:10px; margin-bottom:11px; }
  .comp-preview-img { width:52px; height:52px; object-fit:cover; border-radius:8px; border:1px solid rgba(156,26,53,.3); flex-shrink:0; }
  .comp-preview-meta { flex:1; min-width:0; }
  .comp-preview-fname { font-size:12px; color:#fff; font-weight:600; white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
  .comp-preview-date { font-size:10px; color:rgba(255,255,255,.38); margin-top:2px; }
  .comp-action-row { display:flex; gap:7px; }
  .comp-btn-post { flex:1; padding:10px 18px; background:linear-gradient(135deg,#7a1028,#b01535); color:#fff; font-size:13px; font-weight:700; border:none; border-radius:9px; cursor:pointer; display:flex; align-items:center; justify-content:center; gap:7px; transition:all .22s; box-shadow:0 4px 16px rgba(156,26,53,.38); font-family:inherit; }
  .comp-btn-post:hover:not(:disabled) { transform:translateY(-1px); box-shadow:0 6px 22px rgba(156,26,53,.52); }
  .comp-btn-post:disabled { opacity:.6; cursor:not-allowed; }
  .comp-btn-decline { padding:10px 14px; background:rgba(255,255,255,.05); color:rgba(255,255,255,.5); font-size:13px; font-weight:600; border:1px solid rgba(255,255,255,.1); border-radius:9px; cursor:pointer; transition:all .18s; font-family:inherit; }
  .comp-btn-decline:hover { background:rgba(255,255,255,.09); color:#fff; }
  .comp-success-card { display:flex; align-items:center; gap:12px; padding:15px; background:rgba(22,163,74,.08); border:1px solid rgba(22,163,74,.22); border-radius:11px; color:#4ade80; font-size:13px; }
  .comp-success-card i { font-size:20px; flex-shrink:0; }
  .tool-card-img { margin-top:10px; border-radius:10px; overflow:hidden; }
  .tool-card-img img { width:100%; display:block; border-radius:10px; }
  .tool-action-spinner { width:12px; height:12px; border:2px solid rgba(156,26,53,.15); border-top-color:var(--primary); border-radius:50%; animation:fa-spin .7s linear infinite; flex-shrink:0; }
  .typing-dot {
    width: 7px; height: 7px; border-radius: 50%;
    background: var(--primary-l); animation: typing 1.2s infinite;
  }
  .typing-dot:nth-child(2) { animation-delay: 0.2s; }
  .typing-dot:nth-child(3) { animation-delay: 0.4s; }
  @keyframes typing {
    0%, 60%, 100% { transform: translateY(0); opacity: 0.4; }
    30% { transform: translateY(-5px); opacity: 1; }
  }
  .chat-input-area {
    padding: 12px 24px 18px; border-top: 1px solid var(--card-border);
    background: #ffffff; flex-shrink: 0;
  }
  .chat-input-wrap {
    background: #fafafa; border: 1.5px solid var(--card-border);
    border-radius: 16px; overflow: hidden;
    transition: border-color var(--transition), box-shadow var(--transition), background var(--transition);
    box-shadow: 0 1px 6px rgba(0,0,0,0.04);
  }
  .chat-input-wrap:focus-within {
    background: #ffffff; border-color: rgba(122,16,40,0.3);
    box-shadow: 0 0 0 4px rgba(122,16,40,0.07), 0 2px 8px rgba(0,0,0,0.06);
  }
  #chat-input {
    width: 100%; background: transparent; border: none;
    padding: 14px 18px 10px; color: var(--text); font-size: 14px;
    font-family: inherit; resize: none; outline: none;
    min-height: 52px; max-height: 160px; line-height: 1.6; display: block;
  }
  #chat-input::placeholder { color: var(--text-dim); }
  .chat-input-footer {
    display: flex; align-items: center; gap: 8px;
    padding: 8px 12px; border-top: 1px solid rgba(0,0,0,0.04);
    background: rgba(0,0,0,0.01);
  }
  .attach-btn {
    width: 30px; height: 30px; flex-shrink: 0;
    background: transparent; border: 1px solid transparent;
    border-radius: 8px; font-size: 14px; cursor: pointer;
    display: flex; align-items: center; justify-content: center;
    transition: all 0.15s; color: var(--text-dim);
  }
  .attach-btn:hover { background: var(--primary-dim); border-color: var(--primary-border); color: var(--primary); }
  #voice-btn.listening { background: var(--primary); border-color: var(--primary); color: #fff; animation: voice-pulse 1s ease-in-out infinite; }
  @keyframes voice-pulse { 0%,100%{box-shadow:0 0 0 0 rgba(122,16,40,.4)} 50%{box-shadow:0 0 0 6px rgba(122,16,40,0)} }
  .input-hint { font-size: 11px; color: var(--text-dim); flex: 1; }
  #send-btn {
    width: 36px; height: 36px; border-radius: 10px;
    background: var(--primary);
    border: none; cursor: pointer; color: #fff;
    display: flex; align-items: center; justify-content: center;
    font-size: 14px; transition: all 0.15s; flex-shrink: 0;
    box-shadow: 0 2px 8px rgba(122,16,40,0.3);
  }
  #send-btn:hover { background: var(--primary-l); transform: scale(1.05); box-shadow: 0 4px 16px rgba(122,16,40,0.4); }
  #send-btn:disabled { opacity: 0.35; cursor: not-allowed; transform: none; box-shadow: none; }

  /* ── STOP BUTTON ── */
  #stop-btn { display:none; width:36px; height:36px; border-radius:10px; background:#fff; border:1.5px solid #fecaca; cursor:pointer; color:#dc2626; align-items:center; justify-content:center; font-size:13px; transition:all .15s; flex-shrink:0; }
  #stop-btn.active { display:flex; }
  #stop-btn:hover { background:#fef2f2; border-color:#f87171; }
  /* ── INPUT FOOTER HINT ── */
  .input-footer-hint { font-size:10.5px; color:var(--text-dim); flex:1; text-align:center; }

  /* ── MESSAGE BODY ── */
  .msg-body { position:relative; display:flex; flex-direction:column; min-width:0; flex:1; }
  .msg.user .msg-body { align-items:flex-end; }
  .msg-actions { display:flex; gap:3px; opacity:0; transition:opacity .15s; }
  .msg-actions { margin-top:4px; }
  .msg-action-row { display:flex; gap:3px; opacity:0; transition:opacity .15s; margin:4px 0 -18px 40px; position:relative; z-index:1; }
  .msg:hover .msg-actions { opacity:1; }
  .msg-action-btn { width:26px; height:26px; border-radius:6px; border:1px solid var(--card-border); background:#fff; color:var(--text-muted); cursor:pointer; display:flex; align-items:center; justify-content:center; font-size:11px; transition:all .12s; box-shadow:0 1px 3px rgba(0,0,0,.05); }
  .msg-action-btn:hover { background:var(--primary-dim); border-color:var(--primary-border); color:var(--primary); }
  .msg-action-btn.copied { background:#f0fdf4; border-color:#bbf7d0; color:#16a34a; }
  .msg-timestamp { font-size:9px; color:var(--text-dim); margin-top:5px; opacity:0; transition:opacity .15s; }
  .msg.user .msg-timestamp { text-align:right; }
  .msg:hover .msg-timestamp { opacity:1; }

  /* ── CODE BLOCKS ── */
  .code-block-wrap { position:relative; margin:8px 0; border-radius:10px; overflow:hidden; font-size:0; }
  .code-block-top { display:flex; align-items:center; justify-content:space-between; padding:6px 12px; background:#181825; border-bottom:1px solid rgba(255,255,255,.06); }
  .code-lang-label { font-size:9px; font-weight:700; text-transform:uppercase; letter-spacing:.1em; color:rgba(205,214,244,.45); font-family:'Courier New',monospace; }
  .code-copy-btn { padding:2px 8px; background:rgba(255,255,255,.07); border:1px solid rgba(255,255,255,.1); border-radius:5px; color:rgba(205,214,244,.65); font-size:10px; font-weight:600; cursor:pointer; font-family:inherit; transition:all .15s; }
  .code-copy-btn:hover { background:rgba(255,255,255,.14); color:#cdd6f4; }
  .code-block-wrap pre { background:#1e1e2e; color:#cdd6f4; padding:12px 16px; font-size:12px; overflow-x:auto; font-family:'Courier New',monospace; line-height:1.65; margin:0; }

  /* ── SUGGESTED PROMPTS ── */
  .suggested-prompts { display:flex; flex-wrap:wrap; gap:8px; justify-content:center; margin-top:18px; }
  .prompt-chip { padding:7px 14px; border-radius:20px; border:1px solid var(--card-border); background:#fff; color:var(--text); font-size:12px; cursor:pointer; transition:all .15s; font-family:inherit; box-shadow:0 1px 3px rgba(0,0,0,.05); }
  .prompt-chip:hover { background:var(--primary-dim); border-color:var(--primary-border); color:var(--primary); transform:translateY(-1px); box-shadow:0 3px 8px rgba(122,16,40,.1); }

  /* ── CHAT SEARCH ── */
  .chat-search-bar { display:none; padding:8px 20px; border-bottom:1px solid var(--card-border); background:#fafafa; gap:8px; align-items:center; flex-shrink:0; }
  .chat-search-bar.open { display:flex; }
  .chat-search-input { flex:1; background:#fff; border:1px solid var(--card-border); border-radius:8px; padding:7px 12px; font-size:13px; font-family:inherit; outline:none; color:var(--text); transition:border-color .15s; }
  .chat-search-input:focus { border-color:var(--primary-border); box-shadow:0 0 0 2px rgba(122,16,40,.06); }
  .search-match-count { font-size:11px; color:var(--text-muted); flex-shrink:0; min-width:70px; text-align:right; }
  mark.chat-hl { background:rgba(122,16,40,.15); color:var(--primary); border-radius:3px; padding:0 2px; }
  mark.chat-hl.current { background:rgba(122,16,40,.38); }

  /* ── CHAR COUNTER ── */
  .char-counter { font-size:11px; color:var(--text-dim); transition:color .15s; white-space:nowrap; }
  .char-counter.over { color:#f59e0b; }

  /* ── HEADER ICON BTNS ── */
  .chat-hdr-btn { width:30px; height:30px; border-radius:7px; border:1px solid var(--card-border); background:#fff; color:var(--text-muted); cursor:pointer; display:flex; align-items:center; justify-content:center; font-size:12px; transition:all .15s; flex-shrink:0; }
  .chat-hdr-btn:hover { border-color:var(--primary-border); color:var(--primary); background:var(--primary-dim); }
  .chat-hdr-btn.active { border-color:var(--primary); color:var(--primary); background:var(--primary-dim); }

  /* ── PIN ── */
  .session-pin-btn { opacity:0; font-size:10px; background:none; border:none; cursor:pointer; padding:2px 4px; color:var(--text-dim); transition:all .15s; flex-shrink:0; }
  .chat-session-item:hover .session-pin-btn { opacity:1; }
  .session-pin-btn.pinned { opacity:1 !important; color:var(--primary); }
  .chat-session-item.pinned-item { border-left:2px solid var(--primary); }

  /* ── MODAL ── */
  #modal-overlay {
    display: none;
    position: fixed; inset: 0;
    background: rgba(0,0,0,0.7);
    backdrop-filter: blur(6px);
    z-index: 1000;
    align-items: center;
    justify-content: center;
  }
  #modal-overlay.open { display: flex; }
  .modal {
    background: #ffffff;
    border: 1px solid var(--card-border);
    border-radius: var(--radius);
    width: 600px;
    max-width: 90vw;
    max-height: 80vh;
    display: flex;
    flex-direction: column;
    box-shadow: 0 24px 80px rgba(78,9,23,0.35);
    animation: modalIn 0.2s ease;
  }
  @keyframes modalIn { from { opacity: 0; transform: scale(0.95); } to { opacity: 1; transform: scale(1); } }
  .modal-header {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 18px 22px;
    border-bottom: 1px solid var(--card-border);
  }
  .modal-title { font-size: 15px; font-weight: 700; color: var(--text); }
  .modal-close {
    background: none; border: none; color: var(--text-muted);
    cursor: pointer; font-size: 20px; transition: color var(--transition);
  }
  .modal-close:hover { color: var(--text); }
  .modal-body { padding: 20px 22px; overflow-y: auto; }
  .modal-pre {
    background: rgba(122,16,40,0.04);
    border: 1px solid var(--card-border);
    border-radius: var(--radius-sm);
    padding: 16px;
    font-size: 12px;
    white-space: pre-wrap;
    word-break: break-word;
    color: var(--text-muted);
    font-family: 'Courier New', monospace;
    max-height: 400px;
    overflow-y: auto;
  }

  /* ── POST / EMAIL ACTION MODALS ── */
  .action-modal-overlay {
    display: none; position: fixed; inset: 0;
    background: rgba(0,0,0,0.5); z-index: 1100;
    backdrop-filter: blur(3px);
    align-items: center; justify-content: center;
  }
  .action-modal-overlay.open { display: flex; }
  .action-modal {
    background: #fff;
    border-radius: 20px;
    width: 540px; max-width: 95vw; max-height: 88vh;
    display: flex; flex-direction: column;
    box-shadow: 0 32px 80px rgba(78,9,23,0.3);
    animation: modalIn 0.22s cubic-bezier(0.16,1,0.3,1);
    overflow: hidden;
  }
  /* Wide 2-col variant for FB/IG post modals */
  .action-modal.am-wide {
    width: min(920px, 96vw);
    height: min(600px, 90vh);
  }
  .am-two-col {
    display: flex; flex: 1; overflow: hidden;
  }
  /* LEFT — image panel */
  .am-col-img {
    width: 44%; flex-shrink: 0;
    background: #1a0a0f;
    display: flex; align-items: center; justify-content: center;
    position: relative; overflow: hidden;
  }
  .am-col-img img {
    width: 100%; height: 100%; object-fit: cover;
    display: none;
  }
  .am-col-img img.loaded { display: block; }
  .am-img-empty {
    display: flex; flex-direction: column; align-items: center; gap: 10px;
    color: rgba(255,255,255,0.2); font-size: 12px; letter-spacing: .05em;
    text-transform: uppercase; font-weight: 600;
  }
  .am-img-empty i { font-size: 36px; opacity: .3; }
  /* RIGHT — form panel */
  .am-col-form {
    flex: 1; display: flex; flex-direction: column;
    padding: 22px 24px 0; overflow-y: auto;
  }
  .am-col-form-inner { flex: 1; }
  .am-footer-inline {
    padding: 14px 0 18px;
    border-top: 1px solid var(--card-border);
    display: flex; gap: 10px; justify-content: flex-end;
    flex-shrink: 0; margin-top: auto;
  }
  .am-header {
    display: flex; align-items: center; gap: 12px;
    padding: 18px 24px 14px;
    border-bottom: 1px solid var(--card-border);
    flex-shrink: 0;
  }
  .am-header-icon {
    width: 38px; height: 38px; border-radius: 11px;
    display: flex; align-items: center; justify-content: center;
    font-size: 16px; flex-shrink: 0;
  }
  .am-icon-fb    { background: rgba(24,119,242,0.12); color: #1877f2; }
  .am-icon-ig    { background: rgba(225,48,108,0.1);  color: #e1306c; }
  .am-icon-email { background: var(--primary-dim); color: var(--primary); }
  .am-title { font-size: 16px; font-weight: 700; color: var(--text); flex: 1; }
  .am-close { background: none; border: none; color: var(--text-dim); cursor: pointer; font-size: 20px; line-height: 1; padding: 4px; transition: color .2s; }
  .am-close:hover { color: var(--text); }
  .am-body { padding: 20px 24px; overflow-y: auto; flex: 1; }
  .am-footer { padding: 14px 24px; border-top: 1px solid var(--card-border); display: flex; gap: 10px; justify-content: flex-end; flex-shrink: 0; }
  .am-label { font-size: 10px; font-weight: 700; text-transform: uppercase; letter-spacing: .08em; color: var(--text-muted); margin-bottom: 6px; display: flex; align-items: center; justify-content: space-between; }
  .am-group { margin-bottom: 14px; }
  /* ── CUSTOM DROPDOWN ── */
  .cs-wrap { position: relative; }
  .cs-trigger {
    width: 100%; display: flex; align-items: center; gap: 10px;
    background: #fdfafa; border: 1.5px solid var(--card-border);
    border-radius: 10px; padding: 10px 13px; cursor: pointer;
    text-align: left; font-family: inherit;
    transition: border-color .2s, box-shadow .2s;
  }
  .cs-trigger:hover { border-color: rgba(122,16,40,0.25); }
  .cs-wrap.open .cs-trigger {
    border-color: var(--primary-border);
    box-shadow: 0 0 0 3px rgba(122,16,40,0.07);
    border-radius: 10px 10px 0 0;
  }
  .cs-trigger-icon { color: var(--primary); font-size: 13px; flex-shrink: 0; }
  .cs-trigger-content { flex: 1; overflow: hidden; min-width: 0; }
  .cs-trigger-text {
    font-size: 13px; color: var(--text); font-weight: 500;
    white-space: nowrap; overflow: hidden; text-overflow: ellipsis; display: block;
  }
  .cs-trigger-badge {
    font-size: 9px; font-weight: 700; letter-spacing: .04em;
    background: var(--primary-dim); color: var(--primary);
    border: 1px solid var(--primary-border);
    border-radius: 20px; padding: 2px 8px; white-space: nowrap; flex-shrink: 0;
  }
  .cs-trigger-badge.today-badge { background: rgba(217,119,6,0.1); color: #b45309; border-color: rgba(217,119,6,0.25); }
  .cs-trigger-arrow { color: var(--text-muted); font-size: 10px; flex-shrink: 0; transition: transform .2s; }
  .cs-wrap.open .cs-trigger-arrow { transform: rotate(180deg); }
  .cs-list-wrap {
    position: absolute; top: 100%; left: 0; right: 0;
    background: #fff; border: 1.5px solid var(--primary-border);
    border-top: none; border-radius: 0 0 12px 12px;
    max-height: 260px; overflow-y: auto;
    z-index: 300; display: none;
    box-shadow: 0 16px 36px rgba(78,9,23,0.18);
  }
  .cs-wrap.open .cs-list-wrap { display: block; }
  .cs-list { padding: 6px; }
  .cs-option {
    display: flex; align-items: center; gap: 10px;
    padding: 8px 10px; border-radius: 8px; cursor: pointer;
    transition: background .12s;
  }
  .cs-option:hover { background: var(--primary-dim); }
  .cs-option.cs-selected { background: var(--primary-dim); }
  .cs-option.cs-today { border-bottom: 1px solid var(--card-border); margin-bottom: 4px; padding-bottom: 10px; }
  .cs-today-star { color: #b45309; font-size: 9px; flex-shrink: 0; }
  .cs-opt-dot { width: 5px; height: 5px; border-radius: 50%; background: rgba(122,16,40,0.15); flex-shrink: 0; margin: 0 1px; }
  .cs-opt-date-badge {
    font-size: 10px; font-weight: 700; color: var(--text-muted); min-width: 52px; flex-shrink: 0;
  }
  .cs-option.cs-today .cs-opt-date-badge { color: #b45309; }
  .cs-opt-name {
    font-size: 12px; color: var(--text); flex: 1;
    white-space: nowrap; overflow: hidden; text-overflow: ellipsis;
  }
  .cs-option.cs-selected .cs-opt-name { color: var(--primary); font-weight: 600; }
  /* ── END CUSTOM DROPDOWN ── */
  .am-select, .am-textarea, .am-input {
    width: 100%; background: #fdfafa; border: 1.5px solid var(--card-border);
    border-radius: 10px; padding: 10px 13px; color: var(--text);
    font-size: 13px; font-family: inherit; outline: none;
    transition: border-color .2s, box-shadow .2s;
  }
  .am-select:focus, .am-textarea:focus, .am-input:focus {
    border-color: var(--primary-border);
    box-shadow: 0 0 0 3px rgba(122,16,40,0.07);
  }
  .am-select:hover, .am-textarea:hover { border-color: rgba(122,16,40,0.25); }
  .am-textarea { resize: vertical; line-height: 1.6; }
  .am-caption-row { display: flex; align-items: center; justify-content: space-between; margin-bottom: 6px; }
  .am-gen-btn { font-size: 11px; color: var(--primary); background: var(--primary-dim); border: 1px solid var(--primary-border); border-radius: 20px; padding: 4px 11px; cursor: pointer; font-family: inherit; font-weight: 600; transition: all .2s; display: flex; align-items: center; gap: 5px; white-space: nowrap; }
  .am-gen-btn:hover { background: rgba(122,16,40,0.12); }
  .am-image-chip { display: inline-flex; align-items: center; gap: 6px; background: var(--primary-dim); border: 1px solid var(--primary-border); border-radius: 8px; padding: 6px 12px; font-size: 12px; color: var(--text); margin-bottom: 14px; }
  .am-summary-row { display: flex; align-items: center; gap: 10px; background: var(--primary-dim); border-radius: 10px; padding: 12px 16px; margin-bottom: 10px; }
  .am-summary-num { font-size: 28px; font-weight: 800; color: var(--primary); line-height: 1; }
  .am-summary-label { font-size: 12px; color: var(--text-muted); }
  .am-lead-preview { background: #fdfafa; border: 1px solid var(--card-border); border-radius: 10px; overflow: hidden; }
  .am-lead-row { display: flex; align-items: center; gap: 10px; padding: 9px 14px; font-size: 12px; border-bottom: 1px solid var(--card-border); }
  .am-lead-row:last-child { border-bottom: none; }
  .am-lead-name { font-weight: 600; color: var(--text); flex: 1; }
  .am-lead-step { font-size: 10px; background: var(--primary-dim); color: var(--primary); border-radius: 20px; padding: 2px 8px; font-weight: 600; }

  /* ── TOAST ── */
  #toast {
    position: fixed;
    bottom: 24px;
    right: 24px;
    background: var(--primary-d);
    border: 1px solid rgba(255,255,255,0.1);
    border-radius: var(--radius);
    padding: 14px 20px;
    display: flex;
    align-items: center;
    gap: 10px;
    font-size: 13px;
    font-weight: 500;
    color: #f5ecec;
    box-shadow: 0 8px 32px rgba(78,9,23,0.45);
    transform: translateY(80px);
    opacity: 0;
    transition: all 0.3s cubic-bezier(0.4,0,0.2,1);
    z-index: 2000;
    max-width: 320px;
  }
  #toast.show { transform: translateY(0); opacity: 1; }
  .toast-icon { font-size: 18px; flex-shrink: 0; }

  /* ── LOADING SKELETON ── */
  .skeleton {
    background: linear-gradient(90deg, rgba(255,255,255,0.05) 25%, rgba(255,255,255,0.08) 50%, rgba(255,255,255,0.05) 75%);
    background-size: 200% 100%;
    animation: shimmer 1.5s infinite;
    border-radius: 6px;
    height: 14px;
  }
  @keyframes shimmer { 0% { background-position: -200% 0; } 100% { background-position: 200% 0; } }

  /* ── RESPONSIVE ── */
  @media (max-width: 900px) {
    .stats-grid { grid-template-columns: repeat(3, 1fr); }
    .two-col { grid-template-columns: 1fr; }
  }
  @media (max-width: 640px) {
    #sidebar { position: fixed; left: -220px; height: 100vh; transition: left 0.3s ease; }
    #sidebar.open { left: 0; }
    .stats-grid { grid-template-columns: repeat(3, 1fr); }
    .view { padding: 16px; }
  }

  /* ── DIVIDER ── */
  .divider { height: 1px; background: rgba(255,255,255,0.07); margin: 4px 0; }

  /* ── EMPTY STATE ── */
  .empty-state { text-align: center; padding: 40px; color: var(--text-dim); }
  .empty-state .empty-icon { font-size: 36px; margin-bottom: 12px; }
  .empty-state p { font-size: 13px; }

  /* scrollable table wrapper */
  .table-wrap { overflow-x: auto; }

  /* ── LEAD SCORE BADGES ── */
  .badge-hot  { background:rgba(239,68,68,0.15);  color:#F87171; border:1px solid rgba(239,68,68,0.3); }
  .badge-warm { background:rgba(245,158,11,0.15); color:#FCD34D; border:1px solid rgba(245,158,11,0.3); }
  .badge-cold { background:rgba(96,165,250,0.15); color:#93C5FD; border:1px solid rgba(96,165,250,0.3); }

  /* ── SEARCH BAR ── */
  .search-wrap { position:relative; margin-bottom:14px; }
  .search-wrap input { width:100%; background:#fdfafa; border:1px solid var(--card-border); border-radius:var(--radius-sm); padding:9px 12px 9px 36px; color:var(--text); font-size:13px; font-family:inherit; outline:none; transition:border-color .2s; }
  .search-wrap input:focus { border-color:var(--primary-border); background:#fff; }
  .search-icon { position:absolute; left:11px; top:50%; transform:translateY(-50%); font-size:13px; color:var(--text-dim); pointer-events:none; }

  /* ── INTELLIGENCE VIEW ── */
  #view-intelligence { padding: 0; height: 100%; overflow: hidden; }
  #view-intelligence.active { display: flex !important; flex-direction: column; }
  #view-intelligence .view-header { padding: 20px 24px 0; flex-shrink: 0; }
  .intel-layout { display: flex; flex: 1; overflow: hidden; }
  .intel-sidebar {
    width: 200px; min-width: 200px;
    border-right: 1px solid var(--card-border);
    overflow-y: auto; padding: 8px 0; flex-shrink: 0;
  }
  .intel-list-header { font-size: 9px; font-weight: 700; letter-spacing: .12em; text-transform: uppercase; color: var(--text-dim); padding: 8px 16px 6px; }
  .intel-list-item {
    padding: 9px 14px; cursor: pointer; border-radius: 8px; margin: 1px 6px;
    font-size: 12px; color: var(--text-muted); transition: all .18s;
  }
  .intel-list-item:hover  { background: var(--primary-dim); color: var(--text); }
  .intel-list-item.active { background: var(--primary-dim); color: var(--primary); font-weight: 700; }
  .intel-list-date { display: block; font-weight: 600; font-size: 12px; }
  .intel-list-sub  { display: block; font-size: 10px; color: var(--text-dim); margin-top: 1px; }
  .intel-main { flex: 1; overflow-y: auto; padding: 20px 24px; }
  .intel-header { display:flex; align-items:center; gap:12px; margin-bottom:16px; flex-wrap:wrap; }
  .intel-meta { font-size:12px; color:var(--text-dim); }
  .intel-meta strong { color:var(--text); font-weight:600; }
  .threat-badge { padding:5px 14px; border-radius:20px; font-size:11px; font-weight:700; letter-spacing:.05em; text-transform:uppercase; display:inline-flex; align-items:center; gap:6px; }
  .threat-low    { background:rgba(74,222,128,0.15); color:#2a8a50; border:1px solid rgba(74,222,128,0.3); }
  .threat-medium { background:rgba(245,158,11,0.15); color:#FCD34D; border:1px solid rgba(245,158,11,0.3); }
  .threat-high   { background:rgba(239,68,68,0.15);  color:#F87171; border:1px solid rgba(239,68,68,0.3); }
  .threat-unknown{ background:rgba(255,255,255,0.06); color:var(--text-dim); border:1px solid var(--card-border); }
  .intel-body { font-size:13px; line-height:1.75; color:var(--text-muted); }
  .intel-body h2 { font-size:11px; color:var(--primary); font-weight:700; margin:14px 0 6px; text-transform:uppercase; letter-spacing:0.08em; padding-bottom:5px; border-bottom:1px solid var(--card-border); }
  .intel-body h3 { font-size:13px; color:var(--text); font-weight:700; margin:12px 0 4px; }
  .intel-body strong { color:var(--text); font-weight:600; }
  .intel-body em { color:var(--text-dim); }
  .intel-body ul { padding-left:18px; margin:8px 0; }
  .intel-body li { margin-bottom:4px; }
  .intel-action-box { background:rgba(122,16,40,0.05); border-left:3px solid var(--primary); border-radius:0 8px 8px 0; padding:10px 14px; margin:10px 0; font-size:13px; color:var(--text); }
  .intel-action-box .action-label { font-size:10px; font-weight:700; text-transform:uppercase; letter-spacing:.1em; color:var(--primary); margin-bottom:4px; }
  .news-card { background:#fff; border:1px solid var(--card-border); border-radius:10px; padding:13px 16px; margin-bottom:8px; display:flex; gap:12px; align-items:flex-start; transition:box-shadow .2s; }
  .news-card:hover { box-shadow:0 2px 12px rgba(122,16,40,0.08); }
  .news-date-chip { min-width:54px; text-align:center; background:var(--primary-dim); border-radius:6px; padding:4px 6px; font-size:9px; font-weight:700; color:var(--primary); text-transform:uppercase; letter-spacing:.04em; line-height:1.3; flex-shrink:0; }
  .news-card-body { flex:1; min-width:0; }
  .news-card .news-title { font-size:13px; color:var(--text); font-weight:500; line-height:1.45; margin-bottom:4px; }
  .news-card .news-source { display:inline-block; font-size:10px; font-weight:600; color:var(--primary); background:var(--primary-dim); padding:2px 8px; border-radius:20px; }
  .change-card { background:#fff; border:1px solid rgba(180,40,60,0.18); border-radius:10px; padding:12px 16px; margin-bottom:8px; display:flex; align-items:center; justify-content:space-between; gap:10px; }
  .cc-name { font-size:13px; font-weight:700; color:var(--text); }
  .cc-url-link { font-size:11px; color:var(--primary); text-decoration:none; word-break:break-all; }
  .cc-url-link:hover { text-decoration:underline; }
  .intel-sources-footer { margin-top:20px; background:var(--primary-dim); border:1px solid var(--primary-border); border-radius:10px; padding:14px 18px; }
  .intel-sources-footer .sf-title { font-size:10px; font-weight:700; text-transform:uppercase; letter-spacing:.1em; color:var(--text-dim); margin-bottom:8px; }
  .intel-sources-footer .sf-row { font-size:12px; color:var(--text-muted); display:flex; gap:8px; align-items:baseline; margin-bottom:4px; }
  .intel-sources-footer .sf-row i { color:var(--primary); font-size:10px; flex-shrink:0; }
  .date-selector { display:flex; gap:6px; flex-wrap:wrap; margin-bottom:16px; }
  .date-pill { padding:4px 12px; border-radius:20px; font-size:11px; font-weight:600; cursor:pointer; border:1px solid var(--card-border); color:var(--text-muted); background:transparent; font-family:inherit; transition:all .2s; }
  .date-pill:hover { border-color:var(--primary-border); color:var(--primary); }
  .date-pill.active { background:var(--primary-dim); border-color:var(--primary); color:var(--primary); }

  /* ── AUTOMATIONS VIEW ── */
  .auto-grid { display:flex; flex-direction:column; gap:0; }
  .auto-card { background:#ffffff; border:1px solid var(--card-border); border-radius:var(--radius); padding:18px 20px; display:flex; flex-direction:column; gap:10px; transition:transform .2s,box-shadow .2s; box-shadow:0 1px 4px rgba(122,16,40,0.06); }
  .auto-card:hover { box-shadow:0 4px 16px rgba(122,16,40,0.1); }
  .auto-card:hover { transform:translateY(-2px); }
  .auto-card-head { display:flex; align-items:center; gap:10px; }
  .auto-emoji { font-size:20px; line-height:1; width:24px; flex-shrink:0; text-align:center; }
  .auto-name  { font-size:14px; font-weight:700; color:var(--text); flex:1; min-width:0; }
  .auto-status { margin-left:auto; flex-shrink:0; }
  .auto-detail  { font-size:12px; color:var(--text-muted); }
  .auto-schedule{ font-size:11px; color:var(--text-dim); margin-top:2px; }
  .status-pill { padding:3px 10px; border-radius:20px; font-size:10px; font-weight:700; letter-spacing:.04em; text-transform:uppercase; white-space:nowrap; }
  .status-success  { background:rgba(74,222,128,0.12);  color:#4ADE80; border:1px solid rgba(74,222,128,0.3); }
  .status-warning  { background:rgba(245,158,11,0.12);  color:#FCD34D; border:1px solid rgba(245,158,11,0.3); }
  .status-inactive { background:rgba(255,255,255,0.06); color:var(--text-dim); border:1px solid var(--card-border); }
  .status-scheduled{ background:rgba(99,102,241,0.12);  color:#818CF8; border:1px solid rgba(99,102,241,0.3); }
  .status-pending  { background:rgba(245,158,11,0.12);  color:#FCD34D; border:1px solid rgba(245,158,11,0.3); }
  .status-approved { background:rgba(74,222,128,0.12);  color:#4ADE80; border:1px solid rgba(74,222,128,0.3); }
  .status-none     { background:rgba(255,255,255,0.06); color:var(--text-dim); border:1px solid var(--card-border); }

  /* ── UNIFIED WORKFLOW CARDS ── */
  .wf-unified-grid { display:grid; grid-template-columns:repeat(auto-fill,minmax(360px,1fr)); gap:16px; }
  .wf-unified-card { background:#fff; border:1px solid var(--card-border); border-radius:var(--radius); padding:20px; display:flex; flex-direction:column; gap:14px; box-shadow:0 1px 4px rgba(122,16,40,0.06); transition:transform .2s,box-shadow .2s; }
  .wf-unified-card:hover { transform:translateY(-2px); box-shadow:0 6px 20px rgba(122,16,40,0.10); }
  .wf-uc-top { display:flex; gap:14px; align-items:flex-start; }
  .wf-uc-icon { width:40px; height:40px; border-radius:10px; background:rgba(156,26,53,0.08); display:flex; align-items:center; justify-content:center; font-size:18px; color:var(--primary); flex-shrink:0; }
  .wf-uc-info { flex:1; min-width:0; }
  .wf-uc-name { font-size:14px; font-weight:700; color:var(--text); margin-bottom:4px; }
  .wf-uc-desc { font-size:12px; color:var(--text-muted); line-height:1.5; }
  .wf-uc-mid { display:flex; align-items:center; gap:10px; flex-wrap:wrap; }
  .wf-uc-detail { font-size:12px; color:var(--text-muted); }
  .wf-uc-footer { display:flex; align-items:center; justify-content:space-between; gap:10px; border-top:1px solid var(--card-border); padding-top:12px; margin-top:2px; flex-wrap:wrap; }
  .wf-uc-sched { font-size:11px; color:var(--text-dim); display:flex; align-items:center; gap:5px; }
  .wf-uc-actions { display:flex; align-items:center; gap:10px; }
  .wf-run-btn { background:var(--primary); color:#fff; border:none; border-radius:8px; padding:6px 14px; font-size:12px; font-weight:600; cursor:pointer; display:flex; align-items:center; gap:6px; transition:opacity .2s,transform .15s; }
  .wf-run-btn:hover:not(:disabled) { opacity:.88; transform:translateY(-1px); }
  .wf-run-btn:disabled { opacity:.5; cursor:not-allowed; transform:none; }
  .wf-uc-toggle { display:flex; align-items:center; gap:8px; }
  .wf-enabled-badge { font-size:10px; font-weight:700; letter-spacing:.06em; padding:2px 8px; border-radius:10px; text-transform:uppercase; }
  .wf-badge-on  { background:rgba(74,222,128,0.14); color:#22c55e; border:1px solid rgba(74,222,128,0.35); }
  .wf-badge-off { background:rgba(148,163,184,0.12); color:#94a3b8; border:1px solid rgba(148,163,184,0.3); }
  .wf-switch { position:relative; display:inline-block; width:44px; height:24px; flex-shrink:0; }
  .wf-switch input { opacity:0; width:0; height:0; }
  .wf-slider { position:absolute; cursor:pointer; inset:0; background:#e2e8f0; border-radius:24px; transition:.3s; }
  .wf-slider:before { content:""; position:absolute; height:18px; width:18px; left:3px; bottom:3px; background:#fff; border-radius:50%; transition:.3s; box-shadow:0 1px 3px rgba(0,0,0,.2); }
  .wf-switch input:checked + .wf-slider { background:var(--primary); }
  .wf-switch input:checked + .wf-slider:before { transform:translateX(20px); }
  .wf-switch input:disabled + .wf-slider { opacity:.5; cursor:not-allowed; }

  /* ── CALENDAR VIEW ── */
  #view-calendar { padding: 0; overflow: hidden; }
  #view-calendar.active { display: flex !important; flex-direction: column; height: 100%; }
  .cal-toolbar {
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 24px; border-bottom: 1px solid var(--card-border);
    flex-shrink: 0; background: #ffffff; gap: 16px;
  }
  .cal-month-nav { display: flex; align-items: center; gap: 12px; }
  .cal-month-nav span { font-size: 17px; font-weight: 700; min-width: 150px; text-align: center; color: var(--text); }
  #cal-month-view { flex: 1; min-height: 0; display: flex; flex-direction: column; padding: 12px 20px 16px; overflow: hidden; }
  #cal-list-view { flex: 1; min-height: 0; overflow-y: auto; padding: 16px 20px; }
  .cal-tabs { display:flex; gap:6px; }
  .cal-tab {
    padding:6px 14px; border-radius:8px; font-size:12px; font-weight:600;
    border:1px solid var(--card-border); background:#fff; color:var(--text-muted);
    cursor:pointer; transition:all .15s; font-family:inherit;
    display:flex; align-items:center; gap:6px;
  }
  .cal-tab.active { background:var(--primary-dim); color:var(--primary); border-color:var(--primary-border); }
  .cal-tab:hover { background:var(--primary-dim); color:var(--primary); }
  .cal-btn { background:transparent; border:1px solid var(--card-border); border-radius:8px; color:var(--text-muted); padding:6px 12px; cursor:pointer; font-size:15px; font-family:inherit; transition:all .2s; line-height:1; }
  .cal-btn:hover { border-color:var(--primary-border); color:var(--primary); background:var(--primary-dim); }
  .cal-grid { display:grid; grid-template-columns:repeat(7,1fr); gap:3px; }
  #cal-grid { flex:1; min-height:0; grid-auto-rows:1fr; }
  .cal-day-hdr { text-align:center; font-size:10px; font-weight:700; color:var(--text-dim); padding:6px 0; letter-spacing:.06em; text-transform:uppercase; }
  .cal-cell { background:#ffffff; border:1px solid var(--card-border); border-radius:6px; min-height:0; padding:6px 8px; cursor:default; transition:background .15s; position:relative; overflow:hidden; }
  .cal-cell.has-post { cursor:pointer; border-color:rgba(122,16,40,0.2); }
  .cal-cell.has-post:hover { background:var(--primary-dim); border-color:var(--primary-border); }
  .cal-cell.today { border-color:var(--primary); box-shadow:inset 0 0 0 1px rgba(122,16,40,0.2); }
  .cal-cell.other-month { opacity:.35; }
  .cal-cell.past.has-post { border-color:rgba(74,222,128,0.2); }
  .cal-num { font-size:13px; font-weight:600; color:var(--text-muted); }
  .cal-cell.today .cal-num { color:var(--primary); font-weight:800; }
  .cal-dot { width:6px; height:6px; border-radius:50%; background:var(--primary); position:absolute; top:8px; right:8px; box-shadow:0 0 4px var(--primary-l); }
  .cal-dot.past { background:#4ADE80; box-shadow:0 0 4px #4ADE80; }
  .cal-img-name { font-size:9px; color:var(--text-dim); margin-top:4px; line-height:1.3; overflow:hidden; display:-webkit-box; -webkit-line-clamp:2; -webkit-box-orient:vertical; }

  /* ── LOGOUT BTN ── */
  .logout-btn { display:flex; align-items:center; gap:8px; padding:9px 16px; cursor:pointer; color:rgba(245,236,236,0.45); font-size:12px; font-weight:500; transition:color .2s; background:none; border:none; font-family:inherit; width:100%; text-align:left; border-top:1px solid rgba(255,255,255,0.08); }
  .logout-btn:hover { color:#fca5a5; }
  /* ── VIEW HEADER ── */
  .view-header h2 { color:var(--text); }
  /* ── MSG BUBBLE BR SPACING ── */
  .msg-bubble br { display:block; margin-bottom:5px; content:""; }

  /* ── LEAD DETAIL MODAL ── */
  .lead-modal-overlay { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.45); z-index:2000; align-items:center; justify-content:center; padding:20px; }
  .lead-modal-overlay.open { display:flex; }
  .lead-modal { background:#fff; border-radius:20px; padding:28px 28px 32px; width:100%; max-width:460px; box-shadow:0 12px 50px rgba(0,0,0,0.22); animation:fadeScaleIn .2s ease; }
  @keyframes fadeScaleIn { from{transform:scale(0.94);opacity:0} to{transform:scale(1);opacity:1} }
  .lead-modal-header { display:flex; align-items:center; justify-content:space-between; margin-bottom:20px; }
  .lead-modal-name { font-size:18px; font-weight:700; color:var(--text); }
  .lead-modal-close { background:none; border:none; font-size:20px; cursor:pointer; color:var(--text-dim); line-height:1; padding:4px; }
  .lead-modal-field { display:flex; align-items:flex-start; gap:12px; margin-bottom:14px; }
  .lead-modal-field-icon { width:34px; height:34px; border-radius:8px; background:var(--primary-dim); display:flex; align-items:center; justify-content:center; color:var(--primary); font-size:14px; flex-shrink:0; }
  .lead-modal-field-content { display:flex; flex-direction:column; gap:2px; min-width:0; }
  .lead-modal-field-label { font-size:10px; font-weight:600; color:var(--text-dim); text-transform:uppercase; letter-spacing:.06em; }
  .lead-modal-field-value { font-size:13.5px; color:var(--text); font-weight:500; word-break:break-all; }
  .lead-modal-field-value a { color:var(--primary); text-decoration:none; }
  .lead-modal-field-value a:hover { text-decoration:underline; }
  .drip-bar-wrap { margin-top:4px; }
  .drip-bar-track { background:var(--primary-dim); border-radius:6px; height:8px; overflow:hidden; margin-bottom:4px; }
  .drip-bar-fill { background:linear-gradient(90deg,var(--primary),var(--primary-l)); height:100%; border-radius:6px; transition:width .4s ease; }
  .drip-bar-label { font-size:11px; color:var(--text-muted); }
  .lead-modal-actions { display:flex; gap:10px; margin-top:8px; padding-top:16px; border-top:1px solid var(--card-border); }
  .lm-btn { flex:1; padding:9px 14px; border-radius:10px; font-size:13px; font-weight:600; cursor:pointer; border:none; display:flex; align-items:center; justify-content:center; gap:7px; transition:opacity .2s,transform .15s; }
  .lm-btn:hover:not(:disabled) { opacity:.85; transform:translateY(-1px); }
  .lm-btn:disabled { opacity:.5; cursor:not-allowed; transform:none; }
  .lm-btn-toggle { background:rgba(99,102,241,0.1); color:#6366f1; border:1px solid rgba(99,102,241,0.3); }
  .lm-btn-delete { background:rgba(239,68,68,0.1); color:#ef4444; border:1px solid rgba(239,68,68,0.3); }
  .lm-btn-delete:hover:not(:disabled) { background:rgba(239,68,68,0.18); }

  /* ── POST PREVIEW MODAL ── */
  .pp-overlay { display:none; position:fixed; inset:0; background:rgba(0,0,0,0.55); z-index:2000; align-items:center; justify-content:center; padding:20px; }
  .pp-overlay.open { display:flex; }
  .pp-modal { background:#fff; border-radius:16px; width:100%; max-width:720px; max-height:90vh; overflow:auto; box-shadow:0 20px 60px rgba(0,0,0,0.25); }
  .pp-header { display:flex; align-items:center; justify-content:space-between; padding:18px 22px 14px; border-bottom:1px solid var(--card-border); }
  .pp-header h3 { font-size:15px; font-weight:700; color:var(--text); margin:0; }
  .pp-close { background:none; border:none; font-size:20px; cursor:pointer; color:var(--text-dim); padding:4px; line-height:1; }
  .pp-frames { display:grid; grid-template-columns:1fr 1fr; gap:18px; padding:20px; }
  @media(max-width:520px){ .pp-frames { grid-template-columns:1fr; } }
  .pp-frame { border-radius:12px; overflow:hidden; box-shadow:0 2px 12px rgba(0,0,0,0.1); border:1px solid #e0e0e0; background:#fff; }
  .pp-frame-header { display:flex; align-items:center; gap:8px; padding:10px 12px; }
  .pp-frame-label { font-size:11px; font-weight:700; letter-spacing:.04em; }
  .fb-pp-header { background:#1877f2; }
  .fb-pp-header .pp-frame-label { color:#fff; }
  .ig-pp-header { background:#fff; border-bottom:1px solid #efefef; }
  .ig-pp-header .pp-frame-label { color:#262626; }
  .pp-img-wrap { width:100%; aspect-ratio:1; background:#f5f5f5; display:flex; align-items:center; justify-content:center; overflow:hidden; }
  .pp-img-wrap img { width:100%; height:100%; object-fit:cover; }
  .pp-img-empty { display:flex; flex-direction:column; align-items:center; justify-content:center; gap:8px; color:#bbb; font-size:12px; }
  .pp-caption { padding:10px 12px; font-size:11px; color:#555; line-height:1.5; background:#fafafa; border-top:1px solid #efefef; min-height:48px; }
  .pp-no-image { text-align:center; padding:40px 20px; color:var(--text-muted); font-size:13px; }
  .pp-meta { padding:8px 20px 16px; text-align:center; font-size:11px; color:var(--text-dim); }

  /* ── STRATEGY STUDIO VIEW ── */
  #view-strategy { padding:0; height:100%; overflow:hidden; }
  #view-strategy.active { display:flex !important; flex-direction:column; }
  .strategy-body { display:flex; flex:1; overflow:hidden; gap:0; }
  .strategy-form-panel {
    width:340px; min-width:280px; max-width:380px;
    padding:24px 20px; overflow-y:auto; border-right:1px solid var(--card-border);
    flex-shrink:0; display:flex; flex-direction:column; gap:16px;
  }
  .strategy-output-panel {
    flex:1; display:flex; flex-direction:column; overflow:hidden; padding:24px 28px;
  }
  .strategy-tool-tabs { display:flex; gap:6px; flex-wrap:wrap; }
  .strat-tab {
    padding:7px 14px; border-radius:20px; font-size:12px; font-weight:600;
    border:1px solid var(--card-border); background:transparent; color:var(--text-secondary);
    cursor:pointer; transition:all .15s;
  }
  .strat-tab.active { background:#7a1028; color:#fff; border-color:#7a1028; }
  .strat-tab:hover:not(.active) { border-color:#7a1028; color:#7a1028; }
  .strat-field { display:flex; flex-direction:column; gap:5px; }
  .strat-field label { font-size:11px; font-weight:600; color:var(--text-secondary); text-transform:uppercase; letter-spacing:.5px; }
  .strat-field input, .strat-field textarea, .strat-field select {
    padding:9px 12px; border:1px solid var(--card-border); border-radius:8px;
    background:var(--card-bg); color:var(--text-primary); font-size:13px;
    font-family:inherit; resize:vertical; transition:border-color .15s;
  }
  .strat-field input:focus, .strat-field textarea:focus, .strat-field select:focus {
    outline:none; border-color:#7a1028;
  }
  .strat-generate-btn {
    padding:11px 16px; background:#7a1028; color:#fff; border:none;
    border-radius:10px; font-size:14px; font-weight:600; cursor:pointer;
    transition:background .15s; display:flex; align-items:center; justify-content:center; gap:8px;
  }
  .strat-generate-btn:hover { background:#9b1535; }
  .strat-generate-btn:disabled { background:#999; cursor:not-allowed; }
  .strat-output-header { display:flex; align-items:center; justify-content:space-between; margin-bottom:16px; flex-shrink:0; }
  .strat-output-header h3 { font-size:14px; font-weight:600; color:var(--text-secondary); margin:0; }
  .strat-export-btns { display:flex; gap:8px; }
  .strat-export-btn {
    padding:7px 14px; border-radius:8px; font-size:12px; font-weight:600;
    border:1px solid var(--card-border); background:transparent; color:var(--text-secondary);
    cursor:pointer; display:flex; align-items:center; gap:6px; transition:all .15s;
  }
  .strat-export-btn:hover { border-color:#7a1028; color:#7a1028; }
  .strat-export-btn.pptx:hover { background:#7a1028; color:#fff; border-color:#7a1028; }
  .strat-output-scroll {
    flex:1; overflow-y:auto; background:var(--card-bg); border:1px solid var(--card-border);
    border-radius:12px; padding:24px 28px; font-size:14px; line-height:1.75;
    color:var(--text-primary);
  }
  .strat-output-scroll h1,.strat-output-scroll h2 { color:#7a1028; margin-top:1.5em; }
  body.dark .strat-output-scroll h1, body.dark .strat-output-scroll h2 { color:#e88a9b; }
  .strat-output-scroll h3 { color:var(--text-primary); margin-top:1.2em; }
  .strat-output-scroll table { width:100%; border-collapse:collapse; margin:1em 0; font-size:13px; }
  .strat-output-scroll th { background:#7a1028; color:#fff; padding:8px 12px; text-align:left; }
  .strat-output-scroll td { padding:8px 12px; border-bottom:1px solid var(--card-border); }
  .strat-output-scroll tr:hover td { background:rgba(122,16,40,.05); }
  .strat-empty-state { display:flex; flex-direction:column; align-items:center; justify-content:center;
    height:100%; gap:12px; color:var(--text-secondary); }
  .strat-empty-state i { font-size:48px; opacity:.2; }
  .strat-empty-state p { font-size:14px; margin:0; }
  .strat-streaming-cursor { display:inline-block; width:2px; height:14px; background:#7a1028; margin-left:2px; animation:blink 1s step-end infinite; vertical-align:middle; }
  @keyframes blink { 0%,100%{opacity:1} 50%{opacity:0} }

  /* ── ANALYTICS VIEW ── */
  #view-analytics { padding:0; height:100%; overflow:auto; }
  #view-analytics.active { display:flex !important; flex-direction:column; }
  #view-analytics .view-header { padding:24px 28px 0; flex-shrink:0; }
  .analytics-body { padding:16px 28px 32px; flex:1; overflow:auto; }

  /* KPI Strip */
  .analytics-kpi-row { display:grid; grid-template-columns:repeat(4,1fr); gap:14px; margin-bottom:22px; }
  @media(max-width:900px){ .analytics-kpi-row { grid-template-columns:repeat(2,1fr); } }
  .analytics-kpi {
    background:#fff; border:1px solid var(--card-border); border-radius:14px;
    padding:18px 20px 16px; box-shadow:0 2px 8px rgba(122,16,40,0.05);
    position:relative; overflow:hidden; cursor:default;
    transition:transform .18s ease, box-shadow .18s ease;
  }
  .analytics-kpi:hover { transform:translateY(-2px); box-shadow:0 6px 20px rgba(122,16,40,0.1); }
  .analytics-kpi::before {
    content:''; position:absolute; top:0; left:0; right:0; height:3px;
    background:linear-gradient(90deg, var(--primary), var(--primary-l));
    border-radius:14px 14px 0 0;
  }
  .analytics-kpi-icon { font-size:18px; width:38px; height:38px; border-radius:10px; display:flex; align-items:center; justify-content:center; margin-bottom:10px; }
  .analytics-kpi:nth-child(1) .analytics-kpi-icon { background:rgba(122,16,40,0.1);  color:#7a1028; }
  .analytics-kpi:nth-child(2) .analytics-kpi-icon { background:rgba(37,99,235,0.1);  color:#2563eb; }
  .analytics-kpi:nth-child(3) .analytics-kpi-icon { background:rgba(234,88,12,0.1);  color:#ea580c; }
  .analytics-kpi:nth-child(4) .analytics-kpi-icon { background:rgba(22,163,74,0.1);  color:#16a34a; }
  .analytics-kpi:nth-child(1)::before { background:linear-gradient(90deg,#7a1028,#9b1a35); }
  .analytics-kpi:nth-child(2)::before { background:linear-gradient(90deg,#2563eb,#60a5fa); }
  .analytics-kpi:nth-child(3)::before { background:linear-gradient(90deg,#ea580c,#fb923c); }
  .analytics-kpi:nth-child(4)::before { background:linear-gradient(90deg,#16a34a,#4ade80); }
  .analytics-kpi:nth-child(1) .analytics-kpi-value { color:#7a1028; }
  .analytics-kpi:nth-child(2) .analytics-kpi-value { color:#2563eb; }
  .analytics-kpi:nth-child(3) .analytics-kpi-value { color:#ea580c; }
  .analytics-kpi:nth-child(4) .analytics-kpi-value { color:#16a34a; }
  .analytics-kpi-value {
    font-size:32px; font-weight:900; color:var(--primary); line-height:1;
    margin-bottom:4px; font-variant-numeric:tabular-nums;
    transition:all .4s ease;
  }
  .analytics-kpi-label { font-size:10px; font-weight:700; color:var(--text-dim); text-transform:uppercase; letter-spacing:.08em; }
  .analytics-kpi-delta { font-size:11px; color:#22c55e; font-weight:600; margin-top:4px; }

  /* Chart grid */
  .analytics-charts-row { display:grid; grid-template-columns:1fr 1fr; gap:16px; margin-bottom:16px; }
  @media(max-width:800px){ .analytics-charts-row { grid-template-columns:1fr; } }
  .analytics-chart-card {
    background:#fff; border:1px solid var(--card-border); border-radius:14px;
    padding:20px 22px; box-shadow:0 2px 8px rgba(122,16,40,0.05);
    transition:box-shadow .18s ease;
  }
  .analytics-chart-card:hover { box-shadow:0 6px 24px rgba(122,16,40,0.09); }
  .analytics-chart-header { display:flex; align-items:center; justify-content:space-between; margin-bottom:16px; }
  .analytics-chart-title { font-size:11px; font-weight:800; color:var(--text-dim); text-transform:uppercase; letter-spacing:.08em; }
  .analytics-chart-badge { font-size:10px; font-weight:700; padding:2px 8px; border-radius:20px; background:var(--primary-dim); color:var(--primary); }
  .analytics-chart-canvas-wrap { position:relative; width:100%; }

  /* Section divider label */
  .analytics-section-label {
    font-size:10px; font-weight:800; color:var(--text-dim); text-transform:uppercase;
    letter-spacing:.12em; padding:0 2px; margin:22px 0 12px;
    display:flex; align-items:center; gap:10px;
  }
  .analytics-section-label::after { content:''; flex:1; height:1px; background:var(--card-border); }

  /* Animated number count-up */
  @keyframes kpi-fade-in { from{opacity:0;transform:translateY(8px)} to{opacity:1;transform:none} }
  .analytics-kpi { animation:kpi-fade-in .35s ease both; }
  .analytics-kpi:nth-child(1){animation-delay:.05s}
  .analytics-kpi:nth-child(2){animation-delay:.10s}
  .analytics-kpi:nth-child(3){animation-delay:.15s}
  .analytics-kpi:nth-child(4){animation-delay:.20s}
  @keyframes card-rise { from{opacity:0;transform:translateY(12px)} to{opacity:1;transform:none} }
  .analytics-chart-card { animation:card-rise .4s ease both; }

  /* Donut centre stat */
  .donut-centre-wrap { position:relative; }
  .donut-centre-label {
    position:absolute; top:50%; left:50%; transform:translate(-50%,-50%);
    text-align:center; pointer-events:none;
  }
  .donut-centre-value { font-size:22px; font-weight:900; color:var(--primary); line-height:1; }
  .donut-centre-sub { font-size:9px; font-weight:700; color:var(--text-dim); text-transform:uppercase; letter-spacing:.06em; margin-top:3px; }

  /* ── STRATEGY STUDIO ── */
  #view-strategy { padding:0; display:none; height:100%; overflow:hidden; }
  #view-strategy.active { display:flex !important; }
  .ss-layout { display:flex; height:100%; width:100%; }
  .ss-form-panel {
    width:340px; min-width:280px; flex-shrink:0;
    border-right:1px solid var(--card-border);
    display:flex; flex-direction:column; overflow-y:auto; gap:14px;
    padding:22px 18px 24px; background:var(--bg);
  }
  .ss-form-panel::-webkit-scrollbar { width:4px; }
  .ss-form-panel::-webkit-scrollbar-thumb { background:var(--card-border); border-radius:2px; }
  .ss-output-panel { flex:1; min-width:0; display:flex; flex-direction:column; background:var(--bg-card); }
  .ss-output-header { padding:18px 24px 14px; border-bottom:1px solid var(--card-border); flex-shrink:0; display:flex; align-items:center; justify-content:space-between; }
  .ss-output-header h3 { font-size:13px; font-weight:600; color:var(--text-muted); margin:0; }
  .ss-output-body { flex:1; overflow-y:auto; padding:28px 32px; }
  .ss-output-body::-webkit-scrollbar { width:5px; }
  .ss-output-body::-webkit-scrollbar-thumb { background:var(--card-border); border-radius:3px; }
  /* Type pills */
  .ss-type-pills { display:flex; flex-wrap:wrap; gap:7px; }
  .ss-type-pill {
    padding:6px 13px; border-radius:20px; border:1.5px solid var(--primary-border);
    background:transparent; color:var(--text-muted); font-size:12px; font-weight:500;
    cursor:pointer; transition:all .15s; font-family:inherit;
  }
  .ss-type-pill:hover { border-color:var(--primary); color:var(--primary); }
  .ss-type-pill.ss-active { background:var(--primary); border-color:var(--primary); color:#fff; }
  /* Form fields */
  .ss-field { display:flex; flex-direction:column; gap:5px; }
  .ss-field label { font-size:10.5px; font-weight:700; letter-spacing:.05em; text-transform:uppercase; color:var(--text-dim); }
  .ss-field input, .ss-field textarea {
    width:100%; padding:8px 11px; border-radius:8px;
    border:1.5px solid var(--card-border); background:var(--bg);
    color:var(--text); font-size:12.5px; font-family:inherit;
    transition:border-color .15s; resize:vertical;
  }
  .ss-field input:focus, .ss-field textarea:focus { outline:none; border-color:var(--primary); box-shadow:0 0 0 3px var(--primary-dim); }
  .ss-field textarea { min-height:64px; }
  .ss-field input::placeholder, .ss-field textarea::placeholder { color:var(--text-dim); }
  /* Generate button */
  .ss-generate-btn {
    width:100%; padding:11px; border-radius:10px; background:var(--primary);
    color:#fff; font-size:13px; font-weight:600; border:none; cursor:pointer;
    display:flex; align-items:center; justify-content:center; gap:8px;
    transition:background .15s; margin-top:4px; font-family:inherit;
  }
  .ss-generate-btn:hover:not(:disabled) { background:var(--primary-d); }
  .ss-generate-btn:disabled { opacity:.6; cursor:not-allowed; }
  /* Empty state */
  .ss-empty-state {
    display:flex; flex-direction:column; align-items:center; justify-content:center;
    height:100%; gap:14px; color:var(--text-dim); text-align:center; padding:40px;
  }
  .ss-empty-state i { font-size:44px; color:var(--card-border); }
  .ss-empty-state p { font-size:15px; font-weight:600; color:var(--text-muted); margin:0; }
  .ss-empty-state span { font-size:12px; color:var(--text-dim); }
  /* Streaming cursor */
  .ss-cursor { display:inline-block; width:2px; height:1em; background:var(--primary); border-radius:1px; animation:ss-blink .7s steps(1) infinite; vertical-align:text-bottom; margin-left:2px; }
  @keyframes ss-blink { 0%,100%{opacity:1} 50%{opacity:0} }
  /* Rendered output */
  .ss-doc-output { font-size:13.5px; line-height:1.75; color:var(--text); max-width:780px; }
  .ss-doc-output h1 { font-size:21px; font-weight:800; color:var(--primary); margin:0 0 18px; }
  .ss-doc-output h2 { font-size:15px; font-weight:700; color:var(--text); margin:22px 0 8px; border-bottom:1.5px solid var(--card-border); padding-bottom:5px; }
  .ss-doc-output h3 { font-size:13px; font-weight:700; color:var(--primary-l); margin:14px 0 5px; }
  .ss-doc-output p { margin:0 0 10px; }
  .ss-doc-output ul, .ss-doc-output ol { padding-left:22px; margin:0 0 10px; }
  .ss-doc-output li { margin-bottom:4px; }
  .ss-doc-output strong { color:var(--text); }
  .ss-doc-output table { width:100%; border-collapse:collapse; margin:10px 0 14px; font-size:12.5px; }
  .ss-doc-output th { background:var(--primary-dim); color:var(--primary); padding:7px 12px; text-align:left; font-weight:700; border:1px solid var(--card-border); }
  .ss-doc-output td { padding:6px 12px; border:1px solid var(--card-border); }
  .ss-doc-output tr:nth-child(even) td { background:var(--primary-dim); }
  /* Output actions bar */
  .ss-output-actions { padding:12px 24px; border-top:1px solid var(--card-border); display:flex; gap:8px; flex-shrink:0; flex-wrap:wrap; }
  .ss-action-btn {
    padding:7px 14px; border-radius:8px; font-size:12px; font-weight:500;
    border:1.5px solid var(--card-border); background:transparent; color:var(--text-muted);
    cursor:pointer; display:flex; align-items:center; gap:6px; transition:all .15s; font-family:inherit;
  }
  .ss-action-btn:hover { border-color:var(--primary); color:var(--primary); }
  .ss-action-btn.ss-primary { background:var(--primary); border-color:var(--primary); color:#fff; }
  .ss-action-btn.ss-primary:hover { background:var(--primary-d); }
  /* generating indicator */
  .ss-generating { display:flex; align-items:center; gap:10px; color:var(--text-muted); font-size:13px; padding:0 0 16px; }
  .ss-gen-dots span { display:inline-block; width:6px; height:6px; border-radius:50%; background:var(--primary); opacity:.3; animation:ss-dot 1.2s ease-in-out infinite; margin:0 2px; }
  .ss-gen-dots span:nth-child(2) { animation-delay:.2s; }
  .ss-gen-dots span:nth-child(3) { animation-delay:.4s; }
  @keyframes ss-dot { 0%,80%,100%{opacity:.3;transform:scale(1)} 40%{opacity:1;transform:scale(1.3)} }
  /* dark mode */
  body.dark .ss-form-panel { background:#111118; border-right-color:rgba(255,255,255,.07); }
  body.dark .ss-output-panel { background:#0f0f13; }
  body.dark .ss-output-header { border-bottom-color:rgba(255,255,255,.07); }
  body.dark .ss-field input, body.dark .ss-field textarea { background:#1a1a24; border-color:rgba(255,255,255,.1); color:#e8e8f0; }
  body.dark .ss-field input:focus, body.dark .ss-field textarea:focus { border-color:var(--primary-l); box-shadow:0 0 0 3px rgba(155,26,53,.15); }
  body.dark .ss-field input::placeholder, body.dark .ss-field textarea::placeholder { color:#484860; }
  body.dark .ss-type-pill { border-color:rgba(255,255,255,.12); color:#9090a8; }
  body.dark .ss-output-actions { border-top-color:rgba(255,255,255,.07); }
  body.dark .ss-action-btn { border-color:rgba(255,255,255,.1); color:#9090a8; }
  body.dark .ss-action-btn:hover { border-color:var(--primary-l); color:var(--primary-l); }
  body.dark .ss-doc-output { color:#e0e0ee; }
  body.dark .ss-doc-output h2 { border-bottom-color:rgba(255,255,255,.07); color:#e0e0ee; }
  body.dark .ss-doc-output th { background:rgba(122,16,40,.2); }
  body.dark .ss-doc-output td { border-color:rgba(255,255,255,.07); }
  body.dark .ss-doc-output tr:nth-child(even) td { background:rgba(122,16,40,.08); }
</style>
</head>
<body>

<!-- ═══════════════════════════════════════ SIDEBAR ═══════════════════════════════════════ -->
<nav id="sidebar">
  <div class="sidebar-brand">
    <img src="/assets/Enderun-Colleges-Logo-Normal-600x291.png" alt="Enderun Colleges" class="sidebar-logo">
    <div class="brand-sub">Marketing Hub &middot; AI-Powered</div>
  </div>

  <div class="sidebar-scroll">
    <div class="sidebar-section-label">Navigation</div>
    <div class="nav-item active" onclick="showView('dashboard', this)" data-view="dashboard">
      <span class="nav-icon"><i class="fa-solid fa-comments"></i></span>
      <span class="nav-label">AI Chat</span>
    </div>
    <div class="nav-item" onclick="showView('leads', this)" data-view="leads">
      <span class="nav-icon"><i class="fa-solid fa-users"></i></span>
      <span class="nav-label">Leads</span>
    </div>
    <div class="nav-item" onclick="showView('calendar', this)" data-view="calendar">
      <span class="nav-icon"><i class="fa-solid fa-calendar"></i></span>
      <span class="nav-label">Calendar</span>
    </div>

    <div class="divider"></div>
    <div class="sidebar-section-label">Intelligence</div>
    <div class="nav-item" onclick="showView('intelligence', this)" data-view="intelligence">
      <span class="nav-icon"><i class="fa-solid fa-magnifying-glass"></i></span>
      <span class="nav-label">Competitor Monitor</span>
    </div>
    <div class="nav-item" onclick="showView('automations', this)" data-view="automations">
      <span class="nav-icon"><i class="fa-solid fa-gears"></i></span>
      <span class="nav-label">Automations</span>
    </div>
    <div class="nav-item" onclick="showView('analytics', this)" data-view="analytics">
      <span class="nav-icon"><i class="fa-solid fa-chart-line"></i></span>
      <span class="nav-label">Analytics</span>
    </div>

    <div class="divider"></div>
    <div class="sidebar-section-label">Strategy</div>
    <div class="nav-item" onclick="showView('strategy', this)" data-view="strategy">
      <span class="nav-icon"><i class="fa-solid fa-lightbulb"></i></span>
      <span class="nav-label">Strategy Studio</span>
    </div>

  </div>

  <div class="sidebar-status">
    <div class="status-dot"></div>
    <div class="status-text">All systems operational</div>
  </div>
  <button class="logout-btn" onclick="window.location='/logout'"><i class="fa-solid fa-right-from-bracket"></i> <span>Sign Out</span></button>
</nav>

<!-- ═══════════════════════════════════════ MAIN ═══════════════════════════════════════ -->
<div id="main">
  <!-- TOPBAR -->
  <div id="topbar">
    <div style="flex:1;min-width:0">
      <div class="topbar-title" id="topbar-title">AI Chat</div>
      <div class="topbar-date" id="topbar-date"><i class="fa-regular fa-calendar"></i> Loading...</div>
    </div>
    <div class="topbar-actions">
      <!-- Automate dropdown -->
      <div class="tb-dropdown" id="tb-automate-dd">
        <button class="tb-automate-btn" onclick="toggleTbDd('tb-automate-dd')" title="Automations">
          <i class="fa-solid fa-bolt"></i> Automate <i class="fa-solid fa-chevron-down" style="font-size:9px;opacity:.7"></i>
        </button>
        <div class="tb-dd-menu">
          <div class="tb-dd-item" onclick="runPost();closeTbDd()"><i class="fa-brands fa-facebook-f" style="color:#1877f2"></i> Post to Facebook</div>
          <div class="tb-dd-item" onclick="openIgModal();closeTbDd()"><i class="fa-brands fa-instagram" style="color:#e1306c"></i> Post to Instagram</div>
          <div class="tb-dd-item" onclick="runEmail();closeTbDd()"><i class="fa-solid fa-envelope" style="color:#7a1028"></i> Send Drip Emails</div>
          <div class="tb-dd-divider"></div>
          <div class="tb-dd-item" onclick="openPostPreview();closeTbDd()"><i class="fa-regular fa-eye" style="color:#64748b"></i> Preview Today's Post</div>
          <div class="tb-dd-item" onclick="runSocialListening();closeTbDd()"><i class="fa-solid fa-satellite-dish" style="color:#64748b"></i> Run Social Listening</div>
          <div class="tb-dd-item" onclick="runWeeklyPreview();closeTbDd()"><i class="fa-solid fa-calendar-week" style="color:#64748b"></i> Weekly Campaign Preview</div>
        </div>
      </div>
      <div style="width:1px;height:24px;background:var(--card-border);flex-shrink:0"></div>
      <button class="tb-icon-btn notif-bell" id="notif-bell" onclick="toggleDrawer()" title="Notifications">
        <i class="fa-solid fa-bell"></i>
        <span class="notif-badge" id="notif-badge" style="display:none">0</span>
      </button>
      <button class="tb-icon-btn" id="dark-toggle" onclick="toggleDark()" title="Toggle dark mode"><i class="fa-solid fa-moon"></i></button>
      <button class="tb-icon-btn" onclick="showShortcuts()" title="Keyboard shortcuts (?)"><i class="fa-solid fa-keyboard"></i></button>
      <button class="tb-icon-btn" onclick="window.location='/profile'" title="Profile Settings"><i class="fa-solid fa-user-gear"></i></button>
      <button class="tb-icon-btn" onclick="window.location='/logout'" title="Sign Out"><i class="fa-solid fa-right-from-bracket"></i></button>
    </div>
  </div>

  <!-- NOTIFICATION DRAWER -->
  <div id="notif-overlay" onclick="closeDrawer()"></div>
  <div id="notif-drawer">
    <div class="nd-header">
      <div>
        <div class="nd-title">Notifications</div>
        <div class="nd-subtitle" id="nd-subtitle">All agents</div>
      </div>
      <div style="display:flex;gap:8px;align-items:center">
        <button class="nd-action-btn" onclick="markAllRead()">Mark all read</button>
        <button class="nd-action-btn nd-clear" onclick="clearAll()">Clear</button>
        <button class="nd-close" onclick="closeDrawer()">✕</button>
      </div>
    </div>
    <div class="nd-body" id="nd-body">
      <div class="nd-empty">No notifications yet</div>
    </div>
  </div>

  <!-- CONTENT -->
  <div id="content">

    <!-- ── DASHBOARD ── -->
    <div class="view active" id="view-dashboard" style="padding:0;height:100%;overflow:hidden;">
      <!-- Left sessions sidebar -->
      <div class="chat-sidebar" id="chat-sidebar">
        <div class="sidebar-drag-handle" id="sidebar-drag-handle"></div>
        <div class="chat-sidebar-top">
          <button class="new-chat-btn" onclick="newChat()">
            <i class="fa-solid fa-plus"></i> New Chat
          </button>
          <div class="custom-dd" id="agent-dd">
            <div class="custom-dd-trigger" onclick="toggleDd('agent-dd')">
              <span id="agent-dd-label" style="flex:1;min-width:0;overflow:hidden;text-overflow:ellipsis;white-space:nowrap">— Select Agent —</span>
              <i class="fa-solid fa-chevron-down custom-dd-arrow"></i>
            </div>
            <div class="custom-dd-menu" id="agent-dd-menu">
              <div style="padding:6px 8px;border-bottom:1px solid var(--card-border);position:sticky;top:0;background:#fff;z-index:1">
                <input id="agent-search" type="text" placeholder="Search agents…" autocomplete="off"
                  style="width:100%;padding:5px 10px;border:1px solid var(--card-border);border-radius:6px;font-size:12px;font-family:inherit;color:var(--text);outline:none"
                  oninput="filterAgents(this.value)" onclick="event.stopPropagation()">
              </div>
            </div>
          </div>
        </div>
        <div class="chat-sessions-list" id="chat-sessions-list">
          <div style="padding:24px 12px;text-align:center;font-size:11px;color:var(--text-dim)">
            Select an agent<br>to start chatting
          </div>
        </div>
      </div>
      <!-- Main chat area -->
      <div class="chat-main chat-column">
        <div id="drag-overlay">
          <div class="drag-icon">📎</div>
          <div class="drag-label">Drop image, PDF, or text file</div>
        </div>
        <div class="chat-header">
          <div class="agent-avatar" id="chat-avatar"><img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun"></div>
          <div style="min-width:0;flex:1">
            <div class="chat-agent-name" id="chat-name">Agent Chat</div>
            <div class="chat-agent-role" id="chat-role">Select an agent from the left panel to begin</div>
          </div>
          <div class="chat-header-actions">
            <div class="chat-mode-bar">
              <div class="model-dd custom-dd" id="model-dd">
                <div class="custom-dd-trigger" onclick="toggleDd('model-dd')">
                  <span id="model-dd-label">Opus 4.7</span>
                  <i class="fa-solid fa-chevron-down custom-dd-arrow"></i>
                </div>
                <div class="custom-dd-menu">
                  <div class="custom-dd-opt selected" data-val="claude-opus-4-7" onclick="pickModel('claude-opus-4-7','Opus 4.7')">Opus 4.7</div>
                  <div class="custom-dd-opt" data-val="claude-sonnet-4-6" onclick="pickModel('claude-sonnet-4-6','Sonnet 4.6')">Sonnet 4.6</div>
                  <div class="custom-dd-opt" data-val="claude-haiku-4-5-20251001" onclick="pickModel('claude-haiku-4-5-20251001','Haiku 4.5')">Haiku 4.5</div>
                </div>
              </div>
              <div class="chat-mode-pills">
                <button class="mode-pill active" data-mode="normal"   onclick="setChatMode('normal')">Normal</button>
                <button class="mode-pill"        data-mode="deep"     onclick="setChatMode('deep')">Deep</button>
                <button class="mode-pill"        data-mode="concise"  onclick="setChatMode('concise')">Concise</button>
                <button class="mode-pill"        data-mode="creative" onclick="setChatMode('creative')">Creative</button>
              </div>
            </div>
            <button class="btn btn-ghost btn-sm" onclick="clearCurrentChat()" title="New conversation"><i class="fa-solid fa-rotate-right"></i></button>
            <button class="chat-hdr-btn" onclick="toggleChatSearch()" id="search-toggle-btn" title="Search in chat"><i class="fa-solid fa-magnifying-glass"></i></button>
            <span class="msg-counter" id="msg-counter" style="display:none" title="Messages in conversation"></span>
            <button class="chat-hdr-btn" onclick="exportChat()" title="Export as TXT"><i class="fa-solid fa-file-lines"></i></button>
            <button class="chat-hdr-btn" onclick="exportChatPDF()" title="Export as PDF"><i class="fa-solid fa-file-pdf"></i></button>
          </div>
        </div>
        <div class="chat-search-bar" id="chat-search-bar">
          <i class="fa-solid fa-magnifying-glass" style="color:var(--text-dim);font-size:12px;flex-shrink:0"></i>
          <input type="text" class="chat-search-input" id="chat-search-input" placeholder="Search in conversation…" oninput="searchInChat(this.value)" onkeydown="if(event.key==='Escape')toggleChatSearch()">
          <span class="search-match-count" id="search-match-count"></span>
          <button style="background:none;border:none;color:var(--text-dim);cursor:pointer;font-size:14px;flex-shrink:0" onclick="toggleChatSearch()"><i class="fa-solid fa-xmark"></i></button>
        </div>
        <div class="chat-messages" id="chat-messages" style="position:relative">
          <button class="scroll-pause-btn" id="scroll-pause-btn" onclick="resumeScroll()" title="Resume auto-scroll"><i class="fa-solid fa-arrow-down"></i> New content</button>
          <div class="welcome-box" id="chat-welcome">
            <div class="welcome-hero">
              <div class="welcome-logo-ring">
                <img src="/assets/logos/Enderun-Colleges.png" class="welcome-logo" alt="Enderun">
              </div>
              <div class="welcome-greeting" id="welcome-greeting">Good morning</div>
              <div class="welcome-title">How can I help you today?</div>
              <div class="welcome-sub">27 AI agents trained on Enderun's brand, strategy &amp; data — powered by Claude Opus 4.7</div>
            </div>
            <div class="welcome-cards">
              <button class="wcard" onclick="usePrompt('Write a Facebook post for Enderun Extension showcasing our world-class culinary programs.')">
                <span class="wcard-icon"><i class="fa-brands fa-facebook-f"></i></span>
                <span class="wcard-text"><span class="wcard-title">Social Post</span><span class="wcard-desc">Write a branded Facebook or Instagram post</span></span>
              </button>
              <button class="wcard" onclick="usePrompt('Analyze my current leads and suggest follow-up strategies for each program.')">
                <span class="wcard-icon"><i class="fa-solid fa-chart-bar"></i></span>
                <span class="wcard-text"><span class="wcard-title">Analyze Leads</span><span class="wcard-desc">Get insights and follow-up strategies</span></span>
              </button>
              <button class="wcard" onclick="usePrompt('Check for competitor news and recent social media activity. Give me a full brief.')">
                <span class="wcard-icon"><i class="fa-solid fa-magnifying-glass"></i></span>
                <span class="wcard-text"><span class="wcard-title">Competitor Intel</span><span class="wcard-desc">Monitor DLSU, CCA, ISCAHM and more</span></span>
              </button>
              <button class="wcard" onclick="usePrompt('Generate a designed PDF report on our current marketing performance and lead funnel.')">
                <span class="wcard-icon"><i class="fa-solid fa-file-pdf"></i></span>
                <span class="wcard-text"><span class="wcard-title">Generate Report</span><span class="wcard-desc">Branded PDF with charts and KPI cards</span></span>
              </button>
              <button class="wcard" onclick="usePrompt('What are the best enrollment season marketing strategies for Enderun? Give me a full campaign plan.')">
                <span class="wcard-icon"><i class="fa-solid fa-lightbulb"></i></span>
                <span class="wcard-text"><span class="wcard-title">Campaign Ideas</span><span class="wcard-desc">Enrollment season strategies and tactics</span></span>
              </button>
              <button class="wcard" onclick="usePrompt('Write a personalized drip email for a new hospitality management lead.')">
                <span class="wcard-icon"><i class="fa-solid fa-envelope"></i></span>
                <span class="wcard-text"><span class="wcard-title">Drip Email</span><span class="wcard-desc">Personalized email for a specific lead</span></span>
              </button>
            </div>
          </div>
        </div>
        <div id="quick-chips" style="display:none;padding:6px 14px 2px;overflow-x:auto;white-space:nowrap;scrollbar-width:none">
          <button class="q-chip" onclick="usePrompt('Show me a summary of current leads.')">📊 Leads</button>
          <button class="q-chip" onclick="usePrompt('What images are scheduled for posting this week?')">📅 Schedule</button>
          <button class="q-chip" onclick="sendDirectMsg('Post to Facebook and Instagram')">🖼️ Post FB &amp; IG</button>
          <button class="q-chip" onclick="usePrompt('Check for competitor news and recent social media activity.')">🔍 Competitors</button>
          <button class="q-chip" onclick="usePrompt('Run the weekly analytics report.')">📈 Analytics</button>
          <button class="q-chip" onclick="usePrompt('Show me all my saved memories.')">🧠 Memories</button>
          <button class="q-chip" onclick="usePrompt('Draft a drip email for a new hospitality management lead.')">📧 Drip email</button>
          <button class="q-chip" onclick="usePrompt('Generate content ideas for this week\'s social media.')">💡 Content ideas</button>
        </div>
        <div class="chat-input-area" style="position:relative">
          <!-- @mention dropdown -->
          <div class="mention-dropdown" id="mention-dropdown"></div>
          <!-- URL preview chip -->
          <div class="url-preview-chip" id="url-preview-chip">
            <i class="fa-solid fa-link" style="color:var(--primary);font-size:12px;flex-shrink:0"></i>
            <a id="url-preview-link" href="#" target="_blank" rel="noopener"></a>
            <span onclick="dismissUrlPreview()" style="cursor:pointer;color:var(--text-dim);font-size:14px;flex-shrink:0"><i class="fa-solid fa-xmark"></i></span>
          </div>
          <!-- Multi-file queue -->
          <div class="file-queue" id="file-queue"></div>
          <div id="attach-preview" style="display:none;margin-bottom:8px;padding:8px 12px;background:var(--gold-dim);border:1px solid var(--gold-border);border-radius:8px;font-size:12px;color:var(--text);align-items:center;gap:8px">
            <span style="font-size:16px" id="attach-icon">🖼</span>
            <span id="attach-label" style="flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap"></span>
            <span onclick="clearAttachment()" style="cursor:pointer;color:var(--text-dim);font-size:16px;line-height:1">×</span>
          </div>
          <div class="chat-input-wrap">
            <textarea id="chat-input" rows="1" placeholder="Ask your agent anything… (@ to mention agent)" onkeydown="handleKey(event);_mentionKeydown(event)" oninput="autoResize(this);updateCharCount(this);_mentionOninput(this);_detectUrl(this)"></textarea>
            <div class="chat-input-footer">
              <button class="attach-btn" onclick="document.getElementById('file-input').click()" title="Attach file"><i class="fa-solid fa-paperclip"></i></button>
              <button class="attach-btn" onclick="promptImageUrl()" title="Paste image URL"><i class="fa-solid fa-link"></i></button>
              <button class="attach-btn" id="voice-btn" onclick="toggleVoiceInput()" title="Voice input"><i class="fa-solid fa-microphone"></i></button>
              <span class="input-hint" style="text-align:center">↵ Send &nbsp;·&nbsp; ⇧↵ New line</span>
              <span class="char-counter" id="char-counter"></span>
              <button id="stop-btn" onclick="stopStream()" title="Stop generation"><i class="fa-solid fa-stop"></i></button>
              <button id="send-btn" onclick="sendMsg()" title="Send"><i class="fa-solid fa-paper-plane"></i></button>
            </div>
          </div>
          <input type="file" id="file-input" accept="image/*,.pdf,.csv,.txt" style="display:none" multiple onchange="handleFileQueue(event)">
        </div>
      </div>
    </div>

    <!-- ── LEADS ── -->
    <div class="view" id="view-leads">
      <div class="view-header">
        <h2><i class="fa-solid fa-users"></i> Leads</h2>
        <div style="display:flex;gap:10px;align-items:center">
          <button class="btn btn-ghost" onclick="exportLeadsCSV()"><i class="fa-solid fa-download"></i> Export CSV</button>
          <button class="btn btn-gold" onclick="toggleForm()"><i class="fa-solid fa-plus"></i> Add Lead</button>
        </div>
      </div>

      <!-- Lead Stats -->
      <div class="stats-grid" id="stats-grid" style="margin-bottom:12px">
        <div class="stat-card">
          <span class="stat-icon"><i class="fa-solid fa-users"></i></span>
          <div class="stat-number" id="stat-total">—</div>
          <div class="stat-label">Total Leads</div>
          <div class="stat-sub" id="stat-total-sub">All registered</div>
        </div>
        <div class="stat-card">
          <span class="stat-icon"><i class="fa-solid fa-bolt"></i></span>
          <div class="stat-number" id="stat-active">—</div>
          <div class="stat-label">Active Leads</div>
          <div class="stat-sub" id="stat-active-sub">Receiving emails</div>
        </div>
        <div class="stat-card">
          <span class="stat-icon" style="color:#F87171"><i class="fa-solid fa-fire"></i></span>
          <div class="stat-number" id="stat-hot" style="color:#F87171">—</div>
          <div class="stat-label">Hot Leads</div>
          <div class="stat-sub">10+ emails sent</div>
        </div>
        <div class="stat-card">
          <span class="stat-icon" style="color:#FCD34D"><i class="fa-solid fa-temperature-half"></i></span>
          <div class="stat-number" id="stat-warm" style="color:#FCD34D">—</div>
          <div class="stat-label">Warm Leads</div>
          <div class="stat-sub">5–9 emails sent</div>
        </div>
        <div class="stat-card">
          <span class="stat-icon" style="color:#93C5FD"><i class="fa-solid fa-snowflake"></i></span>
          <div class="stat-number" id="stat-cold" style="color:#93C5FD">—</div>
          <div class="stat-label">Cold Leads</div>
          <div class="stat-sub">Under 5 emails</div>
        </div>
      </div>

      <div class="search-wrap">
        <i class="fa-solid fa-magnifying-glass search-icon"></i>
        <input type="text" id="leads-search" placeholder="Search by name, email, or program..." oninput="filterLeads(this.value)">
      </div>

      <div id="add-form">
        <div class="card-title"><i class="fa-solid fa-user-plus"></i> New Lead</div>
        <div class="form-grid">
          <div class="form-group">
            <label class="form-label">First Name</label>
            <input class="form-input" id="f-fname" placeholder="e.g. Maria" />
          </div>
          <div class="form-group">
            <label class="form-label">Last Name</label>
            <input class="form-input" id="f-lname" placeholder="e.g. Santos" />
          </div>
          <div class="form-group">
            <label class="form-label">Email</label>
            <input class="form-input" id="f-email" type="email" placeholder="email@example.com" />
          </div>
          <div class="form-group">
            <label class="form-label">Program</label>
            <input class="form-input" id="f-program" placeholder="e.g. BS Hospitality Management" />
          </div>
        </div>
        <div class="form-actions">
          <button class="btn btn-gold btn-sm" onclick="addLead()"><i class="fa-solid fa-floppy-disk"></i> Save Lead</button>
          <button class="btn btn-outline btn-sm" onclick="toggleForm()">Cancel</button>
        </div>
      </div>

      <div class="card">
        <div class="table-wrap">
          <table id="leads-table">
            <thead>
              <tr>
                <th class="sortable" id="th-name" onclick="sortLeads('name')">Name <i class="fa-solid fa-sort sort-icon"></i></th>
                <th>Email</th>
                <th class="sortable" id="th-program" onclick="sortLeads('program')">Program <i class="fa-solid fa-sort sort-icon"></i></th>
                <th class="sortable" id="th-emails" onclick="sortLeads('emails')" style="text-align:center">Emails Sent <i class="fa-solid fa-sort sort-icon"></i></th>
                <th class="sortable" id="th-score" onclick="sortLeads('score')">Score <i class="fa-solid fa-sort sort-icon"></i></th>
                <th class="sortable" id="th-status" onclick="sortLeads('status')">Status <i class="fa-solid fa-sort sort-icon"></i></th>
              </tr>
            </thead>
            <tbody id="leads-tbody">
              <tr><td colspan="6" style="text-align:center;color:var(--text-dim);padding:30px">Loading leads...</td></tr>
            </tbody>
          </table>
        </div>
        <div class="table-pagination" id="leads-pagination"></div>
      </div>
    </div>


    <!-- ── CALENDAR ── -->
    <div class="view" id="view-calendar">
      <!-- Persistent toolbar: nav + tabs -->
      <div class="cal-toolbar">
        <div class="cal-month-nav" id="cal-month-nav">
          <button class="cal-btn" onclick="calPrev()">‹</button>
          <span id="cal-month-label">Loading...</span>
          <button class="cal-btn" onclick="calNext()">›</button>
        </div>
        <div class="cal-tabs">
          <button class="cal-tab active" id="tab-month" onclick="calSetTab('month')">
            <i class="fa-solid fa-calendar-days"></i> Month
          </button>
          <button class="cal-tab" id="tab-list" onclick="calSetTab('list')">
            <i class="fa-solid fa-list"></i> Schedule List
          </button>
        </div>
      </div>
      <!-- Month view -->
      <div id="cal-month-view">
        <div class="cal-grid" id="cal-day-headers"></div>
        <div class="cal-grid" id="cal-grid"></div>
      </div>
      <!-- List view -->
      <div id="cal-list-view" style="display:none">
        <div class="table-wrap">
          <table id="schedule-table">
            <thead><tr><th>Date</th><th>Day</th><th>Image File</th><th>Status</th></tr></thead>
            <tbody id="schedule-tbody">
              <tr><td colspan="4" style="text-align:center;color:var(--text-dim);padding:30px">Loading schedule...</td></tr>
            </tbody>
          </table>
        </div>
      </div>
    </div>

    <!-- ── INTELLIGENCE ── -->
    <div class="view" id="view-intelligence">
      <div class="view-header" style="flex-shrink:0">
        <h2><i class="fa-solid fa-magnifying-glass-chart"></i> Competitor Intelligence</h2>
      </div>
      <div class="intel-layout">
        <div class="intel-sidebar" id="intel-date-list">
          <div class="intel-list-header">Reports</div>
        </div>
        <div class="intel-main" id="intel-content">
          <div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading briefing...</p></div>
        </div>
      </div>
    </div>

    <!-- ── AUTOMATIONS ── -->
    <div class="view" id="view-automations">
      <div class="view-header">
        <h2><i class="fa-solid fa-gears"></i> Automations</h2>
      </div>
      <div id="auto-grid" class="auto-grid">
        <div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading...</p></div>
      </div>
    </div>

    <!-- ── ANALYTICS ── -->
    <div class="view" id="view-analytics">
      <div class="view-header"><h2><i class="fa-solid fa-chart-line"></i> Analytics</h2></div>
      <div class="analytics-body">

        <!-- KPI Strip -->
        <div class="analytics-kpi-row" id="analytics-kpi-row">
          <div class="analytics-kpi">
            <div class="analytics-kpi-icon"><i class="fa-solid fa-calendar-days"></i></div>
            <div class="analytics-kpi-value" id="akpi-total">—</div>
            <div class="analytics-kpi-label">Total Posts</div>
          </div>
          <div class="analytics-kpi">
            <div class="analytics-kpi-icon"><i class="fa-solid fa-calendar-check"></i></div>
            <div class="analytics-kpi-value" id="akpi-active">—</div>
            <div class="analytics-kpi-label">This Month</div>
          </div>
          <div class="analytics-kpi">
            <div class="analytics-kpi-icon"><i class="fa-solid fa-rocket"></i></div>
            <div class="analytics-kpi-value" id="akpi-hot">—</div>
            <div class="analytics-kpi-label">Next 7 Days</div>
          </div>
          <div class="analytics-kpi">
            <div class="analytics-kpi-icon"><i class="fa-solid fa-paper-plane"></i></div>
            <div class="analytics-kpi-value" id="akpi-posts">—</div>
            <div class="analytics-kpi-label">Drip Emails Sent</div>
          </div>
        </div>

        <!-- Row 1: Monthly volume + 14-day calendar -->
        <div class="analytics-section-label">Post Performance</div>
        <div class="analytics-charts-row">
          <div class="analytics-chart-card" style="animation-delay:.1s">
            <div class="analytics-chart-header">
              <div class="analytics-chart-title">Monthly Post Volume</div>
              <span class="analytics-chart-badge" id="badge-monthly">Loading…</span>
            </div>
            <div class="analytics-chart-canvas-wrap" style="height:210px"><canvas id="chart-monthly-posts"></canvas></div>
          </div>
          <div class="analytics-chart-card" style="animation-delay:.18s">
            <div class="analytics-chart-header">
              <div class="analytics-chart-title">Upcoming Posts — Next 14 Days</div>
              <span class="analytics-chart-badge" id="badge-upcoming">—</span>
            </div>
            <div class="analytics-chart-canvas-wrap" style="height:210px"><canvas id="chart-upcoming-posts"></canvas></div>
          </div>
        </div>

        <!-- Row 2: Drip bar + Drip donut -->
        <div class="analytics-section-label">Email Campaign Engagement</div>
        <div class="analytics-charts-row" style="grid-template-columns:1.6fr 1fr">
          <div class="analytics-chart-card" style="animation-delay:.26s">
            <div class="analytics-chart-header">
              <div class="analytics-chart-title">Top Leads by Emails Received</div>
              <span class="analytics-chart-badge" id="badge-drip-total">—</span>
            </div>
            <div class="analytics-chart-canvas-wrap" style="height:220px"><canvas id="chart-drip-engagement"></canvas></div>
          </div>
          <div class="analytics-chart-card" style="animation-delay:.34s">
            <div class="analytics-chart-header">
              <div class="analytics-chart-title">Campaign Reach</div>
              <span class="analytics-chart-badge" id="badge-reach">—</span>
            </div>
            <div class="donut-centre-wrap">
              <div class="analytics-chart-canvas-wrap" style="height:220px"><canvas id="chart-drip-donut"></canvas></div>
              <div class="donut-centre-label">
                <div class="donut-centre-value" id="donut-centre-val">—</div>
                <div class="donut-centre-sub">Leads<br>reached</div>
              </div>
            </div>
          </div>
        </div>

      </div>
    </div>

    <!-- ── STRATEGY STUDIO ── -->
    <div class="view" id="view-strategy">
      <div class="ss-layout">

        <!-- Form Panel -->
        <div class="ss-form-panel">
          <div>
            <h2 style="font-size:15px;font-weight:700;color:var(--text);margin:0 0 3px;display:flex;align-items:center;gap:8px">
              <i class="fa-solid fa-lightbulb" style="color:var(--primary);font-size:14px"></i> Strategy Studio
            </h2>
            <p style="font-size:11.5px;color:var(--text-dim);margin:0">AI-generated strategy documents for Enderun</p>
          </div>

          <div class="ss-field">
            <label>Document Type</label>
            <div class="ss-type-pills">
              <button class="ss-type-pill ss-active" onclick="ssSelectType('business-plan', this)">Business Plan</button>
              <button class="ss-type-pill" onclick="ssSelectType('marketing-strategy', this)">Marketing Strategy</button>
              <button class="ss-type-pill" onclick="ssSelectType('strategic-plan', this)">Strategic Plan</button>
            </div>
          </div>

          <div class="ss-field">
            <label>Program / Initiative Name <span style="color:#c62a47">*</span></label>
            <input type="text" id="ss-program" placeholder="e.g. NCLEX Review Center">
          </div>

          <div class="ss-field">
            <label>Brief Description <span style="color:#c62a47">*</span></label>
            <textarea id="ss-desc" placeholder="What is this program about? What problem does it solve?"></textarea>
          </div>

          <div class="ss-field">
            <label>Target Market</label>
            <input type="text" id="ss-market" placeholder="e.g. Filipino nurses, fresh grads, OFW-bound">
          </div>

          <div class="ss-field">
            <label>Formats / Delivery</label>
            <input type="text" id="ss-formats" placeholder="e.g. Online, Face-to-Face, Hybrid">
          </div>

          <div class="ss-field">
            <label>Budget Range <span style="font-size:10px;color:var(--text-dim);text-transform:none;letter-spacing:0">(optional)</span></label>
            <input type="text" id="ss-budget" placeholder="e.g. PHP 1M–2M startup">
          </div>

          <div class="ss-field">
            <label>Timeline</label>
            <input type="text" id="ss-timeline" placeholder="e.g. Launch by Q3 2026">
          </div>

          <div class="ss-field">
            <label>Additional Notes</label>
            <textarea id="ss-notes" placeholder="Anything else relevant — competitors, unique angle, constraints..."></textarea>
          </div>

          <button class="ss-generate-btn" id="ss-generate-btn" onclick="ssGenerate()">
            <i class="fa-solid fa-bolt"></i> Generate
          </button>
        </div>

        <!-- Output Panel -->
        <div class="ss-output-panel">
          <div class="ss-output-header">
            <h3 id="ss-output-title">Output will appear here</h3>
            <div id="ss-word-count" style="font-size:11px;color:var(--text-dim)"></div>
          </div>
          <div class="ss-output-body" id="ss-output-body">
            <div class="ss-empty-state" id="ss-empty-state">
              <i class="fa-regular fa-lightbulb"></i>
              <p>Fill in the form and hit <strong>Generate</strong></p>
              <span>Your strategy document will stream here in real time.</span>
            </div>
            <div id="ss-stream-area" style="display:none">
              <div class="ss-generating" id="ss-generating-indicator">
                <div class="ss-gen-dots"><span></span><span></span><span></span></div>
                <span>Generating your strategy document…</span>
              </div>
              <div class="ss-doc-output" id="ss-doc-output"></div>
            </div>
          </div>
          <div class="ss-output-actions" id="ss-output-actions" style="display:none">
            <button class="ss-action-btn" onclick="ssCopy()"><i class="fa-regular fa-copy"></i> Copy</button>
            <button class="ss-action-btn" onclick="ssExportPDF()"><i class="fa-regular fa-file-pdf"></i> Export PDF</button>
            <button class="ss-action-btn" onclick="ssOpenInChat()"><i class="fa-solid fa-comments"></i> Open in Chat</button>
            <button class="ss-action-btn ss-primary" onclick="ssGenerate()"><i class="fa-solid fa-rotate-right"></i> Regenerate</button>
          </div>
        </div>

      </div>
    </div>

  </div><!-- /content -->
</div><!-- /main -->

<!-- ═══════════════════════════ LEAD DETAIL MODAL ═══════════════════════════ -->
<div class="lead-modal-overlay" id="lead-modal-overlay" onclick="closeLeadModalBg(event)">
  <div class="lead-modal">
    <div class="lead-modal-header">
      <div>
        <div class="lead-modal-name" id="lm-name">—</div>
        <div id="lm-score-badge" style="margin-top:4px"></div>
      </div>
      <button class="lead-modal-close" onclick="closeLeadModal()">✕</button>
    </div>
    <div class="lead-modal-field">
      <div class="lead-modal-field-icon"><i class="fa-solid fa-envelope"></i></div>
      <div class="lead-modal-field-content">
        <div class="lead-modal-field-label">Email</div>
        <div class="lead-modal-field-value"><a id="lm-email" href="#"></a></div>
      </div>
    </div>
    <div class="lead-modal-field">
      <div class="lead-modal-field-icon"><i class="fa-solid fa-graduation-cap"></i></div>
      <div class="lead-modal-field-content">
        <div class="lead-modal-field-label">Program Interest</div>
        <div class="lead-modal-field-value" id="lm-program">—</div>
      </div>
    </div>
    <div class="lead-modal-field">
      <div class="lead-modal-field-icon"><i class="fa-solid fa-circle-check"></i></div>
      <div class="lead-modal-field-content">
        <div class="lead-modal-field-label">Status</div>
        <div class="lead-modal-field-value" id="lm-status">—</div>
      </div>
    </div>
    <div class="lead-modal-field">
      <div class="lead-modal-field-icon"><i class="fa-solid fa-paper-plane"></i></div>
      <div class="lead-modal-field-content">
        <div class="lead-modal-field-label">Drip Email Progress</div>
        <div class="drip-bar-wrap">
          <div class="drip-bar-track"><div class="drip-bar-fill" id="lm-drip-fill" style="width:0%"></div></div>
          <div class="drip-bar-label" id="lm-drip-label">—</div>
        </div>
      </div>
    </div>
    <div class="lead-modal-actions">
      <button class="lm-btn lm-btn-toggle" id="lm-toggle-btn" onclick="toggleLeadStatus()">
        <i class="fa-solid fa-toggle-on"></i> <span id="lm-toggle-label">Mark Inactive</span>
      </button>
      <button class="lm-btn lm-btn-delete" onclick="deleteLeadConfirm()">
        <i class="fa-solid fa-trash"></i> Delete Lead
      </button>
    </div>
  </div>
</div>

<!-- ═══════════════════════════ POST PREVIEW MODAL ═══════════════════════════ -->
<div class="pp-overlay" id="pp-overlay" onclick="closePpModalBg(event)">
  <div class="pp-modal">
    <div class="pp-header">
      <h3><i class="fa-regular fa-eye" style="color:var(--primary)"></i>&nbsp; Today's Post Preview</h3>
      <button class="pp-close" onclick="closePpModal()">✕</button>
    </div>
    <div id="pp-body">
      <div class="empty-state" style="padding:40px"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading preview...</p></div>
    </div>
    <div class="pp-meta" id="pp-meta"></div>
  </div>
</div>

<!-- ═══════════════════════════════════════ MODAL ═══════════════════════════════════════ -->
<div id="modal-overlay" onclick="closeModalBg(event)">
  <div class="modal">
    <div class="modal-header">
      <div class="modal-title" id="modal-title">Output</div>
      <button class="modal-close" onclick="closeModal()">✕</button>
    </div>
    <div class="modal-body">
      <pre class="modal-pre" id="modal-content"></pre>
    </div>
  </div>
</div>

<!-- ══════════════════════════════════ FACEBOOK POST MODAL ══════════════════════════════════ -->
<div id="fb-modal-overlay" class="action-modal-overlay" onclick="closeFbModalBg(event)">
  <div class="action-modal am-wide">
    <div class="am-header">
      <div class="am-header-icon am-icon-fb"><i class="fa-brands fa-facebook-f"></i></div>
      <div class="am-title">Post to Facebook</div>
      <button class="am-close" onclick="closeFbModal()">✕</button>
    </div>
    <div class="am-two-col">
      <!-- LEFT: image preview -->
      <div class="am-col-img">
        <div class="am-img-empty" id="fb-img-empty">
          <i class="fa-solid fa-image"></i>
          <span>Select an image</span>
        </div>
        <img id="fb-image-preview" src="" alt="" onerror="this.classList.remove('loaded')">
      </div>
      <!-- RIGHT: form -->
      <div class="am-col-form">
        <div class="am-col-form-inner">
          <div class="am-group">
            <div class="am-label">Select Image</div>
            <div class="cs-wrap" id="fb-cs-wrap">
              <button type="button" class="cs-trigger" onclick="toggleCsSel('fb-cs-wrap')">
                <i class="fa-solid fa-calendar-days cs-trigger-icon"></i>
                <div class="cs-trigger-content">
                  <span class="cs-trigger-text" id="fb-cs-text">Loading images...</span>
                </div>
                <span class="cs-trigger-badge" id="fb-cs-badge" style="display:none"></span>
                <i class="fa-solid fa-chevron-down cs-trigger-arrow"></i>
              </button>
              <div class="cs-list-wrap">
                <div class="cs-list" id="fb-cs-list"></div>
              </div>
              <select id="fb-image-select" style="display:none"></select>
            </div>
          </div>
          <div class="am-group">
            <div class="am-caption-row">
              <div class="am-label" style="margin:0">Caption</div>
              <button class="am-gen-btn" id="fb-gen-btn" onclick="generateFbCaption()">
                <i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption
              </button>
            </div>
            <textarea id="fb-caption" class="am-textarea" rows="5" placeholder="Write your caption here..."></textarea>
          </div>
          <div class="am-group">
            <div class="am-label">Hashtags</div>
            <textarea id="fb-hashtags" class="am-textarea" rows="2" placeholder="#EnderunExtension #BGC ..."></textarea>
          </div>
        </div>
        <div class="am-footer-inline">
          <button class="btn btn-ghost" onclick="closeFbModal()">Cancel</button>
          <button class="btn btn-gold" id="fb-submit-btn" onclick="submitFbPost()">
            <i class="fa-brands fa-facebook-f"></i> Post to Facebook
          </button>
        </div>
      </div>
    </div>
  </div>
</div>

<!-- ══════════════════════════════════ INSTAGRAM POST MODAL ══════════════════════════════════ -->
<div id="ig-modal-overlay" class="action-modal-overlay" onclick="closeIgModalBg(event)">
  <div class="action-modal am-wide">
    <div class="am-header">
      <div class="am-header-icon am-icon-ig"><i class="fa-brands fa-instagram"></i></div>
      <div class="am-title">Post to Instagram</div>
      <button class="am-close" onclick="closeIgModal()">✕</button>
    </div>
    <div class="am-two-col">
      <!-- LEFT: image preview -->
      <div class="am-col-img">
        <div class="am-img-empty" id="ig-img-empty">
          <i class="fa-solid fa-image"></i>
          <span>Select an image</span>
        </div>
        <img id="ig-image-preview" src="" alt="" onerror="this.classList.remove('loaded')">
      </div>
      <!-- RIGHT: form -->
      <div class="am-col-form">
        <div class="am-col-form-inner">
          <div class="am-group">
            <div class="am-label">Select Image</div>
            <div class="cs-wrap" id="ig-cs-wrap">
              <button type="button" class="cs-trigger" onclick="toggleCsSel('ig-cs-wrap')">
                <i class="fa-solid fa-calendar-days cs-trigger-icon"></i>
                <div class="cs-trigger-content">
                  <span class="cs-trigger-text" id="ig-cs-text">Loading images...</span>
                </div>
                <span class="cs-trigger-badge" id="ig-cs-badge" style="display:none"></span>
                <i class="fa-solid fa-chevron-down cs-trigger-arrow"></i>
              </button>
              <div class="cs-list-wrap">
                <div class="cs-list" id="ig-cs-list"></div>
              </div>
              <select id="ig-image-select" style="display:none"></select>
            </div>
          </div>
          <div class="am-group">
            <div class="am-caption-row">
              <div class="am-label" style="margin:0">Caption</div>
              <button class="am-gen-btn" id="ig-gen-btn" onclick="generateIgCaption()">
                <i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption
              </button>
            </div>
            <textarea id="ig-caption" class="am-textarea" rows="5" placeholder="Short, punchy IG caption..."></textarea>
          </div>
          <div class="am-group">
            <div class="am-label">Hashtags</div>
            <textarea id="ig-hashtags" class="am-textarea" rows="2" placeholder="#EnderunExtension #BGC ..."></textarea>
          </div>
        </div>
        <div class="am-footer-inline">
          <button class="btn btn-ghost" onclick="closeIgModal()">Cancel</button>
          <button class="btn btn-gold" id="ig-submit-btn" onclick="submitIgPost()">
            <i class="fa-brands fa-instagram"></i> Post to Instagram
          </button>
        </div>
      </div>
    </div>
  </div>
</div>

<!-- ════════════════════════════════════ DRIP EMAIL MODAL ════════════════════════════════════ -->
<div id="email-modal-overlay" class="action-modal-overlay" onclick="closeEmailModalBg(event)">
  <div class="action-modal">
    <div class="am-header">
      <div class="am-header-icon am-icon-email"><i class="fa-solid fa-envelope-open-text"></i></div>
      <div class="am-title">Send Drip Emails</div>
      <button class="am-close" onclick="closeEmailModal()">✕</button>
    </div>
    <div class="am-body" id="email-modal-body">
      <div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading leads...</p></div>
    </div>
    <div class="am-footer">
      <button class="btn btn-ghost" onclick="closeEmailModal()">Cancel</button>
      <button class="btn btn-gold" id="email-submit-btn" onclick="submitDripEmail()">
        <i class="fa-solid fa-paper-plane"></i> Send to All Active Leads
      </button>
    </div>
  </div>
</div>

<!-- ═══════════════════════════════════════ TOAST ═══════════════════════════════════════ -->
<div id="toast">
  <span class="toast-icon" id="toast-icon">✅</span>
  <span id="toast-msg">Success!</span>
</div>

<!-- ═══════════════════════════════════════ SCRIPT ═══════════════════════════════════════ -->
<script>
// ── STATE ──────────────────────────────────────────────────────────────
var currentAgent        = null;
var currentAgentIconHTML= '<img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun">';
var isStreaming         = false;
var toastTimer          = null;
var chatSessions        = {};   // { agentId: [{id,title,messages,updatedAt}] }
var currentSessionId    = {};   // { agentId: sessionId|null }
var currentMessages     = [];   // [{role,content}] live history
var chatModel           = 'claude-opus-4-7';
var chatMode            = 'normal';

// ── INIT ──────────────────────────────────────────────────────────────
document.addEventListener('DOMContentLoaded', function() {
  // Clean up corrupted greeting-only sessions from old auto-greeting bug
  try {
    var raw = localStorage.getItem('eva_chat_v2');
    if (raw) {
      var data = JSON.parse(raw);
      Object.keys(data).forEach(function(agentId) {
        data[agentId] = (data[agentId] || []).filter(function(s) {
          var msgs = s.messages || [];
          var hasReal = msgs.some(function(m) {
            var txt = typeof m.content === 'string' ? m.content : '';
            return !(m.role === 'user' && txt === '__GREET__') && m.role !== 'user' ? false : m.role === 'user' && txt !== '__GREET__';
          }) || msgs.some(function(m){ return m.role === 'assistant'; });
          return hasReal;
        });
      });
      localStorage.setItem('eva_chat_v2', JSON.stringify(data));
    }
  } catch(e) {}
  _loadSessions();
  setDate();
  loadStats();
  loadAgents();
  loadNotifications();
  setInterval(loadNotifications, 30000);
  calInit();
  // Warm the composer thumbnail cache in the background
  setTimeout(function(){ fetch('/api/composer/preload', {method:'POST'}); }, 3000);
});

// ── NOTIFICATIONS ──────────────────────────────────────────────────────
var drawerOpen = false;

function toggleTbDd(id) {
  var dd = document.getElementById(id);
  if (!dd) return;
  var open = dd.classList.contains('open');
  closeTbDd();
  if (!open) dd.classList.add('open');
}
function closeTbDd() {
  document.querySelectorAll('.tb-dropdown.open').forEach(function(d){ d.classList.remove('open'); });
}
document.addEventListener('click', function(e) {
  if (!e.target.closest('.tb-dropdown')) closeTbDd();
});

// ── DARK MODE ──────────────────────────────────────────────────────────
function toggleDark() {
  var isDark = document.body.classList.toggle('dark');
  localStorage.setItem('_darkMode', isDark ? '1' : '0');
  var btn = document.getElementById('dark-toggle');
  if (btn) btn.innerHTML = isDark ? '<i class="fa-solid fa-sun"></i>' : '<i class="fa-solid fa-moon"></i>';
  // Update agent search dropdown bg in dark
  var agentSearch = document.querySelector('#agent-dd-menu > div:first-child');
  if (agentSearch) agentSearch.style.background = isDark ? '#1a1a24' : '#fff';
}
// Restore dark mode immediately (body class) — button update deferred to DOMContentLoaded
if (localStorage.getItem('_darkMode') === '1') {
  document.body.classList.add('dark');
  document.addEventListener('DOMContentLoaded', function() {
    var btn = document.getElementById('dark-toggle');
    if (btn) btn.innerHTML = '<i class="fa-solid fa-sun"></i>';
  }, { once: true });
}

// ── KEYBOARD SHORTCUTS MODAL ───────────────────────────────────────────
function showShortcuts() {
  document.getElementById('shortcuts-overlay').classList.add('open');
}
function closeShortcuts() {
  document.getElementById('shortcuts-overlay').classList.remove('open');
}

// ── GLOBAL KEYBOARD SHORTCUTS ─────────────────────────────────────────
document.addEventListener('keydown', function(e) {
  var tag = (e.target.tagName || '').toLowerCase();
  var inInput = tag === 'textarea' || tag === 'input';
  // ? → shortcuts (not in input)
  if (!inInput && e.key === '?' && !e.ctrlKey && !e.metaKey) { showShortcuts(); return; }
  // Escape → close shortcuts / stop stream
  if (e.key === 'Escape') {
    if (document.getElementById('shortcuts-overlay').classList.contains('open')) { closeShortcuts(); return; }
    if (isStreaming) { stopStream(); return; }
  }
  // Ctrl+Shift+D → dark mode
  if ((e.ctrlKey || e.metaKey) && e.shiftKey && e.key === 'D') { e.preventDefault(); toggleDark(); return; }
  // Ctrl+Shift+N → new chat
  if ((e.ctrlKey || e.metaKey) && e.shiftKey && e.key === 'N') { e.preventDefault(); newChat(); return; }
  // Ctrl+L → focus input
  if ((e.ctrlKey || e.metaKey) && e.key === 'l') { e.preventDefault(); var inp = document.getElementById('chat-input'); if (inp) inp.focus(); return; }
  // Ctrl+F → search in chat (only when chat view active)
  if ((e.ctrlKey || e.metaKey) && e.key === 'f' && document.getElementById('view-dashboard').classList.contains('active')) {
    e.preventDefault(); toggleChatSearch(); return;
  }
  // Ctrl+Shift+E → export chat PDF
  if ((e.ctrlKey || e.metaKey) && e.shiftKey && e.key === 'E') { e.preventDefault(); exportChatPDF(); return; }
  // Alt+1/2/3/4 → navigation
  if (e.altKey && !inInput) {
    var navMap = {'1':'dashboard','2':'leads','3':'calendar','4':'intelligence'};
    if (navMap[e.key]) { e.preventDefault(); showView(navMap[e.key]); return; }
  }
});

function toggleDrawer() {
  drawerOpen ? closeDrawer() : openDrawer();
}
function openDrawer() {
  drawerOpen = true;
  document.getElementById('notif-drawer').classList.add('open');
  document.getElementById('notif-overlay').classList.add('open');
}
function closeDrawer() {
  drawerOpen = false;
  document.getElementById('notif-drawer').classList.remove('open');
  document.getElementById('notif-overlay').classList.remove('open');
}

function loadNotifications() {
  fetch('/api/notifications')
    .then(function(r){ return r.json(); })
    .then(function(data) {
      var notifs = data.notifications || [];
      var unread = data.unread || 0;

      // Update badge
      var badge = document.getElementById('notif-badge');
      var bell  = document.getElementById('notif-bell');
      if (unread > 0) {
        badge.textContent = unread > 99 ? '99+' : unread;
        badge.style.display = 'flex';
        bell.classList.add('has-unread');
      } else {
        badge.style.display = 'none';
        bell.classList.remove('has-unread');
      }

      // Update subtitle
      var sub = document.getElementById('nd-subtitle');
      sub.textContent = unread > 0 ? unread + ' unread notification' + (unread > 1 ? 's' : '') : 'All caught up';

      // Render cards
      var body = document.getElementById('nd-body');
      if (!notifs.length) {
        body.innerHTML = '<div class="nd-empty"><i class="fa-solid fa-bell-slash" style="font-size:28px;display:block;margin-bottom:10px;color:var(--text-dim)"></i>No notifications yet<br><span style="font-size:11px;color:var(--text-dim)">Automation activity will appear here</span></div>';
        return;
      }
      body.innerHTML = notifs.map(function(n) {
        return renderCard(n);
      }).join('');
    })
    .catch(function(){});
}

function renderCard(n) {
  var readClass = n.read ? 'read' : 'unread';
  var dot       = n.read ? '' : '<div class="nd-unread-dot"></div>';
  var time      = relativeTime(n.timestamp);
  var reportBtn = n.report_path
    ? '<a class="nd-report-btn" href="#" onclick="return false;" title="' + n.report_path + '"><i class="fa-solid fa-file-lines"></i> Report</a>'
    : '';
  var levelLabel = {
    critical: '<i class="fa-solid fa-circle-exclamation"></i> Critical',
    warning:  '<i class="fa-solid fa-triangle-exclamation"></i> Warning',
    success:  '<i class="fa-solid fa-circle-check"></i> Done',
    info:     '<i class="fa-solid fa-circle-info"></i> Info'
  }[n.level] || n.level;
  return [
    '<div class="nd-card level-' + n.level + ' ' + readClass + '" id="nc-' + n.id + '">',
      dot,
      '<div class="nd-card-emoji">' + (n.agent_emoji || '🤖') + '</div>',
      '<div class="nd-card-body">',
        '<div class="nd-card-agent">' + (n.agent_label || '') + ' &nbsp;·&nbsp; ' + levelLabel + '</div>',
        '<div class="nd-card-title">' + escHtml(n.title) + '</div>',
        '<div class="nd-card-msg">'  + escHtml(n.message) + '</div>',
        '<div class="nd-card-footer">',
          '<span class="nd-card-time">' + time + '</span>',
          '<div class="nd-card-actions">',
            reportBtn,
            n.read ? '' : '<button class="nd-btn-read" onclick="markRead(\'' + n.id + '\')">Mark read</button>',
          '</div>',
        '</div>',
      '</div>',
    '</div>'
  ].join('');
}

function markRead(id) {
  fetch('/api/notifications/read', {
    method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({id: id})
  }).then(function(){ loadNotifications(); });
}
function markAllRead() {
  fetch('/api/notifications/read-all', { method:'POST' })
    .then(function(){ loadNotifications(); });
}
function clearAll() {
  if (!confirm('Clear all notifications?')) return;
  fetch('/api/notifications/clear', { method:'POST' })
    .then(function(){ loadNotifications(); });
}

function relativeTime(iso) {
  if (!iso) return '';
  var diff = Math.floor((Date.now() - new Date(iso).getTime()) / 1000);
  if (diff < 60)   return 'Just now';
  if (diff < 3600) return Math.floor(diff/60) + 'm ago';
  if (diff < 86400)return Math.floor(diff/3600) + 'h ago';
  return Math.floor(diff/86400) + 'd ago';
}
function escHtml(s) {
  return String(s||'').replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
}

function setDate() {
  var now = new Date();
  var opts = { weekday: 'long', year: 'numeric', month: 'long', day: 'numeric' };
  document.getElementById('topbar-date').innerHTML = '<i class="fa-regular fa-calendar"></i> ' + now.toLocaleDateString('en-PH', opts);
}

// ── VIEW SWITCHING ─────────────────────────────────────────────────────
function showView(id, el) {
  document.querySelectorAll('.view').forEach(function(v) { v.classList.remove('active'); v.style.cssText = ''; });
  var target = document.getElementById('view-' + id);
  if (target) target.classList.add('active');

  document.querySelectorAll('.nav-item').forEach(function(n) { n.classList.remove('active'); });
  if (el) el.classList.add('active');

  var titles = {
    dashboard:        'AI Chat',
    leads:            'Leads Management',
    calendar:         'Content Calendar',
    intelligence:     'Competitor Intelligence',
    automations:      'Automation Status',
    analytics:        'Analytics',
    strategy:         'Strategy Studio',
    'agent-workspace':'Agent Workspace'
  };
  document.getElementById('topbar-title').textContent = titles[id] || 'Marketing Hub';

  if (id === 'leads')           loadLeads();
  if (id === 'calendar')        calRender();
  if (id === 'intelligence')    loadIntelligence();
  if (id === 'automations')     loadAutomations();
  if (id === 'analytics')       loadAnalytics();
}

// ── SESSION STORAGE ────────────────────────────────────────────────────
function _saveSessions() {
  try { localStorage.setItem('eva_chat_v2', JSON.stringify(chatSessions)); } catch(e) {}
}
function _loadSessions() {
  try { var r = localStorage.getItem('eva_chat_v2'); if (r) chatSessions = JSON.parse(r); } catch(e) {}
}
function _genId() {
  return Date.now().toString(36) + Math.random().toString(36).slice(2, 6);
}
function _timeAgo(ts) {
  var d = Date.now() - ts;
  if (d < 60000)    return 'just now';
  if (d < 3600000)  return Math.floor(d/60000) + 'm ago';
  if (d < 86400000) return Math.floor(d/3600000) + 'h ago';
  return Math.floor(d/86400000) + 'd ago';
}

function openAgent(agentId, agentName, agentEmoji, agentRole) {
  // Always reset stream state when switching agents
  isStreaming = false;
  document.getElementById('send-btn').disabled = false;
  var stale = document.getElementById('typing-msg'); if (stale) stale.remove();

  _loadSessions();
  currentAgent = agentId;
  var iconHTML = '<img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun">';
  currentAgentIconHTML = iconHTML;
  document.getElementById('chat-avatar').innerHTML = iconHTML;
  document.getElementById('chat-name').textContent  = agentName || 'Agent';
  document.getElementById('chat-role').textContent  = agentRole || 'AI Marketing Specialist';
  var qc = document.getElementById('quick-chips'); if (qc) qc.style.display = 'block';

  var lbl = document.getElementById('agent-dd-label');
  if (lbl) lbl.textContent = agentName || agentId;
  document.querySelectorAll('#agent-dd-menu .custom-dd-opt').forEach(function(o) {
    o.classList.toggle('selected', o.getAttribute('data-val') === agentId);
  });

  // Restore most recent session or show welcome
  var sessions = chatSessions[agentId] || [];
  var lastId = currentSessionId[agentId];
  var session = lastId ? sessions.find(function(s){ return s.id === lastId; }) : null;
  if (!session && sessions.length) session = sessions[sessions.length - 1];

  // Only restore if there's at least one visible (non-__GREET__) message
  var visibleMsgs = session ? (session.messages || []).filter(function(m) {
    var txt = typeof m.content === 'string' ? m.content : '';
    return !(m.role === 'user' && txt === '__GREET__');
  }) : [];

  if (session && visibleMsgs.length) {
    currentSessionId[agentId] = session.id;
    currentMessages = session.messages ? session.messages.slice() : [];
    _restoreMessages(session.messages, iconHTML);
  } else {
    currentSessionId[agentId] = null;
    currentMessages = [];
    _showWelcome(agentName, iconHTML, agentRole);
  }

  renderSessionsList(agentId);
  showView('dashboard', null);
  scrollBottom();
  document.getElementById('chat-input').focus();
}

function _showWelcome(name, iconHTML, role) {
  var msgs = document.getElementById('chat-messages');
  msgs.innerHTML = '';
  var wb = document.createElement('div');
  wb.className = 'welcome-box'; wb.id = 'chat-welcome';
  wb.innerHTML =
    '<div class="welcome-icon"><img src="/assets/logos/Enderun-Colleges.png" class="welcome-logo" alt="Enderun"></div>' +
    '<div class="welcome-title">' + escHtml(name||'Agent') + '</div>' +
    '<div class="welcome-sub">' + escHtml(role||'AI Marketing Specialist') + '</div>';
  msgs.appendChild(wb);
}

function _triggerGreeting() {
  if (!currentAgent) return;
  var greetingAgent = currentAgent; // capture — abort callbacks if agent changes
  isStreaming = true;
  document.getElementById('send-btn').disabled = true;
  var msgs = document.getElementById('chat-messages');
  msgs.innerHTML = '';
  var agentEmoji = '<img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun">';
  var typingEl = document.createElement('div');
  typingEl.className = 'msg'; typingEl.id = 'typing-msg';
  typingEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-bubble"><div class="typing-indicator"><div class="typing-dot"></div><div class="typing-dot"></div><div class="typing-dot"></div></div></div>';
  msgs.appendChild(typingEl);
  scrollBottom();

  fetch('/api/chat', {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify({ agent: greetingAgent, greeting: true, model: chatModel })
  })
  .then(function(r) {
    if (!r.ok) throw new Error('HTTP ' + r.status);
    var reader = r.body.getReader();
    var decoder = new TextDecoder();
    var aiEl = null; var aiBubble = null; var accumulated = '';

    function readChunk() {
      reader.read().then(function(result) {
        if (currentAgent !== greetingAgent) { reader.cancel(); return; } // agent changed — abort
        if (result.done) { finishGreeting(); return; }
        var chunk = decoder.decode(result.value, { stream: true });
        var isDone = false;
        chunk.split('\n').forEach(function(line) {
          if (isDone || !line.startsWith('data: ')) return;
          var data = line.slice(6).trim();
          if (data === '[DONE]') { isDone = true; finishGreeting(); return; }
          try {
            var parsed = JSON.parse(data);
            var token = parsed.token || parsed.content || parsed.text || '';
            if (token) {
              accumulated += token;
              if (!aiEl) {
                var t = document.getElementById('typing-msg'); if (t) t.remove();
                aiEl = document.createElement('div'); aiEl.className = 'msg';
                aiEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-bubble"></div>';
                msgs.appendChild(aiEl); aiBubble = aiEl.querySelector('.msg-bubble');
              }
              aiBubble.innerHTML = md(accumulated); scrollBottom();
            }
          } catch(e) {}
        });
        if (!isDone) readChunk();
      }).catch(function() { finishGreeting(); });
    }

    function finishGreeting() {
      if (currentAgent !== greetingAgent) return; // agent changed — discard result
      var t = document.getElementById('typing-msg'); if (t) t.remove();
      if (!aiEl && accumulated) {
        aiEl = document.createElement('div'); aiEl.className = 'msg';
        aiEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-bubble"></div>';
        msgs.appendChild(aiEl);
      }
      if (aiEl) {
        aiBubble = aiEl.querySelector('.msg-bubble');
        aiBubble.innerHTML = _renderWithThinking(accumulated);
      }
      if (accumulated) {
        currentMessages = [
          { role: 'user',      content: '__GREET__' },
          { role: 'assistant', content: accumulated }
        ];
        _saveCurrentSession();
      }
      isStreaming = false;
      document.getElementById('send-btn').disabled = false;
      scrollBottom();
    }

    readChunk();
  })
  .catch(function() {
    var t = document.getElementById('typing-msg'); if (t) t.remove();
    isStreaming = false;
    document.getElementById('send-btn').disabled = false;
  });
}

function _restoreMessages(messages, iconHTML) {
  var msgs = document.getElementById('chat-messages');
  msgs.innerHTML = '';
  (messages || []).forEach(function(m) {
    var txt = typeof m.content === 'string' ? m.content
      : (Array.isArray(m.content) ? (m.content.filter(function(x){return x.type==='text';})[0]||{}).text||'' : '');
    if (m.role === 'user' && txt === '__GREET__') return;
    var el = document.createElement('div');
    var restoreName = document.getElementById('chat-name').textContent || 'Agent';
    el.className = 'msg' + (m.role === 'user' ? ' user' : ' ai');
    var av = m.role === 'user'
      ? '<div class="msg-avatar user-av"><i class="fa-solid fa-user"></i></div>'
      : '<div class="msg-avatar">' + (iconHTML||'<img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun">') + '</div>';
    var bubble = m.role === 'user'
      ? '<div class="msg-body"><div class="msg-bubble">' + escHtml(txt) + '</div></div>'
      : '<div class="msg-body"><span class="msg-name">' + escHtml(restoreName) + '</span><div class="msg-bubble">' + _renderWithThinking(txt) + '</div></div>';
    el.innerHTML = av + bubble;
    if (m.role !== 'user') { var rb = el.querySelector('.msg-bubble'); if (rb) _processCodeBlocks(rb); }
    msgs.appendChild(el);
    // Restore tool cards saved with this message
    if (m._toolCards && m._toolCards.length) {
      m._toolCards.forEach(function(tc) {
        var card = document.createElement('div');
        card.className = 'tool-action-card ' + (tc.status || 'done');
        if (tc.img_url) {
          card.classList.add('has-image');
          card.innerHTML = '<div class="tool-card-header"><i class="fa-solid fa-circle-check" style="font-size:13px;flex-shrink:0"></i><span>' + escHtml(tc.result||'Image generated!') + '</span></div><div class="tool-card-img"><img src="' + escHtml(tc.img_url) + '" alt="Generated image" onload="scrollBottom()" onclick="openLightbox(this.src)" style="cursor:zoom-in"></div>';
        } else {
          var ico = tc.status === 'failed' ? 'fa-circle-xmark' : 'fa-circle-check';
          card.innerHTML = '<i class="fa-solid ' + ico + '" style="font-size:13px;flex-shrink:0"></i><span>' + escHtml(tc.result||'') + '</span>';
        }
        msgs.appendChild(card);
      });
    }
    // Restore post composer if it was open when session was saved
    if (m._composer) {
      openPostComposer(m._composer.platform || 'facebook');
    }
  });
}

// ── LOAD AGENTS ────────────────────────────────────────────────────────
var AGENT_DEFS = [
  { id: 'general',             emoji: '<i class="fa-solid fa-brain"></i>',                    name: 'Eva AI',                  role: 'General intelligence — any topic',       cat: 'ops',        icon: 'fa-brain',                    iconCls: 'icon-ops',        workspace: 'none',     quickActions: ['Help me plan this week', 'Explain a marketing concept', 'Brainstorm campaign ideas'] },
  { id: 'marketing-manager',   emoji: '<i class="fa-solid fa-bullseye"></i>',                 name: 'Marketing Manager',       role: 'Chief coordinator & orchestrator',       cat: 'core',       icon: 'fa-bullseye',                 iconCls: 'icon-core',       workspace: 'overview', quickActions: ['Review all active campaigns', 'Create this week\'s strategy brief', 'Generate a full marketing status report'] },
  { id: 'social-media',        emoji: '<i class="fa-solid fa-mobile-screen-button"></i>',     name: 'Social Media',            role: 'Daily content — FB, IG, TikTok',        cat: 'core',       icon: 'fa-mobile-screen-button',     iconCls: 'icon-core',       workspace: 'schedule', quickActions: ['Write a Facebook post for Enderun Extension', 'Draft an Instagram caption for BS HM', 'Write a TikTok script for campus life'] },
  { id: 'drip-campaign',       emoji: '<i class="fa-solid fa-envelope-open-text"></i>',       name: 'Drip Campaign',           role: 'Email sequences & lead nurture',         cat: 'core',       icon: 'fa-envelope-open-text',       iconCls: 'icon-core',       workspace: 'drip',     quickActions: ['Draft a follow-up email for an HM lead', 'Write a 5-email nurture sequence for CA leads', 'Create a re-engagement email for cold leads'] },
  { id: 'data-analysis',       emoji: '<i class="fa-solid fa-chart-bar"></i>',                name: 'Data Analysis',           role: 'KPIs, metrics & reporting',              cat: 'core',       icon: 'fa-chart-bar',                iconCls: 'icon-core',       workspace: 'leads',    quickActions: ['Analyze current lead pipeline', 'Generate a PDF report with charts', 'Show hot lead breakdown by program'] },
  { id: 'content-strategy',    emoji: '<i class="fa-solid fa-pen-to-square"></i>',            name: 'Content Strategy',        role: 'Editorial calendar & content pillars',   cat: 'core',       icon: 'fa-pen-to-square',            iconCls: 'icon-core',       workspace: 'schedule', quickActions: ['Build next month\'s editorial calendar', 'Define content pillars for Enderun Colleges', 'Audit current content mix and suggest improvements'] },
  { id: 'pr',                  emoji: '<i class="fa-solid fa-newspaper"></i>',                name: 'PR Agent',                role: 'Press releases & media relations',       cat: 'core',       icon: 'fa-newspaper',                iconCls: 'icon-core',       workspace: 'caps',     quickActions: ['Write a press release about Les Roches affiliation', 'Draft a media pitch for lifestyle editors', 'Write talking points for CEO interview'] },
  { id: 'competitor-analysis', emoji: '<i class="fa-solid fa-magnifying-glass-chart"></i>',  name: 'Competitor Analysis',     role: 'Market intelligence & insights',         cat: 'core',       icon: 'fa-magnifying-glass-chart',   iconCls: 'icon-core',       workspace: 'intel',    quickActions: ['Run competitor analysis vs DLSU', 'What are CCA and ISCAHM doing this month?', 'Identify market gaps we can exploit'] },
  { id: 'designer',            emoji: '<i class="fa-solid fa-palette"></i>',                 name: 'Designer Agent',          role: 'Creative direction & visual briefs',     cat: 'core',       icon: 'fa-palette',                  iconCls: 'icon-core',       workspace: 'caps',     quickActions: ['Create a Canva brief for Enrollment Season campaign', 'Design direction for Instagram Reels thumbnail', 'Write a visual brief for an Open House banner'] },
  { id: 'seo-digital',         emoji: '<i class="fa-solid fa-magnifying-glass"></i>',        name: 'SEO & Digital Ads',       role: 'Google, Meta & TikTok advertising',      cat: 'specialist', icon: 'fa-magnifying-glass',         iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Optimize the Enderun Colleges homepage for SEO', 'Write a Google Search ad for BS HM admissions', 'Recommend Meta audience targeting for enrollment campaign'] },
  { id: 'admissions',          emoji: '<i class="fa-solid fa-graduation-cap"></i>',          name: 'Admissions Agent',        role: 'Lead qualification & campus tours',       cat: 'specialist', icon: 'fa-graduation-cap',           iconCls: 'icon-specialist', workspace: 'admissions',quickActions: ['Reply to a BS Hospitality inquiry', 'Follow up 5 hot leads this week', 'Create campus tour talking points'] },
  { id: 'video-multimedia',    emoji: '<i class="fa-solid fa-clapperboard"></i>',            name: 'Video & Multimedia',      role: 'TikTok, Reels & YouTube scripts',        cat: 'specialist', icon: 'fa-clapperboard',             iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Write a TikTok script for a campus tour', 'Script a 60-sec Reel for BS Culinary Arts', 'Write a YouTube script about Les Roches partnership'] },
  { id: 'events-activations',  emoji: '<i class="fa-solid fa-calendar-check"></i>',          name: 'Events & Activations',    role: 'Open house, webinars & campus tours',    cat: 'specialist', icon: 'fa-calendar-check',           iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Plan next month\'s Open House event', 'Create a webinar script for parents', 'Design a campus tour experience flow'] },
  { id: 'alumni-relations',    emoji: '<i class="fa-solid fa-handshake"></i>',               name: 'Alumni Relations',        role: 'Alumni engagement & referral programs',  cat: 'specialist', icon: 'fa-handshake',                iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Write an alumni spotlight story template', 'Design a referral incentive program', 'Draft a monthly alumni newsletter'] },
  { id: 'influencer-kol',      emoji: '<i class="fa-solid fa-star"></i>',                    name: 'Influencer & KOL',        role: 'Creator outreach & campaign management', cat: 'specialist', icon: 'fa-star',                     iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Identify 5 KOLs for BS Culinary Arts promo', 'Write an influencer outreach pitch', 'Plan a creator campus tour event'] },
  { id: 'marketing-analysis',  emoji: '<i class="fa-solid fa-chart-line"></i>',              name: 'Marketing Analysis',      role: 'Deep analytics & data visualization',    cat: 'specialist', icon: 'fa-chart-line',               iconCls: 'icon-specialist', workspace: 'reports',  quickActions: ['Generate a lead funnel PDF report', 'Create a program performance chart', 'Analyze which program has best lead velocity'] },
  { id: 'business-analyst',    emoji: '<i class="fa-solid fa-briefcase"></i>',               name: 'Business Analyst',        role: 'Revenue models & strategic planning',    cat: 'specialist', icon: 'fa-briefcase',                iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Model ROI of enrollment season ad spend', 'Size the hospitality education market in PH', 'Build a 3-year enrollment growth projection'] },
  { id: 'events-banquetes',    emoji: '<i class="fa-solid fa-champagne-glasses"></i>',       name: 'Events / Banquetes',      role: 'Venue marketing & catering leads',        cat: 'specialist', icon: 'fa-champagne-glasses',        iconCls: 'icon-specialist', workspace: 'events',   quickActions: ['Write a wedding venue promo post', 'Follow up with event inquiry leads', 'Draft a corporate events pitch deck outline'] },
  { id: 'researcher',          emoji: '<i class="fa-solid fa-microscope"></i>',              name: 'Researcher Agent',        role: 'Brand research & fact-checking',         cat: 'specialist', icon: 'fa-microscope',               iconCls: 'icon-specialist', workspace: 'caps',     quickActions: ['Research enrollment trends in PH higher ed', 'Fact-check Les Roches global ranking claims', 'Analyze Enderun brand sentiment online'] },
  { id: 'lead-generation',     emoji: '<i class="fa-solid fa-magnet"></i>',                  name: 'Lead Generation',         role: 'Paid ads & lead acquisition',            cat: 'growth',     icon: 'fa-magnet',                   iconCls: 'icon-growth',     workspace: 'leads',    quickActions: ['Write a Facebook Lead Form ad for BS HM', 'Plan a Google Search campaign for Enderun Extension', 'Create a retargeting audience strategy'] },
  { id: 'community-manager',   emoji: '<i class="fa-solid fa-comments"></i>',                name: 'Community Manager',       role: 'Real-time social engagement & DMs',      cat: 'growth',     icon: 'fa-comments',                 iconCls: 'icon-growth',     workspace: 'caps',     quickActions: ['Draft responses for 5 Facebook comments', 'Write a reply template for admissions DMs', 'Handle a negative Google Review professionally'] },
  { id: 'parent-engagement',   emoji: '<i class="fa-solid fa-people-roof"></i>',             name: 'Parent Engagement',       role: 'Parent-focused content & email drips',   cat: 'growth',     icon: 'fa-people-roof',              iconCls: 'icon-growth',     workspace: 'parents',  quickActions: ['Write a parent-focused Facebook post', 'Draft a parent email about ROI of Enderun', 'Create Open House invite for parents'] },
  { id: 'whatsapp-sms',        emoji: '<i class="fa-brands fa-whatsapp"></i>',               name: 'WhatsApp & SMS',          role: 'Conversational lead nurture & closing',  cat: 'growth',     icon: 'fa-comment-sms',              iconCls: 'icon-growth',     workspace: 'caps',     quickActions: ['Write a WhatsApp follow-up for a hot lead', 'Create a 5-message WhatsApp nurture flow', 'Draft an SMS blast for enrollment deadline'] },
  { id: 'enrollment-tracker',  emoji: '<i class="fa-solid fa-funnel-dollar"></i>',           name: 'Enrollment Tracker',      role: 'Funnel analytics & forecasting',         cat: 'growth',     icon: 'fa-chart-gantt',              iconCls: 'icon-growth',     workspace: 'leads',    quickActions: ['Show today\'s funnel health report', 'Forecast end-of-month enrollment numbers', 'Identify which funnel stage has the most drop-off'] },
  { id: 'blog-seo-content',    emoji: '<i class="fa-solid fa-blog"></i>',                    name: 'Blog & SEO Content',      role: 'SEO articles & organic traffic',         cat: 'growth',     icon: 'fa-blog',                     iconCls: 'icon-growth',     workspace: 'caps',     quickActions: ['Write an SEO article: "Best Hospitality Schools in PH"', 'Create a keyword strategy for Enderun Colleges', 'Draft a blog post about Les Roches partnership'] },
  { id: 'social-listening',    emoji: '<i class="fa-solid fa-satellite-dish"></i>',          name: 'Social Listening',        role: 'Brand mentions & competitor monitoring',  cat: 'ops',        icon: 'fa-satellite-dish',           iconCls: 'icon-ops',        workspace: 'intel',    quickActions: ['Summarize today\'s intelligence briefing', 'What are competitors saying this week?', 'Identify trends in hospitality education news'] },
  { id: 'chief-strategist',    emoji: '<i class="fa-solid fa-chess-king"></i>',              name: 'Chief Strategist',        role: 'Senior advisor & campaign approvals',    cat: 'ops',        icon: 'fa-chess-king',               iconCls: 'icon-ops',        workspace: 'caps',     quickActions: ['Review and approve this week\'s campaign plan', 'What\'s our biggest strategic priority right now?', 'Give me a SWOT analysis of Enderun vs competitors'] },
  { id: 'testing',             emoji: '<i class="fa-solid fa-flask"></i>',                   name: 'Testing Agent',           role: 'Demo posts & preview without side effects', cat: 'ops',     icon: 'fa-flask',                    iconCls: 'icon-ops',        workspace: 'testing',  quickActions: ['Run a demo Facebook post', 'Test Instagram post pipeline', 'Send a preview drip email'] },
  { id: 'crypto-trader',       emoji: '<i class="fa-solid fa-coins"></i>',                   name: 'Crypto Trader',           role: 'Paper trading & market analysis',        cat: 'ops',        icon: 'fa-coins',                    iconCls: 'icon-ops',        workspace: 'trading',  quickActions: ['Scan for buy signals now', 'Show my paper portfolio', 'Analyze BTC price action with SMC'] }
];


function loadAgents() {
  var menu = document.getElementById('agent-dd-menu');
  if (!menu) return;
  // Keep the search box (first child) and rebuild the rest
  var searchWrap = menu.querySelector('div');
  menu.innerHTML = '';
  if (searchWrap) menu.appendChild(searchWrap);
  var shownDiv = false;
  AGENT_DEFS.forEach(function(a) {
    if (!shownDiv && a.id !== 'general') {
      var div = document.createElement('div');
      div.className = 'custom-dd-divider';
      div.textContent = 'Marketing Agents';
      menu.appendChild(div);
      shownDiv = true;
    }
    var d = document.createElement('div');
    d.className = 'custom-dd-opt' + (a.id === 'general' ? ' agent-general-opt' : '');
    d.setAttribute('data-val', a.id);
    d.setAttribute('data-name', a.name.toLowerCase());
    d.textContent = a.name;
    d.onclick = function() { pickAgent(a.id, a.name); };
    menu.appendChild(d);
  });
}
function filterAgents(q) {
  var term = q.trim().toLowerCase();
  document.querySelectorAll('#agent-dd-menu .custom-dd-opt').forEach(function(o) {
    if (o.classList.contains('placeholder-opt')) return;
    var name = o.getAttribute('data-name') || o.textContent.toLowerCase();
    o.style.display = (!term || name.includes(term)) ? '' : 'none';
  });
  document.querySelectorAll('#agent-dd-menu .custom-dd-divider').forEach(function(d) {
    d.style.display = term ? 'none' : '';
  });
}
function toggleDd(id) {
  var dd = document.getElementById(id);
  if (!dd) return;
  var wasOpen = dd.classList.contains('open');
  closeAllDd();
  if (!wasOpen) {
    dd.classList.add('open');
    if (id === 'agent-dd') {
      var s = document.getElementById('agent-search');
      if (s) { s.value = ''; filterAgents(''); setTimeout(function(){ s.focus(); }, 50); }
    }
  }
}
function closeAllDd() {
  document.querySelectorAll('.custom-dd.open').forEach(function(d) { d.classList.remove('open'); });
}
document.addEventListener('click', function(e) {
  if (!e.target.closest('.custom-dd')) closeAllDd();
});
function pickAgent(agentId, agentName) {
  var lbl = document.getElementById('agent-dd-label');
  if (lbl) lbl.textContent = agentName;
  document.querySelectorAll('#agent-dd-menu .custom-dd-opt').forEach(function(o) {
    o.classList.toggle('selected', o.getAttribute('data-val') === agentId);
  });
  closeAllDd();
  switchAgent(agentId);
}
function pickModel(modelId, modelName) {
  chatModel = modelId;
  var lbl = document.getElementById('model-dd-label');
  if (lbl) lbl.textContent = modelName;
  document.querySelectorAll('#model-dd .custom-dd-opt').forEach(function(o) {
    o.classList.toggle('selected', o.getAttribute('data-val') === modelId);
  });
  closeAllDd();
}

function switchAgent(agentId) {
  if (!agentId) return;
  var def = AGENT_DEFS.find(function(d) { return d.id === agentId; });
  if (def) openAgent(def.id, def.name, def.emoji, def.role);
}

// ── SESSION MANAGEMENT ─────────────────────────────────────────────────
function newChat() {
  if (!currentAgent) { toast(TI.warn, 'Select an agent first.'); return; }
  currentMessages = [];
  currentSessionId[currentAgent] = null;
  renderSessionsList(currentAgent);
  document.getElementById('chat-input').focus();
  _triggerGreeting();
}

function startGreeting() {
  _triggerGreeting();
}

function clearCurrentChat() { newChat(); }

function renderSessionsList(agentId) {
  var list = document.getElementById('chat-sessions-list');
  if (!list) return;
  var sessions = (chatSessions[agentId] || []).slice().reverse();
  var activeId = currentSessionId[agentId];
  if (!sessions.length) {
    list.innerHTML = '<div style="padding:20px 10px;text-align:center;font-size:11px;color:var(--text-dim)">No conversations yet.<br>Start chatting to save history.</div>';
    return;
  }
  // Sort pinned first
  sessions.sort(function(a, b) {
    var pa = _pins[agentId + '|' + a.id] ? 1 : 0;
    var pb = _pins[agentId + '|' + b.id] ? 1 : 0;
    return pb - pa;
  });
  list.innerHTML = sessions.map(function(s) {
    var active  = s.id === activeId ? ' active' : '';
    var pinned  = _pins[agentId + '|' + s.id] ? ' pinned-item' : '';
    var pinCls  = _pins[agentId + '|' + s.id] ? ' pinned' : '';
    var count   = s.messages ? Math.ceil(s.messages.length / 2) : 0;
    return '<div class="chat-session-item' + active + pinned + '" onclick="loadSession(\'' + escHtml(agentId) + '\',\'' + s.id + '\')">' +
      '<div class="chat-session-agent-icon"><img src="/assets/logos/Enderun-Colleges.png" alt="Enderun"></div>' +
      '<div class="chat-session-info">' +
        '<div class="chat-session-title">' + escHtml(s.title || 'Chat') + '</div>' +
        '<div class="chat-session-meta">' + _timeAgo(s.updatedAt) + ' · ' + count + ' msg' + (count !== 1 ? 's' : '') + '</div>' +
      '</div>' +
      '<button class="session-pin-btn' + pinCls + '" onclick="togglePin(\'' + escHtml(agentId) + '\',\'' + s.id + '\',this,event)" title="Pin"><i class="fa-solid fa-thumbtack"></i></button>' +
      '<button class="chat-session-del" onclick="deleteSession(event,\'' + escHtml(agentId) + '\',\'' + s.id + '\')" title="Delete"><i class="fa-solid fa-trash"></i></button>' +
    '</div>';
  }).join('');
}

function loadSession(agentId, sessionId) {
  _loadSessions();
  var session = (chatSessions[agentId] || []).find(function(s){ return s.id === sessionId; });
  if (!session) return;
  currentSessionId[agentId] = sessionId;
  currentMessages = session.messages ? session.messages.slice() : [];
  _restoreMessages(session.messages, currentAgentIconHTML);
  renderSessionsList(agentId);
  scrollBottom();
}

function deleteSession(e, agentId, sessionId) {
  e.stopPropagation();
  if (!chatSessions[agentId]) return;
  chatSessions[agentId] = chatSessions[agentId].filter(function(s){ return s.id !== sessionId; });
  _saveSessions();
  if (currentSessionId[agentId] === sessionId) {
    currentSessionId[agentId] = null;
    currentMessages = [];
    var def = AGENT_DEFS.find(function(d){ return d.id === agentId; });
    _showWelcome(def ? def.name : 'Agent', currentAgentIconHTML, def ? def.role : '');
  }
  renderSessionsList(agentId);
}

function _saveCurrentSession() {
  if (!currentAgent) return;
  // Never persist a session that only has the hidden greeting trigger (no real user message)
  var hasRealMsg = currentMessages.some(function(m) {
    return m.role === 'user' && (typeof m.content === 'string' ? m.content : '') !== '__GREET__';
  });
  if (!hasRealMsg) return;
  if (!chatSessions[currentAgent]) chatSessions[currentAgent] = [];
  var sessions = chatSessions[currentAgent];
  var sid = currentSessionId[currentAgent];
  if (sid) {
    var idx = sessions.findIndex(function(s){ return s.id === sid; });
    if (idx >= 0) { sessions[idx].messages = currentMessages.slice(); sessions[idx].updatedAt = Date.now(); }
  } else {
    var firstUser = currentMessages.find(function(m){ return m.role === 'user'; });
    var raw = firstUser ? (typeof firstUser.content === 'string' ? firstUser.content
      : (firstUser.content.filter(function(x){return x.type==='text';})[0]||{}).text||'') : '';
    var title = (raw || 'New Chat').slice(0, 44) + (raw.length > 44 ? '…' : '');
    var ns = { id: _genId(), title: title, messages: currentMessages.slice(), updatedAt: Date.now() };
    sessions.push(ns);
    currentSessionId[currentAgent] = ns.id;
  }
  _saveSessions();
  renderSessionsList(currentAgent);
}

// ── LOAD STATS ─────────────────────────────────────────────────────────
function loadStats() {
  fetch('/api/stats')
    .then(function(r) { return r.json(); })
    .then(function(d) {
      var _set = function(id, val) { var el = document.getElementById(id); if (el) el.textContent = val; };
      _set('stat-total',   d.total_leads  != null ? d.total_leads  : '—');
      _set('stat-active',  d.active_leads != null ? d.active_leads : '—');
      _set('stat-hot',     d.hot_leads    != null ? d.hot_leads    : '—');
      _set('stat-warm',    d.warm_leads   != null ? d.warm_leads   : '—');
      _set('stat-cold',    d.cold_leads   != null ? d.cold_leads   : '—');
      renderProgBars(d.programs || []);
      renderUpcoming(d.upcoming || []);
    })
    .catch(function() {
      ['stat-total','stat-active','stat-hot','stat-warm','stat-cold'].forEach(function(id){ var el=document.getElementById(id); if(el) el.textContent='—'; });
      renderProgBars([]);
    });
}

function renderProgBars(programs) {
  var container = document.getElementById('prog-bars');
  if (!container) return;
  if (!programs.length) {
    container.innerHTML = '<div class="empty-state" style="padding:20px"><div class="empty-icon"><i class="fa-solid fa-chart-bar"></i></div><p>No program data yet</p></div>';
    return;
  }
  var max = Math.max.apply(null, programs.map(function(p) { return p.count || 0; })) || 1;
  container.innerHTML = programs.map(function(p) {
    var pct = Math.round(((p.count || 0) / max) * 100);
    return '<div class="prog-item"><div class="prog-header"><span class="prog-label">' + escHtml(p.name) + '</span><span class="prog-val">' + (p.count || 0) + '</span></div><div class="prog-bar"><div class="prog-fill" data-pct="' + pct + '" style="width:0%"></div></div></div>';
  }).join('');
  setTimeout(function() {
    document.querySelectorAll('.prog-fill').forEach(function(el) {
      el.style.width = (el.dataset.pct || 0) + '%';
    });
  }, 100);
}

function renderUpcoming(items) {
  var container = document.getElementById('upcoming-posts');
  if (!container) return;
  if (!items.length) {
    container.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-calendar-days"></i></div><p>No upcoming posts scheduled</p></div>';
    return;
  }
  var today = new Date().toISOString().slice(0, 10);
  container.innerHTML = items.slice(0, 7).map(function(item) {
    var isToday = item.date === today;
    var cls = isToday ? 'post-item today' : 'post-item';
    var badge = isToday ? '<span class="badge badge-today">TODAY</span>' :
                item.status === 'Posted' ? '<span class="badge badge-posted">Posted</span>' :
                '<span class="badge badge-scheduled">Scheduled</span>';
    var dateObj = new Date(item.date);
    var dayStr = dateObj.toLocaleDateString('en-PH', { weekday: 'short' });
    var dateStr = dateObj.toLocaleDateString('en-PH', { month: 'short', day: 'numeric' });
    return '<div class="' + cls + '"><div class="post-date-badge">' + dateStr + '</div><div class="post-dot"></div><div class="post-info"><div class="post-filename">' + escHtml(item.image || 'post_' + item.date + '.jpg') + '</div><div class="post-day">' + dayStr + '</div></div>' + badge + '</div>';
  }).join('');
}

// ── DRIP SCHEDULE ──────────────────────────────────────────────────────
function loadDripSchedule() {
  var container = document.getElementById('drip-schedule');
  if (!container) return;
  fetch('/api/drip-schedule')
    .then(function(r) { return r.json(); })
    .then(function(leads) {
      if (!leads.length) {
        container.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-envelope"></i></div><p>No active leads in drip campaign</p></div>';
        return;
      }
      var html = '<table style="width:100%;border-collapse:collapse">';
      html += '<thead><tr>';
      html += '<th style="text-align:left;padding:8px 12px;font-size:10px;font-weight:700;letter-spacing:1.5px;text-transform:uppercase;color:var(--text-dim);border-bottom:1px solid var(--card-border)">Lead</th>';
      html += '<th style="text-align:left;padding:8px 12px;font-size:10px;font-weight:700;letter-spacing:1.5px;text-transform:uppercase;color:var(--text-dim);border-bottom:1px solid var(--card-border)">Program</th>';
      html += '<th style="text-align:center;padding:8px 12px;font-size:10px;font-weight:700;letter-spacing:1.5px;text-transform:uppercase;color:var(--text-dim);border-bottom:1px solid var(--card-border)">Sent</th>';
      html += '<th style="text-align:left;padding:8px 12px;font-size:10px;font-weight:700;letter-spacing:1.5px;text-transform:uppercase;color:var(--text-dim);border-bottom:1px solid var(--card-border)">Next Email</th>';
      html += '</tr></thead><tbody>';
      leads.forEach(function(l) {
        var imgShort = l.next_img.replace(/\.[^.]+$/, '');
        html += '<tr style="transition:background .15s" onmouseover="this.style.background=\'var(--gold-dim)\'" onmouseout="this.style.background=\'transparent\'">';
        html += '<td style="padding:10px 12px"><div style="font-weight:600;color:var(--text)">' + escHtml(l.name.trim()) + '</div><div style="font-size:11px;color:var(--text-dim)">' + escHtml(l.email) + '</div></td>';
        html += '<td style="padding:10px 12px;font-size:12px;color:var(--text-muted)">' + escHtml(l.program || '—') + '</td>';
        html += '<td style="padding:10px 12px;text-align:center"><span style="display:inline-flex;align-items:center;justify-content:center;width:28px;height:28px;background:var(--gold-dim);border-radius:50%;font-size:12px;font-weight:700;color:var(--gold)">' + l.sent + '</span></td>';
        html += '<td style="padding:10px 12px"><div style="display:flex;align-items:center;gap:8px"><span style="background:var(--gold);color:var(--navy-d);font-size:9px;font-weight:800;padding:2px 7px;border-radius:10px">#' + l.next_num + '</span><span style="font-size:12px;color:var(--text)">' + escHtml(imgShort) + '</span></div></td>';
        html += '</tr>';
      });
      html += '</tbody></table>';
      container.innerHTML = html;
    })
    .catch(function() {
      container.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-triangle-exclamation"></i></div><p>Could not load drip schedule</p></div>';
    });
}

// ── LOAD LEADS ─────────────────────────────────────────────────────────
var _allLeads = [];
var _filteredLeads = [];
var _leadsPage = 1;
var _leadsPerPage = 15;
var _leadsSortCol = null;
var _leadsSortDir = 1;

function loadLeads() {
  var tbody = document.getElementById('leads-tbody');
  tbody.innerHTML = '<tr><td colspan="6" style="text-align:center;color:var(--text-dim);padding:30px">Loading leads...</td></tr>';
  fetch('/api/leads')
    .then(function(r) { return r.json(); })
    .then(function(leads) {
      _allLeads = leads;
      _filteredLeads = leads.slice();
      _leadsPage = 1;
      _renderLeadsPage();
    })
    .catch(function() {
      tbody.innerHTML = '<tr><td colspan="6" style="text-align:center;color:var(--text-dim);padding:30px">Could not load leads.</td></tr>';
    });
}

function _rowHtml(l) {
  var cnt = parseInt(l.email_count || l.emails_sent || 0);
  var scoreBadge = cnt >= 10
    ? '<span class="badge badge-hot"><i class="fa-solid fa-fire"></i> Hot</span>'
    : cnt >= 5
    ? '<span class="badge badge-warm"><i class="fa-solid fa-temperature-half"></i> Warm</span>'
    : '<span class="badge badge-cold"><i class="fa-solid fa-snowflake"></i> Cold</span>';
  var statusBadge = (l.active === false || l.status === 'inactive')
    ? '<span class="badge badge-inactive">Inactive</span>'
    : '<span class="badge badge-active">Active</span>';
  var leadData = escHtml(JSON.stringify({
    name: ((l.first_name||'') + ' ' + (l.last_name||'')).trim(),
    email: l.email||'',
    program: l.program_interest || l.program || '—',
    status: l.status || 'active',
    email_count: cnt
  }));
  return '<tr style="cursor:pointer" data-lead=\'' + leadData + '\' onclick="openLeadModal(this)">' +
    '<td>' + escHtml((l.first_name||'') + ' ' + (l.last_name||'')) + '</td>' +
    '<td>' + escHtml(l.email||'') + '</td>' +
    '<td>' + escHtml(l.program_interest || l.program || '—') + '</td>' +
    '<td style="text-align:center">' + cnt + '</td>' +
    '<td>' + scoreBadge + '</td>' +
    '<td>' + statusBadge + '</td></tr>';
}

function _renderLeadsPage() {
  var tbody = document.getElementById('leads-tbody');
  var total = _filteredLeads.length;
  if (!total) {
    tbody.innerHTML = '<tr><td colspan="6" style="text-align:center;color:var(--text-dim);padding:30px">No leads found.</td></tr>';
    var pg = document.getElementById('leads-pagination'); if (pg) pg.innerHTML = '';
    return;
  }
  var start = (_leadsPage - 1) * _leadsPerPage;
  var end = Math.min(start + _leadsPerPage, total);
  tbody.innerHTML = _filteredLeads.slice(start, end).map(_rowHtml).join('');
  _renderPagination(total);
}

function _renderPagination(total) {
  var el = document.getElementById('leads-pagination');
  if (!el) return;
  var pages = Math.ceil(total / _leadsPerPage);
  var from = (_leadsPage - 1) * _leadsPerPage + 1;
  var to   = Math.min(_leadsPage * _leadsPerPage, total);
  if (pages <= 1) { el.innerHTML = '<span class="pg-info">' + total + ' lead' + (total !== 1 ? 's' : '') + '</span>'; return; }
  var nums = '';
  var lo = Math.max(1, _leadsPage - 2), hi = Math.min(pages, _leadsPage + 2);
  if (lo > 1) nums += '<button class="pg-btn" onclick="leadsGoPage(1)">1</button>' + (lo > 2 ? '<span style="padding:0 4px;color:var(--text-dim)">…</span>' : '');
  for (var i = lo; i <= hi; i++) {
    nums += '<button class="pg-btn' + (i === _leadsPage ? ' active' : '') + '" onclick="leadsGoPage(' + i + ')">' + i + '</button>';
  }
  if (hi < pages) nums += (hi < pages - 1 ? '<span style="padding:0 4px;color:var(--text-dim)">…</span>' : '') + '<button class="pg-btn" onclick="leadsGoPage(' + pages + ')">' + pages + '</button>';
  el.innerHTML =
    '<span class="pg-info">' + from + '–' + to + ' of ' + total + ' leads</span>' +
    '<div class="pg-btns">' +
      '<button class="pg-btn" onclick="leadsGoPage(' + (_leadsPage - 1) + ')"' + (_leadsPage <= 1 ? ' disabled' : '') + '><i class="fa-solid fa-chevron-left"></i></button>' +
      nums +
      '<button class="pg-btn" onclick="leadsGoPage(' + (_leadsPage + 1) + ')"' + (_leadsPage >= pages ? ' disabled' : '') + '><i class="fa-solid fa-chevron-right"></i></button>' +
    '</div>';
}

function leadsGoPage(p) {
  var pages = Math.ceil(_filteredLeads.length / _leadsPerPage);
  _leadsPage = Math.max(1, Math.min(p, pages));
  _renderLeadsPage();
}

function sortLeads(col) {
  if (_leadsSortCol === col) { _leadsSortDir *= -1; }
  else { _leadsSortCol = col; _leadsSortDir = 1; }
  _filteredLeads.sort(function(a, b) {
    var av, bv;
    if (col === 'name')    { av = ((a.first_name||'') + ' ' + (a.last_name||'')).toLowerCase(); bv = ((b.first_name||'') + ' ' + (b.last_name||'')).toLowerCase(); }
    else if (col === 'program') { av = (a.program_interest||a.program||'').toLowerCase(); bv = (b.program_interest||b.program||'').toLowerCase(); }
    else if (col === 'emails' || col === 'score') { av = parseInt(a.email_count||a.emails_sent||0); bv = parseInt(b.email_count||b.emails_sent||0); }
    else if (col === 'status') { av = (a.active === false || a.status === 'inactive') ? 0 : 1; bv = (b.active === false || b.status === 'inactive') ? 0 : 1; }
    if (av < bv) return -1 * _leadsSortDir;
    if (av > bv) return  1 * _leadsSortDir;
    return 0;
  });
  ['name','program','emails','score','status'].forEach(function(c) {
    var th = document.getElementById('th-' + c);
    if (!th) return;
    var icon = th.querySelector('.sort-icon');
    th.classList.remove('asc','desc');
    if (icon) icon.className = 'fa-solid fa-sort sort-icon';
    if (c === col) {
      th.classList.add(_leadsSortDir === 1 ? 'asc' : 'desc');
      if (icon) icon.className = 'fa-solid ' + (_leadsSortDir === 1 ? 'fa-sort-up' : 'fa-sort-down') + ' sort-icon';
    }
  });
  _leadsPage = 1;
  _renderLeadsPage();
}

function renderLeads(leads) {
  _filteredLeads = leads.slice();
  _leadsPage = 1;
  _renderLeadsPage();
}

function filterLeads(q) {
  var base = _allLeads;
  if (q) {
    var ql = q.toLowerCase();
    base = base.filter(function(l) {
      return ((l.first_name||'') + ' ' + (l.last_name||'')).toLowerCase().includes(ql) ||
             (l.email||'').toLowerCase().includes(ql) ||
             (l.program_interest||l.program||'').toLowerCase().includes(ql);
    });
  }
  _filteredLeads = base;
  _leadsPage = 1;
  if (_leadsSortCol) sortLeads(_leadsSortCol); else _renderLeadsPage();
}
function exportLeadsCSV() {
  var data = _allLeads;
  if (!data.length) { toast(TI.warn,'No leads to export'); return; }
  var keys = Object.keys(data[0]);
  var csv = [keys.join(',')].concat(data.map(function(r){ return keys.map(function(k){ return '"'+(r[k]||'')+'"'; }).join(','); })).join('\n');
  var a = document.createElement('a');
  a.href = 'data:text/csv;charset=utf-8,' + encodeURIComponent(csv);
  a.download = 'leads_' + new Date().toISOString().slice(0,10) + '.csv';
  a.click();
  toast(TI.ok,'Leads exported!');
}

// ── LOAD SCHEDULE ──────────────────────────────────────────────────────
function loadSchedule() {
  var tbody = document.getElementById('schedule-tbody');
  tbody.innerHTML = '<tr><td colspan="4" style="text-align:center;color:var(--text-dim);padding:30px">Loading schedule...</td></tr>';
  fetch('/api/schedule')
    .then(function(r) { return r.json(); })
    .then(function(sched) {
      var keys = Object.keys(sched).sort();
      if (!keys.length) {
        tbody.innerHTML = '<tr><td colspan="4" style="text-align:center;color:var(--text-dim);padding:30px">No schedule data found.</td></tr>';
        return;
      }
      var today = new Date().toISOString().slice(0, 10);
      tbody.innerHTML = keys.map(function(dt) {
        var isToday = dt === today;
        var isPast  = dt < today;
        var dateObj = new Date(dt + 'T00:00:00');
        var dateStr = dateObj.toLocaleDateString('en-PH', { year: 'numeric', month: 'short', day: 'numeric' });
        var dayStr  = dateObj.toLocaleDateString('en-PH', { weekday: 'long' });
        var badge   = isToday ? '<span class="badge badge-today">Today</span>' :
                      isPast  ? '<span class="badge badge-posted">Posted</span>' :
                                '<span class="badge badge-scheduled">Scheduled</span>';
        return '<tr' + (isToday ? ' style="background:var(--gold-dim)"' : '') + '><td>' + dateStr + '</td><td>' + dayStr + '</td><td>' + escHtml(sched[dt] || '—') + '</td><td>' + badge + '</td></tr>';
      }).join('');
    })
    .catch(function() {
      tbody.innerHTML = '<tr><td colspan="4" style="text-align:center;color:var(--text-dim);padding:30px">Could not load schedule.</td></tr>';
    });
}

// ── FORM ───────────────────────────────────────────────────────────────
function toggleForm() {
  var form = document.getElementById('add-form');
  form.classList.toggle('open');
  if (form.classList.contains('open')) {
    document.getElementById('f-fname').focus();
  }
}

function addLead() {
  var fname = document.getElementById('f-fname').value.trim();
  var lname = document.getElementById('f-lname').value.trim();
  var email = document.getElementById('f-email').value.trim();
  var program = document.getElementById('f-program').value.trim();
  if (!fname || !email) { toast(TI.warn, 'First name and email are required.'); return; }
  fetch('/api/leads', {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify({ first_name: fname, last_name: lname, email: email, program: program })
  })
    .then(function(r) { return r.json(); })
    .then(function(d) {
      if (d.success || d.status === 'ok') {
        toast(TI.ok, 'Lead added successfully!');
        document.getElementById('f-fname').value = '';
        document.getElementById('f-lname').value = '';
        document.getElementById('f-email').value = '';
        document.getElementById('f-program').value = '';
        toggleForm();
        loadLeads();
      } else {
        toast(TI.err, d.message || 'Failed to add lead.');
      }
    })
    .catch(function() { toast(TI.err, 'Error saving lead. Check server.'); });
}

// ── FACEBOOK POST MODAL ────────────────────────────────────────────────
function runPost() { openFbModal(); }

function openFbModal() {
  var overlay = document.getElementById('fb-modal-overlay');
  overlay.classList.add('open');
  document.getElementById('fb-caption').value  = '';
  document.getElementById('fb-gen-btn').innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption';
  document.getElementById('fb-gen-btn').disabled = false;
  var sel = document.getElementById('fb-image-select');
  sel.value = '';
  _loadScheduleIntoSelect(sel, 'fb-cs-wrap', fbImageChanged);
}

function _fmtDate(iso) {
  var d = new Date(iso + 'T12:00:00');
  return d.toLocaleDateString('en-PH', { month: 'short', day: 'numeric' });
}

function toggleCsSel(wrapId) {
  var wrap = document.getElementById(wrapId);
  if (!wrap) return;
  var isOpen = wrap.classList.contains('open');
  document.querySelectorAll('.cs-wrap.open').forEach(function(w){ w.classList.remove('open'); });
  if (!isOpen) wrap.classList.add('open');
}

document.addEventListener('click', function(e) {
  if (!e.target.closest('.cs-wrap')) {
    document.querySelectorAll('.cs-wrap.open').forEach(function(w){ w.classList.remove('open'); });
  }
});

function _buildCustomDropdown(csWrapId, entries, today, hiddenSel, onChanged) {
  var wrap = document.getElementById(csWrapId);
  if (!wrap) return;
  var listEl   = wrap.querySelector('.cs-list');
  var textEl   = wrap.querySelector('.cs-trigger-text');
  var badgeEl  = wrap.querySelector('.cs-trigger-badge');
  var todayEntry = entries.find(function(e){ return e[0] === today; });
  var others     = entries.filter(function(e){ return e[0] !== today; });
  var sorted     = todayEntry ? [todayEntry].concat(others) : entries;

  // Populate hidden <select> so .value works correctly
  hiddenSel.innerHTML = entries.map(function(e){
    return '<option value="' + escHtml(e[1]) + '">' + escHtml(e[0] + ' — ' + e[1]) + '</option>';
  }).join('');

  function selectEntry(e) {
    hiddenSel.value = e[1];
    var name = e[1].replace(/\.[^.]+$/, '');
    if (textEl) textEl.textContent = name;
    if (badgeEl) {
      var isToday = e[0] === today;
      badgeEl.textContent  = isToday ? 'TODAY' : _fmtDate(e[0]);
      badgeEl.className    = 'cs-trigger-badge' + (isToday ? ' today-badge' : '');
      badgeEl.style.display = '';
    }
    listEl.querySelectorAll('.cs-option').forEach(function(o){ o.classList.remove('cs-selected'); });
    var target = listEl.querySelector('[data-val="' + CSS.escape(e[1]) + '"]');
    if (target) target.classList.add('cs-selected');
    wrap.classList.remove('open');
    if (onChanged) onChanged();
  }

  listEl.innerHTML = sorted.map(function(e) {
    var isToday   = e[0] === today;
    var shortDate = isToday ? 'TODAY' : _fmtDate(e[0]);
    var name      = e[1].replace(/\.[^.]+$/, '');
    return '<div class="cs-option' + (isToday ? ' cs-today' : '') + '" data-val="' + escHtml(e[1]) + '" data-iso="' + escHtml(e[0]) + '">' +
      (isToday ? '<i class="fa-solid fa-star cs-today-star"></i>' : '<span class="cs-opt-dot"></span>') +
      '<span class="cs-opt-date-badge">' + escHtml(shortDate) + '</span>' +
      '<span class="cs-opt-name">' + escHtml(name) + '</span>' +
    '</div>';
  }).join('');

  listEl.querySelectorAll('.cs-option').forEach(function(el) {
    el.addEventListener('click', function(ev) {
      ev.stopPropagation();
      selectEntry({ 0: el.dataset.iso, 1: el.dataset.val });
    });
  });

  // Auto-select today or first
  if (todayEntry) {
    selectEntry(todayEntry);
  } else if (entries.length) {
    selectEntry(entries[0]);
  }
}

function _loadScheduleIntoSelect(sel, csWrapId, onChanged) {
  fetch('/api/schedule')
    .then(function(r){
      if (r.status === 401) { window.location.href = '/login'; return null; }
      if (!r.ok) throw new Error('HTTP ' + r.status);
      return r.json();
    })
    .then(function(sched){
      if (!sched) return;
      var today   = new Date().toISOString().slice(0,10);
      var entries = Object.entries(sched).sort(function(a,b){ return a[0].localeCompare(b[0]); });
      if (!entries.length) {
        var wrap = document.getElementById(csWrapId);
        if (wrap) { var t = wrap.querySelector('.cs-trigger-text'); if(t) t.textContent = 'No images scheduled'; }
        return;
      }
      _buildCustomDropdown(csWrapId, entries, today, sel, onChanged);
    })
    .catch(function(e){
      var wrap = document.getElementById(csWrapId);
      if (wrap) { var t = wrap.querySelector('.cs-trigger-text'); if(t) t.textContent = 'Error: ' + String(e); }
    });
}

function closeFbModal() {
  document.getElementById('fb-modal-overlay').classList.remove('open');
}
function closeFbModalBg(e) {
  if (e.target === document.getElementById('fb-modal-overlay')) closeFbModal();
}

function fbImageChanged() {
  var img   = document.getElementById('fb-image-select').value;
  var prev  = document.getElementById('fb-image-preview');
  var empty = document.getElementById('fb-img-empty');
  if (img && prev) {
    prev.onload = function() { prev.classList.add('loaded'); if (empty) empty.style.display = 'none'; };
    prev.onerror = function() { prev.classList.remove('loaded'); if (empty) empty.style.display = 'flex'; };
    prev.src = '/api/image-preview/' + encodeURIComponent(img);
  } else {
    if (prev)  { prev.classList.remove('loaded'); prev.src = ''; }
    if (empty) { empty.style.display = 'flex'; }
  }
}

function generateFbCaption() {
  var img = document.getElementById('fb-image-select').value;
  if (!img) { toast(TI.warn, 'Select an image first.'); return; }
  var btn = document.getElementById('fb-gen-btn');
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Generating...';
  btn.disabled = true;
  fetch('/api/generate-caption?image=' + encodeURIComponent(img))
    .then(function(r){ return r.json(); })
    .then(function(d){
      if (d.error) {
        btn.innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption';
        btn.disabled = false;
        toast(TI.err, 'Caption failed: ' + d.error);
        return;
      }
      document.getElementById('fb-caption').value  = d.caption  || '';
      document.getElementById('fb-hashtags').value = d.hashtags || '';
      btn.innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Regenerate';
      btn.disabled = false;
    })
    .catch(function(e){
      btn.innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption';
      btn.disabled = false;
      toast(TI.err, 'Caption generation failed: ' + e);
    });
}

function submitFbPost() {
  var caption  = document.getElementById('fb-caption').value.trim();
  var hashtags = document.getElementById('fb-hashtags').value.trim();
  var image    = document.getElementById('fb-image-select').value;
  if (!caption) { toast(TI.warn, 'Please write or generate a caption first.'); return; }
  var btn = document.getElementById('fb-submit-btn');
  btn.disabled = true;
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Posting...';
  fetch('/api/post-facebook', {
    method: 'POST',
    headers: {'Content-Type':'application/json'},
    body: JSON.stringify({ caption: caption, hashtags: hashtags, image: image })
  })
    .then(function(r){ return r.json(); })
    .then(function(d){
      btn.disabled = false;
      btn.innerHTML = '<i class="fa-brands fa-facebook-f"></i> Post to Facebook';
      closeFbModal();
      if (d.status === 'ok') {
        toast(TI.ok, 'Facebook post sent successfully!');
      } else {
        toast(TI.err, 'Post failed. Check output for details.');
        if (d.output) showModal('Post Output', d.output);
      }
    })
    .catch(function(){
      btn.disabled = false;
      btn.innerHTML = '<i class="fa-brands fa-facebook-f"></i> Post to Facebook';
      toast(TI.err, 'Could not reach server.');
    });
}

// ── DRIP EMAIL MODAL ───────────────────────────────────────────────────
function runEmail() { openEmailModal(); }

function openEmailModal() {
  var overlay = document.getElementById('email-modal-overlay');
  overlay.classList.add('open');
  var body = document.getElementById('email-modal-body');
  body.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading leads...</p></div>';
  fetch('/api/leads')
    .then(function(r){ return r.json(); })
    .then(function(leads){
      var active = leads.filter(function(l){ return (l.status||'').toLowerCase() === 'active'; });
      var hot  = active.filter(function(l){ return parseInt(l.email_count||0) >= 10; }).length;
      var warm = active.filter(function(l){ var c=parseInt(l.email_count||0); return c>=5 && c<10; }).length;
      var cold = active.filter(function(l){ return parseInt(l.email_count||0) < 5; }).length;
      var preview = active.slice(0, 5).map(function(l){
        var step = parseInt(l.email_count||0) + 1;
        return '<div class="am-lead-row">' +
          '<div class="am-lead-name">' + escHtml((l.first_name||'') + ' ' + (l.last_name||'')) + '</div>' +
          '<div style="font-size:11px;color:var(--text-dim);flex:1">' + escHtml(l.program_interest||'') + '</div>' +
          '<div class="am-lead-step">Email #' + step + '</div>' +
        '</div>';
      }).join('');
      body.innerHTML =
        '<div style="display:flex;gap:10px;margin-bottom:16px">' +
          '<div class="am-summary-row" style="flex:1"><div class="am-summary-num">' + active.length + '</div><div class="am-summary-label">Active leads<br>will receive email</div></div>' +
          '<div style="display:flex;flex-direction:column;gap:6px;justify-content:center">' +
            '<div style="font-size:12px;color:var(--text-muted)"><i class="fa-solid fa-fire" style="color:#f87171"></i> Hot: <strong>' + hot + '</strong></div>' +
            '<div style="font-size:12px;color:var(--text-muted)"><i class="fa-solid fa-temperature-half" style="color:#fcd34d"></i> Warm: <strong>' + warm + '</strong></div>' +
            '<div style="font-size:12px;color:var(--text-muted)"><i class="fa-solid fa-snowflake" style="color:#93c5fd"></i> Cold: <strong>' + cold + '</strong></div>' +
          '</div>' +
        '</div>' +
        '<div class="am-label" style="margin-bottom:8px">Next in Queue (preview)</div>' +
        '<div class="am-lead-preview">' + (preview || '<div class="am-lead-row"><div class="am-lead-name" style="color:var(--text-dim)">No active leads</div></div>') + '</div>' +
        (active.length > 5 ? '<div style="font-size:11px;color:var(--text-dim);margin-top:8px;text-align:right">+' + (active.length-5) + ' more leads</div>' : '');
    })
    .catch(function(){
      body.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-triangle-exclamation"></i></div><p>Could not load leads</p></div>';
    });
}

function closeEmailModal() {
  document.getElementById('email-modal-overlay').classList.remove('open');
}
function closeEmailModalBg(e) {
  if (e.target === document.getElementById('email-modal-overlay')) closeEmailModal();
}

function submitDripEmail() {
  var btn = document.getElementById('email-submit-btn');
  btn.disabled = true;
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Sending...';
  fetch('/api/run-email', { method: 'POST' })
    .then(function(r){ return r.json(); })
    .then(function(d){
      btn.disabled = false;
      btn.innerHTML = '<i class="fa-solid fa-paper-plane"></i> Send to All Active Leads';
      closeEmailModal();
      if (d.status === 'ok') {
        toast(TI.ok, 'Drip emails queued — sending in ~1-2 min. Check Telegram for confirmation.');
      } else {
        toast(TI.err, 'Email trigger failed. Check output.');
        if (d.output) showModal('Email Output', d.output);
      }
    })
    .catch(function(){
      btn.disabled = false;
      btn.innerHTML = '<i class="fa-solid fa-paper-plane"></i> Send to All Active Leads';
      toast(TI.err, 'Could not reach server.');
    });
}

// ── SOCIAL LISTENING / WEEKLY PREVIEW ─────────────────────────────────
function runSocialListening() {
  toast(TI.ok, 'Social Listening runs automatically at 7:50 AM PHT via GitHub Actions. Check your email for today\'s briefing.');
}

function runWeeklyPreview() {
  toast(TI.ok, 'Weekly Campaign Preview runs every Sunday at 5:00 PM PHT via GitHub Actions. Check your email for the preview.');
}

// ── ATTACHMENT ─────────────────────────────────────────────────────────
var pendingAttachment = null; // { type: 'base64'|'url', data: '...', mediaType: '...', label: '...' }

function handleFileAttach(evt) {
  var file = evt.target.files[0];
  if (!file) return;
  evt.target.value = '';
  var isText = file.type === 'text/csv' || file.type === 'text/plain' || file.name.endsWith('.csv') || file.name.endsWith('.txt');
  var isPdf  = file.type === 'application/pdf' || file.name.endsWith('.pdf');
  if (isPdf) { handlePdfFile(file); return; }
  if (isText) {
    var r = new FileReader();
    r.onload = function(e) {
      pendingAttachment = { type: 'text', content: e.target.result, name: file.name };
      showAttachPreview('📄', file.name);
    };
    r.readAsText(file);
    return;
  }
  // Image
  var reader = new FileReader();
  reader.onload = function(e) {
    var dataUrl = e.target.result;
    var base64 = dataUrl.split(',')[1];
    var mediaType = file.type || 'image/jpeg';
    pendingAttachment = { type: 'base64', data: base64, mediaType: mediaType, label: file.name };
    showAttachPreview('🖼', file.name);
  };
  reader.readAsDataURL(file);
}

function promptImageUrl() {
  var url = prompt('Paste image URL:');
  if (!url || !url.trim()) return;
  url = url.trim();
  pendingAttachment = { type: 'url', data: url, mediaType: 'image/jpeg', label: url };
  showAttachPreview('🔗', url);
}

function showAttachPreview(icon, label) {
  var p = document.getElementById('attach-preview');
  document.getElementById('attach-icon').textContent = icon;
  document.getElementById('attach-label').textContent = label;
  p.style.display = 'flex';
}

function clearAttachment() {
  pendingAttachment = null;
  var p = document.getElementById('attach-preview');
  p.style.display = 'none';
}

function arrayBufferToBase64(buffer) {
  var bytes = new Uint8Array(buffer), binary = '', chunk = 8192;
  for (var i = 0; i < bytes.length; i += chunk)
    binary += String.fromCharCode.apply(null, bytes.subarray(i, i + chunk));
  return btoa(binary);
}

function handlePdfFile(file) {
  var r = new FileReader();
  r.onload = function(e) {
    pendingAttachment = { type: 'pdf', data: arrayBufferToBase64(e.target.result), name: file.name };
    showAttachPreview('📕', file.name + ' (PDF — AI will read it)');
  };
  r.readAsArrayBuffer(file);
}

function handleDroppedFile(file) {
  var isPdf  = file.type === 'application/pdf' || file.name.endsWith('.pdf');
  var isImg  = file.type.startsWith('image/');
  var isTxt  = file.type === 'text/csv' || file.type === 'text/plain' || file.name.endsWith('.csv') || file.name.endsWith('.txt');
  if (isPdf) { handlePdfFile(file); return; }
  if (isImg) {
    var reader = new FileReader();
    reader.onload = function(e) {
      var dataUrl = e.target.result;
      pendingAttachment = { type: 'base64', data: dataUrl.split(',')[1], mediaType: file.type || 'image/jpeg', label: file.name };
      showAttachPreview('🖼', file.name);
    };
    reader.readAsDataURL(file); return;
  }
  if (isTxt) {
    var r = new FileReader();
    r.onload = function(e) { pendingAttachment = { type: 'text', content: e.target.result, name: file.name }; showAttachPreview('📄', file.name); };
    r.readAsText(file); return;
  }
  toast(TI.warn, 'Drop images, PDFs, CSV or text files.');
}

// ── DRAG & DROP ────────────────────────────────────────────────────────
(function() {
  var col = document.querySelector('.chat-column');
  var overlay = document.getElementById('drag-overlay');
  if (!col || !overlay) return;
  col.addEventListener('dragover', function(e) { e.preventDefault(); overlay.style.display = 'flex'; });
  col.addEventListener('dragleave', function(e) { if (!col.contains(e.relatedTarget)) overlay.style.display = 'none'; });
  col.addEventListener('drop', function(e) {
    e.preventDefault(); overlay.style.display = 'none';
    var files = e.dataTransfer.files;
    if (files && files.length) handleDroppedFile(files[0]);
  });
})();

// ── CHAT ───────────────────────────────────────────────────────────────
function toolLabel(name) {
  var labels = {
    send_drip_emails:        'Sending drip emails to all leads',
    open_post_composer:      'Opening post composer',
    post_to_facebook:        'Posting to Facebook',
    post_to_instagram:       'Posting to Instagram',
    trigger_social_listening:'Running social listening',
    trigger_weekly_analytics:'Generating analytics report',
    trigger_weekly_preview:  'Generating campaign preview',
    get_leads:               'Reading leads data',
    get_posting_schedule:    'Reading posting schedule',
    read_drive_file:         'Reading file from Google Drive',
    list_drive_files:        'Listing Google Drive files',
    read_intelligence_briefing: 'Reading competitor intelligence briefings',
    web_search:              'Searching the web',
    get_calendar:            'Reading Google Calendar',
    add_calendar_event:      'Adding calendar event',
    search_email:            'Searching Gmail inbox',
    remember:                'Saving to memory',
    recall_memories:         'Recalling saved memories',
    forget_memory:           'Deleting memory',
    generate_image:          'Generating image with AI',
    execute_python:          'Running Python code'
  };
  return labels[name] || name;
}

// ── POST COMPOSER ───────────────────────────────────────────────────────────
var _compPlatform = 'facebook';
var _compSelected = null; // {filename, thumbUrl, date}
var _compTab = 'a';
var _compCaptions = {a:'', b:''};

function openPostComposer(platform) {
  _compPlatform = platform || 'facebook';
  _compSelected = null; _compTab = 'a'; _compCaptions = {a:'', b:''};
  var msgs = document.getElementById('chat-messages');
  var ex = document.getElementById('composer-widget'); if (ex) ex.remove();
  var wb = document.getElementById('chat-welcome'); if (wb) wb.remove();
  var icon = {facebook:'<i class="fa-brands fa-facebook"></i>', instagram:'<i class="fa-brands fa-instagram"></i>', both:'<i class="fa-solid fa-share-nodes"></i>'}[_compPlatform] || '';
  var label = {facebook:'Facebook', instagram:'Instagram', both:'Facebook + Instagram'}[_compPlatform] || _compPlatform;
  var el = document.createElement('div');
  el.id = 'composer-widget'; el.className = 'composer-msg';
  el.innerHTML =
    '<div class="composer-card">' +
      '<div class="composer-header">' +
        '<div class="comp-platform-badge">' + icon + ' ' + label + '</div>' +
        '<div class="comp-title">Create Post</div>' +
        '<button class="comp-close-btn" onclick="closeComposer()" title="Close"><i class="fa-solid fa-xmark"></i></button>' +
      '</div>' +
      // Step 1
      '<div class="comp-section" id="comp-sec-images">' +
        '<div class="comp-step-label"><span class="comp-step-num">1</span> Choose an image</div>' +
        '<div class="comp-image-grid" id="comp-image-grid">' +
          '<div class="comp-grid-loading"><div class="typing-indicator"><div class="typing-dot"></div><div class="typing-dot"></div><div class="typing-dot"></div></div><span style="margin-left:8px">Loading images…</span></div>' +
        '</div>' +
        '<label class="comp-upload-label"><i class="fa-solid fa-cloud-arrow-up"></i> Upload your own<input type="file" id="comp-upload-input" accept="image/*" style="display:none" onchange="compHandleUpload(this)"></label>' +
      '</div>' +
      // Step 2 (hidden until image selected)
      '<div class="comp-section" id="comp-sec-caption" style="display:none">' +
        '<div class="comp-step-label"><span class="comp-step-num">2</span> Caption</div>' +
        '<div class="comp-caption-tabs">' +
          '<button id="comp-tab-a" class="comp-tab active" onclick="compSwitchTab(\'a\')">Version A</button>' +
          '<button id="comp-tab-b" class="comp-tab" onclick="compSwitchTab(\'b\')">Version B</button>' +
          '<button class="comp-regen-btn" onclick="compRegenCaption()" title="Regenerate captions"><i class="fa-solid fa-rotate-right"></i> Regen</button>' +
        '</div>' +
        '<div id="comp-caption-loading" class="comp-caption-loading"><div class="typing-indicator"><div class="typing-dot"></div><div class="typing-dot"></div><div class="typing-dot"></div></div><span style="margin-left:8px">Generating captions with AI…</span></div>' +
        '<textarea id="comp-caption-a" class="comp-caption-input" style="display:none" placeholder="Version A" oninput="_compCaptions.a=this.value"></textarea>' +
        '<textarea id="comp-caption-b" class="comp-caption-input" style="display:none" placeholder="Version B" oninput="_compCaptions.b=this.value"></textarea>' +
      '</div>' +
      // Footer (hidden until image selected)
      '<div class="comp-footer" id="comp-footer" style="display:none">' +
        '<div class="comp-preview-row">' +
          '<img id="comp-preview-img" class="comp-preview-img" src="" alt="">' +
          '<div class="comp-preview-meta"><div class="comp-preview-fname" id="comp-preview-fname"></div><div class="comp-preview-date" id="comp-preview-date"></div></div>' +
        '</div>' +
        '<div class="comp-action-row">' +
          '<button class="comp-btn-post" id="comp-btn-post" onclick="compPost()"><i class="fa-solid fa-paper-plane"></i> Post Now</button>' +
          '<button class="comp-btn-decline" onclick="closeComposer()"><i class="fa-solid fa-xmark"></i> Cancel</button>' +
        '</div>' +
      '</div>' +
    '</div>';
  msgs.appendChild(el); scrollBottom();

  // Load images
  fetch('/api/composer/images')
    .then(function(r){ return r.json(); })
    .then(function(images) {
      var grid = document.getElementById('comp-image-grid');
      if (!grid) return;
      if (!images || !images.length) { grid.innerHTML = '<div class="comp-grid-empty">No images found in schedule or staging folder.</div>'; return; }
      grid.innerHTML = images.slice(0, 40).map(function(img) {
        var badge = img.is_today
          ? '<span class="comp-today-badge">Today</span>'
          : (img.date ? '<span class="comp-date-badge">' + escHtml(img.date.slice(5)) + '</span>' : '');
        var imgTag = '<img src="' + escHtml(img.thumb_url) + '" alt="" loading="lazy" data-src="' + escHtml(img.thumb_url) + '" onerror="compThumbRetry(this,4)">';
        return '<div class="comp-img-item"' +
          ' data-filename="' + escHtml(img.filename||'') + '"' +
          ' data-thumb="' + escHtml(img.thumb_url||'') + '"' +
          ' data-date="' + escHtml(img.date||'') + '"' +
          ' onclick="compSelectImage(this.dataset.filename,this.dataset.thumb,this.dataset.date,this)">' +
          '<div class="comp-img-wrap">' + imgTag + badge + '<div class="comp-img-check"><i class="fa-solid fa-circle-check"></i></div></div>' +
          '<div class="comp-img-name">' + escHtml((img.filename||'').replace(/\.[^.]+$/,'').slice(0,30)) + '</div>' +
          '</div>';
      }).join('');
      scrollBottom();
    })
    .catch(function(){ var grid=document.getElementById('comp-image-grid'); if(grid) grid.innerHTML='<div class="comp-grid-empty">Could not load images.</div>'; });
}

function compThumbRetry(img, left) {
  if (!left) { img.style.display = 'none'; img.parentElement.classList.add('comp-img-missing'); return; }
  setTimeout(function() {
    img.style.display = '';
    img.parentElement.classList.remove('comp-img-missing');
    var base = (img.getAttribute('data-src') || img.src).split('?')[0];
    img.onerror = function() { compThumbRetry(img, left - 1); };
    img.src = base + '?r=' + Date.now();
  }, 4000);
}

function compSelectImage(filename, thumbUrl, date, itemEl, b64) {
  _compSelected = {filename: filename, thumbUrl: thumbUrl, date: date, b64: b64};
  document.querySelectorAll('.comp-img-item').forEach(function(e){ e.classList.remove('selected'); });
  itemEl.classList.add('selected');
  var sec = document.getElementById('comp-sec-caption'); if (sec) sec.style.display = 'block';
  var ft  = document.getElementById('comp-footer');       if (ft)  ft.style.display  = 'block';
  var pi  = document.getElementById('comp-preview-img');  if (pi)  pi.src = thumbUrl;
  var pf  = document.getElementById('comp-preview-fname'); if (pf) pf.textContent = filename;
  var pd  = document.getElementById('comp-preview-date');  if (pd) pd.textContent = date ? 'Scheduled: ' + date : 'Unscheduled';
  compGenerateCaptions(filename);
  // Scroll composer card down to show caption section
  setTimeout(function() {
    var captionSec = document.getElementById('comp-sec-caption');
    if (captionSec) captionSec.scrollIntoView({behavior:'smooth', block:'nearest'});
  }, 100);
}

function compGenerateCaptions(filename) {
  var load = document.getElementById('comp-caption-loading');
  var ca   = document.getElementById('comp-caption-a');
  var cb   = document.getElementById('comp-caption-b');
  if (load) load.style.display = 'flex';
  if (ca)   { ca.style.display = 'none'; ca.value = ''; }
  if (cb)   { cb.style.display = 'none'; cb.value = ''; }
  _compCaptions = {a:'', b:''};
  fetch('/api/composer/caption', {
    method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({filename: filename, platform: _compPlatform})
  })
  .then(function(r){ return r.json(); })
  .then(function(data) {
    if (load) load.style.display = 'none';
    _compCaptions.a = data.caption_a || '';
    _compCaptions.b = data.caption_b || '';
    if (ca) { ca.value = _compCaptions.a; }
    if (cb) { cb.value = _compCaptions.b; }
    compSwitchTab(_compTab);
    scrollBottom();
  })
  .catch(function() {
    if (load) load.innerHTML = '<span style="color:#f87171">Could not generate captions. Try again.</span>';
  });
}

function compSwitchTab(tab) {
  _compTab = tab;
  var ta = document.getElementById('comp-tab-a');
  var tb = document.getElementById('comp-tab-b');
  var ca = document.getElementById('comp-caption-a');
  var cb = document.getElementById('comp-caption-b');
  if (ta) ta.classList.toggle('active', tab === 'a');
  if (tb) tb.classList.toggle('active', tab === 'b');
  if (ca) ca.style.display = tab === 'a' ? 'block' : 'none';
  if (cb) cb.style.display = tab === 'b' ? 'block' : 'none';
}

function compRegenCaption() {
  if (!_compSelected) return;
  compGenerateCaptions(_compSelected.filename);
}

function compHandleUpload(input) {
  if (!input.files || !input.files[0]) return;
  var file = input.files[0];
  var url  = URL.createObjectURL(file);
  var reader = new FileReader();
  reader.onload = function(e) {
    var b64 = e.target.result.split(',')[1];
    var grid = document.getElementById('comp-image-grid');
    if (grid) {
      var uploadItem = document.createElement('div');
      uploadItem.className = 'comp-img-item';
      uploadItem.innerHTML = '<div class="comp-img-wrap"><img src="' + url + '" alt=""><span class="comp-today-badge">Upload</span><div class="comp-img-check"><i class="fa-solid fa-circle-check"></i></div></div><div class="comp-img-name">' + escHtml(file.name.slice(0,30)) + '</div>';
      uploadItem.onclick = function() { compSelectImage(file.name, url, '', uploadItem, b64); };
      document.querySelectorAll('.comp-img-item').forEach(function(el){ el.classList.remove('selected'); });
      grid.insertBefore(uploadItem, grid.firstChild);
    }
    compSelectImage(file.name, url, '', grid ? grid.firstChild : document.createElement('div'), b64);
  };
  reader.readAsDataURL(file);
}

function compPost() {
  if (!_compSelected) { toast && toast(TI && TI.warn, 'Select an image first.'); return; }
  var activeEl = _compTab === 'a' ? document.getElementById('comp-caption-a') : document.getElementById('comp-caption-b');
  var caption  = activeEl ? activeEl.value.trim() : '';
  if (!caption) { toast && toast(TI && TI.warn, 'Caption is empty — add text first.'); return; }
  var btn = document.getElementById('comp-btn-post');
  if (btn) { btn.disabled = true; btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Posting…'; }
  fetch('/api/composer/post', {
    method:'POST', headers:{'Content-Type':'application/json'},
    body: JSON.stringify({filename: _compSelected.filename, caption: caption, platform: _compPlatform, b64: _compSelected.b64})
  })
  .then(function(r){ return r.json(); })
  .then(function(data) {
    var widget = document.getElementById('composer-widget');
    if (data.ok) {
      var platLabel = {facebook:'Facebook', instagram:'Instagram', both:'Facebook & Instagram'}[_compPlatform] || _compPlatform;
      if (widget) widget.innerHTML = '<div class="comp-success-card"><i class="fa-solid fa-circle-check"></i><div><strong>Posted to ' + platLabel + '!</strong><div style="font-size:11px;opacity:.65;margin-top:3px">' + escHtml(_compSelected.filename) + '</div></div></div>';
      var msgs = document.getElementById('chat-messages');
      var avatarEl = document.getElementById('chat-avatar');
      var avatarImg = avatarEl && avatarEl.querySelector('img');
      var agentEmoji = avatarImg ? '<img src="' + avatarImg.src + '" class="avatar-logo" alt="Agent">' : '<img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun">';
      var aiEl = document.createElement('div'); aiEl.className = 'msg';
      aiEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-body"><div class="msg-bubble">Your post is now live on ' + platLabel + '! The caption has been published with the selected image.</div><div class="msg-timestamp">' + _ts() + '</div></div>';
      msgs.appendChild(aiEl);
      currentMessages.push({role:'assistant', content:'Post published to ' + platLabel + '.'});
    } else {
      if (btn) { btn.disabled = false; btn.innerHTML = '<i class="fa-solid fa-paper-plane"></i> Retry'; }
      var errMsg = Object.values(data.results || {}).join(' | ');
      if (widget) {
        var errDiv = widget.querySelector('.comp-footer');
        if (errDiv) {
          var errInfo = document.createElement('div');
          errInfo.style.cssText = 'font-size:11px;color:#f87171;margin-bottom:8px;padding:6px 10px;background:rgba(220,38,38,.08);border-radius:6px;border:1px solid rgba(220,38,38,.2)';
          errInfo.textContent = errMsg.slice(0, 200);
          errDiv.insertBefore(errInfo, errDiv.firstChild);
        }
      }
    }
    scrollBottom();
  })
  .catch(function(e) {
    if (btn) { btn.disabled = false; btn.innerHTML = '<i class="fa-solid fa-paper-plane"></i> Retry'; }
  });
}

function closeComposer() {
  var w = document.getElementById('composer-widget'); if (w) w.remove();
}

function handleKey(e) {
  if (e.key === 'Enter' && !e.shiftKey) {
    e.preventDefault();
    sendMsg();
  }
}

// ── VOICE INPUT (Web Speech API) ───────────────────────────────────────
var _voiceRecog = null;
var _voiceActive = false;

function toggleVoiceInput() {
  if (!('webkitSpeechRecognition' in window) && !('SpeechRecognition' in window)) {
    toast(TI.warn, 'Voice input not supported in this browser. Try Chrome.');
    return;
  }
  if (_voiceActive) {
    if (_voiceRecog) _voiceRecog.stop();
    _stopVoice();
    return;
  }
  var SpeechRecognition = window.SpeechRecognition || window.webkitSpeechRecognition;
  _voiceRecog = new SpeechRecognition();
  _voiceRecog.continuous = false;
  _voiceRecog.interimResults = true;
  _voiceRecog.lang = 'en-PH';
  var input = document.getElementById('chat-input');
  var baseText = input.value;
  _voiceRecog.onstart = function() {
    _voiceActive = true;
    var btn = document.getElementById('voice-btn');
    if (btn) btn.classList.add('listening');
    toast(TI.ok, 'Listening…');
  };
  _voiceRecog.onresult = function(e) {
    var transcript = '';
    for (var i = 0; i < e.results.length; i++) {
      transcript += e.results[i][0].transcript;
    }
    input.value = (baseText ? baseText + ' ' : '') + transcript;
    autoResize(input);
  };
  _voiceRecog.onend = function() {
    _stopVoice();
    var txt = input.value.trim();
    if (txt) sendMsg();
  };
  _voiceRecog.onerror = function(e) {
    _stopVoice();
    if (e.error !== 'aborted') toast(TI.warn, 'Voice error: ' + e.error);
  };
  _voiceRecog.start();
}

function _stopVoice() {
  _voiceActive = false;
  _voiceRecog = null;
  var btn = document.getElementById('voice-btn');
  if (btn) btn.classList.remove('listening');
}

function autoResize(el) {
  el.style.height = 'auto';
  el.style.height = Math.min(el.scrollHeight, 140) + 'px';
}

function _ts() {
  var d = new Date(), h = d.getHours(), m = d.getMinutes();
  var ap = h >= 12 ? 'PM' : 'AM'; h = h % 12 || 12;
  return h + ':' + (m < 10 ? '0' : '') + m + ' ' + ap;
}

var _abortCtrl = null;

function sendMsg() {
  if (isStreaming) return;
  if (!currentAgent) { toast(TI.warn, 'Select an agent first.'); return; }
  var input = document.getElementById('chat-input');
  var text  = input.value.trim();
  if (!text && !pendingAttachment) return;
  input.value = ''; input.style.height = 'auto'; updateCharCount(input);
  var attachSnap = pendingAttachment;
  clearAttachment();

  var msgs = document.getElementById('chat-messages');
  var wb   = document.getElementById('chat-welcome');
  if (wb) wb.remove();

  // Build content for API
  var userContent;
  if (attachSnap) {
    if (attachSnap.type === 'text') {
      var fh = '[File: ' + attachSnap.name + ']\n```\n' + attachSnap.content.slice(0, 8000) + '\n```';
      userContent = text ? fh + '\n\n' + text : fh;
    } else if (attachSnap.type === 'pdf') {
      userContent = [];
      userContent.push({ type: 'document', source: { type: 'base64', media_type: 'application/pdf', data: attachSnap.data } });
      if (text) userContent.push({ type: 'text', text: text });
      else userContent.push({ type: 'text', text: 'Please read and analyze this PDF document.' });
    } else {
      userContent = [];
      if (attachSnap.type === 'base64') {
        userContent.push({ type: 'image', source: { type: 'base64', media_type: attachSnap.mediaType, data: attachSnap.data } });
      } else {
        userContent.push({ type: 'image', source: { type: 'url', url: attachSnap.data } });
      }
      if (text) userContent.push({ type: 'text', text: text });
    }
  } else {
    userContent = text;
  }
  var msgIdx = currentMessages.length;
  currentMessages.push({ role: 'user', content: userContent });

  // User bubble
  var userEl = document.createElement('div');
  userEl.className = 'msg user';
  var bc = '';
  if (attachSnap) {
    if (attachSnap.type === 'base64') bc += '<img src="data:' + attachSnap.mediaType + ';base64,' + attachSnap.data + '" style="max-width:220px;max-height:160px;border-radius:8px;display:block;margin-bottom:6px">';
    else if (attachSnap.type === 'url') bc += '<img src="' + escHtml(attachSnap.data) + '" style="max-width:220px;max-height:160px;border-radius:8px;display:block;margin-bottom:6px" onerror="this.style.display=\'none\'">';
    else if (attachSnap.type === 'text') bc += '<div style="font-size:11px;opacity:.7;margin-bottom:4px"><i class="fa-solid fa-paperclip"></i> ' + escHtml(attachSnap.name) + '</div>';
    else if (attachSnap.type === 'pdf') bc += '<div style="font-size:11px;opacity:.85;margin-bottom:4px;display:flex;align-items:center;gap:5px"><i class="fa-solid fa-file-pdf" style="color:#f87171"></i> ' + escHtml(attachSnap.name) + '</div>';
  }
  if (text) bc += escHtml(text);
  userEl.innerHTML =
    '<div class="msg-avatar user-av"><i class="fa-solid fa-user"></i></div>' +
    '<div class="msg-body">' +
      '<div class="msg-bubble">' + bc + '</div>' +
      '<div class="msg-actions">' +
        '<button class="msg-action-btn" onclick="editMsgAt(' + msgIdx + ')" title="Edit"><i class="fa-solid fa-pen-to-square"></i></button>' +
        '<button class="msg-action-btn" onclick="copyMsgText(this)" title="Copy"><i class="fa-solid fa-copy"></i></button>' +
      '</div>' +
      '<div class="msg-timestamp">' + _ts() + '</div>' +
    '</div>';
  msgs.appendChild(userEl);
  _updateMsgCounter();
  scrollBottom();
  _streamResponse();
}

// ── MESSAGE COUNTER ────────────────────────────────────────────────────
function _updateMsgCounter() {
  var ctr = document.getElementById('msg-counter');
  if (!ctr) return;
  var count = currentMessages.filter(function(m){ return m.role === 'user'; }).length;
  if (count > 0) {
    ctr.textContent = count + ' msg' + (count !== 1 ? 's' : '');
    ctr.style.display = '';
  } else {
    ctr.style.display = 'none';
  }
}

// ── AUTO-SCROLL PAUSE ─────────────────────────────────────────────────
var _autoScroll = true;
document.addEventListener('DOMContentLoaded', function() {
  var msgs = document.getElementById('chat-messages');
  if (!msgs) return;
  msgs.addEventListener('scroll', function() {
    if (!isStreaming) return;
    var atBottom = msgs.scrollTop + msgs.clientHeight >= msgs.scrollHeight - 60;
    _autoScroll = atBottom;
    var btn = document.getElementById('scroll-pause-btn');
    if (btn) btn.classList.toggle('visible', !atBottom);
  });
});
function scrollBottom() {
  if (!_autoScroll) return;
  var msgs = document.getElementById('chat-messages');
  if (msgs) setTimeout(function(){ msgs.scrollTop = msgs.scrollHeight; }, 30);
}
function resumeScroll() {
  _autoScroll = true;
  var btn = document.getElementById('scroll-pause-btn');
  if (btn) btn.classList.remove('visible');
  var msgs = document.getElementById('chat-messages');
  if (msgs) msgs.scrollTop = msgs.scrollHeight;
}

// ── AUTO-TITLE FROM FIRST MESSAGE ─────────────────────────────────────
function _autoGenerateTitle(userText) {
  var words = userText.trim().split(/\s+/).slice(0, 7).join(' ');
  return words.length < userText.trim().length ? words + '…' : words;
}

// ── @AGENT MENTION SWITCHER ────────────────────────────────────────────
var _mentionActive = false;
var _mentionIdx = 0;
var _mentionQuery = '';
var _mentionStart = 0;
function _mentionOninput(textarea) {
  var val = textarea.value;
  var pos = textarea.selectionStart;
  var before = val.slice(0, pos);
  var m = before.match(/@(\w*)$/);
  if (!m) { _closeMention(); return; }
  _mentionQuery = m[1].toLowerCase();
  _mentionStart = pos - m[0].length;
  var filtered = AGENT_DEFS.filter(function(a){
    return !_mentionQuery || a.name.toLowerCase().includes(_mentionQuery) || a.id.toLowerCase().includes(_mentionQuery);
  }).slice(0, 8);
  if (!filtered.length) { _closeMention(); return; }
  _mentionActive = true;
  _mentionIdx = 0;
  var dd = document.getElementById('mention-dropdown');
  dd.innerHTML = filtered.map(function(a, i) {
    return '<div class="mention-item' + (i === 0 ? ' active' : '') + '" data-idx="' + i + '" data-id="' + escHtml(a.id) + '" onclick="_selectMention(\'' + escHtml(a.id) + '\',\'' + escHtml(a.name) + '\')">' +
      '<span class="mention-emoji">' + (a.emoji || '🤖') + '</span>' +
      '<span class="mention-name">' + escHtml(a.name) + '</span>' +
      '<span class="mention-role">' + escHtml((a.role || '').slice(0, 35)) + '</span>' +
    '</div>';
  }).join('');
  dd.classList.add('open');
}
function _mentionKeydown(e) {
  if (!_mentionActive) return;
  var items = document.querySelectorAll('#mention-dropdown .mention-item');
  if (e.key === 'ArrowDown') { e.preventDefault(); _mentionIdx = Math.min(_mentionIdx+1, items.length-1); }
  else if (e.key === 'ArrowUp') { e.preventDefault(); _mentionIdx = Math.max(_mentionIdx-1, 0); }
  else if (e.key === 'Enter' || e.key === 'Tab') {
    var active = items[_mentionIdx];
    if (active) { e.preventDefault(); _selectMention(active.dataset.id, active.querySelector('.mention-name').textContent); }
    return;
  }
  else if (e.key === 'Escape') { _closeMention(); return; }
  else return;
  items.forEach(function(it, i){ it.classList.toggle('active', i === _mentionIdx); });
}
function _selectMention(agentId, agentName) {
  var textarea = document.getElementById('chat-input');
  var val = textarea.value;
  var before = val.slice(0, _mentionStart);
  var after  = val.slice(textarea.selectionStart);
  textarea.value = before + agentName + ' ' + after;
  textarea.selectionStart = textarea.selectionEnd = before.length + agentName.length + 1;
  _closeMention();
  pickAgent(agentId, agentName);
  textarea.focus();
}
function _closeMention() {
  _mentionActive = false;
  var dd = document.getElementById('mention-dropdown');
  if (dd) dd.classList.remove('open');
}

// ── MULTI-FILE QUEUE ───────────────────────────────────────────────────
var _fileQueue = []; // [{file, name, dataUrl, type}]
function handleFileQueue(event) {
  var files = Array.from(event.target.files || []);
  event.target.value = ''; // reset
  files.forEach(function(file) {
    if (_fileQueue.length >= 5) { toast(TI.warn, 'Max 5 files at once.'); return; }
    var reader = new FileReader();
    reader.onload = function(ev) {
      _fileQueue.push({ file: file, name: file.name, dataUrl: ev.target.result, type: file.type });
      _renderFileQueue();
    };
    reader.readAsDataURL(file);
    // For the first image, also set the legacy attachment so the old pipeline works
    if (file.type.startsWith('image/') && _fileQueue.length === 0) {
      handleFileAttach({ target: { files: [file] } });
    }
  });
}
function _renderFileQueue() {
  var fq = document.getElementById('file-queue');
  if (!fq) return;
  if (!_fileQueue.length) { fq.classList.remove('open'); fq.innerHTML = ''; return; }
  fq.classList.add('open');
  fq.innerHTML = _fileQueue.map(function(item, i) {
    var icon = item.type.startsWith('image/') ? '<i class="fa-solid fa-image" style="color:var(--primary)"></i>' :
               item.type === 'application/pdf' ? '<i class="fa-solid fa-file-pdf" style="color:#e53935"></i>' :
               '<i class="fa-solid fa-file-lines" style="color:#666"></i>';
    return '<div class="file-queue-item">' + icon +
      '<span class="fq-name" title="' + escHtml(item.name) + '">' + escHtml(item.name) + '</span>' +
      '<button class="fq-remove" onclick="_removeFromQueue(' + i + ')" title="Remove"><i class="fa-solid fa-xmark"></i></button>' +
    '</div>';
  }).join('');
}
function _removeFromQueue(idx) {
  _fileQueue.splice(idx, 1);
  _renderFileQueue();
  if (!_fileQueue.length) clearAttachment();
}

// ── SMART URL PASTE DETECTION ─────────────────────────────────────────
var _detectedUrl = null;
function _detectUrl(textarea) {
  var val = textarea.value;
  var urlRe = /https?:\/\/[^\s]{8,}/gi;
  var m = val.match(urlRe);
  var chip = document.getElementById('url-preview-chip');
  var link = document.getElementById('url-preview-link');
  if (m && m[0] && m[0] !== _detectedUrl) {
    _detectedUrl = m[0];
    if (chip && link) {
      var short = _detectedUrl.replace(/^https?:\/\/(www\.)?/, '').slice(0, 50);
      link.href = _detectedUrl;
      link.textContent = short + (_detectedUrl.length > 53 ? '…' : '');
      chip.classList.add('open');
    }
  } else if (!m) {
    _detectedUrl = null;
    if (chip) chip.classList.remove('open');
  }
}
function dismissUrlPreview() {
  _detectedUrl = null;
  var chip = document.getElementById('url-preview-chip');
  if (chip) chip.classList.remove('open');
}

// ── RESIZABLE CHAT SIDEBAR ────────────────────────────────────────────
document.addEventListener('DOMContentLoaded', function() {
  var handle = document.getElementById('sidebar-drag-handle');
  var sidebar = document.getElementById('chat-sidebar');
  if (!handle || !sidebar) return;
  var dragging = false;
  var startX, startW;
  handle.addEventListener('mousedown', function(e) {
    dragging = true; startX = e.clientX; startW = sidebar.offsetWidth;
    document.body.style.cursor = 'col-resize';
    document.body.style.userSelect = 'none';
  });
  document.addEventListener('mousemove', function(e) {
    if (!dragging) return;
    var w = Math.max(180, Math.min(380, startW + (e.clientX - startX)));
    sidebar.style.width = w + 'px';
    sidebar.style.minWidth = w + 'px';
  });
  document.addEventListener('mouseup', function() {
    if (dragging) { dragging = false; document.body.style.cursor = ''; document.body.style.userSelect = ''; }
  });
});

function _streamResponse() {
  var msgs = document.getElementById('chat-messages');
  var avatarEl = document.getElementById('chat-avatar');
  var avatarImg = avatarEl && avatarEl.querySelector('img');
  var agentEmoji = avatarImg ? '<img src="' + avatarImg.src + '" class="avatar-logo" alt="Agent">' : '<img src="/assets/logos/Enderun-Colleges.png" class="avatar-logo" alt="Enderun">';

  var _msgToolCards = []; // tool cards added during this response
  var _composerOpened = false; // track if composer was opened this turn

  var typingEl = document.createElement('div');
  var _agentDisplayName = document.getElementById('chat-name').textContent || 'Agent';
  typingEl.className = 'msg ai'; typingEl.id = 'typing-msg';
  typingEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-body"><span class="msg-name">' + escHtml(_agentDisplayName) + '</span><div class="msg-bubble"><div class="typing-indicator"><div class="typing-dot"></div><div class="typing-dot"></div><div class="typing-dot"></div></div></div></div>';
  msgs.appendChild(typingEl);
  scrollBottom();

  isStreaming = true;
  document.getElementById('send-btn').disabled = true;
  document.getElementById('stop-btn').classList.add('active');
  _abortCtrl = new AbortController();

  fetch('/api/chat', {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify({ agent: currentAgent, messages: currentMessages, model: chatModel, mode_hint: getModeHint(), composer_selected: _compSelected || null }),
    signal: _abortCtrl.signal
  })
  .then(function(r) {
    if (!r.ok) throw new Error('HTTP ' + r.status);
    var reader = r.body.getReader();
    var decoder = new TextDecoder();
    var aiEl = null; var aiBubble = null; var accumulated = '';
    var _sseBuffer = ''; // buffers incomplete SSE lines across HTTP chunks

    function readChunk() {
      reader.read().then(function(result) {
        if (result.done) { finishStream(); return; }
        var chunk = decoder.decode(result.value, { stream: true });
        _sseBuffer += chunk;
        var lines = _sseBuffer.split('\n');
        _sseBuffer = lines.pop(); // keep any incomplete trailing line
        var isDone = false;
        lines.forEach(function(line) {
          if (isDone || !line.startsWith('data: ')) return;
          var data = line.slice(6).trim();
          if (data === '[DONE]') { isDone = true; finishStream(); return; }
          try {
            var parsed = JSON.parse(data);
            if (parsed.action === 'open_composer') {
              var t0 = document.getElementById('typing-msg'); if (t0) t0.remove();
              var tc0 = document.getElementById('tool-card-' + (parsed.tool_id || 'open_post_composer')); if (tc0) tc0.remove();
              _composerOpened = true;
              openPostComposer(parsed.platform || 'facebook');
              return;
            }
            if (parsed.action === 'clear_text') {
              accumulated = '';
              if (aiEl) { aiEl.remove(); aiEl = null; aiBubble = null; }
              return;
            }
            if (parsed.action === 'tool_start') {
              var t = document.getElementById('typing-msg'); if (t) t.remove();
              var card = document.createElement('div');
              card.className = 'tool-action-card running';
              card.id = 'tool-card-' + (parsed.tool_id || parsed.tool);
              card.innerHTML = '<div class="tool-action-spinner"></div><span>' + toolLabel(parsed.tool) + '…</span>';
              msgs.appendChild(card); _msgToolCards.push(card); scrollBottom(); return;
            }
            if (parsed.action === 'tool_done') {
              var card = document.getElementById('tool-card-' + (parsed.tool_id || parsed.tool));
              var ok = parsed.result && !parsed.result.startsWith('❌');
              if (card) {
                card.className = 'tool-action-card ' + (ok ? 'done' : 'failed');
                if (parsed.img_url) {
                  card.classList.add('has-image');
                  card.innerHTML = '<div class="tool-card-header"><i class="fa-solid ' + (ok ? 'fa-circle-check' : 'fa-circle-xmark') + '" style="font-size:13px;flex-shrink:0"></i><span>' + escHtml((parsed.result||'').slice(0,140)) + '</span></div><div class="tool-card-img"><img src="' + escHtml(parsed.img_url) + '" alt="Generated image" onload="scrollBottom()" onclick="openLightbox(this.src)" style="cursor:zoom-in"></div>';
                } else {
                  card.innerHTML = '<i class="fa-solid ' + (ok ? 'fa-circle-check' : 'fa-circle-xmark') + '" style="font-size:13px;flex-shrink:0"></i><span>' + escHtml((parsed.result||'').slice(0,140)) + '</span>';
                }
                scrollBottom();
              }
              return;
            }
            var token = parsed.token || parsed.content || parsed.text || '';
            if (token) {
              accumulated += token;
              if (!aiEl) {
                var t = document.getElementById('typing-msg'); if (t) t.remove();
                aiEl = document.createElement('div'); aiEl.className = 'msg ai';
                aiEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-body"><span class="msg-name">' + escHtml(_agentDisplayName) + '</span><div class="msg-bubble"></div><div class="msg-timestamp">' + _ts() + '</div></div>';
                msgs.appendChild(aiEl); aiBubble = aiEl.querySelector('.msg-bubble');
              }
              aiBubble.innerHTML = md(accumulated); scrollBottom();
            }
          } catch(e) {
            if (data && data !== '[DONE]') {
              accumulated += data;
              if (!aiEl) {
                var t2 = document.getElementById('typing-msg'); if (t2) t2.remove();
                aiEl = document.createElement('div'); aiEl.className = 'msg ai';
                aiEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-body"><span class="msg-name">' + escHtml(_agentDisplayName) + '</span><div class="msg-bubble"></div><div class="msg-timestamp">' + _ts() + '</div></div>';
                msgs.appendChild(aiEl); aiBubble = aiEl.querySelector('.msg-bubble');
              }
              aiBubble.innerHTML = md(accumulated); scrollBottom();
            }
          }
        });
        if (!isDone) readChunk();
      }).catch(function(err) { if (err.name !== 'AbortError') finishStream(); });
    }

    function finishStream() {
      var t = document.getElementById('typing-msg'); if (t) t.remove();
      if (!aiEl && accumulated) {
        aiEl = document.createElement('div'); aiEl.className = 'msg ai';
        aiEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-body"><span class="msg-name">' + escHtml(_agentDisplayName) + '</span><div class="msg-bubble"></div><div class="msg-timestamp">' + _ts() + '</div></div>';
        msgs.appendChild(aiEl);
      }
      if (aiEl) {
        aiBubble = aiEl.querySelector('.msg-bubble');
        var genMatch = accumulated.match(/\[GENERATE:(pdf|pptx|docx|csv)\]([\s\S]*?)\[\/GENERATE\]/i);
        if (genMatch) {
          _processGenerateBlock(accumulated, genMatch, aiBubble);
        } else {
          aiBubble.innerHTML = _renderWithThinking(accumulated);
          _renderLatex(aiBubble);
          _processCodeBlocks(aiBubble);
        }
        // Feedback + action buttons — after last tool card if present, else inside msg-body
        var fbId = 'fb-' + Date.now();
        var actsHtml =
          '<button class="feedback-btn" id="fb-up-' + fbId + '" onclick="_sendFeedback(this,\'up\')" title="Helpful"><i class="fa-regular fa-thumbs-up"></i></button>' +
          '<button class="feedback-btn" id="fb-dn-' + fbId + '" onclick="_sendFeedback(this,\'down\')" title="Not helpful"><i class="fa-regular fa-thumbs-down"></i></button>' +
          '<button class="msg-action-btn" onclick="copyMsgText(this)" title="Copy"><i class="fa-solid fa-copy"></i></button>' +
          '<button class="msg-action-btn" onclick="retryLast()" title="Regenerate"><i class="fa-solid fa-rotate-right"></i></button>';
        var lastCard = _msgToolCards.length ? _msgToolCards[_msgToolCards.length - 1] : null;
        if (lastCard) {
          var acts = document.createElement('div');
          acts.className = 'msg-action-row';
          acts.innerHTML = actsHtml;
          lastCard.insertAdjacentElement('afterend', acts);
          var hoverGroup = [aiEl].concat(_msgToolCards).filter(Boolean);
          function _showActs() { acts.style.opacity = '1'; }
          function _hideActs(e) {
            var to = e.relatedTarget;
            var inside = acts === to || acts.contains(to) || hoverGroup.some(function(g){ return g === to || g.contains(to); });
            if (!inside) acts.style.opacity = '0';
          }
          hoverGroup.forEach(function(el) { el.addEventListener('mouseenter', _showActs); el.addEventListener('mouseleave', _hideActs); });
          acts.addEventListener('mouseleave', _hideActs);
        } else {
          var body = aiEl.querySelector('.msg-body');
          if (body) {
            var acts = document.createElement('div');
            acts.className = 'msg-actions';
            acts.innerHTML = actsHtml;
            var ts = body.querySelector('.msg-timestamp');
            body.insertBefore(acts, ts || null);
          }
        }
        // Follow-up suggestion chips (simple heuristic from response topics)
        _appendFollowUpChips(accumulated);
      }
      if (accumulated) {
        var _savedCards = _msgToolCards.map(function(card) {
          var imgEl = card.querySelector('img');
          var spanEl = card.querySelector('.tool-card-header span') || card.querySelector('span');
          return {
            tool: (card.id || '').replace('tool-card-', ''),
            result: spanEl ? spanEl.textContent.trim() : '',
            img_url: imgEl ? imgEl.src : null,
            status: card.classList.contains('failed') ? 'failed' : 'done'
          };
        });
        var msg = { role: 'assistant', content: accumulated };
        if (_savedCards.length) msg._toolCards = _savedCards;
        if (_composerOpened) msg._composer = { platform: _compPlatform || 'facebook' };
        currentMessages.push(msg);
        // Auto-generate title for session if this is the first exchange
        _maybeSetAutoTitle();
        _saveCurrentSession();
        _updateMsgCounter();
      }
      isStreaming = false; _abortCtrl = null;
      _autoScroll = true;
      var pauseBtn = document.getElementById('scroll-pause-btn');
      if (pauseBtn) pauseBtn.classList.remove('visible');
      document.getElementById('send-btn').disabled = false;
      document.getElementById('stop-btn').classList.remove('active');
      scrollBottom();
    }

    readChunk();
  })
  .catch(function(err) {
    if (err.name === 'AbortError') return;
    var t = document.getElementById('typing-msg'); if (t) t.remove();
    if (currentMessages.length && currentMessages[currentMessages.length-1].role === 'user') currentMessages.pop();
    var errEl = document.createElement('div'); errEl.className = 'msg';
    errEl.innerHTML = '<div class="msg-avatar">' + agentEmoji + '</div><div class="msg-body"><div class="msg-bubble" style="color:#F87171"><i class="fa-solid fa-circle-xmark"></i> Could not reach agent.<br><small>' + escHtml(String(err)) + '</small></div></div>';
    msgs.appendChild(errEl);
    isStreaming = false; _abortCtrl = null;
    document.getElementById('send-btn').disabled = false;
    document.getElementById('stop-btn').classList.remove('active');
    scrollBottom();
  });
}

// ── STOP STREAM ────────────────────────────────────────────────────────
function stopStream() {
  if (_abortCtrl) { _abortCtrl.abort(); _abortCtrl = null; }
  var t = document.getElementById('typing-msg'); if (t) t.remove();
  isStreaming = false;
  document.getElementById('send-btn').disabled = false;
  document.getElementById('stop-btn').classList.remove('active');
  scrollBottom();
}

// ── COPY MESSAGE TEXT ──────────────────────────────────────────────────
function copyMsgText(btn) {
  // Support buttons inside msg-body OR standalone msg-action-row (after tool cards)
  var body = btn.closest('.msg-body');
  var bubble = body ? body.querySelector('.msg-bubble') : null;
  if (!bubble) {
    // msg-action-row: find the preceding .msg sibling
    var row = btn.closest('.msg-action-row');
    if (row) {
      var prev = row.previousElementSibling;
      while (prev && !prev.classList.contains('msg')) prev = prev.previousElementSibling;
      bubble = prev ? prev.querySelector('.msg-bubble') : null;
    }
  }
  if (!bubble) return;
  navigator.clipboard.writeText(bubble.innerText || bubble.textContent || '').then(function() {
    var icon = btn.querySelector('i');
    if (icon) { icon.className = 'fa-solid fa-check'; btn.classList.add('copied'); }
    setTimeout(function() {
      if (icon) icon.className = 'fa-solid fa-copy';
      btn.classList.remove('copied');
    }, 1600);
  });
}

// ── RETRY LAST RESPONSE ────────────────────────────────────────────────
function retryLast() {
  if (isStreaming) return;
  if (currentMessages.length && currentMessages[currentMessages.length-1].role === 'assistant') currentMessages.pop();
  var msgs = document.getElementById('chat-messages');
  var all = msgs.querySelectorAll('.msg:not(#typing-msg)');
  for (var i = all.length - 1; i >= 0; i--) {
    if (!all[i].classList.contains('user')) { all[i].remove(); break; }
  }
  msgs.querySelectorAll('.tool-action-card').forEach(function(c){ c.remove(); });
  _streamResponse();
}

// ── EDIT MESSAGE ───────────────────────────────────────────────────────
function editMsgAt(idx) {
  if (isStreaming) return;
  var msg = currentMessages[idx];
  if (!msg) return;
  var text = typeof msg.content === 'string' ? msg.content :
    (Array.isArray(msg.content) ? msg.content.filter(function(b){return b.type==='text';}).map(function(b){return b.text;}).join('') : '');
  var input = document.getElementById('chat-input');
  input.value = text; autoResize(input); updateCharCount(input); input.focus();
  currentMessages.splice(idx);
  var msgs = document.getElementById('chat-messages');
  var all = msgs.querySelectorAll('.msg:not(#typing-msg), .tool-action-card');
  for (var i = idx; i < all.length; i++) all[i].remove();
}

// ── EXPORT CHAT ────────────────────────────────────────────────────────
function exportChat() {
  if (!currentMessages.length) { toast(TI.warn, 'No messages to export.'); return; }
  var agent = document.getElementById('chat-name').textContent || 'Agent';
  var lines = ['Enderun Marketing Hub — Chat Export', 'Agent: ' + agent, 'Date: ' + new Date().toLocaleString(), '─'.repeat(50), ''];
  currentMessages.forEach(function(m) {
    var c = typeof m.content === 'string' ? m.content :
      (Array.isArray(m.content) ? m.content.filter(function(b){return b.type==='text';}).map(function(b){return b.text;}).join('') : '');
    lines.push('[' + (m.role === 'user' ? 'You' : agent) + ']');
    lines.push(c); lines.push('');
  });
  var a = document.createElement('a');
  a.href = URL.createObjectURL(new Blob([lines.join('\n')], {type:'text/plain;charset=utf-8'}));
  a.download = 'chat-' + Date.now() + '.txt'; a.click();
  toast(TI.ok, 'Exported as TXT.');
}

function exportChatPDF() {
  if (!currentMessages.length) { toast(TI.warn, 'No messages to export.'); return; }
  var agent = document.getElementById('chat-name').textContent || 'Agent';
  toast(TI.ok, 'Generating PDF…');
  fetch('/api/export-chat-pdf', {
    method: 'POST', headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({ messages: currentMessages, agent_name: agent })
  })
  .then(function(r){ return r.json(); })
  .then(function(d) {
    if (d.status === 'ok') {
      var a = document.createElement('a');
      a.href = '/chat-files/' + encodeURIComponent(d.filename);
      a.download = d.filename; a.click();
      toast(TI.ok, 'PDF exported!');
    } else {
      toast(TI.err, 'PDF export failed: ' + (d.message || 'Unknown error'));
    }
  })
  .catch(function() { toast(TI.err, 'PDF export failed.'); });
}

// ── SEARCH IN CHAT ─────────────────────────────────────────────────────
function toggleChatSearch() {
  var bar = document.getElementById('chat-search-bar');
  var btn = document.getElementById('search-toggle-btn');
  var open = bar.classList.toggle('open');
  btn.classList.toggle('active', open);
  if (open) { document.getElementById('chat-search-input').focus(); }
  else { _clearSearchHL(); document.getElementById('chat-search-input').value = ''; document.getElementById('search-match-count').textContent = ''; }
}
function _clearSearchHL() {
  document.querySelectorAll('mark.chat-hl').forEach(function(m) {
    m.parentNode.replaceChild(document.createTextNode(m.textContent), m);
  });
  document.querySelectorAll('#chat-messages span[data-hl]').forEach(function(s) {
    s.outerHTML = s.innerHTML;
  });
}
function searchInChat(q) {
  _clearSearchHL();
  document.getElementById('search-match-count').textContent = '';
  if (!q.trim()) return;
  var re = new RegExp('(' + q.replace(/[.*+?^${}()|[\]\\]/g,'\\$&') + ')', 'gi');
  var count = 0;
  document.querySelectorAll('#chat-messages .msg-bubble').forEach(function(bubble) {
    var walker = document.createTreeWalker(bubble, NodeFilter.SHOW_TEXT, null, false);
    var nodes = []; var n;
    while ((n = walker.nextNode())) nodes.push(n);
    nodes.forEach(function(node) {
      if (!node.textContent.match(re)) return;
      var span = document.createElement('span');
      span.dataset.hl = '1';
      span.innerHTML = node.textContent.replace(re, function(m){ count++; return '<mark class="chat-hl">' + m + '</mark>'; });
      node.parentNode.replaceChild(span, node);
    });
  });
  var cnt = document.getElementById('search-match-count');
  cnt.textContent = count ? count + ' match' + (count > 1 ? 'es' : '') : 'No matches';
  var first = document.querySelector('mark.chat-hl');
  if (first) first.scrollIntoView({behavior:'smooth', block:'center'});
}

// ── CHAR COUNTER ───────────────────────────────────────────────────────
function updateCharCount(el) {
  var c = document.getElementById('char-counter'); if (!c) return;
  var n = el.value.length;
  c.textContent = n > 0 ? n.toLocaleString() : '';
  c.className = 'char-counter' + (n > 3000 ? ' over' : '');
}

// ── PIN CONVERSATION ───────────────────────────────────────────────────
var _pins = JSON.parse(localStorage.getItem('_chatPins') || '{}');
function togglePin(agentId, sid, btn, ev) {
  ev.stopPropagation();
  var key = agentId + '|' + sid;
  if (_pins[key]) { delete _pins[key]; btn.classList.remove('pinned'); }
  else { _pins[key] = 1; btn.classList.add('pinned'); }
  localStorage.setItem('_chatPins', JSON.stringify(_pins));
  renderSessionsList(agentId);
}

// ── SUGGESTED PROMPT ──────────────────────────────────────────────────
function usePrompt(text) {
  if (!currentAgent) { toast(TI.warn, 'Select an agent first.'); return; }
  var input = document.getElementById('chat-input');
  input.value = text; autoResize(input); updateCharCount(input); input.focus();
}

function sendDirectMsg(text) {
  if (!currentAgent) { toast(TI.warn, 'Select an agent first.'); return; }
  var input = document.getElementById('chat-input');
  input.value = text; autoResize(input); sendMsg();
}

// ── CHAT FILE GENERATION ───────────────────────────────────────────────
function _renderWithThinking(text) {
  var thinkings = [];
  var cleaned = text.replace(/\[THINKING\]([\s\S]*?)\[\/THINKING\]/gi, function(_, inner) {
    var ph = '|||THINK' + thinkings.length + '|||';
    thinkings.push(inner.trim());
    return ph;
  });
  var rendered = md(cleaned);
  thinkings.forEach(function(content, i) {
    rendered = rendered.replace('|||THINK' + i + '|||',
      '<details class="thinking-block"><summary><i class="fa-solid fa-brain"></i> Show thinking process</summary>' +
      '<div class="thinking-content">' + md(content) + '</div></details>');
  });
  return rendered;
}

function _processGenerateBlock(fullText, genMatch, bubbleEl) {
  var docType  = genMatch[1].toLowerCase();
  var content  = genMatch[2].trim();
  var before   = _renderWithThinking(fullText.slice(0, genMatch.index).trim());
  var after    = _renderWithThinking(fullText.slice(genMatch.index + genMatch[0].length).trim());
  var typeLabel = {pdf:'PDF Document', pptx:'Presentation (PPTX)', docx:'Word Document', csv:'CSV Spreadsheet'}[docType] || docType.toUpperCase();
  var iconCls   = {pdf:'fa-file-pdf', pptx:'fa-file-powerpoint', docx:'fa-file-word', csv:'fa-file-csv'}[docType] || 'fa-file';
  var iconColor = {pdf:'#e53935', pptx:'#d84315', docx:'#1565c0', csv:'#2e7d32'}[docType] || '#666';

  bubbleEl.innerHTML = (before ? md(before) + '<br>' : '') +
    '<div class="file-gen-loading"><i class="fa-solid fa-spinner fa-spin"></i> Generating ' + typeLabel + '…</div>' +
    (after ? '<br>' + md(after) : '');
  scrollBottom();

  fetch('/api/generate-doc', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({type: docType, content: content})
  })
  .then(function(r){ return r.json(); })
  .then(function(d) {
    if (d.status === 'ok') {
      bubbleEl.innerHTML = (before ? md(before) + '<br><br>' : '') +
        '<a href="/chat-files/' + encodeURIComponent(d.filename) + '" download class="file-download-card">' +
          '<i class="fa-solid ' + iconCls + '" style="color:' + iconColor + ';font-size:30px;flex-shrink:0"></i>' +
          '<div><div class="file-card-name">' + escHtml(d.title) + '</div>' +
          '<div class="file-card-type">' + typeLabel + ' &nbsp;·&nbsp; Click to download</div></div>' +
        '</a>' +
        (after ? md(after) : '');
    } else {
      bubbleEl.innerHTML = (before ? md(before) + '<br>' : '') +
        '<span style="color:#e53935"><i class="fa-solid fa-circle-xmark"></i> Could not generate file: ' + escHtml(d.message || 'Unknown error') + '</span>';
    }
    scrollBottom();
  })
  .catch(function() {
    bubbleEl.innerHTML = '<span style="color:#e53935"><i class="fa-solid fa-circle-xmark"></i> File generation failed.</span>';
    scrollBottom();
  });
}

// ── MARKDOWN RENDERER ─────────────────────────────────────────────────
function copyCodeBlock(id) {
  var pre = document.getElementById(id);
  if (!pre) return;
  navigator.clipboard.writeText(pre.textContent).then(function() {
    var btn = pre.closest('.code-block-wrap') && pre.closest('.code-block-wrap').querySelector('.code-copy-btn');
    if (btn) { btn.textContent = 'Copied!'; setTimeout(function(){ btn.textContent = 'Copy'; }, 1800); }
  });
}

// ── CODE BLOCK HEADERS ─────────────────────────────────────────────────────
function _processCodeBlocks(el) {
  el.querySelectorAll('pre').forEach(function(pre) {
    if (pre.querySelector('.code-block-header')) return; // already processed
    var codeEl = pre.querySelector('code');
    var lang = '';
    if (codeEl) {
      var cls = codeEl.className || '';
      var m = cls.match(/language-(\w+)/);
      lang = m ? m[1] : '';
      // Apply Prism syntax highlighting
      if (lang && typeof Prism !== 'undefined') {
        if (!codeEl.classList.contains('language-' + lang)) {
          codeEl.className = 'language-' + lang;
        }
        try { Prism.highlightElement(codeEl); } catch(e) {}
      } else if (typeof Prism !== 'undefined' && Prism.plugins && Prism.plugins.autoloader) {
        try { Prism.highlightElement(codeEl); } catch(e) {}
      }
    }
    var hdr = document.createElement('div');
    hdr.className = 'code-block-header';
    hdr.innerHTML = '<span class="code-block-lang">' + (lang || 'code') + '</span>' +
      '<button class="code-copy-btn" onclick="_copyCode(this)"><i class="fa-regular fa-copy"></i> Copy</button>';
    pre.insertBefore(hdr, pre.firstChild);
  });
}
function _copyCode(btn) {
  var pre = btn.closest('pre');
  var code = pre ? (pre.querySelector('code') || pre) : null;
  if (!code) return;
  navigator.clipboard.writeText(code.innerText).then(function() {
    btn.innerHTML = '<i class="fa-solid fa-check"></i> Copied';
    btn.classList.add('copied');
    setTimeout(function() { btn.innerHTML = '<i class="fa-regular fa-copy"></i> Copy'; btn.classList.remove('copied'); }, 2000);
  });
}

// ── FEEDBACK BUTTONS ───────────────────────────────────────────────────
function _sendFeedback(btn, dir) {
  var fbId = btn.id.replace('fb-' + dir + '-', '');
  var upBtn = document.getElementById('fb-up-' + fbId);
  var dnBtn = document.getElementById('fb-dn-' + fbId);
  if (!upBtn || !dnBtn) return;
  // Toggle off if already active
  var isActive = btn.classList.contains('active-' + dir);
  upBtn.classList.remove('active-up'); dnBtn.classList.remove('active-down');
  if (!isActive) {
    btn.classList.add(dir === 'up' ? 'active-up' : 'active-down');
    toast(dir === 'up' ? TI.ok : TI.warn, dir === 'up' ? 'Thanks for the feedback!' : 'Noted — will improve.');
  }
}

// ── FOLLOW-UP CHIPS ────────────────────────────────────────────────────
function _appendFollowUpChips(text) {
  var chips = _suggestChips(text);
  if (!chips.length) return;
  var msgs = document.getElementById('chat-messages');
  var fup = document.createElement('div');
  fup.className = 'followup-chips';
  fup.innerHTML = chips.map(function(c) {
    return '<button class="fup-chip" onclick="usePrompt(' + JSON.stringify(c) + ');this.closest(\'.followup-chips\').remove()">' + escHtml(c) + '</button>';
  }).join('');
  msgs.appendChild(fup);
  scrollBottom();
}
function _suggestChips(text) {
  var chips = [];
  var t = text.toLowerCase();
  if (/facebook|instagram|social|post/.test(t)) chips.push('Can you generate a visual brief for this post?');
  if (/email|drip|lead/.test(t)) chips.push('Write a follow-up email for this lead.');
  if (/competitor|analysis|market/.test(t)) chips.push('Generate a PDF report from this analysis.');
  if (/report|pdf|analytics/.test(t)) chips.push('Show me the data as a bar chart.');
  if (/strategy|plan|campaign/.test(t)) chips.push('Turn this into a content calendar.');
  if (/hospitality|culinary|program/.test(t)) chips.push('Write Instagram captions for this program.');
  return chips.slice(0, 3);
}

// ── AUTO-TITLE FOR SESSIONS ────────────────────────────────────────────
function _maybeSetAutoTitle() {
  if (!currentAgent) return;
  var sid = currentSessionId[currentAgent];
  if (!sid) return;
  var sessions = chatSessions[currentAgent] || [];
  var session = sessions.find(function(s){ return s.id === sid; });
  if (!session) return;
  // Only set title if it's still the default 'Chat' or empty
  if (session.title && session.title !== 'Chat' && session.title !== '') return;
  // Use first user message text
  var firstUser = currentMessages.find(function(m){ return m.role === 'user'; });
  if (!firstUser) return;
  var txt = typeof firstUser.content === 'string' ? firstUser.content :
    (Array.isArray(firstUser.content) ? firstUser.content.filter(function(b){return b.type==='text';}).map(function(b){return b.text;}).join(' ') : '');
  if (!txt || txt === '__GREET__') return;
  session.title = _autoGenerateTitle(txt);
  renderSessionsList(currentAgent);
}

// ── TIME-BASED GREETING ─────────────────────────────────────────────────────
(function() {
  var h = new Date().getHours();
  var g = h < 12 ? 'Good morning' : h < 18 ? 'Good afternoon' : 'Good evening';
  var el = document.getElementById('welcome-greeting');
  if (el) el.textContent = g + ', Eva';
})();

function _renderLatex(el) {
  if (typeof katex === 'undefined') return;
  el.innerHTML = el.innerHTML.replace(/\$\$([\s\S]+?)\$\$/g, function(_, expr) {
    try { return '<div class="katex-display">' + katex.renderToString(expr.trim(), {displayMode:true, throwOnError:false}) + '</div>'; }
    catch(e) { return _; }
  });
  el.innerHTML = el.innerHTML.replace(/\$([^$\n<>]{1,200})\$/g, function(_, expr) {
    try { return katex.renderToString(expr.trim(), {displayMode:false, throwOnError:false}); }
    catch(e) { return _; }
  });
}

function md(text) {
  if (!text) return '';
  var s = text;
  // Strip hashtag-only lines (social media tags have no place in chat)
  s = s.replace(/\n([ \t]*(#[A-Za-z][A-Za-z0-9_]*[ \t]*)+)+$/g, '');
  s = s.replace(/^([ \t]*(#[A-Za-z][A-Za-z0-9_]*[ \t]*)+\n?)+/g, '');
  // Extract fenced code blocks BEFORE HTML escaping
  var codeBlocks = [];
  s = s.replace(/```([a-zA-Z0-9]*)\n?([\s\S]*?)```/g, function(_, lang, code) {
    var ph = '\x01CODE' + codeBlocks.length + '\x01';
    var id = 'cb' + Math.random().toString(36).slice(2, 8);
    var safeCode = code.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;');
    var langLabel = lang || 'code';
    codeBlocks.push('<div class="code-block-wrap"><div class="code-block-top"><span class="code-lang-label">' + langLabel + '</span><button class="code-copy-btn" onclick="copyCodeBlock(\'' + id + '\')">Copy</button></div><pre id="' + id + '">' + safeCode + '</pre></div>');
    return ph;
  });
  // Escape HTML
  s = s.replace(/&/g, '&amp;').replace(/</g, '&lt;').replace(/>/g, '&gt;');
  // Headers
  s = s.replace(/^### (.+)$/gm, '<h3>$1</h3>');
  s = s.replace(/^## (.+)$/gm, '<h2>$1</h2>');
  s = s.replace(/^# (.+)$/gm, '<h2>$1</h2>');
  // Bold & italic
  s = s.replace(/[*][*]([^*\n]+)[*][*]/g, '<strong>$1</strong>');
  s = s.replace(/[*]([^*\n]+)[*]/g, '<em>$1</em>');
  // Inline code
  s = s.replace(/`([^`\n]+)`/g, '<code>$1</code>');
  // Horizontal rule
  s = s.replace(/^---+$/gm, '<hr>');
  // Mark list items with placeholders so we can group them
  s = s.replace(/^[ \t]*[-*] (.+)$/gm, '\x02UL\x03$1\x02/UL\x03');
  s = s.replace(/^[ \t]*[0-9]+\. (.+)$/gm, '\x02OL\x03$1\x02/OL\x03');
  // Group consecutive UL items into <ul>
  s = s.replace(/(\x02UL\x03[\s\S]*?\x02\/UL\x03\n?)+/g, function(m) {
    var items = m.replace(/\x02UL\x03([\s\S]*?)\x02\/UL\x03\n?/g, '<li>$1</li>');
    return '<ul>' + items + '</ul>';
  });
  // Group consecutive OL items into <ol>
  s = s.replace(/(\x02OL\x03[\s\S]*?\x02\/OL\x03\n?)+/g, function(m) {
    var items = m.replace(/\x02OL\x03([\s\S]*?)\x02\/OL\x03\n?/g, '<li>$1</li>');
    return '<ol>' + items + '</ol>';
  });
  // Markdown tables → <table>
  s = s.replace(/((?:[ \t]*\|.+\|\n?)+)/g, function(tableBlock) {
    var rows = tableBlock.trim().split('\n').map(function(r){ return r.trim(); }).filter(Boolean);
    if (rows.length < 2) return tableBlock;
    var isSep = function(r){ return /^\|[\s\-:|]+\|$/.test(r.replace(/\|[-: ]+/g,'|---')); };
    var sepIdx = rows.findIndex(isSep);
    if (sepIdx < 0) return tableBlock;
    var headerRows = rows.slice(0, sepIdx);
    var bodyRows   = rows.slice(sepIdx + 1);
    var parseRow = function(r, tag) {
      var cells = r.replace(/^\||\|$/g,'').split('|');
      return '<tr>' + cells.map(function(c){ return '<' + tag + '>' + c.trim() + '</' + tag + '>'; }).join('') + '</tr>';
    };
    var thead = '<thead>' + headerRows.map(function(r){ return parseRow(r,'th'); }).join('') + '</thead>';
    var tbody = bodyRows.length ? '<tbody>' + bodyRows.map(function(r){ return parseRow(r,'td'); }).join('') + '</tbody>' : '';
    return '<div class="md-table-wrap"><table class="md-table">' + thead + tbody + '</table></div>\n';
  });
  // Split on blank lines — each chunk becomes a <p> unless it's a block element
  var blocks = s.split(/\n{2,}/);
  var out = blocks.map(function(block) {
    block = block.trim();
    if (!block) return '';
    if (/^<(h[123]|ul|ol|hr|div|table)/.test(block)) return block;
    return '<p>' + block.replace(/\n/g, '<br>') + '</p>';
  }).filter(Boolean).join('');
  // Clean up <br> around block tags that ended up inside <p>
  out = out.replace(/<p>(<(ul|ol|h[123]|hr)[^>]*>)/g, '$1');
  out = out.replace(/(<\/(ul|ol|h[123]|hr)>)<\/p>/g, '$1');
  // Restore code blocks
  codeBlocks.forEach(function(block, i) {
    out = out.replace(new RegExp('<p>\\x01CODE' + i + '\\x01<\\/p>', 'g'), block);
    out = out.replace(new RegExp('\\x01CODE' + i + '\\x01', 'g'), block);
  });
  return out;
}

// ── MODAL ──────────────────────────────────────────────────────────────
function showModal(title, content) {
  document.getElementById('modal-title').textContent = title;
  document.getElementById('modal-content').textContent = content;
  document.getElementById('modal-overlay').classList.add('open');
}
function closeModal() {
  document.getElementById('modal-overlay').classList.remove('open');
}
function closeModalBg(e) {
  if (e.target === document.getElementById('modal-overlay')) closeModal();
}

// ── IMAGE LIGHTBOX ─────────────────────────────────────────────────────
function openLightbox(src) {
  var lb = document.getElementById('img-lightbox');
  var lbImg = document.getElementById('lightbox-img');
  if (!lb || !lbImg) return;
  lbImg.src = src;
  lb.style.display = 'flex';
  document.addEventListener('keydown', _lbKey);
}
function closeLightbox() {
  var lb = document.getElementById('img-lightbox');
  if (lb) lb.style.display = 'none';
  document.removeEventListener('keydown', _lbKey);
}
function _lbKey(e) { if (e.key === 'Escape') closeLightbox(); }

// ── TOAST ──────────────────────────────────────────────────────────────
var TI = {
  ok:   '<i class="fa-solid fa-circle-check" style="color:#4ADE80"></i>',
  warn: '<i class="fa-solid fa-triangle-exclamation" style="color:#FCD34D"></i>',
  err:  '<i class="fa-solid fa-circle-xmark" style="color:#F87171"></i>'
};
function toast(icon, msg) {
  var el = document.getElementById('toast');
  document.getElementById('toast-icon').innerHTML = icon;
  document.getElementById('toast-msg').textContent = msg;
  el.classList.add('show');
  if (toastTimer) clearTimeout(toastTimer);
  toastTimer = setTimeout(function() { el.classList.remove('show'); }, 5000);
}

// ── UTILS ──────────────────────────────────────────────────────────────
function escHtml(str) {
  return String(str || '')
    .replace(/&/g, '&amp;')
    .replace(/</g, '&lt;')
    .replace(/>/g, '&gt;')
    .replace(/"/g, '&quot;');
}

// ── CALENDAR ──────────────────────────────────────────────────────────
var calYear, calMonth, calSchedule = {};

function calInit() {
  var now = new Date();
  calYear  = now.getFullYear();
  calMonth = now.getMonth();
  var hdr = document.getElementById('cal-day-headers');
  if (hdr) hdr.innerHTML = ['Sun','Mon','Tue','Wed','Thu','Fri','Sat'].map(function(d){ return '<div class="cal-day-hdr">'+d+'</div>'; }).join('');
  fetch('/api/schedule')
    .then(function(r){ return r.json(); })
    .then(function(s){ calSchedule = s; });
}

function calRender() {
  fetch('/api/schedule')
    .then(function(r){ return r.json(); })
    .then(function(s){ calSchedule = s; calDraw(); });
}

function calDraw() {
  var label = document.getElementById('cal-month-label');
  var grid  = document.getElementById('cal-grid');
  if (!label || !grid) return;
  var months = ['January','February','March','April','May','June','July','August','September','October','November','December'];
  label.textContent = months[calMonth] + ' ' + calYear;
  var today   = new Date().toISOString().slice(0,10);
  var first   = new Date(calYear, calMonth, 1);
  var lastDay = new Date(calYear, calMonth+1, 0).getDate();
  var startDow= first.getDay();
  var cells   = [];
  // Prev month padding
  var prevLast = new Date(calYear, calMonth, 0).getDate();
  for (var i = startDow - 1; i >= 0; i--) {
    cells.push({ day: prevLast - i, other: true, iso: '' });
  }
  // This month
  for (var d = 1; d <= lastDay; d++) {
    var mo = String(calMonth + 1).padStart(2,'0');
    var dd = String(d).padStart(2,'0');
    var iso = calYear + '-' + mo + '-' + dd;
    cells.push({ day: d, other: false, iso: iso, isToday: iso===today, isPast: iso<today, img: calSchedule[iso]||'' });
  }
  // Next month padding
  var remaining = 42 - cells.length;
  for (var n = 1; n <= remaining; n++) cells.push({ day: n, other: true, iso: '' });

  grid.innerHTML = cells.map(function(c) {
    var cls = 'cal-cell';
    if (c.other)   cls += ' other-month';
    if (c.isToday) cls += ' today';
    if (c.isPast)  cls += ' past';
    if (c.img)     cls += ' has-post';
    var dot = c.img ? '<div class="cal-dot' + (c.isPast?' past':'') + '"></div>' : '';
    var imgName = c.img ? '<div class="cal-img-name">' + escHtml(c.img.replace(/\.[^.]+$/,'')) + '</div>' : '';
    var onclick = c.img ? ' onclick="calClick(\'' + escHtml(c.iso) + '\',\'' + escHtml(c.img) + '\')"' : '';
    return '<div class="' + cls + '"' + onclick + '>' + dot + '<div class="cal-num">' + c.day + '</div>' + imgName + '</div>';
  }).join('');
}

function calPrev() { calMonth--; if (calMonth < 0) { calMonth = 11; calYear--; } calDraw(); }
function calNext() { calMonth++; if (calMonth > 11) { calMonth = 0;  calYear++; } calDraw(); }
function calClick(iso, img) { showModal(iso, '📸 ' + img + '\n\nScheduled for: ' + iso); }

// ── INTELLIGENCE ──────────────────────────────────────────────────────
var _intelDates = [];

function loadIntelligence(dateFilter) {
  var el = document.getElementById('intel-content');
  el.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading briefing...</p></div>';
  var url = '/api/intelligence' + (dateFilter ? '?date=' + dateFilter : '');
  fetch(url)
    .then(function(r){ if (!r.ok) throw new Error('not found'); return r.json(); })
    .then(function(d){
      _intelDates = d.all_dates || [];
      renderDateList(_intelDates, d.date);
      renderIntelligence(d);
    })
    .catch(function(){
      el.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-magnifying-glass"></i></div><p>No intelligence briefings found.<br><span style="font-size:11px">Run Social Listening to generate a briefing.</span></p></div>';
      renderDateList([], null);
    });
}

function renderDateList(dates, activeDate) {
  var today = new Date().toISOString().slice(0, 10);
  var list  = document.getElementById('intel-date-list');
  if (!dates.length) {
    list.innerHTML = '<div class="intel-list-header">Reports</div><div style="padding:12px 16px;font-size:11px;color:var(--text-dim)">No reports yet</div>';
    return;
  }
  list.innerHTML = '<div class="intel-list-header">Reports</div>' +
    dates.map(function(d) {
      var isActive = d === activeDate;
      var isToday  = d === today;
      var sub = isToday ? 'Today' : '';
      return '<div class="intel-list-item' + (isActive ? ' active' : '') + '" onclick="loadIntelligence(\'' + d + '\')">' +
        '<span class="intel-list-date">' + escHtml(d) + '</span>' +
        (sub ? '<span class="intel-list-sub">' + sub + '</span>' : '') +
      '</div>';
    }).join('');
}

function renderIntelligence(d) {
  var threatClass = { low:'threat-low', medium:'threat-medium', high:'threat-high', unknown:'threat-unknown' }[d.threat] || 'threat-unknown';
  var threatLabel = { low:'Low Threat', medium:'Medium', high:'High Alert', unknown:'Unknown' }[d.threat] || d.threat;

  // Build news cards — plain text, no icons
  var newsHtml = '';
  if (d.news && d.news.length) {
    var newCards = d.news.map(function(n) {
      var m = n.match(/^([^—]+?)\s*—\s*(.+?)\s*\(([^)]+)\)$/);
      if (m) {
        var dateParts = m[1].trim().split(' ');
        var shortDate = dateParts.length >= 2 ? dateParts[1] + ' ' + dateParts[2] : m[1].trim();
        return '<div class="news-card">' +
          '<div class="news-date-chip">' + escHtml(shortDate) + '</div>' +
          '<div class="news-card-body">' +
            '<div class="news-title">' + escHtml(m[2].trim()) + '</div>' +
            '<span class="news-source">' + escHtml(m[3].trim()) + '</span>' +
          '</div></div>';
      }
      return '<div class="news-card"><div class="news-card-body"><div class="news-title">' + escHtml(n) + '</div></div></div>';
    }).join('');
    newsHtml = '<div class="card" style="margin-bottom:12px">' +
      '<div class="card-title" style="margin-bottom:12px">News Coverage &nbsp;<span style="font-weight:400;color:var(--text-dim);font-size:10px">' + d.news.length + ' articles</span></div>' +
      newCards + '</div>';
  }

  // Build changes cards — plain text, no icons
  var changesHtml = '';
  if (d.changes && d.changes.length) {
    var changeCards = d.changes.map(function(c) {
      return '<div class="change-card">' +
        '<div>' +
          '<div class="cc-name">' + escHtml(c.competitor) + '</div>' +
          '<div style="font-size:11px;color:var(--text-dim);margin-top:2px">Website change detected</div>' +
        '</div>' +
        '<a href="' + escHtml(c.url) + '" target="_blank" class="cc-url-link" style="flex-shrink:0;font-size:12px;color:var(--primary);text-decoration:none">View site</a>' +
      '</div>';
    }).join('');
    changesHtml = '<div class="card" style="margin-bottom:12px">' +
      '<div class="card-title" style="margin-bottom:12px">Website Changes Detected</div>' +
      changeCards + '</div>';
  }

  // Sources footer — plain text, no icons
  var sfRows = '';
  sfRows += '<div class="sf-row"><span>Source file: <strong>' + escHtml(d.source_file || d.date + '_intelligence_briefing.txt') + '</strong></span></div>';
  sfRows += '<div class="sf-row"><span>Generated by Social Listening Agent — daily 7:50 AM PHT</span></div>';
  if (d.sources && d.sources.length) {
    sfRows += '<div class="sf-row"><span>News sources: ' + d.sources.map(function(s){ return '<strong>' + escHtml(s) + '</strong>'; }).join(', ') + '</span></div>';
  }
  if (d.changes && d.changes.length) {
    sfRows += '<div class="sf-row"><span>Monitored: ' + d.changes.map(function(c){ return '<a href="'+escHtml(c.url)+'" target="_blank" style="color:var(--primary);text-decoration:none">'+escHtml(c.competitor)+'</a>'; }).join(', ') + '</span></div>';
  }
  var sourcesFooter = '<div class="intel-sources-footer"><div class="sf-title">Data Sources &amp; References</div>' + sfRows + '</div>';

  document.getElementById('intel-content').innerHTML =
    '<div class="intel-header">' +
      '<div class="intel-meta">Briefing: <strong>' + d.date + '</strong></div>' +
      '<span class="threat-badge ' + threatClass + '">' + threatLabel + '</span>' +
    '</div>' +
    '<div class="card" style="margin-bottom:12px">' +
      '<div class="card-title" style="margin-bottom:10px">AI Intelligence Analysis</div>' +
      '<div class="intel-body" id="intel-analysis"></div>' +
    '</div>' +
    changesHtml +
    newsHtml +
    sourcesFooter;

  // Render analysis — strip any leftover ---, →, and Action: labels, no icons
  var raw = md(d.analysis || '');
  // strip any --- / <hr> that slipped through
  raw = raw.replace(/<hr\s*\/?>/gi, '');
  raw = raw.replace(/(<br>)\s*-{3,}\s*(<br>)/g, '$1');
  raw = raw.replace(/^-{3,}<br>/gm, '');
  raw = raw.replace(/^-{3,}$/gm, '');
  raw = raw.replace(/-{3,}/g, '');
  // collapse excessive blank lines (3+ consecutive <br>) to max 1
  raw = raw.replace(/(<br>\s*){3,}/g, '<br>');
  // strip → arrow
  raw = raw.replace(/→\s*/g, '');
  // style Action: as a plain callout box (no icon)
  raw = raw.replace(
    /(<strong>)?Action:(<\/strong>)?\s*/g,
    '</p><div class="intel-action-box"><div class="action-label">Recommended Action</div>'
  );
  raw = raw.replace(/(<h2>)/g, '</div>$1');
  raw = raw.replace(/^<\/div>/, '');
  // Close any unclosed intel-action-box (happens when Action: is in the last section)
  raw += '</div>';
  var analysisEl = document.getElementById('intel-analysis');
  if (analysisEl) analysisEl.innerHTML = raw;
}

// ── CALENDAR TABS ─────────────────────────────────────────────────────
function calSetTab(tab) {
  document.getElementById('cal-month-view').style.display = tab === 'month' ? '' : 'none';
  document.getElementById('cal-list-view').style.display  = tab === 'list'  ? '' : 'none';
  document.getElementById('tab-month').classList.toggle('active', tab === 'month');
  document.getElementById('tab-list').classList.toggle('active', tab === 'list');
  var nav = document.getElementById('cal-month-nav');
  if (nav) nav.style.visibility = tab === 'month' ? 'visible' : 'hidden';
  if (tab === 'list') loadSchedule();
}

// ── CHAT MODEL + MODE ─────────────────────────────────────────────────
function chatModelChanged() { /* model set via pickModel() */ }
function setChatMode(m) {
  chatMode = m;
  document.querySelectorAll('.mode-pill').forEach(function(p) {
    p.classList.toggle('active', p.dataset.mode === m);
  });
}
function getModeHint() {
  if (chatMode === 'deep')     return 'Think step by step. Be comprehensive and thorough in your reasoning.';
  if (chatMode === 'concise')  return 'Be brief and actionable. Max 3-5 bullet points. Skip filler words.';
  if (chatMode === 'creative') return 'Think creatively and propose unconventional or surprising angles.';
  return '';
}

// ── INSTAGRAM MODAL ───────────────────────────────────────────────────
function openIgModal() {
  var overlay = document.getElementById('ig-modal-overlay');
  overlay.classList.add('open');
  document.getElementById('ig-caption').value  = '';
  document.getElementById('ig-hashtags').value = '';
  document.getElementById('ig-gen-btn').innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption';
  document.getElementById('ig-gen-btn').disabled = false;
  var sel = document.getElementById('ig-image-select');
  sel.value = '';
  _loadScheduleIntoSelect(sel, 'ig-cs-wrap', igImageChanged);
}
function closeIgModal() {
  document.getElementById('ig-modal-overlay').classList.remove('open');
}
function closeIgModalBg(e) {
  if (e.target === document.getElementById('ig-modal-overlay')) closeIgModal();
}
function igImageChanged() {
  var img   = document.getElementById('ig-image-select').value;
  var prev  = document.getElementById('ig-image-preview');
  var empty = document.getElementById('ig-img-empty');
  if (img && prev) {
    prev.onload = function() { prev.classList.add('loaded'); if (empty) empty.style.display = 'none'; };
    prev.onerror = function() { prev.classList.remove('loaded'); if (empty) empty.style.display = 'flex'; };
    prev.src = '/api/image-preview/' + encodeURIComponent(img);
  } else {
    if (prev)  { prev.classList.remove('loaded'); prev.src = ''; }
    if (empty) { empty.style.display = 'flex'; }
  }
}
function generateIgCaption() {
  var img = document.getElementById('ig-image-select').value;
  if (!img) { toast(TI.warn, 'Select an image first.'); return; }
  var btn = document.getElementById('ig-gen-btn');
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Generating...';
  btn.disabled = true;
  fetch('/api/generate-ig-caption?image=' + encodeURIComponent(img))
    .then(function(r){ return r.json(); })
    .then(function(d){
      if (d.error) {
        btn.innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption';
        btn.disabled = false;
        toast(TI.err, 'Caption failed: ' + d.error);
        return;
      }
      document.getElementById('ig-caption').value  = d.caption  || '';
      document.getElementById('ig-hashtags').value = d.hashtags || '';
      btn.innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Regenerate';
      btn.disabled = false;
    })
    .catch(function(e){
      btn.innerHTML = '<i class="fa-solid fa-wand-magic-sparkles"></i> Generate AI Caption';
      btn.disabled = false;
      toast(TI.err, 'Caption generation failed: ' + e);
    });
}
function _pollIgJob(jobId) {
  fetch('/api/job/' + jobId)
    .then(function(r){ return r.json(); })
    .then(function(d){
      if (d.status === 'running') {
        setTimeout(function(){ _pollIgJob(jobId); }, 4000);
        return;
      }
      if (d.status === 'ok') {
        toast(TI.ok, 'Instagram post sent successfully!');
      } else {
        toast(TI.err, 'Instagram post failed. Check output.');
        if (d.output) showModal('Post Output', d.output);
      }
    })
    .catch(function(){ toast(TI.err, 'Could not check Instagram job status.'); });
}

function submitIgPost() {
  var caption  = document.getElementById('ig-caption').value.trim();
  var hashtags = document.getElementById('ig-hashtags').value.trim();
  var image    = document.getElementById('ig-image-select').value;
  if (!caption) { toast(TI.warn, 'Please write or generate a caption first.'); return; }
  var btn = document.getElementById('ig-submit-btn');
  btn.disabled = true;
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Posting...';
  fetch('/api/post-instagram', {
    method: 'POST',
    headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({ caption: caption, hashtags: hashtags, image: image })
  })
    .then(function(r){ return r.json(); })
    .then(function(d){
      btn.disabled = false;
      btn.innerHTML = '<i class="fa-brands fa-instagram"></i> Post to Instagram';
      closeIgModal();
      if (d.status === 'started') {
        toast(TI.ok, 'Instagram posting... (downloading image & uploading, may take ~20 sec)');
        _pollIgJob(d.job_id);
      } else if (d.status === 'ok') {
        toast(TI.ok, 'Instagram post sent successfully!');
      } else {
        toast(TI.err, 'Post failed. Check output for details.');
        if (d.output) showModal('Post Output', d.output);
      }
    })
    .catch(function(){
      btn.disabled = false;
      btn.innerHTML = '<i class="fa-brands fa-instagram"></i> Post to Instagram';
      toast(TI.err, 'Could not reach server.');
    });
}

// ── AUTOMATIONS ──────────────────────────────────────────────────────
function loadAutomations() {
  var grid = document.getElementById('auto-grid');
  grid.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading...</p></div>';
  fetch('/api/workflows/status')
    .then(function(r){ return r.json(); })
    .then(function(d){ renderAutomations(d); })
    .catch(function(){ grid.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-triangle-exclamation"></i></div><p>Could not load automations</p></div>'; });
}

function toggleWorkflow(wfId, enable, toggleEl) {
  toggleEl.disabled = true;
  var card = toggleEl.closest('.wf-unified-card');
  var badge = card ? card.querySelector('.wf-enabled-badge') : null;
  fetch('/api/workflows/toggle', {
    method: 'POST', headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({workflow_id: wfId, enable: enable})
  })
  .then(function(r){ return r.json(); })
  .then(function(d){
    toggleEl.disabled = false;
    if (d.success) {
      if (badge) {
        badge.textContent = enable ? 'ON' : 'OFF';
        badge.className = 'wf-enabled-badge ' + (enable ? 'wf-badge-on' : 'wf-badge-off');
      }
      toast('<i class="fa-solid fa-circle-check"></i>', enable ? 'Schedule enabled — will run automatically.' : 'Schedule disabled — manual only.');
    } else {
      toggleEl.checked = !enable;
      toast('<i class="fa-solid fa-triangle-exclamation"></i>', 'Error: ' + (d.error || 'Unknown'));
    }
  })
  .catch(function(){
    toggleEl.disabled = false;
    toggleEl.checked = !enable;
    toast('<i class="fa-solid fa-triangle-exclamation"></i>', 'Could not reach server.');
  });
}

function runWorkflowNow(wfId, btn) {
  btn.disabled = true;
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Starting…';
  fetch('/api/workflows/run', {
    method: 'POST', headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({workflow_id: wfId})
  })
  .then(function(r){ return r.json(); })
  .then(function(d){
    btn.disabled = false;
    btn.innerHTML = '<i class="fa-solid fa-play"></i> Run Now';
    if (d.success) {
      toast('<i class="fa-solid fa-circle-check"></i>', 'Workflow triggered! Check GitHub Actions for progress.');
    } else {
      toast('<i class="fa-solid fa-triangle-exclamation"></i>', 'Error: ' + (d.error || 'Unknown'));
    }
  })
  .catch(function(){
    btn.disabled = false;
    btn.innerHTML = '<i class="fa-solid fa-play"></i> Run Now';
    toast('<i class="fa-solid fa-triangle-exclamation"></i>', 'Could not reach server.');
  });
}

function renderAutomations(workflows) {
  var grid = document.getElementById('auto-grid');
  if (!Array.isArray(workflows) || workflows.error) {
    grid.innerHTML = '<div class="empty-state"><div class="empty-icon"><i class="fa-solid fa-triangle-exclamation"></i></div><p>Could not load workflows</p></div>';
    return;
  }

  var statusMeta = {
    success:   { cls: 'status-success',   label: '<i class="fa-solid fa-circle-check"></i> Ran today' },
    warning:   { cls: 'status-warning',   label: '<i class="fa-solid fa-triangle-exclamation"></i> Check' },
    inactive:  { cls: 'status-inactive',  label: '<i class="fa-regular fa-circle"></i> Idle' },
    scheduled: { cls: 'status-scheduled', label: '<i class="fa-regular fa-clock"></i> Scheduled' },
    pending:   { cls: 'status-pending',   label: '<i class="fa-regular fa-hourglass-half"></i> Pending' },
    none:      { cls: 'status-none',      label: '<i class="fa-regular fa-circle"></i> —' },
  };

  var cards = workflows.map(function(wf) {
    var sm     = statusMeta[wf.live_status] || statusMeta.none;
    var enCls  = wf.enabled ? 'wf-badge-on' : 'wf-badge-off';
    var enLabel= wf.enabled ? 'ON' : 'OFF';
    var chk    = wf.enabled ? 'checked' : '';
    return (
      '<div class="wf-unified-card">' +
        '<div class="wf-uc-top">' +
          '<div class="wf-uc-icon"><i class="' + wf.emoji + '"></i></div>' +
          '<div class="wf-uc-info">' +
            '<div class="wf-uc-name">' + escHtml(wf.name) + '</div>' +
            '<div class="wf-uc-desc">' + escHtml(wf.description || '') + '</div>' +
          '</div>' +
        '</div>' +
        '<div class="wf-uc-mid">' +
          '<span class="status-pill ' + sm.cls + '">' + sm.label + '</span>' +
          (wf.live_detail ? '<span class="wf-uc-detail">' + escHtml(wf.live_detail) + '</span>' : '') +
        '</div>' +
        '<div class="wf-uc-footer">' +
          '<span class="wf-uc-sched"><i class="fa-regular fa-clock"></i> ' + escHtml(wf.schedule) + '</span>' +
          '<div class="wf-uc-actions">' +
            '<button class="wf-run-btn" onclick="runWorkflowNow(\'' + escHtml(wf.id) + '\', this)">' +
              '<i class="fa-solid fa-play"></i> Run Now' +
            '</button>' +
            '<div class="wf-uc-toggle">' +
              '<span class="wf-enabled-badge ' + enCls + '">' + enLabel + '</span>' +
              '<label class="wf-switch">' +
                '<input type="checkbox" ' + chk + ' onchange="toggleWorkflow(\'' + escHtml(wf.id) + '\', this.checked, this)">' +
                '<span class="wf-slider"></span>' +
              '</label>' +
            '</div>' +
          '</div>' +
        '</div>' +
      '</div>'
    );
  }).join('');

  grid.innerHTML = '<div class="wf-unified-grid">' + cards + '</div>';
}

// ── LEAD DETAIL MODAL ─────────────────────────────────────────────────
var _currentLeadEmail = '';
var _currentLeadActive = true;

function openLeadModal(trEl) {
  var raw = trEl.getAttribute('data-lead');
  if (!raw) return;
  var l;
  try { l = JSON.parse(raw); } catch(e) { return; }
  var cnt = parseInt(l.email_count || 0);
  _currentLeadEmail  = l.email || '';
  _currentLeadActive = (l.status || '').toLowerCase() !== 'inactive';
  document.getElementById('lm-name').textContent = l.name || '—';
  var scoreBadge = cnt >= 10
    ? '<span class="badge badge-hot"><i class="fa-solid fa-fire"></i> Hot</span>'
    : cnt >= 5
    ? '<span class="badge badge-warm"><i class="fa-solid fa-temperature-half"></i> Warm</span>'
    : '<span class="badge badge-cold"><i class="fa-solid fa-snowflake"></i> Cold</span>';
  document.getElementById('lm-score-badge').innerHTML = scoreBadge;
  var emailEl = document.getElementById('lm-email');
  emailEl.textContent = l.email || '—';
  emailEl.href = l.email ? 'mailto:' + l.email : '#';
  document.getElementById('lm-program').textContent = l.program || '—';
  document.getElementById('lm-status').innerHTML = _currentLeadActive
    ? '<span class="badge badge-active">Active</span>'
    : '<span class="badge badge-inactive">Inactive</span>';
  var pct = Math.min(100, Math.round(cnt / 30 * 100));
  document.getElementById('lm-drip-fill').style.width = pct + '%';
  document.getElementById('lm-drip-label').textContent = 'Email #' + cnt + ' sent';
  document.getElementById('lm-toggle-label').textContent = _currentLeadActive ? 'Mark Inactive' : 'Mark Active';
  document.getElementById('lm-toggle-btn').querySelector('i').className = _currentLeadActive ? 'fa-solid fa-toggle-off' : 'fa-solid fa-toggle-on';
  document.getElementById('lead-modal-overlay').classList.add('open');
}

function toggleLeadStatus() {
  if (!_currentLeadEmail) return;
  var btn       = document.getElementById('lm-toggle-btn');
  var newStatus = _currentLeadActive ? 'inactive' : 'active';
  btn.disabled  = true;
  fetch('/api/update-lead-status', {
    method: 'POST', headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({email: _currentLeadEmail, status: newStatus})
  })
  .then(function(r){ return r.json(); })
  .then(function(d){
    btn.disabled = false;
    if (d.status === 'updated') {
      _currentLeadActive = !_currentLeadActive;
      document.getElementById('lm-status').innerHTML = _currentLeadActive
        ? '<span class="badge badge-active">Active</span>'
        : '<span class="badge badge-inactive">Inactive</span>';
      document.getElementById('lm-toggle-label').textContent = _currentLeadActive ? 'Mark Inactive' : 'Mark Active';
      btn.querySelector('i').className = _currentLeadActive ? 'fa-solid fa-toggle-off' : 'fa-solid fa-toggle-on';
      toast('<i class="fa-solid fa-circle-check"></i>', d.message);
      setTimeout(function(){ loadLeads(); }, 500);
    } else {
      toast('<i class="fa-solid fa-triangle-exclamation"></i>', d.error || 'Failed to update.');
    }
  })
  .catch(function(){ btn.disabled = false; toast('<i class="fa-solid fa-triangle-exclamation"></i>', 'Could not reach server.'); });
}

function deleteLeadConfirm() {
  if (!_currentLeadEmail) return;
  var name = document.getElementById('lm-name').textContent;
  if (!confirm('Delete ' + name + ' (' + _currentLeadEmail + ')?\n\nThis cannot be undone.')) return;
  fetch('/api/delete-lead', {
    method: 'POST', headers: {'Content-Type': 'application/json'},
    body: JSON.stringify({email: _currentLeadEmail})
  })
  .then(function(r){ return r.json(); })
  .then(function(d){
    if (d.status === 'deleted') {
      closeLeadModal();
      toast('<i class="fa-solid fa-trash"></i>', name + ' deleted.');
      setTimeout(function(){ loadLeads(); }, 500);
    } else {
      toast('<i class="fa-solid fa-triangle-exclamation"></i>', d.error || 'Failed to delete.');
    }
  })
  .catch(function(){ toast('<i class="fa-solid fa-triangle-exclamation"></i>', 'Could not reach server.'); });
}

function closeLeadModal() {
  document.getElementById('lead-modal-overlay').classList.remove('open');
}
function closeLeadModalBg(e) {
  if (e.target === document.getElementById('lead-modal-overlay')) closeLeadModal();
}

// ── POST PREVIEW MODAL ────────────────────────────────────────────────
function openPostPreview() {
  var overlay = document.getElementById('pp-overlay');
  overlay.classList.add('open');
  document.getElementById('pp-body').innerHTML = '<div class="empty-state" style="padding:40px"><div class="empty-icon"><i class="fa-solid fa-spinner fa-spin"></i></div><p>Loading preview...</p></div>';
  document.getElementById('pp-meta').textContent = '';
  fetch('/api/post-preview')
    .then(function(r){ return r.json(); })
    .then(function(d) {
      if (!d.has_image) {
        document.getElementById('pp-body').innerHTML = '<div class="pp-no-image"><i class="fa-solid fa-calendar-xmark" style="font-size:32px;color:var(--text-dim);margin-bottom:12px"></i><br>No post scheduled for today.</div>';
        document.getElementById('pp-meta').textContent = '';
        return;
      }
      var imgTag = '<img src="' + escHtml(d.image_url) + '" alt="Post image">';
      var caption = d.filename ? 'Today\'s post · ' + escHtml(d.date) : '';
      document.getElementById('pp-body').innerHTML =
        '<div class="pp-frames">' +
          '<div class="pp-frame">' +
            '<div class="pp-frame-header fb-pp-header"><i class="fa-brands fa-facebook-f" style="color:#fff;font-size:14px"></i><span class="pp-frame-label">Facebook</span></div>' +
            '<div class="pp-img-wrap">' + imgTag + '</div>' +
            '<div class="pp-caption">Enderun Colleges &mdash; ' + escHtml(d.date) + '</div>' +
          '</div>' +
          '<div class="pp-frame">' +
            '<div class="pp-frame-header ig-pp-header"><i class="fa-brands fa-instagram" style="color:#d62976;font-size:14px"></i><span class="pp-frame-label">Instagram</span></div>' +
            '<div class="pp-img-wrap">' + imgTag + '</div>' +
            '<div class="pp-caption">@enderuncolleges &mdash; ' + escHtml(d.date) + '</div>' +
          '</div>' +
        '</div>';
      var meta = d.days_until === 0 ? 'Scheduled for today' : d.days_until === 1 ? 'Scheduled for tomorrow' : 'Scheduled in ' + d.days_until + ' days';
      document.getElementById('pp-meta').textContent = escHtml(d.filename) + ' · ' + meta;
    })
    .catch(function() {
      document.getElementById('pp-body').innerHTML = '<div class="pp-no-image">Could not load post preview.</div>';
    });
}
function closePpModal() {
  document.getElementById('pp-overlay').classList.remove('open');
}
function closePpModalBg(e) {
  if (e.target === document.getElementById('pp-overlay')) closePpModal();
}

// ── ANALYTICS ─────────────────────────────────────────────────────────
var _analyticsCharts = {};

function _countUp(el, target, duration) {
  if (!el) return;
  var start = 0, step = target / (duration / 16);
  var timer = setInterval(function() {
    start = Math.min(start + step, target);
    el.textContent = Math.round(start);
    if (start >= target) clearInterval(timer);
  }, 16);
}

function loadAnalytics() {
  var today = new Date().toISOString().slice(0,10);
  var thisMonth = today.slice(0,7);

  fetch('/api/stats')
    .then(function(r){ return r.json(); })
    .then(function(stats) {
      var allPosts = stats.upcoming || [];
      var todayD = new Date(today + 'T00:00:00');

      // KPIs with count-up animation
      var totalPosts = allPosts.length;
      var thisMonthPosts = allPosts.filter(function(p){ return p.date.startsWith(thisMonth); }).length;
      var next7 = allPosts.filter(function(p) {
        var diff = (new Date(p.date + 'T00:00:00') - todayD) / 86400000;
        return diff >= 0 && diff < 7;
      }).length;
      _countUp(document.getElementById('akpi-total'),  totalPosts,      700);
      _countUp(document.getElementById('akpi-active'), thisMonthPosts,  700);
      _countUp(document.getElementById('akpi-hot'),    next7,           700);

      // Monthly volume badge
      var badgeM = document.getElementById('badge-monthly');
      if (badgeM) badgeM.textContent = totalPosts + ' total';

      // Monthly post volume bar chart
      var monthCounts = {};
      allPosts.forEach(function(p) {
        var m = p.date.slice(0,7);
        monthCounts[m] = (monthCounts[m] || 0) + 1;
      });
      var monthKeys = Object.keys(monthCounts).sort();
      if (_analyticsCharts.monthly) _analyticsCharts.monthly.destroy();
      _analyticsCharts.monthly = new Chart(document.getElementById('chart-monthly-posts'), {
        type: 'bar',
        data: {
          labels: monthKeys.map(function(m) {
            return new Date(m + '-01T00:00:00').toLocaleString('en', {month:'short', year:'2-digit'});
          }),
          datasets: [{
            data: monthKeys.map(function(m){ return monthCounts[m]; }),
            backgroundColor: monthKeys.map(function(m){
              return m < thisMonth ? 'rgba(122,16,40,0.22)' : m === thisMonth ? '#7a1028' : '#9b1a35';
            }),
            borderRadius: 7,
            borderSkipped: false
          }]
        },
        options: {
          animation: { duration: 800, easing: 'easeOutQuart' },
          plugins: { legend:{ display:false }, tooltip:{ callbacks:{ title: function(i){ return i[0].label; }, label: function(ctx){ return ctx.raw + ' post' + (ctx.raw !== 1 ? 's' : ''); } } } },
          scales: {
            x: { grid:{ display:false }, ticks:{ font:{size:10}, color:'#b08090' } },
            y: { grid:{ color:'rgba(122,16,40,0.05)' }, ticks:{ stepSize:1, color:'#b08090', font:{size:10} }, border:{display:false} }
          },
          maintainAspectRatio: false
        }
      });

      // Upcoming 14-day calendar chart
      var days14 = [];
      for (var i = 0; i < 14; i++) {
        var d = new Date(todayD.getTime() + i * 86400000);
        var ds = d.toISOString().slice(0,10);
        var post = allPosts.find(function(p){ return p.date === ds; });
        days14.push({ label: d.toLocaleString('en', {weekday:'short'}) + '\n' + d.toLocaleString('en', {month:'short', day:'numeric'}), has: post ? 1 : 0, ds: ds, img: post ? post.image : '' });
      }
      var upcoming14count = days14.filter(function(d){ return d.has; }).length;
      var badgeU = document.getElementById('badge-upcoming');
      if (badgeU) badgeU.textContent = upcoming14count + ' posts';

      if (_analyticsCharts.upcoming) _analyticsCharts.upcoming.destroy();
      _analyticsCharts.upcoming = new Chart(document.getElementById('chart-upcoming-posts'), {
        type: 'bar',
        data: {
          labels: days14.map(function(d){ return d.label.split('\n'); }),
          datasets: [{
            data: days14.map(function(d){ return 1; }),
            backgroundColor: days14.map(function(d){
              if (d.ds === today) return '#f39c12';
              return d.has ? '#7a1028' : 'rgba(122,16,40,0.07)';
            }),
            borderRadius: 6,
            borderSkipped: false
          }]
        },
        options: {
          animation: { duration: 900, easing: 'easeOutBounce' },
          plugins: {
            legend: { display:false },
            tooltip: { callbacks: { title: function(i){ return days14[i[0].dataIndex].label.replace('\n',' '); }, label: function(ctx){ var d = days14[ctx.dataIndex]; return d.ds === today ? 'Today' + (d.has ? ' · ' + d.img : ' · No post') : d.has ? '✓ ' + d.img : 'No post scheduled'; } } }
          },
          scales: {
            x: { grid:{ display:false }, ticks:{ font:{size:9}, color: function(ctx){ return days14[ctx.index] && days14[ctx.index].ds === today ? '#f39c12' : '#b08090'; }, maxRotation:0 } },
            y: { display:false, max:1.5 }
          },
          maintainAspectRatio: false
        }
      });
    });

  // Drip engagement charts
  fetch('/api/leads')
    .then(function(r){ return r.json(); })
    .then(function(leads) {
      var totalSent = leads.reduce(function(acc, l){ return acc + parseInt(l.email_count || 0); }, 0);
      var totalReached = leads.filter(function(l){ return parseInt(l.email_count || 0) > 0; }).length;
      var neverReached = leads.filter(function(l){ return parseInt(l.email_count || 0) === 0; }).length;
      var active = leads.filter(function(l){ return (l.status||'').toLowerCase() === 'active'; }).length;

      _countUp(document.getElementById('akpi-posts'), totalSent, 800);
      var badgeDT = document.getElementById('badge-drip-total');
      if (badgeDT) badgeDT.textContent = totalSent + ' total';
      var badgeR = document.getElementById('badge-reach');
      if (badgeR) badgeR.textContent = totalReached + ' reached';

      // Donut centre
      var centreEl = document.getElementById('donut-centre-val');
      if (centreEl) _countUp(centreEl, totalReached, 700);

      // Horizontal bar: top 12 leads by emails received
      var sorted = leads
        .filter(function(l){ return parseInt(l.email_count || 0) > 0; })
        .sort(function(a,b){ return parseInt(b.email_count||0) - parseInt(a.email_count||0); })
        .slice(0, 12);

      var maxEmails = sorted.length ? parseInt(sorted[0].email_count || 0) : 1;

      if (_analyticsCharts.drip) _analyticsCharts.drip.destroy();
      _analyticsCharts.drip = new Chart(document.getElementById('chart-drip-engagement'), {
        type: 'bar',
        data: {
          labels: sorted.map(function(l){
            return ((l.first_name || '').trim() || l.email.split('@')[0]).slice(0, 15);
          }),
          datasets: [{
            label: 'Emails Sent', data: sorted.map(function(l){ return parseInt(l.email_count || 0); }),
            backgroundColor: sorted.map(function(l){
              var v = parseInt(l.email_count || 0) / maxEmails;
              return 'rgba(122,16,40,' + (0.3 + v * 0.7).toFixed(2) + ')';
            }),
            borderRadius: 5, borderSkipped: false
          }]
        },
        options: {
          indexAxis: 'y',
          animation: { duration: 1000, easing: 'easeOutQuart' },
          plugins: { legend:{ display:false }, tooltip:{ callbacks:{ label: function(ctx){ return ctx.raw + ' emails sent'; } } } },
          scales: {
            x: { grid:{ color:'rgba(122,16,40,0.05)' }, ticks:{ font:{size:10}, color:'#b08090' }, border:{display:false} },
            y: { grid:{ display:false }, ticks:{ font:{size:11}, color:'#7a4050' } }
          },
          maintainAspectRatio: false
        }
      });

      // Donut: lead engagement breakdown
      var highEng  = leads.filter(function(l){ return parseInt(l.email_count||0) >= 10; }).length;
      var midEng   = leads.filter(function(l){ var c=parseInt(l.email_count||0); return c>=3&&c<10; }).length;
      var lowEng   = leads.filter(function(l){ var c=parseInt(l.email_count||0); return c>=1&&c<3; }).length;
      if (_analyticsCharts.donut) _analyticsCharts.donut.destroy();
      _analyticsCharts.donut = new Chart(document.getElementById('chart-drip-donut'), {
        type: 'doughnut',
        data: {
          labels: ['High (10+)', 'Mid (3–9)', 'Low (1–2)', 'Not reached'],
          datasets: [{
            data: [highEng, midEng, lowEng, neverReached],
            backgroundColor: ['#7a1028','#c0392b','#e8a0a8','rgba(122,16,40,0.12)'],
            borderWidth: 0, hoverOffset: 8, borderRadius: 4
          }]
        },
        options: {
          cutout: '68%',
          animation: { animateRotate: true, duration: 1000, easing: 'easeOutQuart' },
          plugins: {
            legend: { position:'bottom', labels:{ font:{size:10}, padding:10, color:'#7a4050', boxWidth:10, borderRadius:3 } },
            tooltip: { callbacks: { label: function(ctx){ return ctx.label + ': ' + ctx.raw + ' leads'; } } }
          },
          maintainAspectRatio: false
        }
      });
    });
}

// ── STRATEGY STUDIO ──────────────────────────────────────────────────────────
var _ssType = 'business-plan';
var _ssRawText = '';

function ssSelectType(type, el) {
  _ssType = type;
  document.querySelectorAll('.ss-type-pill').forEach(function(p){ p.classList.remove('ss-active'); });
  el.classList.add('ss-active');
}

function ssGenerate() {
  var program = (document.getElementById('ss-program').value || '').trim();
  var desc    = (document.getElementById('ss-desc').value || '').trim();
  if (!program || !desc) {
    document.getElementById('ss-program').style.borderColor = program ? '' : '#c62a47';
    document.getElementById('ss-desc').style.borderColor    = desc    ? '' : '#c62a47';
    return;
  }
  document.getElementById('ss-program').style.borderColor = '';
  document.getElementById('ss-desc').style.borderColor    = '';

  var market   = (document.getElementById('ss-market').value   || '').trim();
  var formats  = (document.getElementById('ss-formats').value  || '').trim();
  var budget   = (document.getElementById('ss-budget').value   || '').trim();
  var timeline = (document.getElementById('ss-timeline').value || '').trim();
  var notes    = (document.getElementById('ss-notes').value    || '').trim();

  // Reset output
  _ssRawText = '';
  document.getElementById('ss-empty-state').style.display    = 'none';
  document.getElementById('ss-stream-area').style.display    = '';
  document.getElementById('ss-output-actions').style.display = 'none';
  document.getElementById('ss-generating-indicator').style.display = '';
  document.getElementById('ss-doc-output').innerHTML = '';
  document.getElementById('ss-output-title').textContent = 'Generating…';
  document.getElementById('ss-word-count').textContent = '';

  var btn = document.getElementById('ss-generate-btn');
  btn.disabled = true;
  btn.innerHTML = '<i class="fa-solid fa-spinner fa-spin"></i> Generating…';

  fetch('/api/strategy/generate', {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify({
      type: _ssType, program: program, description: desc,
      target_market: market, formats: formats, budget: budget,
      timeline: timeline, notes: notes
    })
  }).then(function(res) {
    if (!res.ok || !res.body) throw new Error('Stream unavailable');
    var reader = res.body.getReader();
    var decoder = new TextDecoder();
    var buf = '';
    function read() {
      reader.read().then(function(r) {
        if (r.done) {
          _ssOnStreamDone();
          return;
        }
        buf += decoder.decode(r.value, { stream: true });
        var lines = buf.split('\n');
        buf = lines.pop();
        lines.forEach(function(line) {
          if (!line.startsWith('data: ')) return;
          var payload = line.slice(6);
          if (payload === '[DONE]') { _ssOnStreamDone(); return; }
          try {
            var d = JSON.parse(payload);
            if (d.text) {
              _ssRawText += d.text;
              document.getElementById('ss-doc-output').innerHTML = md(_ssRawText) + '<span class="ss-cursor"></span>';
              document.getElementById('ss-output-body').scrollTop = 9999;
            }
          } catch(e) {}
        });
        read();
      }).catch(function(){ _ssOnStreamDone(); });
    }
    read();
  }).catch(function(err) {
    document.getElementById('ss-doc-output').innerHTML = '<p style="color:#c62a47"><i class="fa-solid fa-circle-xmark"></i> Error: ' + err.message + '</p>';
    _ssOnStreamDone();
  });
}

function _ssOnStreamDone() {
  document.getElementById('ss-generating-indicator').style.display = 'none';
  var el = document.getElementById('ss-doc-output');
  el.innerHTML = _ssRawText ? md(_ssRawText) : el.innerHTML.replace('<span class="ss-cursor"></span>', '');
  var words = _ssRawText.trim().split(/\s+/).filter(Boolean).length;
  document.getElementById('ss-word-count').textContent = words ? words + ' words' : '';
  document.getElementById('ss-output-title').textContent = 'Generated Document';
  if (_ssRawText) document.getElementById('ss-output-actions').style.display = '';
  var btn = document.getElementById('ss-generate-btn');
  btn.disabled = false;
  btn.innerHTML = '<i class="fa-solid fa-bolt"></i> Generate';
}

function ssCopy() {
  if (!_ssRawText) return;
  navigator.clipboard.writeText(_ssRawText).then(function() {
    var btn = event.currentTarget;
    var orig = btn.innerHTML;
    btn.innerHTML = '<i class="fa-solid fa-check"></i> Copied!';
    setTimeout(function(){ btn.innerHTML = orig; }, 1800);
  });
}

function ssExportPDF() {
  if (!_ssRawText) return;
  var typeLabel = { 'business-plan': 'Business Plan', 'marketing-strategy': 'Marketing Strategy', 'strategic-plan': 'Strategic Plan' }[_ssType] || 'Strategy Document';
  var program = (document.getElementById('ss-program').value || 'Document').trim();
  fetch('/api/generate-doc', {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify({ type: 'pdf', content: '# ' + program + ' — ' + typeLabel + '\n\n' + _ssRawText })
  }).then(function(r){ return r.json(); }).then(function(d) {
    if (d.status === 'ok') {
      var a = document.createElement('a');
      a.href = '/chat-files/' + d.filename;
      a.download = d.filename;
      a.click();
    } else {
      alert('PDF export failed: ' + (d.message || 'unknown error'));
    }
  }).catch(function(e){ alert('Export error: ' + e.message); });
}

function ssOpenInChat() {
  if (!_ssRawText) return;
  var program = (document.getElementById('ss-program').value || '').trim();
  var typeLabel = { 'business-plan': 'Business Plan', 'marketing-strategy': 'Marketing Strategy', 'strategic-plan': 'Strategic Plan' }[_ssType] || 'Strategy Document';
  showView('dashboard', document.querySelector('[data-view="dashboard"]'));
  var inp = document.getElementById('chat-input');
  if (inp) {
    inp.value = 'Here is the generated ' + typeLabel + ' for ' + program + ':\n\n' + _ssRawText.slice(0, 2000) + ((_ssRawText.length > 2000) ? '\n\n[...document truncated for chat. Full version available in Strategy Studio.]' : '');
    inp.dispatchEvent(new Event('input'));
    inp.focus();
  }
}
</script>
<!-- Image lightbox -->
<div id="img-lightbox" onclick="closeLightbox()">
  <button id="lightbox-close" onclick="closeLightbox()">&#x2715;</button>
  <img id="lightbox-img" src="" alt="Preview" onclick="event.stopPropagation()">
</div>
<!-- ── KEYBOARD SHORTCUTS MODAL ── -->
<div class="shortcuts-overlay" id="shortcuts-overlay" onclick="if(event.target===this)closeShortcuts()">
  <div class="shortcuts-modal">
    <div class="sc-modal-title"><i class="fa-solid fa-keyboard" style="color:var(--primary)"></i> Keyboard Shortcuts</div>
    <button class="sc-close" onclick="closeShortcuts()"><i class="fa-solid fa-xmark"></i></button>
    <div class="sc-group-title">Chat</div>
    <div class="sc-row"><span class="sc-desc">Send message</span><div class="sc-keys"><span class="sc-key">Enter</span></div></div>
    <div class="sc-row"><span class="sc-desc">New line</span><div class="sc-keys"><span class="sc-key">Shift</span><span class="sc-key">Enter</span></div></div>
    <div class="sc-row"><span class="sc-desc">New conversation</span><div class="sc-keys"><span class="sc-key">Ctrl</span><span class="sc-key">Shift</span><span class="sc-key">N</span></div></div>
    <div class="sc-row"><span class="sc-desc">Stop generation</span><div class="sc-keys"><span class="sc-key">Escape</span></div></div>
    <div class="sc-row"><span class="sc-desc">Focus input</span><div class="sc-keys"><span class="sc-key">Ctrl</span><span class="sc-key">L</span></div></div>
    <div class="sc-row"><span class="sc-desc">Search in chat</span><div class="sc-keys"><span class="sc-key">Ctrl</span><span class="sc-key">F</span></div></div>
    <div class="sc-group-title">Interface</div>
    <div class="sc-row"><span class="sc-desc">Toggle dark mode</span><div class="sc-keys"><span class="sc-key">Ctrl</span><span class="sc-key">Shift</span><span class="sc-key">D</span></div></div>
    <div class="sc-row"><span class="sc-desc">Show this panel</span><div class="sc-keys"><span class="sc-key">?</span></div></div>
    <div class="sc-row"><span class="sc-desc">Export as PDF</span><div class="sc-keys"><span class="sc-key">Ctrl</span><span class="sc-key">Shift</span><span class="sc-key">E</span></div></div>
    <div class="sc-group-title">Navigation</div>
    <div class="sc-row"><span class="sc-desc">Go to Dashboard</span><div class="sc-keys"><span class="sc-key">Alt</span><span class="sc-key">1</span></div></div>
    <div class="sc-row"><span class="sc-desc">Go to Leads</span><div class="sc-keys"><span class="sc-key">Alt</span><span class="sc-key">2</span></div></div>
    <div class="sc-row"><span class="sc-desc">Go to Schedule</span><div class="sc-keys"><span class="sc-key">Alt</span><span class="sc-key">3</span></div></div>
    <div class="sc-row"><span class="sc-desc">Go to Intelligence</span><div class="sc-keys"><span class="sc-key">Alt</span><span class="sc-key">4</span></div></div>
  </div>
</div>
</body>
</html>"""

if __name__ == "__main__":
    HTML_FILE.write_text(HTML, encoding="utf-8")
    print("=" * 60)
    print("  Enderun Marketing Hub")
    print("=" * 60)
    print("  Open: http://localhost:8080")
    print("  Stop: Ctrl+C")
    print("=" * 60)
    port = int(os.environ.get("PORT", 8080))
    app.run(host="0.0.0.0", port=port, debug=False)
